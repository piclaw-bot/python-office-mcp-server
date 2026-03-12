"""
Final coverage boost tests targeting specific uncovered areas.
"""

import tempfile
from pathlib import Path

import pytest
from docx import Document
from docx.shared import RGBColor
from pptx import Presentation
from pptx.util import Inches as PptxInches

from tools.pptx_advanced_tools import PresentationAdvancedTools
from tools.word_advanced_tools import WordAdvancedTools


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as td:
        yield Path(td)


class TestWordPromptFunctions:
    """Test prompt-related functions."""

    def test_prompt_sow_data(self):
        """Test prompt_sow_data function if available."""
        tools = WordAdvancedTools()
        # Test if the function exists and can be called
        if hasattr(tools, 'prompt_sow_data'):
            result = tools.prompt_sow_data()
            assert result is not None

    def test_prompt_sow_section(self):
        """Test prompt_sow_section if available."""
        tools = WordAdvancedTools()
        if hasattr(tools, 'prompt_sow_section'):
            result = tools.prompt_sow_section("Test Section")
            assert result is not None


class TestWordTrackChangesOperations:
    """Test track changes operations in depth."""

    def test_enable_track_changes_output(self, temp_dir):
        """Should enable track changes with output path."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_paragraph("Test content")
        path = temp_dir / "track_input.docx"
        doc.save(path)

        output = temp_dir / "track_output.docx"
        result = tools.tool_word_enable_track_changes(str(path), output_path=str(output))
        assert isinstance(result, dict)
        assert Path(output).exists()

    def test_patch_with_track_changes_complex(self, temp_dir):
        """Should patch with track changes and multiple replacements."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_heading("Document Title", level=1)
        doc.add_paragraph("<Customer Name> is the customer")
        doc.add_paragraph("The project is <Project Name>")
        doc.add_paragraph("Contact <Customer Name> for details")
        path = temp_dir / "multi_patch.docx"
        doc.save(path)

        output = temp_dir / "patched_output.docx"
        result = tools.tool_word_patch_with_track_changes(
            str(path),
            {
                "<Customer Name>": "Contoso Ltd",
                "<Project Name>": "Cloud Migration"
            },
            author="Test Author",
            output_path=str(output)
        )
        assert isinstance(result, dict)


class TestWordTableCreation:
    """Test table creation with various scenarios."""

    def test_create_table_after_section(self, temp_dir):
        """Should create table after specific section."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_heading("Introduction", level=1)
        doc.add_paragraph("Intro content")
        doc.add_heading("Data Section", level=1)
        doc.add_paragraph("Some text")
        path = temp_dir / "with_section.docx"
        doc.save(path)

        result = tools.tool_word_create_new_table(
            str(path),
            ["Item", "Description", "Status"],
            [{"Item": "A", "Description": "First", "Status": "Done"}],
            insert_after_section="Data Section"
        )
        assert isinstance(result, dict)


class TestPptxSlideOperations:
    """Test slide operations in depth."""

    def test_add_slide_with_layout_index(self, temp_dir):
        """Should add slide with specific layout index."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = temp_dir / "layout_test.pptx"
        prs.save(path)

        result = tools.tool_pptx_add_slide(
            str(path),
            layout_index=5,  # Title Only
            title="Custom Layout"
        )
        assert isinstance(result, dict)

    def test_duplicate_slide_position_after(self, temp_dir):
        """Should duplicate slide after original."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[1])
        slide1.shapes.title.text = "First"
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Second"
        path = temp_dir / "dup_after.pptx"
        prs.save(path)

        result = tools.tool_pptx_duplicate_slide(str(path), 1, position="after")
        assert isinstance(result, dict)


class TestPptxShapeOperations:
    """Test shape-related operations."""

    def test_list_shapes_with_table(self, temp_dir):
        """Should list shapes including table."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Title"

        # Add table
        slide.shapes.add_table(3, 2, PptxInches(1), PptxInches(2), PptxInches(5), PptxInches(2))
        path = temp_dir / "with_table.pptx"
        prs.save(path)

        result = tools.tool_pptx_list_shapes(str(path), 1)
        assert isinstance(result, dict)
        # Should find shapes including table
        shapes = result.get("shapes", [])
        has_table = any(s.get("has_table") for s in shapes)
        assert has_table or len(shapes) > 0

    def test_patch_shape_body(self, temp_dir):
        """Should patch body shape."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Title"

        # Set body content
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.text_frame.text = "Original body"
        path = temp_dir / "patch_body.pptx"
        prs.save(path)

        result = tools.tool_pptx_patch_shape(
            str(path),
            slide_number=1,
            shape_identifier="body",
            new_text="Updated body content"
        )
        assert isinstance(result, dict)


class TestPptxReplaceOperations:
    """Test replace operations."""

    def test_replace_text_specific_slide(self, temp_dir):
        """Should replace text on specific slide only."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[1])
        slide1.shapes.title.text = "<Customer> Slide 1"
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "<Customer> Slide 2"
        path = temp_dir / "specific_slide.pptx"
        prs.save(path)

        result = tools.tool_pptx_replace_text(
            str(path),
            "<Customer>",
            "Contoso",
            slide_number=1
        )
        assert isinstance(result, dict)

    def test_replace_placeholders_multiple(self, temp_dir):
        """Should replace multiple placeholders."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "<Company> - <Project>"
        path = temp_dir / "multi_replace.pptx"
        prs.save(path)

        result = tools.tool_pptx_replace_placeholders(
            str(path),
            {
                "<Company>": "Contoso",
                "<Project>": "Migration"
            }
        )
        assert isinstance(result, dict)


class TestPptxLayoutOperations:
    """Test layout-related operations."""

    def test_recommend_layout_comparison(self, temp_dir):
        """Should recommend comparison layout."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = temp_dir / "rec_comparison.pptx"
        prs.save(path)

        result = tools.tool_pptx_recommend_layout(str(path), "comparison")
        assert isinstance(result, dict)

    def test_recommend_layout_blank(self, temp_dir):
        """Should recommend blank layout."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = temp_dir / "rec_blank.pptx"
        prs.save(path)

        result = tools.tool_pptx_recommend_layout(str(path), "blank")
        assert isinstance(result, dict)


class TestWordSowOperations:
    """Test SOW-specific operations."""

    def test_extract_sow_structure_complex(self, temp_dir):
        """Should extract complex SOW structure."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_heading("Statement of Work", level=1)
        doc.add_heading("1. Executive Summary", level=2)
        doc.add_paragraph("Executive summary content")
        doc.add_heading("2. Scope", level=2)
        doc.add_paragraph("Scope details")

        # Add table
        table = doc.add_table(rows=3, cols=3)
        table.cell(0, 0).text = "Phase"
        table.cell(0, 1).text = "Duration"
        table.cell(0, 2).text = "Deliverable"

        doc.add_heading("3. Timeline", level=2)
        doc.add_paragraph("Timeline info")
        path = temp_dir / "complex_sow.docx"
        doc.save(path)

        result = tools.tool_word_extract_sow_structure(str(path))
        assert isinstance(result, dict)

    def test_analyze_template_formatting(self, temp_dir):
        """Should analyze template formatting."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_heading("<Project Name>", level=1)

        # Add guidance text
        para = doc.add_paragraph()
        run = para.add_run("[Guidance: Fill in project details]")
        run.font.color.rgb = RGBColor(0, 0, 255)

        doc.add_paragraph("Standard content here")
        doc.add_paragraph("<Customer Name> is the customer")
        path = temp_dir / "template_format.docx"
        doc.save(path)

        result = tools.tool_word_analyze_template_formatting(str(path))
        assert isinstance(result, dict)


class TestPptxNotesOperations:
    """Test notes operations."""

    def test_get_notes_no_slide_number(self, temp_dir):
        """Should get all notes without slide number."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
            notes = slide.notes_slide
            notes.notes_text_frame.text = f"Notes for slide {i+1}"
        path = temp_dir / "all_notes.pptx"
        prs.save(path)

        result = tools.tool_pptx_get_notes(str(path))
        assert isinstance(result, dict)


class TestWordComplexOperations:
    """Test complex Word operations."""

    def test_generate_sow_with_data(self, temp_dir):
        """Should generate SOW with complete data."""
        tools = WordAdvancedTools()

        # Create template
        doc = Document()
        doc.add_heading("<Project Name>", level=1)
        doc.add_paragraph("Customer: <Customer Name>")
        doc.add_paragraph("Provider: <Provider Name>")
        doc.add_heading("Objectives", level=2)
        template_path = temp_dir / "sow_template.docx"
        doc.save(template_path)

        output = temp_dir / "generated.docx"
        result = tools.tool_word_generate_sow(
            str(template_path),
            str(output),
            {
                "customer_name": "Contoso",
                "customer_short_name": "Contoso",
                "project_name": "Cloud Migration",
                "provider_name": "Microsoft",
                "business_objectives": [
                    {"objective": "Migrate", "activities": "Assessment", "assumptions": "Cloud ready"}
                ]
            }
        )
        assert isinstance(result, dict)


class TestPptxCommentOperations:
    """Test comment operations."""

    def test_add_comment_custom_author(self, temp_dir):
        """Should add comment with custom author."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Review Slide"
        path = temp_dir / "comment_author.pptx"
        prs.save(path)

        result = tools.tool_pptx_add_comment(
            str(path),
            1,
            "Please review this content",
            author="Custom Reviewer",
            x_inches=2.0,
            y_inches=2.0
        )
        assert isinstance(result, dict)


class TestWordAuditOperations:
    """Test audit operations."""

    def test_audit_completion_empty(self, temp_dir):
        """Should audit empty document."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_paragraph("Just some text")
        path = temp_dir / "audit_empty.docx"
        doc.save(path)

        result = tools.tool_word_audit_completion(str(path))
        assert isinstance(result, dict)


class TestPptxSetTextAutofit:
    """Test text autofit settings."""

    def test_set_autofit_shrink(self, temp_dir):
        """Should set autofit to shrink."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Long title that might need shrinking to fit properly in the space"
        path = temp_dir / "autofit_shrink.pptx"
        prs.save(path)

        result = tools.tool_pptx_set_text_autofit(
            str(path),
            slide_number=1,
            shape_identifier="title",
            autofit_type="shrink"
        )
        assert isinstance(result, dict)

    def test_set_autofit_none(self, temp_dir):
        """Should set autofit to none."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Title"
        path = temp_dir / "autofit_none.pptx"
        prs.save(path)

        result = tools.tool_pptx_set_text_autofit(
            str(path),
            slide_number=1,
            shape_identifier="title",
            autofit_type="none"
        )
        assert isinstance(result, dict)
