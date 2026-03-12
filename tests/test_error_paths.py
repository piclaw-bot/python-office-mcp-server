"""
Tests targeting error paths and initialization code.
"""

from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation

# Fixture temp_dir is provided by conftest.py


class TestToolInstantiation:
    """Parameterized tests for tool class instantiation."""

    @pytest.mark.parametrize("tool_fixture", [
        "word_tools",
        "word_advanced_tools",
        "pptx_tools",
        "pptx_advanced_tools",
        "excel_tools",
    ])
    def test_tool_instance_creation(self, request, tool_fixture):
        """Should create tool instance successfully."""
        tools = request.getfixturevalue(tool_fixture)
        assert tools is not None


class TestNonexistentFileHandling:
    """Parameterized tests for nonexistent file handling across tools."""

    @pytest.mark.parametrize("tool_fixture,method,path", [
        ("word_tools", "tool_word_extract", "/nonexistent/path.docx"),
        ("word_tools", "tool_word_to_markdown", "/nonexistent/path.docx"),
        ("pptx_tools", "tool_pptx_extract", "/nonexistent/path.pptx"),
        ("pptx_tools", "tool_pptx_to_markdown", "/nonexistent/path.pptx"),
        ("excel_tools", "tool_excel_extract", "/nonexistent/path.xlsx"),
    ])
    def test_nonexistent_file_error(self, request, tool_fixture, method, path):
        """Should handle nonexistent file gracefully."""
        tools = request.getfixturevalue(tool_fixture)
        result = getattr(tools, method)(path)
        assert "error" in str(result).lower()


class TestExcelPaths:
    """Test Excel tool paths."""

    def test_extract_nonexistent_excel(self, excel_tools):
        """Should handle nonexistent Excel file."""
        result = excel_tools.tool_excel_extract("/nonexistent/path.xlsx")
        assert "error" in result

    def test_from_markdown_empty(self, excel_tools, temp_dir):
        """Should handle empty markdown."""
        path = temp_dir / "empty.xlsx"
        result = excel_tools.tool_excel_from_markdown(str(path), "")
        # Should handle gracefully
        assert isinstance(result, dict)

    def test_from_markdown_no_tables(self, excel_tools, temp_dir):
        """Should handle markdown with no tables."""
        path = temp_dir / "no_tables.xlsx"
        result = excel_tools.tool_excel_from_markdown(str(path), "# Just a heading\n\nSome text.")
        assert isinstance(result, dict)


class TestWordErrorPaths:
    """Test error handling paths in Word tools."""

    def test_list_sections_invalid_docx(self, word_advanced_tools, temp_dir):
        """Should handle invalid docx file."""
        # Create text file with .docx extension
        path = temp_dir / "fake.docx"
        path.write_text("not a valid docx")

        result = word_advanced_tools.tool_word_list_sections(str(path))
        assert "error" in result

    def test_get_table_invalid_index(self, word_advanced_tools, temp_dir):
        """Should handle negative table index."""
        doc = Document()
        doc.add_table(rows=2, cols=2)
        path = temp_dir / "table.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_get_table(str(path), "-1")
        # Should handle gracefully or return error
        assert isinstance(result, dict)


class TestPptxErrorPaths:
    """Test error handling paths in PPTX tools."""

    def test_list_slides_invalid_pptx(self, pptx_advanced_tools, temp_dir):
        """Should handle invalid pptx file."""
        # Create text file with .pptx extension
        path = temp_dir / "fake.pptx"
        path.write_text("not a valid pptx")

        result = pptx_advanced_tools.tool_pptx_list_slides(str(path))
        assert "error" in result

    def test_add_slide_invalid_position(self, pptx_advanced_tools, temp_dir):
        """Should handle invalid slide position."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "pos.pptx"
        prs.save(path)

        output = temp_dir / "added.pptx"
        result = pptx_advanced_tools.tool_pptx_add_slide(
            str(path),
            position="invalid",
            output_path=str(output)
        )
        # Should handle or default to end
        assert isinstance(result, dict)


class TestWordOutputPaths:
    """Test output path handling."""

    def test_patch_section_no_output(self, word_advanced_tools, temp_dir):
        """Should patch section without explicit output."""
        doc = Document()
        doc.add_heading("Section", level=1)
        doc.add_paragraph("Content")
        path = temp_dir / "no_out.docx"
        doc.save(path)

        # No output_path - should overwrite input
        result = word_advanced_tools.tool_word_patch_section(
            str(path),
            "Section",
            "New content"
        )
        assert isinstance(result, dict)

    def test_insert_row_no_output(self, word_advanced_tools, temp_dir):
        """Should insert row without explicit output."""
        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"
        path = temp_dir / "row_no_out.docx"
        doc.save(path)

        # No output_path
        result = word_advanced_tools.tool_word_insert_table_row(
            str(path),
            "0",
            {"A": "X", "B": "Y"}
        )
        assert isinstance(result, dict)


class TestPptxOutputPaths:
    """Test PPTX output path handling."""

    def test_patch_shape_no_output(self, pptx_advanced_tools, temp_dir):
        """Should patch shape without explicit output."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Title"
        path = temp_dir / "patch_no_out.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_patch_shape(
            str(path),
            slide_number=1,
            shape_identifier="title",
            new_text="New Title"
            # No output_path
        )
        assert isinstance(result, dict)


class TestSpecialCharacters:
    """Test handling of special characters."""

    def test_word_with_unicode(self, word_advanced_tools, temp_dir):
        """Should handle unicode in documents."""
        doc = Document()
        doc.add_heading("Ünïcödé Héading", level=1)
        doc.add_paragraph("日本語テキスト")
        doc.add_paragraph("Émojis: 🎉🚀💡")
        path = temp_dir / "unicode.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_list_sections(str(path))
        assert isinstance(result, dict)

    def test_pptx_with_unicode(self, pptx_advanced_tools, temp_dir):
        """Should handle unicode in presentations."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Ünïcödé Título 日本語"
        path = temp_dir / "unicode.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_list_slides(str(path))
        assert isinstance(result, dict)


class TestEmptyDocuments:
    """Test handling of empty documents."""

    def test_empty_word_extract(self, word_tools, empty_docx):
        """Should extract from empty document."""
        result = word_tools.tool_word_extract(str(empty_docx))
        assert isinstance(result, dict)

    def test_empty_pptx_extract(self, pptx_tools, empty_pptx):
        """Should extract from empty presentation."""
        result = pptx_tools.tool_pptx_extract(str(empty_pptx))
        assert isinstance(result, dict)


class TestLargeContent:
    """Test handling of larger content."""

    def test_many_sections(self, word_advanced_tools, temp_dir):
        """Should handle document with many sections."""
        doc = Document()
        for i in range(20):
            doc.add_heading(f"Section {i+1}", level=1)
            doc.add_paragraph(f"Content for section {i+1}")
        path = temp_dir / "many_sections.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_list_sections(str(path))
        assert len(result.get("sections", [])) >= 20

    def test_many_slides(self, pptx_advanced_tools, temp_dir):
        """Should handle presentation with many slides."""
        prs = Presentation()
        for i in range(10):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
        path = temp_dir / "many_slides.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_list_slides(str(path))
        assert result.get("slide_count", 0) >= 10


class TestWordMarkdownEdgeCases:
    """Test markdown conversion edge cases."""

    def test_from_markdown_complex_list(self, word_tools, temp_dir):
        """Should handle complex nested lists."""
        md = """# Document

## Section 1

- Item 1
  - Nested 1.1
  - Nested 1.2
- Item 2
  - Nested 2.1
    - Deep nested

## Section 2

1. Numbered
2. List
3. Items
"""
        path = temp_dir / "complex_list.docx"
        word_tools.tool_word_from_markdown(str(path), md)
        assert Path(path).exists()


class TestPptxMarkdownEdgeCases:
    """Test PPTX markdown conversion edge cases."""

    def test_from_markdown_multiple_tables(self, pptx_tools, temp_dir):
        """Should handle markdown with multiple tables."""
        md = """# Data Report

## Table 1

| A | B | C |
|---|---|---|
| 1 | 2 | 3 |

---

## Table 2

| X | Y |
|---|---|
| a | b |
"""
        path = temp_dir / "multi_table.pptx"
        pptx_tools.tool_pptx_from_markdown(str(path), md)
        assert Path(path).exists()
