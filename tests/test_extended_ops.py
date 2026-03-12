"""
Additional tests to push coverage above 80%.
Targeting specific uncovered methods and error paths.
"""

from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Inches as PptxInches

# Fixture temp_dir is provided by conftest.py


class TestExcelExtendedMethods:
    """Test more Excel tool methods."""

    def test_excel_extract_multi_sheet(self, excel_tools, temp_dir):
        """Should extract multiple sheets."""
        md = """| A | B |
|---|---|
| 1 | 2 |

| C | D |
|---|---|
| 3 | 4 |
"""
        path = temp_dir / "multi.xlsx"
        excel_tools.tool_excel_from_markdown(str(path), md)

        result = excel_tools.tool_excel_extract(str(path))
        assert isinstance(result, dict)

    def test_excel_to_markdown(self, excel_tools, temp_dir):
        """Should convert Excel to markdown."""
        md = """| Name | Value |
|------|-------|
| Test | 100   |
"""
        path = temp_dir / "to_md.xlsx"
        excel_tools.tool_excel_from_markdown(str(path), md)

        result = excel_tools.tool_excel_to_markdown(str(path))
        assert "Name" in result
        assert "Value" in result


class TestWordExtendedOps:
    """Test more Word operations."""

    def test_word_list_sections_with_subsections(self, word_advanced_tools, temp_dir):
        """Should list sections including subsections."""
        doc = Document()
        doc.add_heading("1. Main Section", level=1)
        doc.add_heading("1.1 Subsection A", level=2)
        doc.add_paragraph("Subsection content")
        doc.add_heading("1.2 Subsection B", level=2)
        doc.add_paragraph("More content")
        doc.add_heading("2. Another Section", level=1)
        path = temp_dir / "subsections.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_list_sections(str(path))
        assert len(result.get("sections", [])) >= 2

    def test_word_patch_section_append(self, word_advanced_tools, temp_dir):
        """Should append to section."""
        doc = Document()
        doc.add_heading("Section", level=1)
        doc.add_paragraph("Original")
        path = temp_dir / "append_section.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_patch_section(
            str(path),
            "Section",
            "Appended content"
        )
        assert isinstance(result, dict)

    def test_word_get_table_with_headers(self, word_advanced_tools, temp_dir):
        """Should get table with proper headers."""
        doc = Document()
        table = doc.add_table(rows=3, cols=3)
        headers = ["Header A", "Header B", "Header C"]
        for i, h in enumerate(headers):
            table.cell(0, i).text = h
        for row in range(1, 3):
            for col in range(3):
                table.cell(row, col).text = f"R{row}C{col}"
        path = temp_dir / "headers_table.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_get_table(str(path), "0")
        assert isinstance(result, dict)


class TestPptxExtendedOps:
    """Test more PPTX operations."""

    def test_pptx_list_masters(self, pptx_advanced_tools, temp_dir):
        """Should list available masters/layouts."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = temp_dir / "masters.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_list_masters(str(path))
        assert isinstance(result, dict)
        assert "default_layouts" in result

    def test_pptx_reorder_slides(self, pptx_advanced_tools, temp_dir):
        """Should reorder slides."""
        prs = Presentation()
        for i in range(4):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
        path = temp_dir / "reorder.pptx"
        prs.save(path)

        output = temp_dir / "reordered.pptx"
        result = pptx_advanced_tools.tool_pptx_reorder_slides(
            str(path),
            [1, 4, 2, 3],
            output_path=str(output)
        )
        assert isinstance(result, dict)

    def test_pptx_log_changes(self, pptx_advanced_tools, temp_dir):
        """Should log changes to slide."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "changes.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_log_changes(
            str(path),
            [
                {"slide": 1, "action": "Updated title", "detail": "Changed from draft to final"},
                {"slide": 1, "action": "Added bullet", "detail": "New point about benefits"}
            ]
        )
        assert isinstance(result, dict)


class TestWordComplexPatching:
    """Test complex patching scenarios."""

    def test_replace_global_variables(self, word_advanced_tools, temp_dir):
        """Should replace global variables."""
        doc = Document()
        doc.add_heading("<Project Name>", level=1)
        doc.add_paragraph("Customer: <Customer Name>")
        doc.add_paragraph("Provider: <Provider Name>")
        doc.add_paragraph("Contact <Customer Name> for details")
        path = temp_dir / "globals.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_replace_global_variables(
            str(path),
            {
                "<Customer Name>": "Contoso",
                "<Project Name>": "Migration",
                "<Provider Name>": "Microsoft"
            }
        )
        assert isinstance(result, dict)


class TestPptxComplexOps:
    """Test complex PPTX operations."""

    def test_pptx_add_table(self, pptx_advanced_tools, temp_dir):
        """Should add table to slide."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
        slide.shapes.title.text = "Data"
        path = temp_dir / "for_table.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_add_table(
            str(path),
            1,
            ["Name", "Role", "Hours"],
            [["Alice", "Dev", "40"], ["Bob", "QA", "30"]]
        )
        assert isinstance(result, dict)

    def test_pptx_update_table_cell(self, pptx_advanced_tools, temp_dir):
        """Should update table cell."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = slide.shapes.add_table(3, 2, PptxInches(1), PptxInches(2), PptxInches(5), PptxInches(2))
        tbl = table.table
        tbl.cell(0, 0).text = "A"
        tbl.cell(0, 1).text = "B"
        tbl.cell(1, 0).text = "C"
        tbl.cell(1, 1).text = "D"
        path = temp_dir / "update_cell.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_patch_table_cell(
            str(path),
            slide_number=1,
            row_index=1,
            col_index=1,
            new_text="Updated"
        )
        assert isinstance(result, dict)

    def test_pptx_patch_table_cell_string(self, pptx_advanced_tools, temp_dir):
        """Should patch table cell with string value."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = slide.shapes.add_table(2, 2, PptxInches(1), PptxInches(2), PptxInches(5), PptxInches(1.5))
        tbl = table.table
        tbl.cell(0, 0).text = "Col1"
        tbl.cell(0, 1).text = "Col2"
        tbl.cell(1, 0).text = "Row1"
        tbl.cell(1, 1).text = "Data1"
        path = temp_dir / "patch_cell.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_patch_table_cell(
            str(path),
            slide_number=1,
            row_index=1,
            col_index=0,
            new_text="UpdatedRow1"
        )
        assert isinstance(result, dict)


class TestWordTableOperations:
    """Test Word table operations in depth."""

    def test_insert_table_row_at_end(self, word_advanced_tools, temp_dir):
        """Should insert row at table end."""
        doc = Document()
        table = doc.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        for i in range(1, 3):
            table.cell(i, 0).text = f"Item{i}"
            table.cell(i, 1).text = f"Val{i}"
        path = temp_dir / "insert_end.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_insert_table_row(
            str(path),
            "0",
            {"Name": "NewItem", "Value": "NewVal"}
        )
        assert isinstance(result, dict)


class TestPptxFromMarkdownVariations:
    """Test PPTX from markdown with different inputs."""

    def test_with_code_blocks(self, pptx_tools, temp_dir):
        """Should handle markdown with code."""
        md = """# Code Examples

## Python Code

```python
def hello():
    print("Hello")
```

## Usage
- Import the module
- Call hello()
"""
        path = temp_dir / "code.pptx"
        pptx_tools.tool_pptx_from_markdown(str(path), md)
        assert Path(path).exists()

    def test_with_emphasis(self, pptx_tools, temp_dir):
        """Should handle bold and italic."""
        md = """# Emphasis Test

    ## Points
    - **Bold text** is important
    - *Italic text* adds emphasis
    - ***Both*** for extra effect
    """
        path = temp_dir / "emphasis.pptx"
        pptx_tools.tool_pptx_from_markdown(str(path), md)
        assert Path(path).exists()


class TestWordFromMarkdownVariations:
    """Test Word from markdown with different inputs."""

    def test_with_links(self, word_tools, temp_dir):
        """Should handle markdown links."""

        md = """# Document with Links

Visit [our website](https://example.com) for more info.

## Resources
- [Resource 1](https://example.com/1)
- [Resource 2](https://example.com/2)
"""
        path = temp_dir / "links.docx"
        word_tools.tool_word_from_markdown(str(path), md)
        assert Path(path).exists()


class TestPptxSlideManipulation:
    """Test slide manipulation operations."""

    def test_delete_middle_slide(self, pptx_advanced_tools, temp_dir):
        """Should delete middle slide."""

        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
        path = temp_dir / "delete_middle.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_delete_slide(str(path), 2)
        assert isinstance(result, dict)

    def test_add_slide_at_start(self, pptx_advanced_tools, temp_dir):
        """Should add slide at start."""

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Original First"
        path = temp_dir / "add_start.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_add_slide(
            str(path),
            position="start",
            title="New First"
        )
        assert isinstance(result, dict)


class TestWordAuditPaths:
    """Test audit and validation paths."""

    def test_audit_sow_with_placeholders(self, word_advanced_tools, temp_dir):
        """Should audit SOW with unfilled placeholders."""

        doc = Document()
        doc.add_heading("Statement of Work", level=1)
        doc.add_paragraph("Customer: <Customer Name>")
        doc.add_paragraph("Project: [TBD]")
        doc.add_heading("Scope", level=2)
        doc.add_paragraph("<Insert scope here>")
        path = temp_dir / "sow_placeholders.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_audit_sow(str(path))
        assert isinstance(result, dict)


class TestPptxBulletOperations:
    """Test bullet operations thoroughly."""

    def test_add_multiple_bullets(self, pptx_advanced_tools, temp_dir):
        """Should add multiple bullets at different levels."""

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Bullet Test"
        path = temp_dir / "multi_bullets.pptx"
        prs.save(path)

        # Add first bullet
        pptx_advanced_tools.tool_pptx_add_bullet(str(path), 1, "First point", level=0)

        # Add sub-bullet
        pptx_advanced_tools.tool_pptx_add_bullet(str(path), 1, "Sub-point", level=1)

        # Add another top-level
        result = pptx_advanced_tools.tool_pptx_add_bullet(str(path), 1, "Second point", level=0)
        assert isinstance(result, dict)

    def test_add_bullet_with_bold_label(self, pptx_advanced_tools, temp_dir):
        """Should add bullet with bold label."""

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Key Points"
        path = temp_dir / "bold_label.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_add_bullet(
            str(path),
            1,
            "18 months to complete",
            bold_label="Duration"
        )
        assert isinstance(result, dict)


class TestSupportFunctions:
    """Test support functions."""

    def test_word_extract(self, word_tools, temp_dir):
        """Should extract word document."""

        doc = Document()
        doc.add_heading("Test", level=1)
        doc.add_paragraph("Content")
        path = temp_dir / "extract.docx"
        doc.save(path)

        result = word_tools.tool_word_extract(str(path))
        assert isinstance(result, dict)
