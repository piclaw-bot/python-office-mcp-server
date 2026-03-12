"""
Tests for pptx_advanced_tools.py - Advanced PowerPoint manipulation

Tests cover:
- Slide listing and introspection
- Shape manipulation
- Notes and comments
- Table operations
- Template operations
"""

import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from pptx import Presentation
from pptx.util import Inches, Mm

from tools.pptx_advanced_tools import (
    _classify_layout,
    _find_shape_by_identifier,
    _find_table_in_slide,
    _get_layout_recommendations,
    _get_shape_info,
)

# Fixtures temp_dir and pptx_advanced_tools are provided by conftest.py


@pytest.fixture
def sample_pptx(temp_dir):
    """Create a test presentation with multiple slides."""
    prs = Presentation()
    prs.slide_width = Mm(338.67)
    prs.slide_height = Mm(190.5)

    # Title slide
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Test Presentation"

    # Content slide with bullets
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Slide 2"
    for shape in slide.shapes:
        if shape.has_text_frame and shape.placeholder_format.idx == 1:
            shape.text_frame.paragraphs[0].text = "Bullet 1"

    # Slide with table
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Table Slide"
    table_shape = slide.shapes.add_table(
        3, 2, Inches(1), Inches(2), Inches(10), Inches(3)
    )
    table = table_shape.table
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(1, 0).text = "Data 1"
    table.cell(1, 1).text = "Data 2"

    path = temp_dir / "advanced_test.pptx"
    prs.save(path)
    return path


class TestHelperFunctions:
    """Tests for helper functions."""

    def test_get_shape_info(self, sample_pptx):
        """Should extract shape information."""
        prs = Presentation(sample_pptx)
        slide = prs.slides[0]
        for shape in slide.shapes:
            info = _get_shape_info(shape)
            assert "name" in info
            assert "shape_id" in info
            assert "is_placeholder" in info

    def test_find_shape_by_title(self, sample_pptx):
        """Should find title placeholder."""
        prs = Presentation(sample_pptx)
        slide = prs.slides[0]
        shape = _find_shape_by_identifier(slide, "title")
        assert shape is not None

    def test_find_shape_by_body(self, sample_pptx):
        """Should find body placeholder."""
        prs = Presentation(sample_pptx)
        slide = prs.slides[1]  # Content slide
        shape = _find_shape_by_identifier(slide, "body")
        assert shape is not None

    def test_find_shape_by_index(self, sample_pptx):
        """Should find shape by numeric index."""
        prs = Presentation(sample_pptx)
        slide = prs.slides[0]
        shape = _find_shape_by_identifier(slide, "0")
        assert shape is not None

    def test_find_table_in_slide(self, sample_pptx):
        """Should find table in slide."""
        prs = Presentation(sample_pptx)
        slide = prs.slides[2]  # Table slide
        table_shape, table = _find_table_in_slide(slide)
        assert table is not None
        assert table_shape is not None

    def test_classify_layout(self):
        """Should classify layouts correctly."""
        assert _classify_layout(['CENTER_TITLE (3)', 'SUBTITLE (4)'], 'Title Slide') == 'title_slide'
        assert _classify_layout(['TITLE (1)', 'BODY (2)'], 'Title and Content') == 'title_and_content'
        assert _classify_layout([], 'Blank') == 'blank'

    def test_get_layout_recommendations(self):
        """Should return layout recommendations."""
        recs = _get_layout_recommendations('title_slide')
        assert isinstance(recs, list)
        assert len(recs) > 0


class TestListSlides:
    """Tests for tool_pptx_list_slides."""

    def test_lists_all_slides(self, pptx_advanced_tools, sample_pptx):
        """Should list all slides in presentation."""
        result = pptx_advanced_tools.tool_pptx_list_slides(str(sample_pptx))
        assert result.get("slide_count") == 3
        assert len(result.get("slides", [])) == 3

    def test_includes_slide_titles(self, pptx_advanced_tools, sample_pptx):
        """Should include slide titles."""
        result = pptx_advanced_tools.tool_pptx_list_slides(str(sample_pptx))
        titles = [s.get("title") for s in result.get("slides", [])]
        assert "Test Presentation" in titles

    def test_file_not_found(self, pptx_advanced_tools):
        """Should handle missing files."""
        result = pptx_advanced_tools.tool_pptx_list_slides("/nonexistent.pptx")
        assert "error" in result


class TestListShapes:
    """Tests for tool_pptx_list_shapes."""

    def test_lists_shapes(self, pptx_advanced_tools, sample_pptx):
        """Should list shapes on a slide."""
        result = pptx_advanced_tools.tool_pptx_list_shapes(str(sample_pptx), 1)
        assert "shapes" in result
        assert len(result["shapes"]) > 0

    def test_invalid_slide_number(self, pptx_advanced_tools, sample_pptx):
        """Should handle invalid slide numbers."""
        result = pptx_advanced_tools.tool_pptx_list_shapes(str(sample_pptx), 100)
        assert "error" in result


class TestGetSlide:
    """Tests for tool_pptx_get_slide."""

    def test_gets_slide_content(self, pptx_advanced_tools, sample_pptx):
        """Should get slide content."""
        result = pptx_advanced_tools.tool_pptx_get_slide(str(sample_pptx), 1)
        # Result has 'slide' key with nested content
        slide_data = result.get("slide", result)
        assert slide_data.get("title") == "Test Presentation" or "Test Presentation" in str(result)

    def test_gets_table_content(self, pptx_advanced_tools, sample_pptx):
        """Should get table content from slide."""
        result = pptx_advanced_tools.tool_pptx_get_slide(str(sample_pptx), 3)
        # Tables are nested in slide data
        slide_data = result.get("slide", result)
        assert "tables" in slide_data


class TestPatchShape:
    """Tests for tool_pptx_patch_shape."""

    def test_updates_title(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should update title text."""
        output = temp_dir / "patched.pptx"
        result = pptx_advanced_tools.tool_pptx_patch_shape(
            str(sample_pptx), 1, "title",
            new_text="New Title",
            output_path=str(output)
        )
        assert result.get("success") is True

        # Verify change
        prs = Presentation(output)
        assert prs.slides[0].shapes.title.text == "New Title"

    def test_updates_body(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should update body content."""
        output = temp_dir / "patched_body.pptx"
        result = pptx_advanced_tools.tool_pptx_patch_shape(
            str(sample_pptx), 2, "body",
            new_text="Updated bullet",
            output_path=str(output)
        )
        assert result.get("success") is True


class TestNotes:
    """Tests for speaker notes."""

    def test_set_notes(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should set speaker notes."""
        output = temp_dir / "notes.pptx"
        result = pptx_advanced_tools.tool_pptx_set_notes(
            str(sample_pptx), 1,
            "These are speaker notes",
            output_path=str(output)
        )
        assert result.get("success") is True

    def test_get_notes(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should get speaker notes."""
        # First set notes
        output = temp_dir / "notes2.pptx"
        pptx_advanced_tools.tool_pptx_set_notes(
            str(sample_pptx), 1,
            "Test notes content",
            output_path=str(output)
        )

        # Then get notes
        result = pptx_advanced_tools.tool_pptx_get_notes(str(output), 1)
        assert "Test notes content" in result.get("notes", "")


class TestComments:
    """Tests for PPTX comments."""

    def test_add_comment(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should add comment to slide."""
        output = temp_dir / "commented.pptx"
        result = pptx_advanced_tools.tool_pptx_add_comment(
            str(sample_pptx), 1,
            "This is a test comment",
            x_inches=1.0,
            y_inches=1.0,
            output_path=str(output)
        )
        assert result.get("success") is True
        assert result.get("position", {}).get("x_inches") == 1.0
        assert result.get("position", {}).get("y_inches") == 1.0

    def test_comment_default_position(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Comment should use default position if not specified."""
        output = temp_dir / "comment_default.pptx"
        result = pptx_advanced_tools.tool_pptx_add_comment(
            str(sample_pptx), 1,
            "Test comment",
            output_path=str(output)
        )
        assert result.get("success") is True
        # Default is 1.0, 1.0
        assert result.get("position", {}).get("x_inches") == 1.0
        assert result.get("position", {}).get("y_inches") == 1.0

    def test_get_comments(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should retrieve comments from slide."""
        output = temp_dir / "get_comments.pptx"
        # Add a comment first
        pptx_advanced_tools.tool_pptx_add_comment(
            str(sample_pptx), 1,
            "Retrievable comment",
            output_path=str(output)
        )

        # Get comments
        result = pptx_advanced_tools.tool_pptx_get_comments(str(output), 1)
        assert result.get("total_comments") >= 1

    def test_get_comments_from_fixture_comments_pptx(self, pptx_advanced_tools, temp_dir):
        """Should read comments from the real comments.pptx test fixture."""
        fixture_path = (
            Path(__file__).resolve().parent /
            "_templates" / "testdata" / "pptx" / "comments.pptx"
        )
        test_file = temp_dir / "fixture_comments.pptx"
        test_file.write_bytes(fixture_path.read_bytes())

        result = pptx_advanced_tools.tool_pptx_get_comments(str(test_file), 1)
        assert "error" not in result
        assert result.get("total_comments", 0) >= 1

        comments = result.get("comments", {}).get(1, [])
        assert comments
        first = comments[0]
        assert first.get("format") == "modern"
        assert first.get("author") == "Rui Carmo"
        assert first.get("text") == "I don't remember it having yellow eyes"

    def test_add_comment_with_none_author(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should default author when None is provided."""
        output = temp_dir / "comment_none_author.pptx"
        result = pptx_advanced_tools.tool_pptx_add_comment(
            str(sample_pptx), 1,
            "Author fallback",
            author=None,
            output_path=str(output)
        )
        assert "error" not in result

    def test_delete_comment_by_index(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should delete a single comment by index from a slide."""
        output = temp_dir / "delete_comment_one.pptx"
        pptx_advanced_tools.tool_pptx_add_comment(
            str(sample_pptx), 1,
            "Delete one",
            output_path=str(output)
        )

        comments = pptx_advanced_tools.tool_pptx_get_comments(str(output), 1)
        assert comments.get("total_comments", 0) >= 1
        idx = comments["comments"][1][0]["index"]

        result = pptx_advanced_tools.tool_pptx_delete_comment(
            str(output),
            1,
            comment_index=int(idx),
        )
        assert result.get("success") is True

    def test_delete_all_comments_on_slide(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should delete all comments from a slide when index is omitted."""
        output = temp_dir / "delete_comment_all.pptx"
        pptx_advanced_tools.tool_pptx_add_comment(
            str(sample_pptx), 1,
            "Delete all",
            output_path=str(output)
        )

        result = pptx_advanced_tools.tool_pptx_delete_comment(str(output), 1)
        assert result.get("success") is True

    def test_get_modern_comments(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should read Office 365 modern comments from modernComment parts."""
        output = temp_dir / "modern_comments.pptx"

        # Start from a valid deck and inject modern-comment parts/relationships.
        with zipfile.ZipFile(sample_pptx, "r") as zf_in:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                tmp_path = tmp.name

            with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zf_out:
                for item in zf_in.namelist():
                    if item == "ppt/slides/_rels/slide1.xml.rels":
                        rels_root = ET.fromstring(zf_in.read(item))
                        rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
                        ET.SubElement(
                            rels_root,
                            f"{{{rel_ns}}}Relationship",
                            {
                                "Id": "rId999",
                                "Type": "http://schemas.microsoft.com/office/2018/10/relationships/comments",
                                "Target": "../comments/modernComment_7FFFFFFF_0.xml",
                            },
                        )
                        zf_out.writestr(item, ET.tostring(rels_root, encoding="utf-8", xml_declaration=True))
                    else:
                        zf_out.writestr(item, zf_in.read(item))

                zf_out.writestr(
                    "ppt/authors.xml",
                    (
                        "<?xml version=\"1.0\" encoding=\"UTF-8\" standalone=\"yes\"?>"
                        "<p188:authorLst xmlns:p188=\"http://schemas.microsoft.com/office/powerpoint/2018/8/main\">"
                        "<p188:author id=\"{AUTHOR-1}\" name=\"GLASSON, Emma\" initials=\"GE\" userId=\"emma@example.com\" providerId=\"AD\"/>"
                        "</p188:authorLst>"
                    ),
                )
                zf_out.writestr(
                    "ppt/comments/modernComment_7FFFFFFF_0.xml",
                    (
                        "<?xml version=\"1.0\" encoding=\"UTF-8\" standalone=\"yes\"?>"
                        "<p188:cmLst"
                        " xmlns:a=\"http://schemas.openxmlformats.org/drawingml/2006/main\""
                        " xmlns:ac=\"http://schemas.microsoft.com/office/drawing/2013/main/command\""
                        " xmlns:pc=\"http://schemas.microsoft.com/office/powerpoint/2013/main/command\""
                        " xmlns:p188=\"http://schemas.microsoft.com/office/powerpoint/2018/8/main\">"
                        "<p188:cm id=\"{CM-1}\" authorId=\"{AUTHOR-1}\" created=\"2026-03-06T09:21:41.409\">"
                        "<ac:deMkLst><pc:sldMk cId=\"0\" sldId=\"2147483648\"/><ac:spMk id=\"5\"/></ac:deMkLst>"
                        "<p188:pos x=\"914400\" y=\"1828800\"/>"
                        "<p188:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>modern comment text</a:t></a:r></a:p></p188:txBody>"
                        "</p188:cm>"
                        "</p188:cmLst>"
                    ),
                )

        Path(tmp_path).replace(output)

        result = pptx_advanced_tools.tool_pptx_get_comments(str(output), 1)
        assert result.get("total_comments", 0) >= 1
        comments = result.get("comments", {}).get(1, [])
        assert comments
        first = comments[0]
        assert first.get("author") == "GLASSON, Emma"
        assert first.get("text") == "modern comment text"
        assert first.get("format") == "modern"


class TestTables:
    """Tests for table operations."""

    def test_get_table(self, pptx_advanced_tools, sample_pptx):
        """Should get table content."""
        result = pptx_advanced_tools.tool_pptx_get_table(str(sample_pptx), 3)
        # Check for 'header' key instead of 'headers'
        assert "header" in result or "data" in result

    def test_add_table(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should add table to slide."""
        output = temp_dir / "new_table.pptx"
        result = pptx_advanced_tools.tool_pptx_add_table(
            str(sample_pptx), 2,
            headers=["Col1", "Col2", "Col3"],
            rows=[["A", "B", "C"], ["D", "E", "F"]],
            output_path=str(output)
        )
        assert result.get("success") is True
        assert result.get("rows") == 3  # header + 2 data rows
        assert result.get("columns") == 3

    def test_table_positioning(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Table should be positioned within slide bounds."""
        output = temp_dir / "table_pos.pptx"
        result = pptx_advanced_tools.tool_pptx_add_table(
            str(sample_pptx), 2,
            headers=["A", "B"],
            left=1.0, top=2.0, width=11.0, height=3.0,
            output_path=str(output)
        )
        assert result.get("success") is True

        # Verify positioning
        prs = Presentation(output)
        slide = prs.slides[1]
        for shape in slide.shapes:
            if shape.has_table:
                # 16:9 is 13.333" x 7.5"
                assert shape.left + shape.width <= prs.slide_width
                assert shape.top + shape.height <= prs.slide_height


class TestSlideManipulation:
    """Tests for slide add/delete/reorder."""

    def test_add_slide(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should add a new slide."""
        output = temp_dir / "added.pptx"
        result = pptx_advanced_tools.tool_pptx_add_slide(
            str(sample_pptx),
            title="New Slide",
            output_path=str(output)
        )
        assert result.get("success") is True

        prs = Presentation(output)
        assert len(prs.slides) == 4  # Original 3 + 1 new

    def test_delete_slide(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should delete a slide."""
        output = temp_dir / "deleted.pptx"
        result = pptx_advanced_tools.tool_pptx_delete_slide(
            str(sample_pptx), 2,
            output_path=str(output)
        )
        # Check for success or remaining slides count
        assert result.get("success") is True or "remaining" in result or output.exists()

    def test_duplicate_slide(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should duplicate a slide."""
        output = temp_dir / "duplicated.pptx"
        result = pptx_advanced_tools.tool_pptx_duplicate_slide(
            str(sample_pptx), 1,
            output_path=str(output)
        )
        assert result.get("success") is True

        prs = Presentation(output)
        assert len(prs.slides) == 4  # Original 3 + 1 duplicate


class TestBullets:
    """Tests for bullet operations."""

    def test_add_bullet(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should add bullet to slide."""
        output = temp_dir / "bullet.pptx"
        result = pptx_advanced_tools.tool_pptx_add_bullet(
            str(sample_pptx), 2,
            "New bullet point",
            output_path=str(output)
        )
        assert result.get("success") is True

    def test_clear_bullets(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should clear bullets from slide."""
        output = temp_dir / "cleared.pptx"
        result = pptx_advanced_tools.tool_pptx_clear_bullets(
            str(sample_pptx), 2,
            output_path=str(output)
        )
        assert result.get("success") is True


class TestLayoutAnalysis:
    """Tests for layout analysis tools."""

    def test_list_masters(self, pptx_advanced_tools, sample_pptx):
        """Should list slide masters/layouts."""
        result = pptx_advanced_tools.tool_pptx_list_masters(str(sample_pptx))
        assert "default_layouts" in result
        assert len(result["default_layouts"]) > 0

    def test_analyze_layouts(self, pptx_advanced_tools, sample_pptx):
        """Should analyze layout placeholders."""
        result = pptx_advanced_tools.tool_pptx_analyze_layouts(str(sample_pptx))
        assert "layouts" in result

    def test_recommend_layout(self, pptx_advanced_tools, sample_pptx):
        """Should recommend layouts for content types."""
        result = pptx_advanced_tools.tool_pptx_recommend_layout(
            str(sample_pptx), "bullets"
        )
        assert "recommended" in result or "layout_index" in result


class TestPlaceholders:
    """Tests for placeholder operations."""

    def test_audit_placeholders(self, pptx_advanced_tools, temp_dir):
        """Should find unfilled placeholders."""
        # Create presentation with placeholder
        prs = Presentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "<Customer Name>"

        path = temp_dir / "placeholders.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_audit_placeholders(str(path))
        # Check for various possible keys
        assert (
            result.get("total_placeholders", 0) >= 1 or
            "status" in result or
            "findings" in result
        )

    def test_replace_placeholders(self, pptx_advanced_tools, temp_dir):
        """Should replace placeholders with values."""
        # Create presentation with placeholder
        prs = Presentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "<Customer Name>"

        path = temp_dir / "replace.pptx"
        prs.save(path)

        output = temp_dir / "replaced.pptx"
        result = pptx_advanced_tools.tool_pptx_replace_placeholders(
            str(path),
            {"<Customer Name>": "Contoso Corp"},
            output_path=str(output)
        )
        assert result.get("success") is True


class TestReplaceText:
    """Tests for text replacement."""

    def test_replace_text(self, pptx_advanced_tools, sample_pptx, temp_dir):
        """Should find and replace text."""
        output = temp_dir / "replaced_text.pptx"
        result = pptx_advanced_tools.tool_pptx_replace_text(
            str(sample_pptx),
            "Test Presentation",
            "Updated Presentation",
            output_path=str(output)
        )
        assert result.get("success") is True

        # Verify replacement
        prs = Presentation(output)
        assert prs.slides[0].shapes.title.text == "Updated Presentation"
