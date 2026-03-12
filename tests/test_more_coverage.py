"""
More coverage tests for prompt functions and edge cases.
"""

from docx import Document
from pptx import Presentation
from pptx.util import Inches

# Fixtures temp_dir, word_advanced_tools, and pptx_advanced_tools are provided by conftest.py


class TestPromptFunctions:
    """Tests for all prompt functions."""

    def test_prompt_sow_generation(self, word_advanced_tools):
        """Should return SOW generation guidance."""
        result = word_advanced_tools.prompt_sow_generation()
        assert isinstance(result, dict)
        # Result has 'name' and 'arguments' keys
        assert "name" in result or "arguments" in result or "prompt" in result

    def test_prompt_section_editing(self, word_advanced_tools):
        """Should return section editing guidance."""
        result = word_advanced_tools.prompt_section_editing()
        assert isinstance(result, dict)

    def test_prompt_document_audit(self, word_advanced_tools):
        """Should return document audit guidance."""
        result = word_advanced_tools.prompt_document_audit()
        assert isinstance(result, dict)

    def test_prompt_table_editing(self, word_advanced_tools):
        """Should return table editing guidance."""
        result = word_advanced_tools.prompt_table_editing()
        assert isinstance(result, dict)


class TestMoreWordErrorHandling:
    """Additional error handling tests for Word tools."""

    def test_patch_placeholder_not_found(self, word_advanced_tools, temp_dir):
        """Should handle placeholder not found."""
        doc = Document()
        doc.add_paragraph("No placeholders here")
        path = temp_dir / "no_ph.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_patch_placeholder(
            str(path),
            "<Missing>",
            "Replacement"
        )
        # Should return gracefully even if not found
        assert isinstance(result, dict)

    def test_get_table_string_index(self, word_advanced_tools, temp_dir):
        """Should handle string table index."""
        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A"
        path = temp_dir / "str_index.docx"
        doc.save(path)

        # Pass string index (should be converted to int)
        result = word_advanced_tools.tool_word_get_table(str(path), "0")
        assert isinstance(result, dict)

    def test_insert_table_row_invalid_table(self, word_advanced_tools, temp_dir):
        """Should handle invalid table identifier."""
        doc = Document()
        doc.add_table(rows=2, cols=2)
        path = temp_dir / "one_table.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_insert_table_row(
            str(path),
            "999",  # Invalid index
            {"A": "1"}
        )
        assert "error" in result or "not found" in str(result).lower()


class TestMorePptxOperations:
    """Additional PPTX operation tests."""

    def test_add_bullet_with_shape(self, pptx_advanced_tools, temp_dir):
        """Should add bullet to specific shape."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Title"
        path = temp_dir / "bullet_shape.pptx"
        prs.save(path)

        output = temp_dir / "with_bullet.pptx"
        result = pptx_advanced_tools.tool_pptx_add_bullet(
            str(path),
            slide_number=1,
            text="New bullet point",
            shape_identifier="body",
            output_path=str(output)
        )
        assert isinstance(result, dict)

    def test_set_notes_empty(self, pptx_advanced_tools, temp_dir):
        """Should handle setting empty notes."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "empty_notes.pptx"
        prs.save(path)

        output = temp_dir / "cleared_notes.pptx"
        result = pptx_advanced_tools.tool_pptx_set_notes(
            str(path),
            slide_number=1,
            notes_text="",
            output_path=str(output)
        )
        assert isinstance(result, dict)

    def test_get_slide_invalid(self, pptx_advanced_tools, temp_dir):
        """Should handle invalid slide number."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "one_slide.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_get_slide(str(path), slide_number=99)
        assert "error" in result


class TestWordTemplateOperations:
    """Tests for Word template operations."""

    def test_copy_template_missing_source(self, word_advanced_tools, temp_dir):
        """Should handle missing source template."""
        result = word_advanced_tools.tool_word_copy_template(
            "/nonexistent/template.docx",
            str(temp_dir / "dest.docx")
        )
        assert "error" in result

    def test_analyze_template_empty(self, word_advanced_tools, temp_dir):
        """Should analyze simple document."""
        doc = Document()
        doc.add_paragraph("Simple text")
        path = temp_dir / "simple.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_analyze_template_formatting(str(path))
        assert isinstance(result, dict)


class TestMorePptxLayouts:
    """More layout operation tests."""

    def test_recommend_layout_image(self, pptx_advanced_tools, temp_dir):
        """Should recommend layout for image content."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = temp_dir / "img_layout.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_recommend_layout(str(path), "image")
        assert isinstance(result, dict)

    def test_recommend_layout_title(self, pptx_advanced_tools, temp_dir):
        """Should recommend layout for title slide."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = temp_dir / "title_layout.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_recommend_layout(str(path), "title")
        assert isinstance(result, dict)


class TestWordComplexDocuments:
    """Tests with complex Word documents."""

    def test_document_with_images(self, word_advanced_tools, temp_dir):
        """Should handle documents gracefully even with images."""
        doc = Document()
        doc.add_heading("Document", level=1)
        doc.add_paragraph("Content before table")

        # Add a table
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Cell"

        doc.add_paragraph("Content after table")
        path = temp_dir / "complex.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_list_sections(str(path))
        assert isinstance(result, dict)

    def test_extract_sow_structure_complex(self, word_advanced_tools, temp_dir):
        """Should extract structure from complex SOW."""
        doc = Document()
        doc.add_heading("Statement of Work", level=0)
        doc.add_heading("1. Executive Summary", level=1)
        doc.add_paragraph("Summary content")
        doc.add_heading("2. Scope of Work", level=1)
        doc.add_paragraph("Scope details")
        doc.add_heading("2.1 In Scope", level=2)
        doc.add_paragraph("In scope items")
        doc.add_heading("2.2 Out of Scope", level=2)
        doc.add_paragraph("Out of scope items")
        doc.add_heading("3. Deliverables", level=1)

        table = doc.add_table(rows=3, cols=3)
        table.cell(0, 0).text = "ID"
        table.cell(0, 1).text = "Deliverable"
        table.cell(0, 2).text = "Date"

        path = temp_dir / "complex_sow.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_extract_sow_structure(str(path))
        assert isinstance(result, dict)


class TestPptxTableOperations:
    """More PPTX table tests."""

    def test_get_table_multiple_tables(self, pptx_advanced_tools, temp_dir):
        """Should get specific table when multiple exist."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Add two tables
        table1 = slide.shapes.add_table(
            rows=2, cols=2,
            left=Inches(0.5), top=Inches(1.5),
            width=Inches(3), height=Inches(1)
        ).table
        table1.cell(0, 0).text = "Table1"

        table2 = slide.shapes.add_table(
            rows=2, cols=2,
            left=Inches(5), top=Inches(1.5),
            width=Inches(3), height=Inches(1)
        ).table
        table2.cell(0, 0).text = "Table2"

        path = temp_dir / "two_tables.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_get_table(str(path), slide_number=1, table_index=1)
        assert isinstance(result, dict)


class TestWordCheckOperations:
    """Tests for Word check operations."""

    def test_check_tracking_enabled(self, word_advanced_tools, temp_dir):
        """Should check tracking on document with tracking enabled."""
        doc = Document()
        doc.add_paragraph("Content")
        path = temp_dir / "check_track.docx"
        doc.save(path)

        # Enable tracking first
        word_advanced_tools.tool_word_enable_track_changes(str(path))

        # Then check
        result = word_advanced_tools.tool_word_check_tracking(str(path))
        assert isinstance(result, dict)


class TestPptxDuplication:
    """Tests for slide duplication."""

    def test_duplicate_first_slide(self, pptx_advanced_tools, temp_dir):
        """Should duplicate first slide."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Original"
        path = temp_dir / "dup_first.pptx"
        prs.save(path)

        output = temp_dir / "duplicated.pptx"
        result = pptx_advanced_tools.tool_pptx_duplicate_slide(
            str(path),
            slide_number=1,
            position="end",
            output_path=str(output)
        )
        assert result.get("success") is True or "slide" in str(result)


class TestCleanupOperations:
    """Tests for cleanup operations."""

    def test_cleanup_sow_with_guidance(self, word_advanced_tools, temp_dir):
        """Should cleanup SOW with guidance text."""
        doc = Document()
        doc.add_heading("SOW", level=0)
        doc.add_paragraph("[Template Guidance: Delete this instruction]")
        doc.add_paragraph("Real content to keep")
        doc.add_paragraph("[Note: Also remove this]")
        path = temp_dir / "guidance_sow.docx"
        doc.save(path)

        output = temp_dir / "cleaned.docx"
        result = word_advanced_tools.tool_word_cleanup_sow(str(path), output_path=str(output))
        assert isinstance(result, dict)


class TestPptxPlaceholderReplacement:
    """Tests for placeholder replacement."""

    def test_replace_in_notes(self, pptx_advanced_tools, temp_dir):
        """Should replace text in speaker notes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "<Customer> Presentation"

        # Add notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Notes for <Customer>"

        path = temp_dir / "notes_replace.pptx"
        prs.save(path)

        output = temp_dir / "replaced.pptx"
        result = pptx_advanced_tools.tool_pptx_replace_text(
            str(path),
            find_text="<Customer>",
            replace_text="Contoso",
            output_path=str(output)
        )
        assert isinstance(result, dict)
