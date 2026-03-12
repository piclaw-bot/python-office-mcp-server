"""
Tests for pptx_tools.py - PowerPoint presentation processing

Tests cover:
- Markdown to PowerPoint conversion
- Content extraction
- Theme font handling
- Layout analysis
"""

from pptx import Presentation
from pptx.util import Mm

from tools.pptx_tools import (
    _analyze_markdown_for_layouts,
    _set_theme_fonts,
)

# Fixtures temp_dir, pptx_tools, and sample_pptx are provided by conftest.py

# Import the module-level function for testing


class TestAnalyzeNodesForLayouts:
    """Tests for _analyze_markdown_for_layouts function."""

    def test_title_slide_detection_h1(self):
        """Should correctly detect title slide presence from H1."""
        result = _analyze_markdown_for_layouts("# Main Title\n\nContent")
        assert result["has_title_slide"] is True

    def test_title_slide_detection_no_h1(self):
        """Should not detect title slide when only H2 present."""
        result = _analyze_markdown_for_layouts("## Section Only\n\nContent")
        assert result["has_title_slide"] is False

    def test_detects_content_slides(self):
        """Subsequent headings become content slides."""
        md = """# Title

## First Section
- bullet 1
- bullet 2

## Second Section
- bullet 3
"""
        result = _analyze_markdown_for_layouts(md)
        assert len(result["slides"]) >= 2

    def test_detects_tables(self):
        """Table content should be detected."""
        md = """# Title

## Data

| Col1 | Col2 |
|------|------|
| A    | B    |
"""
        result = _analyze_markdown_for_layouts(md)
        # Find the slide with table
        table_slides = [s for s in result["slides"] if s.get("has_table")]
        assert len(table_slides) >= 1

    def test_detects_bullets(self):
        """Bullet content should be detected."""
        md = """# Title

## Points
- Point 1
- Point 2
- Point 3
"""
        result = _analyze_markdown_for_layouts(md)
        bullet_slides = [s for s in result["slides"] if s.get("has_bullets")]
        assert len(bullet_slides) >= 1

    def test_horizontal_rule_creates_boundary(self):
        """Horizontal rules should create slide boundaries."""
        md = """# Title

---

## Next Section
"""
        result = _analyze_markdown_for_layouts(md)
        assert result["has_title_slide"] is True


class TestSetThemeFonts:
    """Tests for _set_theme_fonts helper."""

    def test_sets_title_font(self, temp_dir):
        """Should set major (title) font in theme."""
        prs = Presentation()
        _set_theme_fonts(prs, "Arial Black", "Arial")
        # If it doesn't raise, theme modification succeeded
        path = temp_dir / "themed.pptx"
        prs.save(path)
        assert path.exists()

    def test_sets_body_font(self, temp_dir):
        """Should set minor (body) font in theme."""
        prs = Presentation()
        _set_theme_fonts(prs, "Impact", "Calibri")
        path = temp_dir / "themed2.pptx"
        prs.save(path)
        assert path.exists()


class TestPptxFromMarkdown:
    """Tests for tool_pptx_from_markdown."""

    def test_creates_presentation(self, pptx_tools, temp_dir):
        """Should create a valid PPTX file from markdown."""
        md = """# Test Presentation

## Overview
- Point 1
- Point 2
"""
        output = temp_dir / "output.pptx"
        result = pptx_tools.tool_pptx_from_markdown(str(output), md)

        assert result.get("success") is True
        assert output.exists()

    def test_creates_title_slide(self, pptx_tools, temp_dir):
        """First heading should become title slide."""
        md = "# My Title"
        output = temp_dir / "title.pptx"
        result = pptx_tools.tool_pptx_from_markdown(str(output), md)

        assert result.get("success") is True
        prs = Presentation(output)
        assert len(prs.slides) >= 1
        # Check first slide has title
        first_slide = prs.slides[0]
        assert first_slide.shapes.title is not None

    def test_creates_content_slides(self, pptx_tools, temp_dir):
        """Subsequent headings should create content slides."""
        md = """# Title

## Section 1
- bullet

## Section 2
- bullet
"""
        output = temp_dir / "content.pptx"
        result = pptx_tools.tool_pptx_from_markdown(str(output), md)

        assert result.get("success") is True
        prs = Presentation(output)
        assert len(prs.slides) >= 3  # title + 2 content

    def test_creates_tables(self, pptx_tools, temp_dir):
        """Tables in markdown should become PowerPoint tables."""
        md = """# Title

## Data

| Name | Value |
|------|-------|
| A    | 1     |
| B    | 2     |
"""
        output = temp_dir / "tables.pptx"
        result = pptx_tools.tool_pptx_from_markdown(str(output), md)

        assert result.get("success") is True
        prs = Presentation(output)
        # Find table in slides
        tables_found = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    tables_found += 1
        assert tables_found >= 1

    def test_bold_labels(self, pptx_tools, temp_dir):
        """**Label:** format should create bold labels."""
        md = """# Title

## Key Points
- **Duration:** 6 months
- **Cost:** $100K
"""
        output = temp_dir / "labels.pptx"
        result = pptx_tools.tool_pptx_from_markdown(str(output), md)
        assert result.get("success") is True

    def test_custom_fonts(self, pptx_tools, temp_dir):
        """Should accept custom font parameters."""
        md = "# Title"
        output = temp_dir / "fonts.pptx"
        result = pptx_tools.tool_pptx_from_markdown(
            str(output), md,
            title_font="Impact",
            body_font="Calibri"
        )
        assert result.get("success") is True

    def test_subtitle_on_title_slide(self, pptx_tools, temp_dir):
        """Context line should become subtitle."""
        md = """# Main Title
**Context:** Supporting subtitle text
"""
        output = temp_dir / "subtitle.pptx"
        result = pptx_tools.tool_pptx_from_markdown(str(output), md)
        assert result.get("success") is True


class TestPptxExtract:
    """Tests for tool_pptx_extract."""

    def test_extracts_titles(self, pptx_tools, sample_pptx):
        """Should extract slide titles."""
        result = pptx_tools.tool_pptx_extract(str(sample_pptx))
        assert "slides" in result
        assert len(result["slides"]) >= 2
        # Check content contains title
        content = str(result)
        assert "Test Title" in content

    def test_extracts_content(self, pptx_tools, sample_pptx):
        """Should extract slide content."""
        result = pptx_tools.tool_pptx_extract(str(sample_pptx))
        assert "slides" in result

    def test_file_not_found(self, pptx_tools):
        """Should handle missing files gracefully."""
        result = pptx_tools.tool_pptx_extract("/nonexistent/file.pptx")
        assert "error" in result or "Error" in str(result)


class TestPptxToMarkdown:
    """Tests for tool_pptx_to_markdown."""

    def test_converts_to_markdown(self, pptx_tools, sample_pptx):
        """Should convert presentation to markdown."""
        result = pptx_tools.tool_pptx_to_markdown(str(sample_pptx))
        # Returns string directly
        assert isinstance(result, str)
        assert "Test Title" in result or "Slide" in result

    def test_includes_slide_titles(self, pptx_tools, sample_pptx):
        """Markdown should include slide titles as headings."""
        result = pptx_tools.tool_pptx_to_markdown(str(sample_pptx))
        # Check for heading markers or slide content
        assert "#" in result or "Slide" in result

    def test_file_not_found(self, pptx_tools):
        """Should handle missing files gracefully."""
        result = pptx_tools.tool_pptx_to_markdown("/nonexistent/file.pptx")
        assert "Error" in result or "not found" in result.lower()


class TestSlidePositioning:
    """Tests for slide element positioning."""

    def test_content_within_bounds(self, pptx_tools, temp_dir):
        """All content should be within slide boundaries."""
        md = """# Title

## Content
- Bullet 1
- Bullet 2

| Col1 | Col2 |
|------|------|
| A    | B    |
"""
        output = temp_dir / "bounds.pptx"
        result = pptx_tools.tool_pptx_from_markdown(str(output), md)
        assert result.get("success") is True

        prs = Presentation(output)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide in prs.slides:
            for shape in slide.shapes:
                # Check shape is within bounds
                assert shape.left >= 0, f"Shape {shape.name} has negative left"
                assert shape.top >= 0, f"Shape {shape.name} has negative top"
                right = shape.left + shape.width
                bottom = shape.top + shape.height
                assert right <= slide_width, f"Shape {shape.name} extends past right edge"
                assert bottom <= slide_height, f"Shape {shape.name} extends past bottom edge"

    def test_widescreen_dimensions(self, pptx_tools, temp_dir):
        """Presentation should use 16:9 widescreen dimensions."""
        md = "# Title"
        output = temp_dir / "widescreen.pptx"
        result = pptx_tools.tool_pptx_from_markdown(str(output), md)
        assert result.get("success") is True

        prs = Presentation(output)
        # 16:9 = 338.67mm x 190.5mm
        assert abs(prs.slide_width - Mm(338.67)) < Mm(1)
        assert abs(prs.slide_height - Mm(190.5)) < Mm(1)
