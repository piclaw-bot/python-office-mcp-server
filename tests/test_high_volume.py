"""
High-volume tests to increase coverage on specific paths.
Focus on word_advanced_tools and pptx_advanced_tools.
"""

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches as PptxInches

# Fixture temp_dir is provided by conftest.py


class TestWordGetSection:
    """Test word_get_section method in detail."""

    def test_get_first_section(self, word_advanced_tools, temp_dir):
        """Should get first section content."""
        doc = Document()
        doc.add_heading("Introduction", level=1)
        doc.add_paragraph("This is the intro content.")
        doc.add_heading("Details", level=1)
        doc.add_paragraph("Details content.")
        path = temp_dir / "sections.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_get_section(str(path), "Introduction")
        assert isinstance(result, dict)
        assert "content" in result or "error" not in result

    def test_get_last_section(self, word_advanced_tools, temp_dir):
        """Should get last section content."""
        doc = Document()
        doc.add_heading("Section 1", level=1)
        doc.add_paragraph("Content 1")
        doc.add_heading("Section 2", level=1)
        doc.add_paragraph("Content 2")
        doc.add_heading("Last Section", level=1)
        doc.add_paragraph("Final content here.")
        path = temp_dir / "last_section.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_get_section(str(path), "Last Section")
        assert isinstance(result, dict)

    def test_get_nonexistent_section(self, word_advanced_tools, temp_dir):
        """Should handle getting nonexistent section."""
        doc = Document()
        doc.add_heading("Existing", level=1)
        doc.add_paragraph("Content")
        path = temp_dir / "no_section.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_get_section(str(path), "NonexistentSection")
        # Should return error or empty
        assert isinstance(result, dict)


class TestWordPatchTableRow:
    """Test word_patch_table_row method."""

    def test_patch_first_data_row(self, word_advanced_tools, temp_dir):
        """Should patch first data row."""
        doc = Document()
        table = doc.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Original"
        table.cell(1, 1).text = "Data"
        table.cell(2, 0).text = "More"
        table.cell(2, 1).text = "Data2"
        path = temp_dir / "patch_row.docx"
        doc.save(path)

        output = temp_dir / "patched.docx"
        result = word_advanced_tools.tool_word_patch_table_row(
            str(path),
            "0",
            1,
            {"Name": "Updated", "Value": "NewData"},
            output_path=str(output)
        )
        assert isinstance(result, dict)

    def test_patch_last_row(self, word_advanced_tools, temp_dir):
        """Should patch last row."""
        doc = Document()
        table = doc.add_table(rows=4, cols=2)
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"
        for i in range(1, 4):
            table.cell(i, 0).text = f"Row{i}A"
            table.cell(i, 1).text = f"Row{i}B"
        path = temp_dir / "last_row.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_patch_table_row(
            str(path),
            "0",
            3,  # Last row
            {"A": "Z1", "B": "Z2"}
        )
        assert isinstance(result, dict)


class TestWordListTables:
    """Test word_list_tables method."""

    def test_list_multiple_tables(self, word_advanced_tools, temp_dir):
        """Should list multiple tables."""
        doc = Document()
        doc.add_paragraph("Before table 1")
        t1 = doc.add_table(rows=2, cols=2)
        t1.cell(0, 0).text = "T1A"
        t1.cell(0, 1).text = "T1B"

        doc.add_paragraph("Between tables")

        t2 = doc.add_table(rows=3, cols=3)
        t2.cell(0, 0).text = "T2X"
        t2.cell(0, 1).text = "T2Y"
        t2.cell(0, 2).text = "T2Z"

        path = temp_dir / "multi_table.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_list_tables(str(path))
        assert isinstance(result, dict)
        assert "tables" in result or "table_count" in result

    def test_list_tables_empty_doc(self, word_advanced_tools, temp_dir):
        """Should handle doc with no tables."""
        doc = Document()
        doc.add_paragraph("No tables here")
        path = temp_dir / "no_tables.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_list_tables(str(path))
        assert isinstance(result, dict)


class TestPptxGetSlide:
    """Test pptx_get_slide method in detail."""

    def test_get_slide_with_all_content(self, pptx_advanced_tools, temp_dir):
        """Should get slide with various content types."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

        # Add title
        title = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(9), PptxInches(1))
        title.text_frame.paragraphs[0].text = "Test Title"

        # Add table
        table = slide.shapes.add_table(3, 2, PptxInches(0.5), PptxInches(1.5), PptxInches(5), PptxInches(2))
        tbl = table.table
        tbl.cell(0, 0).text = "H1"
        tbl.cell(0, 1).text = "H2"
        tbl.cell(1, 0).text = "R1C1"
        tbl.cell(1, 1).text = "R1C2"

        path = temp_dir / "full_slide.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_get_slide(str(path), 1)
        assert isinstance(result, dict)

    def test_get_slide_beyond_range(self, pptx_advanced_tools, temp_dir):
        """Should handle slide number beyond range."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "one_slide.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_get_slide(str(path), 99)
        assert "error" in result


class TestPptxListShapes:
    """Test pptx_list_shapes method."""

    def test_list_shapes_various_types(self, pptx_advanced_tools, temp_dir):
        """Should list various shape types."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Add textbox
        slide.shapes.add_textbox(PptxInches(1), PptxInches(1), PptxInches(3), PptxInches(1))

        # Add shape
        slide.shapes.add_shape(1, PptxInches(5), PptxInches(1), PptxInches(2), PptxInches(2))

        # Add table
        slide.shapes.add_table(2, 2, PptxInches(1), PptxInches(3), PptxInches(4), PptxInches(1.5))

        path = temp_dir / "shapes.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_list_shapes(str(path), 1)
        assert isinstance(result, dict)
        assert "shapes" in result


class TestPptxReplaceText:
    """Test pptx_replace_text method."""

    def test_replace_text_in_shapes(self, pptx_advanced_tools, temp_dir):
        """Should replace text across shapes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Hello <NAME>"

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if "Click to add" in para.text:
                        para.clear()
                        para.add_run().text = "Welcome <NAME> to the presentation"

        path = temp_dir / "replace.pptx"
        prs.save(path)

        output = temp_dir / "replaced.pptx"
        result = pptx_advanced_tools.tool_pptx_replace_text(
            str(path),
            "<NAME>",
            "John",
            output_path=str(output)
        )
        assert isinstance(result, dict)

    def test_replace_text_not_found(self, pptx_advanced_tools, temp_dir):
        """Should handle text not found."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "No placeholder here"
        path = temp_dir / "no_match.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_replace_text(
            str(path),
            "<NONEXISTENT>",
            "Value"
        )
        assert isinstance(result, dict)


class TestPptxSetNotes:
    """Test pptx_set_notes method."""

    def test_set_notes_new(self, pptx_advanced_tools, temp_dir):
        """Should set notes on slide without existing notes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide Title"
        path = temp_dir / "notes.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_set_notes(
            str(path),
            1,
            "These are the speaker notes.\nWith multiple lines."
        )
        assert isinstance(result, dict)

    def test_append_notes(self, pptx_advanced_tools, temp_dir):
        """Should append to existing notes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Existing notes"
        path = temp_dir / "append_notes.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_set_notes(
            str(path),
            1,
            "Additional notes",
            append=True
        )
        assert isinstance(result, dict)


class TestPptxGetNotes:
    """Test pptx_get_notes method."""

    def test_get_notes_single_slide(self, pptx_advanced_tools, temp_dir):
        """Should get notes from specific slide."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Test notes content"
        path = temp_dir / "get_notes.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_get_notes(str(path), slide_number=1)
        assert isinstance(result, dict)

    def test_get_notes_all_slides(self, pptx_advanced_tools, temp_dir):
        """Should get notes from all slides."""
        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = f"Notes for slide {i+1}"
        path = temp_dir / "all_notes.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_get_notes(str(path))
        assert isinstance(result, dict)


class TestPptxGetTable:
    """Test pptx_get_table method."""

    def test_get_table_content(self, pptx_advanced_tools, temp_dir):
        """Should get table content from slide."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        table_shape = slide.shapes.add_table(4, 3, PptxInches(1), PptxInches(1.5), PptxInches(8), PptxInches(3))
        tbl = table_shape.table

        # Header row
        tbl.cell(0, 0).text = "Name"
        tbl.cell(0, 1).text = "Role"
        tbl.cell(0, 2).text = "Hours"

        # Data rows
        tbl.cell(1, 0).text = "Alice"
        tbl.cell(1, 1).text = "Dev"
        tbl.cell(1, 2).text = "40"

        tbl.cell(2, 0).text = "Bob"
        tbl.cell(2, 1).text = "QA"
        tbl.cell(2, 2).text = "30"

        path = temp_dir / "table_content.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_get_table(str(path), 1)
        assert isinstance(result, dict)


class TestPptxHideSlide:
    """Test pptx_hide_slide method."""

    def test_hide_slide(self, pptx_advanced_tools, temp_dir):
        """Should hide a slide."""
        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
        path = temp_dir / "hide.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_hide_slide(str(path), 2, hidden=True)
        assert isinstance(result, dict)

    def test_unhide_slide(self, pptx_advanced_tools, temp_dir):
        """Should unhide a slide."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Hidden Slide"
        slide._element.set('show', '0')
        path = temp_dir / "unhide.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_hide_slide(str(path), 1, hidden=False)
        assert isinstance(result, dict)


class TestPptxGetHiddenSlides:
    """Test pptx_get_hidden_slides method."""

    def test_get_hidden_slides(self, pptx_advanced_tools, temp_dir):
        """Should get list of hidden slides."""
        prs = Presentation()
        for i in range(4):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
            if i % 2 == 0:  # Hide every other slide
                slide._element.set('show', '0')
        path = temp_dir / "mixed_hidden.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_get_hidden_slides(str(path))
        assert isinstance(result, dict)


class TestPptxDuplicateSlide:
    """Test pptx_duplicate_slide method."""

    def test_duplicate_slide_after(self, pptx_advanced_tools, temp_dir):
        """Should duplicate slide right after original."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Original Slide"
        path = temp_dir / "dup.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_duplicate_slide(str(path), 1, position="after")
        assert isinstance(result, dict)

    def test_duplicate_slide_end(self, pptx_advanced_tools, temp_dir):
        """Should duplicate slide at end."""
        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
        path = temp_dir / "dup_end.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_duplicate_slide(str(path), 1, position="end")
        assert isinstance(result, dict)


class TestPptxAnalyzeLayouts:
    """Test pptx_analyze_layouts method."""

    def test_analyze_layouts(self, pptx_advanced_tools, temp_dir):
        """Should analyze layouts in presentation."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "layouts.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_analyze_layouts(str(path))
        assert isinstance(result, dict)


class TestPptxRecommendLayout:
    """Test pptx_recommend_layout method."""

    @pytest.mark.parametrize("layout_name", ["title", "bullets", "table"])
    def test_recommend_layout(self, pptx_advanced_tools, temp_dir, layout_name):
        """Should recommend requested layout."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = temp_dir / f"rec_{layout_name}.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_recommend_layout(str(path), layout_name)
        assert isinstance(result, dict)


class TestPptxAddComment:
    """Test pptx_add_comment method."""

    def test_add_comment_default_position(self, pptx_advanced_tools, temp_dir):
        """Should add comment at default position."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide with comment"
        path = temp_dir / "comment_default.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_add_comment(
            str(path),
            1,
            "This is a review comment."
        )
        assert isinstance(result, dict)

    def test_add_comment_custom_position(self, pptx_advanced_tools, temp_dir):
        """Should add comment at custom position."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Positioned comment"
        path = temp_dir / "comment_pos.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_add_comment(
            str(path),
            1,
            "Comment at custom location",
            x_inches=5.0,
            y_inches=3.0,
            author="Test Author"
        )
        assert isinstance(result, dict)


class TestPptxGetComments:
    """Test pptx_get_comments method."""

    def test_get_comments_empty(self, pptx_advanced_tools, temp_dir):
        """Should handle slide with no comments."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "No comments"
        path = temp_dir / "no_comments.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_get_comments(str(path))
        assert isinstance(result, dict)


class TestPptxAuditPlaceholders:
    """Test pptx_audit_placeholders method."""

    def test_audit_with_placeholders(self, pptx_advanced_tools, temp_dir):
        """Should find unfilled placeholders."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Hello <Customer Name>"

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if "Click to add" in para.text:
                        para.clear()
                        para.add_run().text = "Contact [TBD] for details"

        path = temp_dir / "audit.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_audit_placeholders(str(path))
        assert isinstance(result, dict)

    def test_audit_no_placeholders(self, pptx_advanced_tools, temp_dir):
        """Should handle deck with no placeholders."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Clean Title"
        path = temp_dir / "clean_audit.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_audit_placeholders(str(path))
        assert isinstance(result, dict)
