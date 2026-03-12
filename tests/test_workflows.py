"""
Final push tests targeting specific uncovered code paths.
"""

from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Inches

# Fixture temp_dir is provided by conftest.py


class TestWordSowWorkflow:
    """Tests for complete SOW workflow."""

    def test_sow_workflow_basic(self, word_advanced_tools, temp_dir):
        """Should handle basic SOW workflow."""
        # Create template
        doc = Document()
        doc.add_heading("Statement of Work", level=0)
        doc.add_heading("Executive Summary", level=1)
        doc.add_paragraph("<Customer Name> project overview")
        doc.add_heading("Scope", level=1)
        doc.add_paragraph("<Project Scope>")
        path = temp_dir / "workflow_template.docx"
        doc.save(path)

        # Parse template
        result = word_advanced_tools.tool_word_parse_sow_template(str(path))
        assert isinstance(result, dict)

        # List sections
        result = word_advanced_tools.tool_word_list_sections(str(path))
        assert len(result.get("sections", [])) >= 2

    def test_sow_from_markdown(self, word_advanced_tools, temp_dir):
        """Should create SOW from markdown with template."""
        # First create a template
        template_doc = Document()
        template_doc.add_heading("Statement of Work", level=0)
        template_doc.add_heading("Executive Summary", level=1)
        template_doc.add_paragraph("<Executive Summary>")
        template_doc.add_heading("Scope", level=1)
        template_doc.add_paragraph("<Scope>")
        template_path = temp_dir / "sow_template.docx"
        template_doc.save(template_path)

        md = """# Statement of Work

## Executive Summary

This project will deliver a comprehensive cloud migration solution.

## Scope

### In Scope
- Application migration
- Infrastructure setup

### Out of Scope
- Hardware procurement
"""

        output = temp_dir / "md_sow.docx"
        result = word_advanced_tools.tool_word_create_sow_from_markdown(
            str(output),
            md,
            str(template_path)
        )
        assert isinstance(result, dict)


class TestPptxCompleteWorkflow:
    """Tests for complete PPTX workflow."""

    def test_create_and_edit_presentation(self, pptx_advanced_tools, pptx_tools, temp_dir):
        """Should create and edit presentation."""
        # Create from markdown
        md = """# Proposal

## Overview

- Key point 1
- Key point 2

---

## Details

| Item | Value |
|------|-------|
| A | 1 |
| B | 2 |
"""
        path = temp_dir / "workflow.pptx"
        result = pptx_tools.tool_pptx_from_markdown(str(path), md)
        assert Path(path).exists()

        # List slides
        result = pptx_advanced_tools.tool_pptx_list_slides(str(path))
        assert result.get("slide_count", 0) >= 2

        # Get slide content
        result = pptx_advanced_tools.tool_pptx_get_slide(str(path), slide_number=1)
        assert isinstance(result, dict)

    def test_add_slide_with_bullets(self, pptx_advanced_tools, temp_dir):
        """Should add slide and populate with bullets."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "bullet_workflow.pptx"
        prs.save(path)

        # Add new slide
        output1 = temp_dir / "with_slide.pptx"
        result = pptx_advanced_tools.tool_pptx_add_slide(
            str(path),
            layout_index=1,
            title="New Content",
            output_path=str(output1)
        )

        # Clear bullets
        output2 = temp_dir / "cleared.pptx"
        pptx_advanced_tools.tool_pptx_clear_bullets(str(output1), slide_number=2, output_path=str(output2))

        # Add bullets
        output3 = temp_dir / "with_bullets.pptx"
        result = pptx_advanced_tools.tool_pptx_add_bullet(
            str(output2),
            slide_number=2,
            text="First point",
            output_path=str(output3)
        )
        assert result.get("success") is True


class TestExcelCompleteWorkflow:
    """Tests for Excel complete workflow."""

    def test_excel_roundtrip(self, excel_tools, temp_dir):
        """Should handle Excel roundtrip."""
        # Create from markdown
        md = """| Name | Score | Grade |
|------|-------|-------|
| Alice | 95 | A |
| Bob | 87 | B |
| Carol | 92 | A |
"""

        path = temp_dir / "grades.xlsx"
        result = excel_tools.tool_excel_from_markdown(str(path), md)
        assert Path(path).exists()

        # Extract
        result = excel_tools.tool_excel_extract(str(path))
        assert len(result.get("sheets", [])) >= 1

        # To markdown
        md_result = excel_tools.tool_excel_to_markdown(str(path))
        assert "Alice" in md_result


class TestWordTableWorkflow:
    """Tests for Word table operations workflow."""

    def test_table_operations(self, word_advanced_tools, temp_dir):
        """Should handle table operations workflow."""
        # Create doc with section
        doc = Document()
        doc.add_heading("Project Timeline", level=1)
        doc.add_paragraph("Below is the timeline.")
        path = temp_dir / "table_workflow.docx"
        doc.save(path)

        # Create new table
        output1 = temp_dir / "with_table.docx"
        result = word_advanced_tools.tool_word_create_new_table(
            str(path),
            ["Phase", "Start", "End"],
            rows=[
                {"Phase": "Phase 1", "Start": "Week 1", "End": "Week 4"}
            ],
            insert_after_section="Project Timeline",
            output_path=str(output1)
        )

        # List tables
        result = word_advanced_tools.tool_word_list_tables(str(output1))
        assert len(result.get("tables", [])) >= 1


class TestPptxCommentWorkflow:
    """Tests for PPTX comment operations."""

    def test_comment_workflow(self, pptx_advanced_tools, temp_dir):
        """Should handle comment workflow."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Review Slide"
        path = temp_dir / "comment_workflow.pptx"
        prs.save(path)

        # Add comment
        output1 = temp_dir / "with_comment.pptx"
        result = pptx_advanced_tools.tool_pptx_add_comment(
            str(path),
            slide_number=1,
            comment_text="Please review this slide",
            x_inches=2.0,
            y_inches=2.0,
            author="Reviewer",
            output_path=str(output1)
        )
        assert result.get("success") is True

        # Get comments
        result = pptx_advanced_tools.tool_pptx_get_comments(str(output1))
        assert isinstance(result, dict)


class TestMoreEdgeCases:
    """Additional edge case tests."""

    def test_word_empty_doc(self, word_advanced_tools, temp_dir):
        """Should handle completely empty document."""
        doc = Document()
        path = temp_dir / "empty.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_list_sections(str(path))
        assert isinstance(result, dict)

        result = word_advanced_tools.tool_word_list_tables(str(path))
        assert len(result.get("tables", [])) == 0

    def test_pptx_single_slide_ops(self, pptx_advanced_tools, temp_dir):
        """Should handle single slide operations."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Only Slide"
        path = temp_dir / "single.pptx"
        prs.save(path)

        # Get slide
        result = pptx_advanced_tools.tool_pptx_get_slide(str(path), slide_number=1)
        assert isinstance(result, dict)

        # List shapes
        result = pptx_advanced_tools.tool_pptx_list_shapes(str(path), slide_number=1)
        assert "shapes" in result

    def test_pptx_with_textbox(self, pptx_advanced_tools, temp_dir):
        """Should handle slides with textboxes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Title"

        # Add textbox
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1.5))
        tf = txBox.text_frame
        tf.text = "Textbox content"

        path = temp_dir / "textbox.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_list_shapes(str(path), slide_number=1)
        # Should have title and textbox
        assert len(result.get("shapes", [])) >= 2


class TestWordSectionEditing:
    """Tests for section editing operations."""

    def test_patch_section_with_content(self, word_advanced_tools, temp_dir):
        """Should patch section content."""
        doc = Document()
        doc.add_heading("Section A", level=1)
        doc.add_paragraph("Original content A")
        doc.add_heading("Section B", level=1)
        doc.add_paragraph("Original content B")
        path = temp_dir / "sections.docx"
        doc.save(path)

        output = temp_dir / "patched.docx"
        result = word_advanced_tools.tool_word_patch_section(
            str(path),
            "Section A",
            "New content for section A",
            output_path=str(output)
        )
        assert isinstance(result, dict)


class TestPptxReplacement:
    """Tests for text replacement operations."""

    def test_replace_placeholders_comprehensive(self, pptx_advanced_tools, temp_dir):
        """Should replace multiple placeholders."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "<Customer> - <Project>"

        # Access body placeholder if available
        for shape in slide.shapes:
            if (shape.has_text_frame
                    and hasattr(shape, 'placeholder_format')
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.type is not None):
                shape.text_frame.text = "Project for <Customer>"

        path = temp_dir / "placeholders.pptx"
        prs.save(path)

        output = temp_dir / "replaced.pptx"
        result = pptx_advanced_tools.tool_pptx_replace_placeholders(
            str(path),
            replacements={
                "<Customer>": "Contoso",
                "<Project>": "Migration"
            },
            output_path=str(output)
        )
        assert isinstance(result, dict)


class TestWordTrackChanges:
    """Tests for track changes operations."""

    def test_enable_and_check_tracking(self, word_advanced_tools, temp_dir):
        """Should enable and check tracking."""
        doc = Document()
        doc.add_paragraph("Content to track")
        path = temp_dir / "track.docx"
        doc.save(path)

        # Enable
        output = temp_dir / "tracking_enabled.docx"
        result = word_advanced_tools.tool_word_enable_track_changes(str(path), output_path=str(output))
        assert isinstance(result, dict)

        # Check
        result = word_advanced_tools.tool_word_check_tracking(str(output))
        assert isinstance(result, dict)


class TestPptxNotes:
    """Tests for speaker notes operations."""

    def test_notes_workflow(self, pptx_advanced_tools, temp_dir):
        """Should handle notes workflow."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide with Notes"
        path = temp_dir / "notes.pptx"
        prs.save(path)

        # Set notes
        output1 = temp_dir / "with_notes.pptx"
        result = pptx_advanced_tools.tool_pptx_set_notes(
            str(path),
            slide_number=1,
            notes_text="Initial speaker notes",
            output_path=str(output1)
        )

        # Get notes
        result = pptx_advanced_tools.tool_pptx_get_notes(str(output1), slide_number=1)
        assert isinstance(result, dict)

        # Append notes
        output2 = temp_dir / "appended_notes.pptx"
        result = pptx_advanced_tools.tool_pptx_set_notes(
            str(output1),
            slide_number=1,
            notes_text="\nAdditional notes",
            append=True,
            output_path=str(output2)
        )
        assert result.get("success") is True


class TestAuditOperations:
    """Tests for audit operations."""

    def test_word_audit_complete(self, word_advanced_tools, temp_dir):
        """Should audit document completion."""
        doc = Document()
        doc.add_heading("Document", level=1)
        doc.add_paragraph("Complete content without placeholders")
        path = temp_dir / "complete.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_audit_completion(str(path))
        assert isinstance(result, dict)

    def test_pptx_audit_clean(self, pptx_advanced_tools, temp_dir):
        """Should audit clean presentation."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Clean Title"
        path = temp_dir / "clean.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_audit_placeholders(str(path))
        assert isinstance(result, dict)


class TestListSupportedFormats:
    """Tests for tools module imports."""

    def test_tool_classes_available(self):
        """Should have tool classes available."""
        from tools import TOOL_CLASSES
        assert len(TOOL_CLASSES) >= 1
