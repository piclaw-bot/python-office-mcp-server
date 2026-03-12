"""
Additional tests for tools/__init__.py and core functionality.
"""

from pathlib import Path

# Fixture temp_dir is provided by conftest.py


class TestExcelToolsMoreCoverage:
    """Extended tests for ExcelTools."""

    def test_extract_empty_workbook(self, excel_tools, temp_dir):
        """Should handle empty workbook."""
        from openpyxl import Workbook

        wb = Workbook()
        path = temp_dir / "empty.xlsx"
        wb.save(path)

        result = excel_tools.tool_excel_extract(str(path))
        assert isinstance(result, dict)

    def test_extract_multiple_sheets(self, excel_tools, temp_dir):
        """Should extract all sheets."""
        from openpyxl import Workbook

        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1["A1"] = "Data 1"

        ws2 = wb.create_sheet("Sheet2")
        ws2["A1"] = "Data 2"

        path = temp_dir / "multi.xlsx"
        wb.save(path)

        result = excel_tools.tool_excel_extract(str(path))
        assert len(result.get("sheets", [])) >= 2

    def test_to_markdown_with_numbers(self, excel_tools, temp_dir):
        """Should handle numeric values."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Amount"
        ws["B1"] = "Percentage"
        ws["A2"] = 1000.50
        ws["B2"] = 0.15

        path = temp_dir / "numbers.xlsx"
        wb.save(path)

        result = excel_tools.tool_excel_to_markdown(str(path))
        assert "Amount" in result
        assert "1000" in result or "1,000" in result


class TestWordToolsMoreCoverage:
    """Extended tests for WordTools."""

    def test_extract_with_lists(self, word_tools, temp_dir):
        """Should extract documents with bullet lists."""
        from docx import Document

        doc = Document()
        doc.add_heading("List Test", level=1)
        for item in ["Item A", "Item B", "Item C"]:
            doc.add_paragraph(item, style='List Bullet')

        path = temp_dir / "lists.docx"
        doc.save(path)

        result = word_tools.tool_word_extract(str(path))
        assert "Item A" in str(result)

    def test_to_markdown_with_headers(self, word_tools, temp_dir):
        """Should convert headers correctly."""
        from docx import Document

        doc = Document()
        doc.add_heading("H1 Title", level=0)
        doc.add_heading("H2 Section", level=1)
        doc.add_paragraph("Content")
        doc.add_heading("H3 Subsection", level=2)

        path = temp_dir / "headers.docx"
        doc.save(path)

        result = word_tools.tool_word_to_markdown(str(path))
        assert "# " in result or "## " in result


class TestPowerPointToolsMoreCoverage:
    """Extended tests for PowerPointTools."""

    def test_extract_with_shapes(self, pptx_tools, temp_dir):
        """Should extract presentations with various shapes."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Shape Test"

        # Add a textbox
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        tf = txBox.text_frame
        tf.text = "Textbox content"

        path = temp_dir / "shapes.pptx"
        prs.save(path)

        result = pptx_tools.tool_pptx_extract(str(path))
        assert "Shape Test" in str(result) or "Textbox" in str(result)

    def test_from_markdown_complex(self, pptx_tools, temp_dir):
        """Should handle complex markdown."""
        md = """# Main Presentation

## First Section

- **Point 1:** Detail about first point
- **Point 2:** Detail about second point
- Sub-point here

---

## Second Section

| Col A | Col B | Col C |
|-------|-------|-------|
| 1 | 2 | 3 |
| 4 | 5 | 6 |

---

## Summary

Final thoughts and conclusions.
"""
        path = temp_dir / "complex.pptx"
        result = pptx_tools.tool_pptx_from_markdown(str(path), md)
        assert Path(path).exists()
        # Result uses "slides" key, not "slide_count"
        assert result.get("slides", 0) >= 3 or result.get("slide_count", 0) >= 3
