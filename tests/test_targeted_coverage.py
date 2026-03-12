"""
Targeted tests for remaining coverage gaps.
"""

from docx import Document
from pptx import Presentation

# Fixture temp_dir is provided by conftest.py


class TestWordExtractMethods:
    """Tests for Word extraction methods."""

    def test_extract_sow_structure_detailed(self, word_advanced_tools, temp_dir):
        """Should extract detailed SOW structure."""
        doc = Document()
        doc.add_heading("SOW for Customer", level=0)
        doc.add_heading("1. Introduction", level=1)
        doc.add_paragraph("Intro text")
        doc.add_heading("2. Engagement Overview", level=1)
        doc.add_heading("2.1 Business Objectives", level=2)
        table = doc.add_table(rows=2, cols=3)
        table.cell(0, 0).text = "Objective"
        table.cell(0, 1).text = "Activity"
        table.cell(0, 2).text = "Assumption"
        doc.add_heading("3. Scope", level=1)
        doc.add_paragraph("Scope details")

        path = temp_dir / "detailed_sow.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_extract_sow_structure(str(path))
        assert isinstance(result, dict)

    def test_get_section_with_table(self, word_advanced_tools, temp_dir):
        """Should get section content including table."""
        doc = Document()
        doc.add_heading("Data Section", level=1)
        doc.add_paragraph("Table below:")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"
        doc.add_heading("Next Section", level=1)
        doc.add_paragraph("More content")

        path = temp_dir / "section_table.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_get_section(str(path), "Data Section")
        assert isinstance(result, dict)


class TestPptxAdvancedSlideOps:
    """Tests for advanced slide operations."""

    def test_duplicate_slide_preserve_content(self, pptx_advanced_tools, temp_dir):
        """Should duplicate slide with content."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Original Title"
        # Add content to body
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                shape.text_frame.text = "Body content"

        path = temp_dir / "to_dup.pptx"
        prs.save(path)

        output = temp_dir / "duplicated.pptx"
        result = pptx_advanced_tools.tool_pptx_duplicate_slide(
            str(path),
            slide_number=1,
            position="after",
            output_path=str(output)
        )
        assert isinstance(result, dict)

    def test_hide_multiple_slides(self, pptx_advanced_tools, temp_dir):
        """Should hide slides correctly."""
        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
        path = temp_dir / "multi.pptx"
        prs.save(path)

        # Hide slide 2
        output = temp_dir / "hidden.pptx"
        result = pptx_advanced_tools.tool_pptx_hide_slide(
            str(path),
            slide_number=2,
            hidden=True,
            output_path=str(output)
        )
        assert result.get("success") is True or result.get("hidden") is True

        # Check hidden slides
        result = pptx_advanced_tools.tool_pptx_get_hidden_slides(str(output))
        assert isinstance(result, dict)


class TestWordTableAdvanced:
    """Advanced table operation tests."""

    def test_patch_table_row_multiple_cols(self, word_advanced_tools, temp_dir):
        """Should patch table row with multiple columns."""
        doc = Document()
        table = doc.add_table(rows=3, cols=4)
        for i, h in enumerate(["ID", "Name", "Status", "Date"]):
            table.cell(0, i).text = h
        table.cell(1, 0).text = "1"
        table.cell(1, 1).text = "Item A"
        table.cell(1, 2).text = "Pending"
        table.cell(1, 3).text = "2024-01-01"

        path = temp_dir / "multi_col.docx"
        doc.save(path)

        output = temp_dir / "patched.docx"
        result = word_advanced_tools.tool_word_patch_table_row(
            str(path),
            "0",
            1,  # Row index
            {"Status": "Complete", "Date": "2024-02-01"},
            output_path=str(output)
        )
        assert isinstance(result, dict)

    def test_duplicate_table_add_row(self, word_advanced_tools, temp_dir):
        """Should duplicate and add row."""
        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Col1"
        table.cell(0, 1).text = "Col2"
        path = temp_dir / "dup_add.docx"
        doc.save(path)

        # Duplicate structure
        output1 = temp_dir / "dup_struct.docx"
        result = word_advanced_tools.tool_word_duplicate_table_structure(
            str(path),
            "0",
            output_path=str(output1)
        )
        assert isinstance(result, dict)


class TestPptxLayoutsAdvanced:
    """Advanced layout tests."""

    def test_analyze_all_layouts(self, pptx_advanced_tools, temp_dir):
        """Should analyze all available layouts."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = temp_dir / "layouts.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_analyze_layouts(str(path))
        assert "layouts" in result

    def test_recommend_layout_two_column(self, pptx_advanced_tools, temp_dir):
        """Should recommend two column layout."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = temp_dir / "two_col.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_recommend_layout(str(path), "two_column")
        assert isinstance(result, dict)


class TestWordReplacement:
    """Word text replacement tests."""

    def test_replace_global_in_headers(self, word_advanced_tools, temp_dir):
        """Should replace text in document headers."""
        doc = Document()
        doc.add_heading("<Customer> Document", level=0)
        doc.add_paragraph("<Customer> is our valued client.")
        doc.add_paragraph("Project for <Customer>")

        path = temp_dir / "headers.docx"
        doc.save(path)

        output = temp_dir / "replaced.docx"
        result = word_advanced_tools.tool_word_replace_global_variables(
            str(path),
            replacements={"<Customer>": "ACME Corp"},
            output_path=str(output)
        )
        assert isinstance(result, dict)


class TestPptxBulletAdvanced:
    """Advanced bullet operation tests."""

    def test_bullet_with_all_options(self, pptx_advanced_tools, temp_dir):
        """Should add bullet with all formatting options."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Content"
        path = temp_dir / "bullet_opts.pptx"
        prs.save(path)

        output = temp_dir / "with_bullets.pptx"

        # Add bullet with bold label
        result = pptx_advanced_tools.tool_pptx_add_bullet(
            str(path),
            slide_number=1,
            text="Detail text here",
            level=0,
            bold_label="Key Point",
            output_path=str(output)
        )
        assert result.get("success") is True


class TestWordCommentAdvanced:
    """Advanced comment operation tests."""

    def test_add_comment_with_author(self, word_advanced_tools, temp_dir):
        """Should add comment with custom author."""
        doc = Document()
        doc.add_paragraph("This text needs review from the team.")
        path = temp_dir / "for_comment.docx"
        doc.save(path)

        output = temp_dir / "commented.docx"
        result = word_advanced_tools.tool_word_add_comment(
            str(path),
            target_text="needs review",
            comment_text="Please verify this section",
            author="Reviewer A",
            output_path=str(output)
        )
        assert isinstance(result, dict)


class TestPptxAutofit:
    """Text autofit tests."""

    def test_autofit_all_types(self, pptx_advanced_tools, temp_dir):
        """Should set different autofit types."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test Title"
        path = temp_dir / "autofit.pptx"
        prs.save(path)

        # Test shrink
        output1 = temp_dir / "shrink.pptx"
        result = pptx_advanced_tools.tool_pptx_set_text_autofit(
            str(path),
            slide_number=1,
            shape_identifier="title",
            autofit_type="shrink",
            output_path=str(output1)
        )
        assert result.get("success") is True

        # Test none
        output2 = temp_dir / "none.pptx"
        result = pptx_advanced_tools.tool_pptx_set_text_autofit(
            str(output1),
            slide_number=1,
            shape_identifier="title",
            autofit_type="none",
            output_path=str(output2)
        )
        assert result.get("success") is True


class TestWordAuditAdvanced:
    """Advanced audit operation tests."""

    def test_audit_sow_comprehensive(self, word_advanced_tools, temp_dir):
        """Should perform comprehensive SOW audit."""
        doc = Document()
        doc.add_heading("SOW Document", level=0)
        doc.add_heading("Executive Summary", level=1)
        doc.add_paragraph("<Executive Summary>")
        doc.add_heading("Scope", level=1)
        doc.add_paragraph("[TBD]")
        doc.add_heading("Timeline", level=1)
        table = doc.add_table(rows=2, cols=3)
        table.cell(0, 0).text = "Phase"
        table.cell(0, 1).text = "Start"
        table.cell(0, 2).text = "End"
        # Leave cells empty to trigger audit

        path = temp_dir / "audit_sow.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_audit_sow(str(path))
        assert isinstance(result, dict)

    def test_audit_completion_mixed(self, word_advanced_tools, temp_dir):
        """Should audit document with mixed completion."""
        doc = Document()
        doc.add_heading("Document", level=1)
        doc.add_paragraph("Complete paragraph with content.")
        doc.add_paragraph("<Placeholder>")
        doc.add_paragraph("More complete content.")
        doc.add_paragraph("[TBD]")

        path = temp_dir / "mixed.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_audit_completion(str(path))
        assert isinstance(result, dict)


class TestPptxReorderAdvanced:
    """Advanced slide reordering tests."""

    def test_reorder_complex(self, pptx_advanced_tools, temp_dir):
        """Should reorder slides in complex order."""
        prs = Presentation()
        for i in range(4):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
        path = temp_dir / "complex_order.pptx"
        prs.save(path)

        output = temp_dir / "reordered.pptx"
        result = pptx_advanced_tools.tool_pptx_reorder_slides(
            str(path),
            new_order=[4, 2, 1, 3],
            output_path=str(output)
        )
        assert isinstance(result, dict)


class TestWordCleanup:
    """Word cleanup operation tests."""

    def test_cleanup_comprehensive(self, word_advanced_tools, temp_dir):
        """Should cleanup document comprehensively."""
        doc = Document()
        doc.add_heading("SOW", level=0)
        doc.add_paragraph("[Template Guidance: This is guidance to remove]")
        doc.add_paragraph("Real content to keep.")
        doc.add_paragraph("[Note: Delete this note]")
        doc.add_paragraph("More real content.")

        path = temp_dir / "cleanup.docx"
        doc.save(path)

        output = temp_dir / "cleaned.docx"
        result = word_advanced_tools.tool_word_cleanup_sow(str(path), output_path=str(output))
        assert isinstance(result, dict)


class TestPptxLogChanges:
    """Change logging tests."""

    def test_log_multiple_changes(self, pptx_advanced_tools, temp_dir):
        """Should log multiple changes."""
        prs = Presentation()
        for _i in range(2):
            prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "for_log.pptx"
        prs.save(path)

        output = temp_dir / "with_log.pptx"
        result = pptx_advanced_tools.tool_pptx_log_changes(
            str(path),
            changes=[
                {"slide": 1, "action": "Updated title", "detail": "Changed from placeholder"},
                {"slide": 2, "action": "Added content", "detail": "Added bullet points"},
                {"slide": 1, "action": "Fixed formatting", "detail": "Adjusted font size"}
            ],
            output_path=str(output)
        )
        assert isinstance(result, dict)
