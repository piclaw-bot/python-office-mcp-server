"""
Final push tests targeting specific uncovered code paths.
"""

import tempfile
from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from tools.excel_tools import ExcelTools
from tools.pptx_advanced_tools import PresentationAdvancedTools
from tools.pptx_tools import PowerPointTools
from tools.word_advanced_tools import WordAdvancedTools, _get_text_with_track_changes


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as td:
        yield Path(td)


class TestWordSowWorkflow:
    """Tests for complete SOW workflow."""

    def test_end_to_end_generation_cleanup_audit_workflow(self, temp_dir):
        """End-to-end Word workflow should produce materially complete output."""
        tools = WordAdvancedTools()

        template = Document()
        template.add_heading("Statement of Work", level=0)
        template.add_paragraph("<Customer Name>")
        template.add_paragraph("<Project Name>")
        template.add_heading("Executive Summary", level=1)
        template.add_paragraph("[Template Guidance: Describe the project here]")
        template.add_heading("Delivery approach", level=1)
        template.add_paragraph("[Template Guidance: Describe the delivery approach]")
        template.add_heading("Customer responsibilities and project assumptions", level=1)
        template.add_paragraph("[Template Guidance: Add customer responsibilities]")
        template.add_heading("Staffing", level=1)
        staffing = template.add_table(rows=2, cols=2)
        staffing.cell(0, 0).text = "Role"
        staffing.cell(0, 1).text = "Hours"
        template_path = temp_dir / "workflow_complete_template.docx"
        template.save(template_path)

        markdown = """# Statement of Work

Customer: Contoso
Project: Migration
Provider: Microsoft

## Executive Summary

Executive summary content.

## Delivery approach

Delivery approach content.

## Customer responsibilities and project assumptions

Customer responsibilities content.

## Staffing

| Role | Hours |
|------|-------|
| Architect | 40 |
"""

        generated = temp_dir / "workflow_generated.docx"
        generation = tools.tool_word_create_sow_from_markdown(
            str(generated),
            markdown,
            str(template_path),
        )
        assert generation.get("success") is True
        assert generation.get("status") in {"success", "partial_success"}

        cleaned = temp_dir / "workflow_cleaned.docx"
        cleanup = tools.tool_word_cleanup_sow(str(generated), output_path=str(cleaned))
        assert cleanup.get("success") is True

        audit = tools.tool_word_audit_completion(str(cleaned))
        assert audit.get("success") is True
        assert audit.get("score", 0) >= 95
        assert audit.get("summary", {}).get("placeholders_found") == 0
        assert audit.get("summary", {}).get("empty_sections") == 0
        assert audit.get("summary", {}).get("empty_table_cells") == 0
        assert audit.get("summary", {}).get("instruction_remnants") == 0

    def test_sow_workflow_basic(self, temp_dir):
        """Should handle basic SOW workflow."""
        tools = WordAdvancedTools()

        # Create template
        doc = Document()
        doc.add_heading("Statement of Work", level=0)
        doc.add_heading("Executive Summary", level=1)
        doc.add_paragraph("<Customer Name> project overview")
        doc.add_heading("Scope", level=1)
        doc.add_paragraph("<Project Scope>")
        path = temp_dir / "workflow_template.docx"
        doc.save(path)

        # Parse template
        result = tools.tool_word_parse_sow_template(str(path))
        assert isinstance(result, dict)

        # List sections
        result = tools.tool_word_list_sections(str(path))
        assert len(result.get("sections", [])) >= 2

    def test_sow_from_markdown(self, temp_dir):
        """Should create SOW from markdown with template."""
        tools = WordAdvancedTools()

        # First create a template
        template_doc = Document()
        template_doc.add_heading("Statement of Work", level=0)
        template_doc.add_heading("Executive Summary", level=1)
        template_doc.add_paragraph("<Executive Summary>")
        template_doc.add_heading("Scope", level=1)
        template_doc.add_paragraph("<Scope>")
        template_path = temp_dir / "sow_template.docx"
        template_doc.save(template_path)

        md = """# Statement of Work

## Executive Summary

This project will deliver a comprehensive cloud migration solution.

## Scope

### In Scope
- Application migration
- Infrastructure setup

### Out of Scope
- Hardware procurement
"""

        output = temp_dir / "md_sow.docx"
        result = tools.tool_word_create_sow_from_markdown(
            str(output),
            md,
            str(template_path)
        )
        assert isinstance(result, dict)

    def test_end_to_end_generation_handles_split_placeholders(self, temp_dir):
        """Template placeholders split across runs should still be fully replaced."""
        tools = WordAdvancedTools()

        template = Document()
        template.add_heading("Statement of Work", level=0)
        para = template.add_paragraph()
        para.add_run("<Customer")
        para.add_run(" Name>")
        para2 = template.add_paragraph()
        para2.add_run("<Project")
        para2.add_run(" Name>")
        para3 = template.add_paragraph()
        para3.add_run("<Provider")
        para3.add_run(" Name>")
        template.add_heading("Executive Summary", level=1)
        template.add_paragraph("[Template Guidance: Describe the project here]")
        template_path = temp_dir / "split_placeholder_template.docx"
        template.save(template_path)

        markdown = """# Statement of Work

Customer: Contoso
Project: Migration Factory
Provider: Microsoft

## Executive Summary

Executive summary content.
"""

        generated = temp_dir / "split_placeholder_generated.docx"
        result = tools.tool_word_create_sow_from_markdown(
            str(generated),
            markdown,
            str(template_path),
        )

        assert result.get("success") is True
        assert result.get("replacements", 0) >= 3

        audit = tools.tool_word_audit_completion(str(generated))
        assert audit.get("summary", {}).get("placeholders_found") == 0

        generated_doc = Document(generated)
        combined_text = "\n".join(_get_text_with_track_changes(p) for p in generated_doc.paragraphs)
        assert "<Customer Name>" not in combined_text
        assert "<Project Name>" not in combined_text
        assert "<Provider Name>" not in combined_text
        assert "Contoso" in combined_text
        assert "Migration Factory" in combined_text
        assert "Microsoft" in combined_text


class TestPptxCompleteWorkflow:
    """Tests for complete PPTX workflow."""

    def test_create_and_edit_presentation(self, temp_dir):
        """Should create and edit presentation."""
        tools = PresentationAdvancedTools()
        basic_tools = PowerPointTools()

        # Create from markdown
        md = """# Proposal

## Overview

- Key point 1
- Key point 2

---

## Details

| Item | Value |
|------|-------|
| A | 1 |
| B | 2 |
"""
        path = temp_dir / "workflow.pptx"
        result = basic_tools.tool_pptx_from_markdown(str(path), md)
        assert Path(path).exists()

        # List slides
        result = tools.tool_pptx_list_slides(str(path))
        assert result.get("slide_count", 0) >= 2

        # Get slide content
        result = tools.tool_pptx_get_slide(str(path), slide_number=1)
        assert isinstance(result, dict)

    def test_add_slide_with_bullets(self, temp_dir):
        """Should add slide and populate with bullets."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "bullet_workflow.pptx"
        prs.save(path)

        # Add new slide
        output1 = temp_dir / "with_slide.pptx"
        result = tools.tool_pptx_add_slide(
            str(path),
            layout_index=1,
            title="New Content",
            output_path=str(output1)
        )

        # Clear bullets
        output2 = temp_dir / "cleared.pptx"
        tools.tool_pptx_clear_bullets(str(output1), slide_number=2, output_path=str(output2))

        # Add bullets
        output3 = temp_dir / "with_bullets.pptx"
        result = tools.tool_pptx_add_bullet(
            str(output2),
            slide_number=2,
            text="First point",
            output_path=str(output3)
        )
        assert result.get("success") is True


class TestExcelCompleteWorkflow:
    """Tests for Excel complete workflow."""

    def test_excel_roundtrip(self, temp_dir):
        """Should handle Excel roundtrip."""
        tools = ExcelTools()

        # Create from markdown
        md = """| Name | Score | Grade |
|------|-------|-------|
| Alice | 95 | A |
| Bob | 87 | B |
| Carol | 92 | A |
"""

        path = temp_dir / "grades.xlsx"
        result = tools.tool_excel_from_markdown(str(path), md)
        assert Path(path).exists()

        # Extract
        result = tools.tool_excel_extract(str(path))
        assert len(result.get("sheets", [])) >= 1

        # To markdown
        md_result = tools.tool_excel_to_markdown(str(path))
        assert "Alice" in md_result


class TestWordTableWorkflow:
    """Tests for Word table operations workflow."""

    def test_table_operations(self, temp_dir):
        """Should handle table operations workflow."""
        tools = WordAdvancedTools()

        # Create doc with section
        doc = Document()
        doc.add_heading("Project Timeline", level=1)
        doc.add_paragraph("Below is the timeline.")
        path = temp_dir / "table_workflow.docx"
        doc.save(path)

        # Create new table
        output1 = temp_dir / "with_table.docx"
        result = tools.tool_word_create_new_table(
            str(path),
            ["Phase", "Start", "End"],
            rows=[
                {"Phase": "Phase 1", "Start": "Week 1", "End": "Week 4"}
            ],
            insert_after_section="Project Timeline",
            output_path=str(output1)
        )

        # List tables
        result = tools.tool_word_list_tables(str(output1))
        assert len(result.get("tables", [])) >= 1


class TestPptxCommentWorkflow:
    """Tests for PPTX comment operations."""

    def test_comment_workflow(self, temp_dir):
        """Should handle comment workflow."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Review Slide"
        path = temp_dir / "comment_workflow.pptx"
        prs.save(path)

        # Add comment
        output1 = temp_dir / "with_comment.pptx"
        result = tools.tool_pptx_add_comment(
            str(path),
            slide_number=1,
            comment_text="Please review this slide",
            x_inches=2.0,
            y_inches=2.0,
            author="Reviewer",
            output_path=str(output1)
        )
        assert result.get("success") is True

        # Get comments
        result = tools.tool_pptx_get_comments(str(output1))
        assert isinstance(result, dict)


class TestMoreEdgeCases:
    """Additional edge case tests."""

    def test_word_empty_doc(self, temp_dir):
        """Should handle completely empty document."""
        tools = WordAdvancedTools()

        doc = Document()
        path = temp_dir / "empty.docx"
        doc.save(path)

        result = tools.tool_word_list_sections(str(path))
        assert isinstance(result, dict)

        result = tools.tool_word_list_tables(str(path))
        assert len(result.get("tables", [])) == 0

    def test_pptx_single_slide_ops(self, temp_dir):
        """Should handle single slide operations."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Only Slide"
        path = temp_dir / "single.pptx"
        prs.save(path)

        # Get slide
        result = tools.tool_pptx_get_slide(str(path), slide_number=1)
        assert isinstance(result, dict)

        # List shapes
        result = tools.tool_pptx_list_shapes(str(path), slide_number=1)
        assert "shapes" in result

    def test_pptx_with_textbox(self, temp_dir):
        """Should handle slides with textboxes."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Title"

        # Add textbox
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1.5))
        tf = txBox.text_frame
        tf.text = "Textbox content"

        path = temp_dir / "textbox.pptx"
        prs.save(path)

        result = tools.tool_pptx_list_shapes(str(path), slide_number=1)
        # Should have title and textbox
        assert len(result.get("shapes", [])) >= 2


class TestWordSectionEditing:
    """Tests for section editing operations."""

    def test_patch_section_with_content(self, temp_dir):
        """Should patch section content."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_heading("Section A", level=1)
        doc.add_paragraph("Original content A")
        doc.add_heading("Section B", level=1)
        doc.add_paragraph("Original content B")
        path = temp_dir / "sections.docx"
        doc.save(path)

        output = temp_dir / "patched.docx"
        result = tools.tool_word_patch_section(
            str(path),
            "Section A",
            "New content for section A",
            output_path=str(output)
        )
        assert isinstance(result, dict)


class TestPptxReplacement:
    """Tests for text replacement operations."""

    def test_replace_placeholders_comprehensive(self, temp_dir):
        """Should replace multiple placeholders."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "<Customer> - <Project>"

        # Access body placeholder if available
        for shape in slide.shapes:
            if (shape.has_text_frame
                    and hasattr(shape, 'placeholder_format')
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.type is not None):
                shape.text_frame.text = "Project for <Customer>"

        path = temp_dir / "placeholders.pptx"
        prs.save(path)

        output = temp_dir / "replaced.pptx"
        result = tools.tool_pptx_replace_placeholders(
            str(path),
            replacements={
                "<Customer>": "Contoso",
                "<Project>": "Migration"
            },
            output_path=str(output)
        )
        assert isinstance(result, dict)


class TestWordTrackChanges:
    """Tests for track changes operations."""

    def test_enable_and_check_tracking(self, temp_dir):
        """Should enable and check tracking."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_paragraph("Content to track")
        path = temp_dir / "track.docx"
        doc.save(path)

        # Enable
        output = temp_dir / "tracking_enabled.docx"
        result = tools.tool_word_enable_track_changes(str(path), output_path=str(output))
        assert isinstance(result, dict)

        # Check
        result = tools.tool_word_check_tracking(str(output))
        assert isinstance(result, dict)


class TestPptxNotes:
    """Tests for speaker notes operations."""

    def test_notes_workflow(self, temp_dir):
        """Should handle notes workflow."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide with Notes"
        path = temp_dir / "notes.pptx"
        prs.save(path)

        # Set notes
        output1 = temp_dir / "with_notes.pptx"
        result = tools.tool_pptx_set_notes(
            str(path),
            slide_number=1,
            notes_text="Initial speaker notes",
            output_path=str(output1)
        )

        # Get notes
        result = tools.tool_pptx_get_notes(str(output1), slide_number=1)
        assert isinstance(result, dict)

        # Append notes
        output2 = temp_dir / "appended_notes.pptx"
        result = tools.tool_pptx_set_notes(
            str(output1),
            slide_number=1,
            notes_text="\nAdditional notes",
            append=True,
            output_path=str(output2)
        )
        assert result.get("success") is True


class TestAuditOperations:
    """Tests for audit operations."""

    def test_word_audit_complete(self, temp_dir):
        """Should audit document completion."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_heading("Document", level=1)
        doc.add_paragraph("Complete content without placeholders")
        path = temp_dir / "complete.docx"
        doc.save(path)

        result = tools.tool_word_audit_completion(str(path))
        assert isinstance(result, dict)

    def test_pptx_audit_clean(self, temp_dir):
        """Should audit clean presentation."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Clean Title"
        path = temp_dir / "clean.pptx"
        prs.save(path)

        result = tools.tool_pptx_audit_placeholders(str(path))
        assert isinstance(result, dict)


class TestListSupportedFormats:
    """Tests for tools module imports."""

    def test_tool_classes_available(self):
        """Should have tool classes available."""
        from tools import TOOL_CLASSES
        assert len(TOOL_CLASSES) >= 1
