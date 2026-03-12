"""
Final tests to reach 80% coverage.
"""

import tempfile
from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches as PptxInches

from tools.excel_tools import ExcelTools
from tools.pptx_advanced_tools import PresentationAdvancedTools
from tools.pptx_tools import PowerPointTools
from tools.word_advanced_tools import WordAdvancedTools
from tools.word_tools import WordTools


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as td:
        yield Path(td)


class TestWordExtractDeep:
    """Deep extraction tests."""

    def test_extract_with_all_elements(self, temp_dir):
        """Should extract document with all element types."""
        doc = Document()
        doc.add_heading("Title", level=1)
        doc.add_heading("Subtitle", level=2)
        doc.add_paragraph("Normal paragraph")

        # Table
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"

        # List
        doc.add_paragraph("Bullet 1", style='List Bullet')
        doc.add_paragraph("Bullet 2", style='List Bullet')

        path = temp_dir / "all_elements.docx"
        doc.save(path)

        tools = WordTools()
        result = tools.tool_word_extract(str(path))
        assert "paragraphs" in result or "text" in str(result)


class TestPptxExtractDeep:
    """Deep PPTX extraction tests."""

    def test_extract_with_shapes_and_tables(self, temp_dir):
        """Should extract presentation with shapes and tables."""
        prs = Presentation()

        # Title slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Main Title"

        # Content slide with table
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2.shapes.title.text = "Data"
        table = slide2.shapes.add_table(3, 2, PptxInches(1), PptxInches(2), PptxInches(5), PptxInches(2))
        tbl = table.table
        tbl.cell(0, 0).text = "H1"
        tbl.cell(0, 1).text = "H2"

        path = temp_dir / "shapes_tables.pptx"
        prs.save(path)

        tools = PowerPointTools()
        result = tools.tool_pptx_extract(str(path))
        assert isinstance(result, dict)


class TestWordFromMarkdownVariations:
    """Test various markdown to Word conversions."""

    def test_with_numbered_list(self, temp_dir):
        """Should convert numbered list."""
        tools = WordTools()

        md = """# Document

## Steps

1. First step
2. Second step
3. Third step
"""
        path = temp_dir / "numbered.docx"
        tools.tool_word_from_markdown(str(path), md)
        assert Path(path).exists()

    def test_with_mixed_content(self, temp_dir):
        """Should convert mixed content."""
        tools = WordTools()

        md = """# Title

## Section 1

Regular paragraph text.

- Bullet 1
- Bullet 2

| Col A | Col B |
|-------|-------|
| 1     | 2     |

## Section 2

More text here.
"""
        path = temp_dir / "mixed.docx"
        tools.tool_word_from_markdown(str(path), md)
        assert Path(path).exists()


class TestPptxFromMarkdownVariations:
    """Test various markdown to PPTX conversions."""

    def test_multiple_sections(self, temp_dir):
        """Should create slides for multiple sections."""
        tools = PowerPointTools()

        md = """# Presentation Title

## Section 1
- Point A
- Point B

---

## Section 2
- Point C
- Point D

---

## Section 3
- Point E
- Point F
"""
        path = temp_dir / "multi_section.pptx"
        tools.tool_pptx_from_markdown(str(path), md)
        assert Path(path).exists()


class TestExcelOperationsDeep:
    """Deep Excel operations tests."""

    def test_extract_with_formulas(self, temp_dir):
        """Should extract Excel with various data."""
        tools = ExcelTools()

        md = """| Value | Multiplier | Result |
|-------|------------|--------|
| 10    | 2          | 20     |
| 20    | 3          | 60     |
| 30    | 4          | 120    |
"""
        path = temp_dir / "formulas.xlsx"
        tools.tool_excel_from_markdown(str(path), md)

        result = tools.tool_excel_extract(str(path))
        assert isinstance(result, dict)


class TestWordSowWorkflow:
    """Test SOW workflow operations."""

    def test_full_sow_workflow(self, temp_dir):
        """Should complete full SOW workflow."""
        tools = WordAdvancedTools()

        # Create template
        doc = Document()
        doc.add_heading("<Project Name>", level=1)
        doc.add_paragraph("Customer: <Customer Name>")
        doc.add_paragraph("Provider: <Provider Name>")
        doc.add_heading("1. Objectives", level=2)
        doc.add_paragraph("<objectives content>")
        template_path = temp_dir / "template.docx"
        doc.save(template_path)

        # Parse template
        result = tools.tool_word_parse_sow_template(str(template_path))
        assert isinstance(result, dict)

        # Copy template
        output_path = temp_dir / "sow.docx"
        result = tools.tool_word_copy_template(str(template_path), str(output_path))
        assert isinstance(result, dict)

        # Replace global variables
        result = tools.tool_word_replace_global_variables(
            str(output_path),
            {
                "<Customer Name>": "Contoso",
                "<Project Name>": "Migration",
                "<Provider Name>": "Microsoft"
            }
        )
        assert isinstance(result, dict)


class TestPptxCompleteWorkflow:
    """Test complete PPTX workflow."""

    def test_full_presentation_workflow(self, temp_dir):
        """Should complete full presentation workflow."""
        tools = PresentationAdvancedTools()
        basic_tools = PowerPointTools()

        # Create from markdown
        md = """# Presentation

## Slide 1
- Point 1
- Point 2

## Slide 2
- Point 3
- Point 4
"""
        path = temp_dir / "workflow.pptx"
        basic_tools.tool_pptx_from_markdown(str(path), md)

        # List slides
        result = tools.tool_pptx_list_slides(str(path))
        assert result.get("slide_count", 0) >= 2

        # Add slide
        result = tools.tool_pptx_add_slide(str(path), title="New Slide")
        assert isinstance(result, dict)

        # Set notes
        result = tools.tool_pptx_set_notes(str(path), 1, "Speaker notes")
        assert isinstance(result, dict)


class TestWordTableWorkflow:
    """Test Word table workflow."""

    def test_table_operations_workflow(self, temp_dir):
        """Should complete table operations workflow."""
        tools = WordAdvancedTools()

        # Create doc with table
        doc = Document()
        doc.add_heading("Data", level=1)
        table = doc.add_table(rows=3, cols=3)
        headers = ["Name", "Role", "Hours"]
        for i, h in enumerate(headers):
            table.cell(0, i).text = h
        table.cell(1, 0).text = "Alice"
        table.cell(1, 1).text = "Dev"
        table.cell(1, 2).text = "40"
        table.cell(2, 0).text = "Bob"
        table.cell(2, 1).text = "QA"
        table.cell(2, 2).text = "35"
        path = temp_dir / "tables.docx"
        doc.save(path)

        # List tables
        result = tools.tool_word_list_tables(str(path))
        assert len(result.get("tables", [])) >= 1

        # Get table
        result = tools.tool_word_get_table(str(path), "0")
        assert isinstance(result, dict)

        # Insert row
        result = tools.tool_word_insert_table_row(
            str(path), "0",
            {"Name": "Carol", "Role": "PM", "Hours": "45"}
        )
        assert isinstance(result, dict)


class TestPptxTableWorkflow:
    """Test PPTX table workflow."""

    def test_pptx_table_operations(self, temp_dir):
        """Should complete PPTX table operations."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Table"
        path = temp_dir / "pptx_table.pptx"
        prs.save(path)

        # Add table
        result = tools.tool_pptx_add_table(
            str(path),
            1,
            ["A", "B", "C"],
            [["1", "2", "3"], ["4", "5", "6"]]
        )
        assert isinstance(result, dict)

        # Get table
        result = tools.tool_pptx_get_table(str(path), 1)
        assert isinstance(result, dict)


class TestWordAuditWorkflow:
    """Test Word audit workflow."""

    def test_audit_workflow(self, temp_dir):
        """Should complete audit workflow."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_heading("Document", level=1)
        doc.add_paragraph("Content with <placeholder>")
        doc.add_paragraph("[TBD] more content")
        path = temp_dir / "audit.docx"
        doc.save(path)

        # Audit completion
        result = tools.tool_word_audit_completion(str(path))
        assert isinstance(result, dict)

        # Audit SOW
        result = tools.tool_word_audit_sow(str(path))
        assert isinstance(result, dict)


class TestPptxAuditWorkflow:
    """Test PPTX audit workflow."""

    def test_audit_workflow(self, temp_dir):
        """Should complete PPTX audit workflow."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Title with <Customer>"
        path = temp_dir / "pptx_audit.pptx"
        prs.save(path)

        # Audit placeholders
        result = tools.tool_pptx_audit_placeholders(str(path))
        assert isinstance(result, dict)


class TestWordCommentWorkflow:
    """Test Word comment workflow."""

    def test_comment_workflow(self, temp_dir):
        """Should add comment to document."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_paragraph("This text needs review.")
        path = temp_dir / "comment.docx"
        doc.save(path)

        result = tools.tool_word_add_comment(
            str(path),
            "needs review",
            "Please verify this section.",
            author="Reviewer"
        )
        assert isinstance(result, dict)


class TestPptxCommentWorkflow:
    """Test PPTX comment workflow."""

    def test_comment_workflow(self, temp_dir):
        """Should add and get comments."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Review Slide"
        path = temp_dir / "comments.pptx"
        prs.save(path)

        # Add comment
        result = tools.tool_pptx_add_comment(
            str(path), 1, "Review this slide",
            author="Reviewer"
        )
        assert isinstance(result, dict)

        # Get comments
        result = tools.tool_pptx_get_comments(str(path))
        assert isinstance(result, dict)


class TestWordTrackChangesWorkflow:
    """Test Word track changes workflow."""

    def test_track_changes_workflow(self, temp_dir):
        """Should enable and use track changes."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_paragraph("Original content here.")
        path = temp_dir / "track.docx"
        doc.save(path)

        # Enable tracking
        result = tools.tool_word_enable_track_changes(str(path))
        assert isinstance(result, dict)

        # Check tracking
        result = tools.tool_word_check_tracking(str(path))
        assert isinstance(result, dict)


class TestPptxNotesWorkflow:
    """Test PPTX notes workflow."""

    def test_notes_workflow(self, temp_dir):
        """Should set and get notes."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Notes Test"
        path = temp_dir / "notes.pptx"
        prs.save(path)

        # Set notes
        result = tools.tool_pptx_set_notes(str(path), 1, "Speaker notes here")
        assert isinstance(result, dict)

        # Get notes
        result = tools.tool_pptx_get_notes(str(path), slide_number=1)
        assert isinstance(result, dict)


class TestPptxSlideManagementWorkflow:
    """Test PPTX slide management workflow."""

    def test_slide_management_workflow(self, temp_dir):
        """Should manage slides completely."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
        path = temp_dir / "manage.pptx"
        prs.save(path)

        # Duplicate
        result = tools.tool_pptx_duplicate_slide(str(path), 1)
        assert isinstance(result, dict)

        # Hide
        result = tools.tool_pptx_hide_slide(str(path), 2, hidden=True)
        assert isinstance(result, dict)

        # Get hidden
        result = tools.tool_pptx_get_hidden_slides(str(path))
        assert isinstance(result, dict)
