"""Tests for SVG and PNG image insertion into Office documents.

Validates office_image and related helpers across Word, Excel, and PowerPoint.
Covers PNG insertion (supported), SVG insertion (currently blocked at validation),
dimension parsing, targeting, and error paths.
"""

from pathlib import Path

import pytest

try:
    from PIL import Image as PILImage
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    import pptx  # noqa: F401
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from docx import Document

from tools import TOOL_CLASSES

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


def _create_combined_tools():
    class CombinedTools(*TOOL_CLASSES):
        pass
    return CombinedTools()


@pytest.fixture
def tools():
    return _create_combined_tools()


@pytest.fixture
def png_image(temp_dir):
    """Create a minimal 100x80 red PNG image."""
    if not HAS_PIL:
        pytest.skip("Pillow not installed")
    img = PILImage.new("RGB", (100, 80), color="red")
    path = temp_dir / "test_image.png"
    img.save(path, dpi=(96, 96))
    return path


@pytest.fixture
def png_image_hires(temp_dir):
    """Create a 300 DPI PNG to test DPI-aware dimension calculation."""
    if not HAS_PIL:
        pytest.skip("Pillow not installed")
    img = PILImage.new("RGB", (300, 300), color="blue")
    path = temp_dir / "hires.png"
    img.save(path, dpi=(300, 300))
    return path


@pytest.fixture
def svg_image_with_dims(temp_dir):
    """Create an SVG with explicit width/height attributes (px)."""
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="192" height="96">'
        '<rect width="192" height="96" fill="green"/>'
        "</svg>"
    )
    path = temp_dir / "dims.svg"
    path.write_text(svg)
    return path


@pytest.fixture
def svg_image_viewbox_only(temp_dir):
    """Create an SVG with viewBox but no width/height attributes."""
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 480 240">'
        '<rect width="480" height="240" fill="orange"/>'
        "</svg>"
    )
    path = temp_dir / "viewbox.svg"
    path.write_text(svg)
    return path


@pytest.fixture
def svg_image_inches(temp_dir):
    """Create an SVG with dimensions in inches."""
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="3in" height="2in">'
        '<rect width="100%" height="100%" fill="purple"/>'
        "</svg>"
    )
    path = temp_dir / "inches.svg"
    path.write_text(svg)
    return path


@pytest.fixture
def svg_image_mm(temp_dir):
    """Create an SVG with dimensions in millimetres."""
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="50.8mm" height="25.4mm">'
        '<rect width="100%" height="100%" fill="teal"/>'
        "</svg>"
    )
    path = temp_dir / "mm.svg"
    path.write_text(svg)
    return path


@pytest.fixture
def svg_image_percent(temp_dir):
    """Create an SVG with percentage dimensions (unresolvable)."""
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="100%" height="100%"'
        ' viewBox="0 0 960 480">'
        '<rect width="960" height="480" fill="gray"/>'
        "</svg>"
    )
    path = temp_dir / "percent.svg"
    path.write_text(svg)
    return path


@pytest.fixture
def sample_docx(temp_dir):
    """Create a Word document with sections for targeted insertion."""
    doc = Document()
    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph("Summary content goes here.")
    doc.add_heading("Architecture", level=1)
    doc.add_paragraph("Architecture details here.")
    path = temp_dir / "doc_for_images.docx"
    doc.save(path)
    return path


@pytest.fixture
def sample_xlsx(temp_dir):
    """Create an Excel workbook with two sheets."""
    if not HAS_OPENPYXL:
        pytest.skip("openpyxl not installed")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Main"
    ws["A1"] = "Data"
    ws2 = wb.create_sheet("Charts")
    ws2["A1"] = "Chart Area"
    path = temp_dir / "wb_for_images.xlsx"
    wb.save(path)
    return path


@pytest.fixture
def sample_pptx(temp_dir):
    """Create a two-slide presentation."""
    if not HAS_PPTX:
        pytest.skip("python-pptx not installed")
    from pptx import Presentation as Prs
    from pptx.util import Mm

    prs = Prs()
    prs.slide_width = Mm(338.67)
    prs.slide_height = Mm(190.5)
    layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(layout)
    if slide1.shapes.title:
        slide1.shapes.title.text = "Title Slide"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    if slide2.shapes.title:
        slide2.shapes.title.text = "Content Slide"
    path = temp_dir / "deck_for_images.pptx"
    prs.save(path)
    return path


# ===================================================================
# PNG insertion – Word
# ===================================================================


class TestPngInsertionWord:
    """PNG image insertion into Word documents."""

    def test_insert_png_at_end(self, tools, sample_docx, png_image, temp_dir):
        """Insert PNG at end of document (default behaviour)."""
        out = temp_dir / "out.docx"
        result = tools.tool_office_image(
            file_path=str(sample_docx),
            image_path=str(png_image),
            output_path=str(out),
        )
        assert result.get("status") == "success"
        assert Path(out).exists()

    def test_insert_png_after_section(self, tools, sample_docx, png_image, temp_dir):
        """Insert PNG after a named section heading."""
        out = temp_dir / "out.docx"
        result = tools.tool_office_image(
            file_path=str(sample_docx),
            image_path=str(png_image),
            target="after:Executive Summary",
            output_path=str(out),
        )
        assert result.get("status") == "success"
        assert result.get("location") == "after:Executive Summary"

    def test_insert_png_section_not_found(self, tools, sample_docx, png_image, temp_dir):
        """Error when the target section does not exist."""
        out = temp_dir / "out.docx"
        result = tools.tool_office_image(
            file_path=str(sample_docx),
            image_path=str(png_image),
            target="after:Nonexistent Section",
            output_path=str(out),
        )
        assert "error" in result

    def test_insert_png_with_width(self, tools, sample_docx, png_image, temp_dir):
        """Insert PNG with explicit width; height auto-scaled."""
        out = temp_dir / "out.docx"
        result = tools.tool_office_image(
            file_path=str(sample_docx),
            image_path=str(png_image),
            width_inches=3.0,
            output_path=str(out),
        )
        assert result.get("status") == "success"

    def test_insert_png_with_both_dims(self, tools, sample_docx, png_image, temp_dir):
        """Insert PNG with explicit width and height."""
        out = temp_dir / "out.docx"
        result = tools.tool_office_image(
            file_path=str(sample_docx),
            image_path=str(png_image),
            width_inches=4.0,
            height_inches=2.5,
            output_path=str(out),
        )
        assert result.get("status") == "success"

    def test_insert_png_overwrites_input(self, tools, sample_docx, png_image):
        """When output_path is omitted, the input file is overwritten."""
        result = tools.tool_office_image(
            file_path=str(sample_docx),
            image_path=str(png_image),
        )
        assert result.get("status") == "success"
        assert result.get("file") == str(sample_docx)


# ===================================================================
# PNG insertion – Excel
# ===================================================================


class TestPngInsertionExcel:
    """PNG image insertion into Excel workbooks."""

    @pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
    def test_insert_png_default_cell(self, tools, sample_xlsx, png_image, temp_dir):
        """Insert PNG at default cell A1 on active sheet."""
        out = temp_dir / "out.xlsx"
        result = tools.tool_office_image(
            file_path=str(sample_xlsx),
            image_path=str(png_image),
            output_path=str(out),
        )
        assert result.get("status") == "success"
        assert result.get("cell") == "A1"

    @pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
    def test_insert_png_specific_cell(self, tools, sample_xlsx, png_image, temp_dir):
        """Insert PNG anchored at a specific cell."""
        out = temp_dir / "out.xlsx"
        result = tools.tool_office_image(
            file_path=str(sample_xlsx),
            image_path=str(png_image),
            target="C5",
            output_path=str(out),
        )
        assert result.get("status") == "success"
        assert result.get("cell") == "C5"

    @pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
    def test_insert_png_specific_sheet(self, tools, sample_xlsx, png_image, temp_dir):
        """Insert PNG on a named sheet at a specified cell."""
        out = temp_dir / "out.xlsx"
        result = tools.tool_office_image(
            file_path=str(sample_xlsx),
            image_path=str(png_image),
            target="Charts!B3",
            output_path=str(out),
        )
        assert result.get("status") == "success"
        assert result.get("sheet") == "Charts"
        assert result.get("cell") == "B3"

    @pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
    def test_insert_png_with_dimensions(self, tools, sample_xlsx, png_image, temp_dir):
        """Insert PNG with explicit pixel-scale dimensions."""
        out = temp_dir / "out.xlsx"
        result = tools.tool_office_image(
            file_path=str(sample_xlsx),
            image_path=str(png_image),
            target="A1",
            width_inches=2.0,
            height_inches=1.5,
            output_path=str(out),
        )
        assert result.get("status") == "success"


# ===================================================================
# PNG insertion – PowerPoint
# ===================================================================


class TestPngInsertionPptx:
    """PNG image insertion into PowerPoint presentations."""

    @pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
    def test_insert_png_first_slide(self, tools, sample_pptx, png_image, temp_dir):
        """Insert PNG on slide 1 (default)."""
        out = temp_dir / "out.pptx"
        result = tools.tool_office_image(
            file_path=str(sample_pptx),
            image_path=str(png_image),
            output_path=str(out),
        )
        assert result.get("status") == "success"
        assert result.get("slide") == 1

    @pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
    def test_insert_png_specific_slide(self, tools, sample_pptx, png_image, temp_dir):
        """Insert PNG on slide 2."""
        out = temp_dir / "out.pptx"
        result = tools.tool_office_image(
            file_path=str(sample_pptx),
            image_path=str(png_image),
            target="slide:2",
            output_path=str(out),
        )
        assert result.get("status") == "success"
        assert result.get("slide") == 2

    @pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
    def test_insert_png_invalid_slide(self, tools, sample_pptx, png_image, temp_dir):
        """Error when slide number is out of range."""
        out = temp_dir / "out.pptx"
        result = tools.tool_office_image(
            file_path=str(sample_pptx),
            image_path=str(png_image),
            target="slide:99",
            output_path=str(out),
        )
        assert "error" in result

    @pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
    def test_insert_png_invalid_slide_format(self, tools, sample_pptx, png_image, temp_dir):
        """Error when slide target is not a number."""
        out = temp_dir / "out.pptx"
        result = tools.tool_office_image(
            file_path=str(sample_pptx),
            image_path=str(png_image),
            target="slide:abc",
            output_path=str(out),
        )
        assert "error" in result

    @pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
    def test_insert_png_with_width_only(self, tools, sample_pptx, png_image, temp_dir):
        """Width specified, height should auto-scale proportionally."""
        out = temp_dir / "out.pptx"
        result = tools.tool_office_image(
            file_path=str(sample_pptx),
            image_path=str(png_image),
            target="slide:1",
            width_inches=5.0,
            output_path=str(out),
        )
        assert result.get("status") == "success"

    @pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
    def test_insert_png_with_height_only(self, tools, sample_pptx, png_image, temp_dir):
        """Height specified, width should auto-scale proportionally."""
        out = temp_dir / "out.pptx"
        result = tools.tool_office_image(
            file_path=str(sample_pptx),
            image_path=str(png_image),
            target="slide:1",
            height_inches=3.0,
            output_path=str(out),
        )
        assert result.get("status") == "success"

    @pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
    def test_insert_png_centered(self, tools, sample_pptx, png_image, temp_dir):
        """Image should be centered on the slide (check position metadata)."""
        out = temp_dir / "out.pptx"
        result = tools.tool_office_image(
            file_path=str(sample_pptx),
            image_path=str(png_image),
            target="slide:1",
            width_inches=3.0,
            height_inches=2.0,
            output_path=str(out),
        )
        assert result.get("status") == "success"
        pos = result.get("position", {})
        # Verify that position metadata is returned
        assert pos, "Expected position metadata for centering verification"


# ===================================================================
# SVG validation gate – office_image currently rejects SVG
# ===================================================================


class TestSvgInsertion:
    """SVG image insertion via OOXML 2016 SVG extension.

    Word and PowerPoint support direct SVG embedding with a 1px PNG fallback.
    Excel does not support SVG (openpyxl limitation).
    """

    def test_svg_succeeds_word(self, tools, sample_docx, svg_image_with_dims, temp_dir):
        """SVG should be embedded in Word via OOXML SVG extension."""
        out = temp_dir / "out.docx"
        result = tools.tool_office_image(
            file_path=str(sample_docx),
            image_path=str(svg_image_with_dims),
            output_path=str(out),
        )
        assert result.get("status") == "success"
        assert Path(out).exists()
        # Verify SVG part is in the package
        from zipfile import ZipFile
        with ZipFile(out) as z:
            svg_parts = [n for n in z.namelist() if n.endswith(".svg")]
            assert svg_parts, "SVG part should be in the docx package"

    def test_svg_word_with_target(self, tools, sample_docx, svg_image_with_dims, temp_dir):
        """SVG insertion after a named section in Word."""
        out = temp_dir / "out.docx"
        result = tools.tool_office_image(
            file_path=str(sample_docx),
            image_path=str(svg_image_with_dims),
            target="after:Executive Summary",
            width_inches=3.0,
            output_path=str(out),
        )
        assert result.get("status") == "success"
        assert result.get("location") == "after:Executive Summary"

    @pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
    def test_svg_rejected_excel(self, tools, sample_xlsx, svg_image_with_dims, temp_dir):
        """SVG should be explicitly rejected for Excel."""
        out = temp_dir / "out.xlsx"
        result = tools.tool_office_image(
            file_path=str(sample_xlsx),
            image_path=str(svg_image_with_dims),
            output_path=str(out),
        )
        assert "error" in result
        assert "not supported" in result["error"].lower()

    @pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
    def test_svg_succeeds_pptx(self, tools, sample_pptx, svg_image_with_dims, temp_dir):
        """SVG should be embedded in PowerPoint via OOXML SVG extension."""
        out = temp_dir / "out.pptx"
        result = tools.tool_office_image(
            file_path=str(sample_pptx),
            image_path=str(svg_image_with_dims),
            target="slide:1",
            width_inches=2.0,
            height_inches=1.0,
            output_path=str(out),
        )
        assert result.get("status") == "success"
        assert Path(out).exists()
        # Verify SVG part is in the package
        from zipfile import ZipFile
        with ZipFile(out) as z:
            svg_parts = [n for n in z.namelist() if n.endswith(".svg")]
            assert svg_parts, "SVG part should be in the pptx package"

    @pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
    def test_svg_pptx_auto_dimensions(self, tools, sample_pptx, svg_image_with_dims, temp_dir):
        """SVG on PPTX with no explicit dimensions should auto-size."""
        out = temp_dir / "out.pptx"
        result = tools.tool_office_image(
            file_path=str(sample_pptx),
            image_path=str(svg_image_with_dims),
            target="slide:2",
            output_path=str(out),
        )
        assert result.get("status") == "success"
        pos = result.get("position", {})
        assert pos.get("width_inches", 0) > 0


# ===================================================================
# _get_image_dimensions – PNG
# ===================================================================


class TestImageDimensionsPng:
    """Dimension parsing for raster (PNG) images."""

    @pytest.mark.skipif(not HAS_PIL, reason="Pillow not installed")
    def test_png_dimensions_96dpi(self, tools, png_image):
        """100x80 px at 96 DPI ≈ 1.042 x 0.833 inches."""
        w, h = tools._get_image_dimensions(str(png_image))
        assert abs(w - 100 / 96) < 0.01
        assert abs(h - 80 / 96) < 0.01

    @pytest.mark.skipif(not HAS_PIL, reason="Pillow not installed")
    def test_png_dimensions_300dpi(self, tools, png_image_hires):
        """300x300 px at 300 DPI = 1.0 x 1.0 inches."""
        w, h = tools._get_image_dimensions(str(png_image_hires))
        assert abs(w - 1.0) < 0.01
        assert abs(h - 1.0) < 0.01


# ===================================================================
# _get_image_dimensions – SVG
# ===================================================================


class TestImageDimensionsSvg:
    """Dimension parsing for SVG images."""

    def test_svg_with_pixel_dims(self, tools, svg_image_with_dims):
        """192x96 px at 96 DPI = 2.0 x 1.0 inches."""
        w, h = tools._get_image_dimensions(str(svg_image_with_dims))
        assert abs(w - 2.0) < 0.01
        assert abs(h - 1.0) < 0.01

    def test_svg_viewbox_fallback(self, tools, svg_image_viewbox_only):
        """viewBox 0 0 480 240 → 5.0 x 2.5 inches at 96 DPI."""
        w, h = tools._get_image_dimensions(str(svg_image_viewbox_only))
        assert abs(w - 5.0) < 0.01
        assert abs(h - 2.5) < 0.01

    def test_svg_inches_unit(self, tools, svg_image_inches):
        """Explicit 3in x 2in dimensions."""
        w, h = tools._get_image_dimensions(str(svg_image_inches))
        assert abs(w - 3.0) < 0.01
        assert abs(h - 2.0) < 0.01

    def test_svg_mm_unit(self, tools, svg_image_mm):
        """50.8mm x 25.4mm = 2.0 x 1.0 inches."""
        w, h = tools._get_image_dimensions(str(svg_image_mm))
        assert abs(w - 2.0) < 0.01
        assert abs(h - 1.0) < 0.01

    def test_svg_percent_falls_back_to_viewbox(self, tools, svg_image_percent):
        """Percentage dims should fall back to viewBox (960x480 → 10 x 5 in)."""
        w, h = tools._get_image_dimensions(str(svg_image_percent))
        assert abs(w - 10.0) < 0.01
        assert abs(h - 5.0) < 0.01

    def test_svg_no_dims_no_viewbox(self, tools, temp_dir):
        """SVG with no dimensions or viewBox uses default 6x4 inches."""
        svg = '<svg xmlns="http://www.w3.org/2000/svg"><circle r="10"/></svg>'
        path = temp_dir / "bare.svg"
        path.write_text(svg)
        w, h = tools._get_image_dimensions(str(path))
        assert abs(w - 6.0) < 0.01
        assert abs(h - 4.0) < 0.01


# ===================================================================
# Common error paths
# ===================================================================


class TestImageInsertionErrors:
    """Error handling for office_image."""

    def test_image_file_not_found(self, tools, sample_docx, temp_dir):
        """Error when image path does not exist."""
        result = tools.tool_office_image(
            file_path=str(sample_docx),
            image_path="/nonexistent/missing.png",
        )
        assert "error" in result
        assert "not found" in result["error"].lower()

    def test_unsupported_image_format(self, tools, sample_docx, temp_dir):
        """Error for an image type not in the allow-list."""
        bmp = temp_dir / "image.bmp"
        if HAS_PIL:
            PILImage.new("RGB", (10, 10)).save(bmp)
        else:
            bmp.write_bytes(b"\x00" * 100)
        result = tools.tool_office_image(
            file_path=str(sample_docx),
            image_path=str(bmp),
        )
        assert "error" in result
        assert "unsupported" in result["error"].lower()

    def test_unsupported_document_format(self, tools, png_image, temp_dir):
        """Error when the document type is not recognised."""
        txt = temp_dir / "notes.txt"
        txt.write_text("hello")
        result = tools.tool_office_image(
            file_path=str(txt),
            image_path=str(png_image),
        )
        assert "error" in result

    def test_document_file_not_found(self, tools, png_image):
        """Error when the document itself does not exist."""
        result = tools.tool_office_image(
            file_path="/nonexistent/doc.docx",
            image_path=str(png_image),
        )
        # Should error - either at file-not-found or during open
        assert "error" in result
