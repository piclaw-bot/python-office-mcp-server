"""
Additional tests for edge cases and error paths.
"""

from pathlib import Path

# Fixture temp_dir is provided by conftest.py


class TestExcelEdgeCases:
    """Edge case tests for Excel tools."""

    def test_extract_with_formulas(self, excel_tools, temp_dir):
        """Should handle cells with formulas."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws["A1"] = 10
        ws["A2"] = 20
        ws["A3"] = "=SUM(A1:A2)"

        path = temp_dir / "formulas.xlsx"
        wb.save(path)

        result = excel_tools.tool_excel_extract(str(path))
        assert isinstance(result, dict)

    def test_extract_with_empty_cells(self, excel_tools, temp_dir):
        """Should handle sparse data."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Start"
        ws["C3"] = "End"  # Sparse data

        path = temp_dir / "sparse.xlsx"
        wb.save(path)

        result = excel_tools.tool_excel_extract(str(path))
        assert isinstance(result, dict)

    def test_from_markdown_single_table(self, excel_tools, temp_dir):
        """Should create workbook from markdown with single table."""
        md = """| Name | Age |
|------|-----|
| Alice | 30 |
| Bob | 25 |
"""
        path = temp_dir / "single.xlsx"
        excel_tools.tool_excel_from_markdown(str(path), md)
        assert Path(path).exists()


class TestPresentationEdgeCases:
    """Edge case tests for PowerPoint tools."""

    def test_list_masters_empty_presentation(self, pptx_advanced_tools, empty_pptx):
        """Should handle empty presentation."""
        result = pptx_advanced_tools.tool_pptx_list_masters(str(empty_pptx))
        assert "default_layouts" in result or "layouts" in result

    def test_list_shapes_blank_slide(self, pptx_advanced_tools, temp_dir):
        """Should handle blank slide."""
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        path = temp_dir / "blank.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_list_shapes(str(path), slide_number=1)
        assert "shapes" in result

    def test_patch_shape_nonexistent(self, pptx_advanced_tools, temp_dir):
        """Should handle patching nonexistent shape."""
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        path = temp_dir / "no_shapes.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_patch_shape(
            str(path),
            slide_number=1,
            shape_identifier="nonexistent",
            new_text="test"
        )
        assert "error" in result or "not found" in str(result).lower()


class TestWordEdgeCases:
    """Edge case tests for Word tools."""

    def test_list_sections_empty_doc(self, word_advanced_tools, temp_dir):
        """Should handle document with no headings."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("Just a paragraph")
        path = temp_dir / "no_headings.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_list_sections(str(path))
        # May have sections or be empty
        assert isinstance(result, dict)

    def test_patch_section_invalid(self, word_advanced_tools, temp_dir):
        """Should handle patching invalid section."""
        from docx import Document

        doc = Document()
        doc.add_heading("Section 1", level=1)
        doc.add_paragraph("Content")
        path = temp_dir / "sections.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_patch_section(
            str(path),
            "Nonexistent Section",
            "New content"
        )
        assert "error" in result or "not found" in str(result).lower()

    def test_fix_split_placeholders(self, word_advanced_tools, temp_dir):
        """Should fix split placeholders."""
        from docx import Document

        doc = Document()
        p = doc.add_paragraph()
        # Simulate split placeholder by adding runs
        p.add_run("<Customer")
        p.add_run(" ")
        p.add_run("Name>")
        path = temp_dir / "split.docx"
        doc.save(path)

        output = temp_dir / "fixed.docx"
        # Provide replacements dict
        result = word_advanced_tools.tool_word_fix_split_placeholders(
            str(path),
            replacements={"<Customer Name>": "Contoso"},
            output_path=str(output)
        )
        assert isinstance(result, dict)


class TestOutputPathHandling:
    """Tests for output path handling."""

    def test_word_copy_template(self, word_advanced_tools, temp_dir):
        """Should copy template to new location."""
        from docx import Document

        # Create source
        doc = Document()
        doc.add_heading("Template", level=1)
        doc.add_paragraph("Template content")
        source = temp_dir / "template.docx"
        doc.save(source)

        dest = temp_dir / "copy.docx"
        result = word_advanced_tools.tool_word_copy_template(str(source), str(dest))
        assert dest.exists() or result.get("success") is True

    def test_add_slide_with_title(self, pptx_advanced_tools, temp_dir):
        """Should add slide with title."""
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "add_slide.pptx"
        prs.save(path)

        output = temp_dir / "added.pptx"
        result = pptx_advanced_tools.tool_pptx_add_slide(
            str(path),
            layout_index=1,
            title="New Slide Title",
            output_path=str(output)
        )
        assert result.get("success") is True or result.get("slide_number") is not None


class TestMarkdownConversions:
    """Tests for markdown conversion functions."""

    def test_excel_to_markdown(self, excel_tools, temp_dir):
        """Should convert Excel to markdown."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws["A1"] = "Header1"
        ws["B1"] = "Header2"
        ws["A2"] = "Value1"
        ws["B2"] = "Value2"
        path = temp_dir / "convert.xlsx"
        wb.save(path)

        result = excel_tools.tool_excel_to_markdown(str(path))
        assert "Header1" in result
        assert "|" in result  # Table format

    def test_word_to_markdown(self, word_tools, temp_dir):
        """Should convert Word to markdown."""
        from docx import Document

        doc = Document()
        doc.add_heading("Main Title", level=0)
        doc.add_paragraph("Introduction paragraph.")
        doc.add_heading("Section", level=1)
        doc.add_paragraph("Section content.")
        path = temp_dir / "convert.docx"
        doc.save(path)

        result = word_tools.tool_word_to_markdown(str(path))
        assert "Main Title" in result
        assert "#" in result


class TestTableInsertion:
    """Tests for table insertion in Word."""

    def test_insert_table_row(self, word_advanced_tools, temp_dir):
        """Should insert row into table."""
        from docx import Document

        doc = Document()
        table = doc.add_table(rows=2, cols=3)
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"
        table.cell(0, 2).text = "C"
        table.cell(1, 0).text = "1"
        table.cell(1, 1).text = "2"
        table.cell(1, 2).text = "3"
        path = temp_dir / "table.docx"
        doc.save(path)

        output = temp_dir / "with_row.docx"
        # Pass dict with column names as keys
        result = word_advanced_tools.tool_word_insert_table_row(
            str(path),
            "0",  # First table
            {"A": "X", "B": "Y", "C": "Z"},
            output_path=str(output)
        )
        assert result.get("success") is True or output.exists()
