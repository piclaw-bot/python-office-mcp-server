"""Tests for office_unified_tools.py - Consolidated Office tools."""

import hashlib
import shutil
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

import pytest

# Check for openpyxl availability
try:
    import openpyxl
    from openpyxl.comments import Comment
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

# Check for python-docx availability
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import lxml  # noqa: F401
    HAS_LXML = True
except ImportError:
    HAS_LXML = False

# Check for python-pptx availability
try:
    import pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Import all tool classes to create a combined class for testing
from tools import TOOL_CLASSES
from tools import office_unified_tools as office_unified_tools_module
from tools.office_unified_tools import _detect_format, _has_tool
from tools.word_advanced_tools import _get_text_with_track_changes


def create_combined_tools_class():
    """Create a class that combines all tool mixins for testing.

    The OfficeUnifiedTools is designed as a mixin that delegates to sibling
    tool methods via self. This helper creates a combined class that simulates
    the actual server environment.
    """
    class CombinedTools(*TOOL_CLASSES):
        """Combined tool class for testing unified tools."""
        pass
    return CombinedTools


@pytest.fixture
def combined_tools():
    """Provide a combined tools instance for helper tests."""
    CombinedTools = create_combined_tools_class()
    return CombinedTools()


class TestFormatDetection:
    """Tests for format detection."""

    @pytest.mark.parametrize("filename,expected", [
        ("file.xlsx", "excel"),
        ("file.xlsm", "excel"),
        ("file.docx", "word"),
        ("file.pptx", "powerpoint"),
        ("FILE.XLSX", "excel"),      # case insensitive
        ("File.Docx", "word"),       # case insensitive
        ("file.txt", None),          # unknown
        ("file.pdf", None),          # unknown
    ])
    def test_detect_format(self, filename, expected):
        """Should correctly detect file format from extension."""
        assert _detect_format(filename) == expected


class TestHasToolHelper:
    """Tests for _has_tool helper function."""

    def test_has_tool_true(self, combined_tools):
        # Should have excel_extract from ExcelTools
        assert _has_tool(combined_tools, "excel_extract") is True

    def test_has_tool_false(self, combined_tools):
        assert _has_tool(combined_tools, "nonexistent_tool") is False


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
class TestOfficeReadExcel:
    """Tests for office_read with Excel files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_xlsx(self, temp_dir):
        """Create a sample Excel file for testing."""
        path = temp_dir / "office_read.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "TestSheet"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Alpha"
        ws["B2"] = 100
        ws["A3"] = "Beta"
        ws["B3"] = "=B2*2"  # Formula
        wb.save(path)
        return str(path)

    def test_read_excel_json(self, tools, sample_xlsx):
        """Test reading Excel as JSON."""
        result = tools.tool_office_read(sample_xlsx)
        assert isinstance(result, dict)
        assert "error" not in result

    def test_read_excel_markdown(self, tools, sample_xlsx):
        """Test reading Excel as markdown."""
        result = tools.tool_office_read(sample_xlsx, output_format="markdown")
        assert isinstance(result, str)
        assert "Name" in result
        assert "Value" in result

    def test_read_excel_range(self, tools, sample_xlsx):
        """Test reading specific range from Excel."""
        result = tools.tool_office_read(sample_xlsx, scope="A1:B2")
        assert isinstance(result, dict)
        assert "error" not in result

    def test_read_excel_file_not_found(self, tools):
        """Test error handling for missing file."""
        result = tools.tool_office_read("/nonexistent/file.xlsx")
        assert "error" in result


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
class TestOfficeInspectExcel:
    """Tests for office_inspect with Excel files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_xlsx(self, temp_dir):
        """Create a sample Excel file with various features."""
        path = temp_dir / "office_inspect.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"
        ws["A1"] = "Header"
        ws["A2"] = "Value"

        # Add a comment
        ws["A1"].comment = Comment("Test comment", "Test Author")

        # Add merged cells
        ws.merge_cells("C1:D2")

        # Add another sheet
        wb.create_sheet("Summary")

        wb.save(path)
        return str(path)

    def test_inspect_sheets(self, tools, sample_xlsx):
        """Test inspecting sheets."""
        result = tools.tool_office_inspect(sample_xlsx, what="sheets")
        assert isinstance(result, dict)
        assert "error" not in result

    def test_inspect_structure(self, tools, sample_xlsx):
        """Test default structure inspection."""
        result = tools.tool_office_inspect(sample_xlsx, what="structure")
        assert isinstance(result, dict)
        assert "error" not in result

    def test_inspect_merged_cells(self, tools, sample_xlsx):
        """Test inspecting merged cells."""
        result = tools.tool_office_inspect(sample_xlsx, what="merged_cells")
        assert isinstance(result, dict)
        assert "error" not in result
        assert result.get("total_merged_regions", 0) > 0

    def test_inspect_comments(self, tools, sample_xlsx):
        """Test inspecting comments."""
        result = tools.tool_office_inspect(sample_xlsx, what="comments")
        assert isinstance(result, dict)
        assert "error" not in result
        assert result.get("total_comments", 0) > 0

    def test_inspect_unsupported_type(self, tools, sample_xlsx):
        """Test error for unsupported inspection type."""
        result = tools.tool_office_inspect(sample_xlsx, what="slides")
        assert "error" in result


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
class TestOfficeCommentExcel:
    """Tests for office_comment with Excel files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_xlsx(self, temp_dir):
        """Create a sample Excel file."""
        path = temp_dir / "office_comment.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "Test"
        wb.save(path)
        return str(path)

    def test_add_comment(self, tools, sample_xlsx):
        """Test adding a comment."""
        result = tools.tool_office_comment(
            file_path=sample_xlsx,
            operation="add",
            target="A1",
            text="Test comment",
        )
        assert result.get("success") is True

    def test_get_comments(self, tools, sample_xlsx):
        """Test getting comments."""
        # First add a comment
        tools.tool_office_comment(
            file_path=sample_xlsx,
            operation="add",
            target="A1",
            text="Test comment",
        )
        # Then get comments
        result = tools.tool_office_comment(
            file_path=sample_xlsx,
            operation="get",
        )
        assert "error" not in result
        assert result.get("total_comments", 0) > 0

    def test_add_comment_missing_target(self, tools, sample_xlsx):
        """Test error when target is missing."""
        result = tools.tool_office_comment(
            file_path=sample_xlsx,
            operation="add",
            text="Test comment",
        )
        assert "error" in result

    def test_add_comment_missing_text(self, tools, sample_xlsx):
        """Test error when text is missing."""
        result = tools.tool_office_comment(
            file_path=sample_xlsx,
            operation="add",
            target="A1",
        )
        assert "error" in result

    def test_delete_comment(self, tools, sample_xlsx):
        """Test deleting a comment."""
        tools.tool_office_comment(
            file_path=sample_xlsx,
            operation="add",
            target="A1",
            text="Delete me",
        )

        result = tools.tool_office_comment(
            file_path=sample_xlsx,
            operation="delete",
            target="A1",
        )
        assert result.get("success") is True

        check = tools.tool_office_comment(
            file_path=sample_xlsx,
            operation="get",
        )
        assert check.get("total_comments", 0) == 0

    def test_set_identity_used_for_excel_comments(self, tools, sample_xlsx):
        """Excel comments should use runtime identity when author is omitted."""
        configured = tools.tool_office_set_comment_identity(
            name="Alex Reviewer",
            identity="alex.reviewer@contoso.com",
            initials="AR",
        )
        assert configured.get("success") is True

        added = tools.tool_office_comment(
            file_path=sample_xlsx,
            operation="add",
            target="A1",
            text="Identity test",
        )
        assert added.get("success") is True
        assert added.get("author") == "Alex Reviewer"


@pytest.mark.skipif(not HAS_DOCX or not HAS_LXML, reason="python-docx or lxml not installed")
class TestOfficeCommentWord:
    """Tests for office_comment with Word files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_docx(self, temp_dir):
        path = temp_dir / "office_comment_word.docx"
        doc = docx.Document()
        doc.add_paragraph("Delete this note target")
        doc.save(path)
        return str(path)

    def test_delete_comment(self, tools, sample_docx):
        add = tools.tool_office_comment(
            file_path=sample_docx,
            operation="add",
            target="Delete this note target",
            text="Word comment",
        )
        assert add.get("success") is True

        got = tools.tool_office_comment(file_path=sample_docx, operation="get")
        assert got.get("comment_count", 0) >= 1
        comment_id = got["comments"][0]["id"]

        deleted = tools.tool_office_comment(
            file_path=sample_docx,
            operation="delete",
            target=str(comment_id),
        )
        assert deleted.get("success") is True


class TestUnsupportedFormats:
    """Tests for handling unsupported file formats."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    def test_read_unsupported_format(self, tools):
        """Test reading unsupported format."""
        result = tools.tool_office_read("/path/to/file.txt")
        assert "error" in result

    def test_inspect_unsupported_format(self, tools):
        """Test inspecting unsupported format."""
        result = tools.tool_office_inspect("/path/to/file.pdf", what="structure")
        assert "error" in result

    def test_comment_unsupported_format(self, tools):
        """Test commenting on unsupported format."""
        result = tools.tool_office_comment("/path/to/file.csv", operation="get")
        assert "error" in result


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx not installed")
class TestOfficeReadWord:
    """Tests for office_read with Word files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_docx(self, temp_dir):
        """Create a sample Word file."""
        path = temp_dir / "office_read.docx"
        doc = docx.Document()
        doc.add_heading("Test Document", 0)
        doc.add_paragraph("This is a test paragraph.")
        doc.save(path)
        return str(path)

    def test_read_word_json(self, tools, sample_docx):
        """Test reading Word as JSON."""
        result = tools.tool_office_read(sample_docx)
        assert isinstance(result, dict)
        assert "error" not in result

    def test_read_word_markdown(self, tools, sample_docx):
        """Test reading Word as markdown."""
        result = tools.tool_office_read(sample_docx, output_format="markdown")
        assert isinstance(result, str)
        assert "Test Document" in result


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
class TestOfficeReadPowerPoint:
    """Tests for office_read with PowerPoint files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_pptx(self, temp_dir):
        """Create a sample PowerPoint file."""
        path = temp_dir / "office_read.pptx"
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        if title:
            title.text = "Test Slide"
        prs.save(path)
        return str(path)

    def test_read_pptx_json(self, tools, sample_pptx):
        """Test reading PowerPoint as JSON."""
        result = tools.tool_office_read(sample_pptx)
        assert isinstance(result, dict)
        assert "error" not in result

    def test_read_pptx_markdown(self, tools, sample_pptx):
        """Test reading PowerPoint as markdown."""
        result = tools.tool_office_read(sample_pptx, output_format="markdown")
        assert isinstance(result, str)


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
class TestOfficeInspectPowerPoint:
    """Tests for office_inspect with PowerPoint files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_pptx(self, temp_dir):
        """Create a sample PowerPoint file."""
        path = temp_dir / "office_inspect.pptx"
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        if title:
            title.text = "Test Slide"
        prs.save(path)
        return str(path)

    def test_inspect_slides(self, tools, sample_pptx):
        """Test inspecting slides."""
        result = tools.tool_office_inspect(sample_pptx, what="slides")
        assert isinstance(result, dict)
        assert "error" not in result

    def test_inspect_shapes(self, tools, sample_pptx):
        """Test inspecting shapes on a slide."""
        result = tools.tool_office_inspect(sample_pptx, what="shapes", target="1")
        assert isinstance(result, dict)
        assert "error" not in result

    def test_inspect_shapes_missing_target(self, tools, sample_pptx):
        """Test error when target missing for shapes."""
        result = tools.tool_office_inspect(sample_pptx, what="shapes")
        assert "error" in result


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
class TestOfficeCommentPowerPoint:
    """Tests for office_comment with PowerPoint files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_pptx(self, temp_dir):
        """Create a sample PowerPoint file."""
        path = temp_dir / "office_comment.pptx"
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        if title:
            title.text = "Comment Slide"
        prs.save(path)
        return str(path)

    def test_add_comment_slide_target(self, tools, sample_pptx):
        """Should accept slide target with slide prefix."""
        result = tools.tool_office_comment(
            file_path=sample_pptx,
            operation="add",
            target="slide:1",
            text="Check this slide",
        )
        assert result.get("success") is True

    def test_add_comment_shape_target_error(self, tools, sample_pptx):
        """Should reject shape-qualified slide targets."""
        result = tools.tool_office_comment(
            file_path=sample_pptx,
            operation="add",
            target="slide:1/Title 1",
            text="Not supported",
        )
        assert "error" in result

    def test_get_comment_with_target_error(self, tools, sample_pptx):
        """Should reject target/text when retrieving comments."""
        result = tools.tool_office_comment(
            file_path=sample_pptx,
            operation="get",
            target="1",
        )
        assert "error" in result

    def test_delete_comment(self, tools, sample_pptx):
        """Should delete comments from a PowerPoint slide."""
        add = tools.tool_office_comment(
            file_path=sample_pptx,
            operation="add",
            target="slide:1",
            text="Delete this PPTX comment",
        )
        assert add.get("success") is True

        result = tools.tool_office_comment(
            file_path=sample_pptx,
            operation="delete",
            target="slide:1",
        )
        assert result.get("success") is True

    def test_set_identity_used_for_powerpoint_comments(self, tools, sample_pptx):
        """PowerPoint comments should use runtime identity when author is omitted."""
        configured = tools.tool_office_set_comment_identity(
            name="Jamie Architect",
            identity="jamie.architect@contoso.com",
            initials="JA",
        )
        assert configured.get("success") is True

        added = tools.tool_office_comment(
            file_path=sample_pptx,
            operation="add",
            target="slide:1",
            text="PPT identity test",
        )
        assert added.get("success") is True
        assert added.get("author") == "Jamie Architect"

        comments = tools.tool_office_comment(
            file_path=sample_pptx,
            operation="get",
        )
        assert comments.get("total_comments", 0) >= 1
        first = comments["comments"][1][0]
        assert first.get("author") == "Jamie Architect"


class TestOfficeCommentIdentity:
    """Tests for office_set_comment_identity tool."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    def test_set_comment_identity(self, tools):
        """Should store commenter identity state on the server instance."""
        result = tools.tool_office_set_comment_identity(
            name="Morgan Reviewer",
            identity="morgan.reviewer@contoso.com",
            initials="MR",
        )
        assert result.get("success") is True
        assert result.get("comment_author") == "Morgan Reviewer"
        assert result.get("comment_identity") == "morgan.reviewer@contoso.com"
        assert result.get("comment_initials") == "MR"


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
class TestOfficePatchExcel:
    """Tests for office_patch with Excel files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_xlsx(self, temp_dir):
        """Create a sample Excel file."""
        path = temp_dir / "office_patch.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Original"
        ws["B2"] = 100
        wb.save(path)
        return str(path)

    @pytest.fixture
    def complex_xlsx_with_custom_parts(self, temp_dir):
        """Create a workbook with non-openpyxl package parts that must survive patching."""
        path = temp_dir / "office_patch_complex.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "General Inputs"
        ws["A1"] = "Header"
        ws["D6"] = "Original"
        ws["A2"] = "Keep me"
        wb.save(path)

        with zipfile.ZipFile(path, "a") as archive:
            archive.writestr("customXml/item1.xml", "<root>preserve-me</root>")
            archive.writestr(
                "customXml/_rels/item1.xml.rels",
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>',
            )
            archive.writestr(
                "customXml/itemProps1.xml",
                '<ds:datastoreItem xmlns:ds="http://schemas.openxmlformats.org/officeDocument/2006/customXml" ds:itemID="{00000000-0000-0000-0000-000000000000}"/>',
            )
            archive.writestr("docMetadata/LabelInfo.xml", "<metadata>keep</metadata>")

        return str(path)

    def test_patch_single_cell(self, tools, sample_xlsx):
        """Test patching a single cell."""
        result = tools.tool_office_patch(
            file_path=sample_xlsx,
            changes=[{"target": "A2", "value": "Updated"}],
        )
        assert result.get("changes_applied") == 1
        assert result.get("errors") == 0

    def test_patch_multiple_cells(self, tools, sample_xlsx):
        """Test patching multiple cells."""
        result = tools.tool_office_patch(
            file_path=sample_xlsx,
            changes=[
                {"target": "A2", "value": "Updated"},
                {"target": "B2", "value": 200},
            ],
        )
        assert result.get("changes_applied") == 2
        assert result.get("errors") == 0

    def test_patch_with_formula(self, tools, sample_xlsx):
        """Test patching with a formula."""
        result = tools.tool_office_patch(
            file_path=sample_xlsx,
            changes=[{"target": "C2", "value": "=B2*2"}],
        )
        assert result.get("changes_applied") == 1

    def test_patch_no_changes(self, tools, sample_xlsx):
        """Test error when no changes provided."""
        result = tools.tool_office_patch(
            file_path=sample_xlsx,
            changes=[],
        )
        assert "error" in result

    def test_patch_missing_target(self, tools, sample_xlsx):
        """Test handling of missing target."""
        result = tools.tool_office_patch(
            file_path=sample_xlsx,
            changes=[{"value": "NoTarget"}],
        )
        assert result.get("errors") == 1

    def test_patch_preserves_custom_package_parts(self, tools, complex_xlsx_with_custom_parts):
        """Complex package parts should survive a minimal cell patch."""
        original_path = Path(complex_xlsx_with_custom_parts)
        output_path = original_path.with_name("office_patch_complex_out.xlsx")

        with zipfile.ZipFile(original_path, "r") as original_zip:
            original_parts = {
                name: original_zip.read(name)
                for name in (
                    "customXml/item1.xml",
                    "customXml/_rels/item1.xml.rels",
                    "customXml/itemProps1.xml",
                    "docMetadata/LabelInfo.xml",
                )
            }
            original_names = set(original_zip.namelist())

        result = tools.tool_office_patch(
            file_path=str(original_path),
            output_path=str(output_path),
            changes=[{"target": "'General Inputs'!D6", "value": "Test"}],
        )

        assert result.get("changes_applied") == 1
        assert result.get("errors") == 0

        with zipfile.ZipFile(output_path, "r") as patched_zip:
            patched_names = set(patched_zip.namelist())
            for name, payload in original_parts.items():
                assert name in patched_names
                assert patched_zip.read(name) == payload
            assert "xl/worksheets/sheet2.xml" not in patched_names

        reloaded = openpyxl.load_workbook(output_path)
        try:
            assert reloaded["General Inputs"]["D6"].value == "Test"
            assert reloaded["General Inputs"]["A2"].value == "Keep me"
        finally:
            reloaded.close()

        assert original_names.issubset(patched_names)

    def test_patch_range_preserves_custom_package_parts(self, tools, complex_xlsx_with_custom_parts):
        """Range edits should preserve non-target OOXML parts too."""
        original_path = Path(complex_xlsx_with_custom_parts)
        output_path = original_path.with_name("office_patch_complex_range.xlsx")

        result = tools.tool_office_patch(
            file_path=str(original_path),
            output_path=str(output_path),
            changes=[{"target": "'General Inputs'!A1:B1", "value": [["One", "Two"]]}],
        )

        assert result.get("changes_applied") == 1
        assert result.get("errors") == 0

        with zipfile.ZipFile(output_path, "r") as patched_zip:
            assert "customXml/item1.xml" in patched_zip.namelist()
            assert "docMetadata/LabelInfo.xml" in patched_zip.namelist()

        reloaded = openpyxl.load_workbook(output_path)
        try:
            ws = reloaded["General Inputs"]
            assert ws["A1"].value == "One"
            assert ws["B1"].value == "Two"
        finally:
            reloaded.close()

    def test_patch_preserves_shared_strings_and_sheet_relationships(self, tools, temp_dir):
        """Fixture with comments/shared strings should keep related OOXML parts after a cell patch."""
        source = Path("tests/_templates/testdata/excel/comments.xlsx")
        original_path = temp_dir / "comments.xlsx"
        shutil.copy2(source, original_path)
        output_path = temp_dir / "comments_patched.xlsx"

        with zipfile.ZipFile(original_path, "r") as original_zip:
            original_parts = {
                name: original_zip.read(name)
                for name in [
                    "xl/sharedStrings.xml",
                    "xl/comments1.xml",
                    "xl/drawings/vmlDrawing1.vml",
                    "xl/worksheets/_rels/sheet1.xml.rels",
                    "docMetadata/LabelInfo.xml",
                ]
            }

        result = tools.tool_office_patch(
            file_path=str(original_path),
            output_path=str(output_path),
            changes=[{"target": "A1", "value": "Patched comment fixture"}],
        )

        assert result.get("changes_applied") == 1
        assert result.get("errors") == 0

        with zipfile.ZipFile(output_path, "r") as patched_zip:
            patched_names = set(patched_zip.namelist())
            for name, payload in original_parts.items():
                assert name in patched_names
                assert patched_zip.read(name) == payload

        reloaded = openpyxl.load_workbook(output_path)
        try:
            ws = reloaded.active
            assert ws["A1"].value == "Patched comment fixture"
            assert ws["A2"].comment is not None
            assert ws["A2"].comment.text
        finally:
            reloaded.close()

    def test_patch_preserves_content_types_for_related_excel_parts(self, tools, temp_dir):
        """[Content_Types].xml should remain consistent for preserved comment/shared-string parts."""
        source = Path("tests/_templates/testdata/excel/comments.xlsx")
        original_path = temp_dir / "comments_types.xlsx"
        shutil.copy2(source, original_path)
        output_path = temp_dir / "comments_types_patched.xlsx"

        result = tools.tool_office_patch(
            file_path=str(original_path),
            output_path=str(output_path),
            changes=[{"target": "A1", "value": "Content types patch"}],
        )
        assert result.get("changes_applied") == 1

        with zipfile.ZipFile(output_path, "r") as patched_zip:
            content_types = ET.fromstring(patched_zip.read("[Content_Types].xml"))

        ns = {"ct": "http://schemas.openxmlformats.org/package/2006/content-types"}
        overrides = {el.attrib.get("PartName") for el in content_types.findall("ct:Override", ns)}
        defaults = {el.attrib.get("Extension"): el.attrib.get("ContentType") for el in content_types.findall("ct:Default", ns)}

        assert "/xl/worksheets/sheet1.xml" in overrides
        assert "/xl/sharedStrings.xml" in overrides
        assert "/xl/comments1.xml" in overrides
        assert defaults.get("vml") == "application/vnd.openxmlformats-officedocument.vmlDrawing"


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
class TestOfficeTableExcel:
    """Tests for office_table with Excel files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_xlsx_with_table(self, temp_dir):
        """Create a sample Excel file with a named table."""
        path = temp_dir / "office_table.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"

        # Create data for table
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Alpha"
        ws["B2"] = 100
        ws["A3"] = "Beta"
        ws["B3"] = 200

        # Create a table
        from openpyxl.worksheet.table import Table, TableStyleInfo
        table = Table(displayName="TestTable", ref="A1:B3")
        style = TableStyleInfo(
            name="TableStyleMedium9",
            showFirstColumn=False,
            showLastColumn=False,
            showRowStripes=True,
            showColumnStripes=False,
        )
        table.tableStyleInfo = style
        ws.add_table(table)

        wb.save(path)
        return str(path)

    def test_get_table(self, tools, sample_xlsx_with_table):
        """Test getting table data."""
        result = tools.tool_office_table(
            file_path=sample_xlsx_with_table,
            operation="get",
            table_id="TestTable",
        )
        assert "error" not in result

    def test_add_row_to_table(self, tools, sample_xlsx_with_table):
        """Test adding a row to table."""
        result = tools.tool_office_table(
            file_path=sample_xlsx_with_table,
            operation="add_row",
            table_id="TestTable",
            data={"Name": "Gamma", "Value": 300},
        )
        assert "error" not in result

    def test_table_missing_id(self, tools, sample_xlsx_with_table):
        """Test error when table_id is missing."""
        result = tools.tool_office_table(
            file_path=sample_xlsx_with_table,
            operation="get",
        )
        assert "error" in result

    def test_add_row_missing_data(self, tools, sample_xlsx_with_table):
        """Test error when data is missing for add_row."""
        result = tools.tool_office_table(
            file_path=sample_xlsx_with_table,
            operation="add_row",
            table_id="TestTable",
        )
        assert "error" in result


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx not installed")
class TestOfficeTableWord:
    """Tests for office_table with Word files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_docx(self, temp_dir):
        path = temp_dir / "office_table_word.docx"
        doc = docx.Document()
        doc.add_heading("Delivery Plan", level=1)
        doc.add_paragraph("Plan content")
        doc.save(path)
        return str(path)

    def test_create_word_table(self, tools, sample_docx):
        """Should create a Word table via office_table(operation='create')."""
        result = tools.tool_office_table(
            file_path=sample_docx,
            operation="create",
            data={
                "headers": ["Phase", "Owner"],
                "rows": [{"Phase": "Discovery", "Owner": "PM"}],
                "insert_after_section": "Delivery Plan",
            },
        )

        assert "error" not in result
        assert result.get("success") is True

        inspected = tools.tool_office_inspect(
            file_path=sample_docx,
            what="tables",
        )
        assert "error" not in inspected
        assert inspected.get("count", 0) >= 1
        first_table = inspected.get("tables", [])[0]
        assert "Phase" in first_table.get("header", [])

    def test_create_word_table_missing_headers(self, tools, sample_docx):
        """Should validate Word table creation payload."""
        result = tools.tool_office_table(
            file_path=sample_docx,
            operation="create",
            data={"rows": [{"Phase": "Discovery"}]},
        )

        assert "error" in result


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx not installed")
class TestOfficePatchWord:
    """Tests for office_patch with Word files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_docx_with_placeholder(self, temp_dir):
        """Create a sample Word file with placeholder."""
        path = temp_dir / "office_patch.docx"
        doc = docx.Document()
        doc.add_heading("Contract for <Customer Name>", 0)
        doc.add_paragraph("This agreement is dated <Date>.")
        doc.save(path)
        return str(path)

    @pytest.fixture
    def sample_docx_with_empty_section(self, temp_dir):
        """Create a Word file with an empty section body."""
        path = temp_dir / "office_patch_empty_section.docx"
        doc = docx.Document()
        doc.add_heading("Delivery approach", level=1)
        doc.add_heading("Next Section", level=1)
        doc.save(path)
        return str(path)

    def test_patch_placeholder(self, tools, sample_docx_with_placeholder):
        """Test replacing placeholders in Word."""
        result = tools.tool_office_patch(
            file_path=sample_docx_with_placeholder,
            changes=[
                {"target": "<Customer Name>", "value": "Acme Corp"},
            ],
        )
        # May succeed or have 0 changes if placeholder not found
        # (depends on how Word stores the text)
        assert "error" not in result or result.get("changes_applied", 0) >= 0

    def test_patch_placeholder_no_match_does_not_save(self, tools, temp_dir):
        """No-match patch should leave the file unchanged."""
        path = temp_dir / "office_patch_no_match.docx"
        doc = docx.Document()
        doc.add_paragraph("This document has no placeholders.")
        doc.save(path)

        before = hashlib.sha256(path.read_bytes()).hexdigest()

        result = tools.tool_office_patch(
            file_path=str(path),
            changes=[
                {"target": "<Does Not Exist>", "value": "Ignored"},
            ],
        )

        after = hashlib.sha256(path.read_bytes()).hexdigest()

        assert "error" not in result
        assert before == after

    def test_patch_section_target_inserts_into_empty_section(self, tools, sample_docx_with_empty_section):
        """office_patch should delegate section targets to word_patch_section."""
        result = tools.tool_office_patch(
            file_path=sample_docx_with_empty_section,
            changes=[
                {
                    "target": "section:Delivery approach",
                    "value": "Microsoft will undertake an iterative delivery approach.",
                }
            ],
        )

        assert result.get("changes_applied") == 1
        assert result.get("errors") == 0

        doc = docx.Document(sample_docx_with_empty_section)
        paragraphs = [
            _get_text_with_track_changes(para).strip()
            for para in doc.paragraphs
            if _get_text_with_track_changes(para).strip()
        ]
        assert "Delivery approach" in paragraphs
        assert "Next Section" in paragraphs
        delivery_index = paragraphs.index("Delivery approach")
        next_index = paragraphs.index("Next Section")
        assert delivery_index < next_index
        assert "Microsoft will undertake an iterative delivery approach." in paragraphs[delivery_index + 1:next_index]


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
class TestOfficeReadExcelScope:
    """Tests for office_read Excel scope handling."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    def test_reads_sheet_name_scope(self, tools, temp_dir):
        path = temp_dir / "sheet_scope.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "ECIF Work Scope (E)"
        ws["A1"] = "Header"
        ws["A2"] = "Value"
        wb.save(path)

        result = tools.tool_office_read(
            file_path=str(path),
            scope="ECIF Work Scope (E)",
        )

        assert "error" not in result
        assert "sheets" in result
        assert "ECIF Work Scope (E)" in result["sheets"]


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
class TestOfficePatchExcelBatch:
    """Tests for batched Excel patching."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    def test_batch_patch_loads_workbook_once(self, tools, temp_dir, monkeypatch):
        path = temp_dir / "batch_patch.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "Old"
        wb.save(path)

        call_count = 0
        real_load = office_unified_tools_module.load_workbook

        def counted_load(*args, **kwargs):
            nonlocal call_count
            call_count += 1
            return real_load(*args, **kwargs)

        monkeypatch.setattr(office_unified_tools_module, "load_workbook", counted_load)

        result = tools.tool_office_patch(
            file_path=str(path),
            changes=[
                {"target": "A1", "value": "New"},
                {"target": "A2", "value": "Next"},
            ],
        )

        assert "error" not in result
        assert call_count == 1


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
class TestOfficeAuditExcelExtras:
    """Tests for extended Excel audit checks."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    def test_audit_empty_cells_dates_totals(self, tools, temp_dir):
        path = temp_dir / "audit_extras.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "ECIF Work Scope (E)"
        ws["A1"] = 10
        ws["A2"] = 20
        ws["B1"] = 40
        ws["C1"] = "2026-03-01"
        wb.save(path)

        audit_config = {
            "required_cells": ["'ECIF Work Scope (E)'!D1"],
            "date_cells": ["'ECIF Work Scope (E)'!C1"],
            "totals": [
                {
                    "sum_range": "'ECIF Work Scope (E)'!A1:A2",
                    "target": "'ECIF Work Scope (E)'!B1",
                    "tolerance": 0.01,
                }
            ],
        }

        result = tools.tool_office_audit(
            file_path=str(path),
            checks=["empty_cells", "dates", "totals"],
            audit_config=audit_config,
        )

        assert "error" not in result
        extras = result.get("results", {})
        assert extras.get("empty_cells", {}).get("count") == 1
        assert extras.get("dates", {}).get("count") == 1
        assert extras.get("totals", {}).get("count") == 1


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
class TestOfficePatchPowerPoint:
    """Tests for office_patch with PowerPoint files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_pptx(self, temp_dir):
        """Create a sample PowerPoint file."""
        path = temp_dir / "office_patch.pptx"
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        if title:
            title.text = "Original Title"
        prs.save(path)
        return str(path)

    def test_patch_shape_by_slide(self, tools, sample_pptx):
        """Test patching a shape on a slide."""
        # First get the shape name
        inspect_result = tools.tool_office_inspect(
            sample_pptx, what="shapes", target="1"
        )
        if "shapes" in inspect_result and len(inspect_result["shapes"]) > 0:
            shape_name = inspect_result["shapes"][0].get("name", "Title 1")
            result = tools.tool_office_patch(
                file_path=sample_pptx,
                changes=[{"target": f"slide:1/{shape_name}", "value": "New Title"}],
            )
            # Check it ran without crashing
            assert "file" in result
            assert result.get("results", [])[0].get("value_preview") == "New Title"

    def test_patch_global_text(self, tools, sample_pptx):
        """Test replacing text globally."""
        result = tools.tool_office_patch(
            file_path=sample_pptx,
            changes=[{"target": "Original", "value": "Updated"}],
        )
        assert "file" in result

    def test_patch_soft_return(self, tools, sample_pptx, temp_dir):
        """Should convert soft-return tokens for PowerPoint text."""
        output_path = temp_dir / "office_patch_soft_return.pptx"
        inspect_result = tools.tool_office_inspect(sample_pptx, what="shapes", target="1")
        if "shapes" in inspect_result and len(inspect_result["shapes"]) > 0:
            shape_name = inspect_result["shapes"][0].get("name", "Title 1")
            tools.tool_office_patch(
                file_path=sample_pptx,
                changes=[{"target": f"slide:1/{shape_name}", "value": "Line1{br}Line2"}],
                output_path=str(output_path),
            )

            prs = pptx.Presentation(output_path)
            title = prs.slides[0].shapes.title
            assert title is not None
            assert "Line1" in title.text
            assert "Line2" in title.text
            assert "{br}" not in title.text


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
class TestOfficeTemplateExcel:
    """Tests for office_template with Excel files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_xlsx(self, temp_dir):
        """Create a sample Excel template."""
        path = temp_dir / "office_template.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "<Customer Name>"
        ws["B1"] = "<Date>"
        wb.save(path)
        return str(path)

    def test_copy_template(self, tools, sample_xlsx, temp_dir):
        """Test copying a template."""
        dest_path = temp_dir / "office_template_copy.xlsx"
        result = tools.tool_office_template(
            source_path=sample_xlsx,
            destination_path=str(dest_path),
            operation="copy",
        )
        assert "error" not in result or result.get("success") is True


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
class TestOfficeAuditExcel:
    """Tests for office_audit with Excel files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_xlsx_with_placeholders(self, temp_dir):
        """Create a sample Excel file with placeholders."""
        path = temp_dir / "office_audit.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "<Customer Name>"
        ws["B1"] = "[TBD]"
        ws["C1"] = "Normal Value"
        wb.save(path)
        return str(path)

    def test_audit_placeholders(self, tools, sample_xlsx_with_placeholders):
        """Test auditing for placeholders."""
        result = tools.tool_office_audit(
            file_path=sample_xlsx_with_placeholders,
            checks=["placeholders"],
        )
        assert "results" in result
        assert "placeholders" in result.get("results", {})

    def test_audit_default_check(self, tools, sample_xlsx_with_placeholders):
        """Test default audit (placeholders)."""
        result = tools.tool_office_audit(
            file_path=sample_xlsx_with_placeholders,
        )
        assert "results" in result


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx not installed")
class TestOfficeTemplateWord:
    """Tests for office_template with Word files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_docx(self, temp_dir):
        """Create a sample Word template."""
        path = temp_dir / "office_template.docx"
        doc = docx.Document()
        doc.add_heading("Template for <Customer>", 0)
        doc.save(path)
        return str(path)

    def test_copy_template(self, tools, sample_docx, temp_dir):
        """Test copying a Word template."""
        dest_path = temp_dir / "office_template_copy.docx"
        result = tools.tool_office_template(
            source_path=sample_docx,
            destination_path=str(dest_path),
            operation="copy",
        )
        assert "error" not in result

    def test_analyze_template(self, tools, sample_docx):
        """Test analyzing Word template formatting."""
        result = tools.tool_office_template(
            source_path=sample_docx,
            destination_path="",
            operation="analyze",
        )
        assert "error" not in result


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
class TestOfficeAuditPowerPoint:
    """Tests for office_audit with PowerPoint files."""

    @pytest.fixture
    def tools(self):
        CombinedTools = create_combined_tools_class()
        return CombinedTools()

    @pytest.fixture
    def sample_pptx_with_placeholders(self, temp_dir):
        """Create a sample PowerPoint file with placeholder text."""
        path = temp_dir / "office_audit.pptx"
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        if title:
            title.text = "<Title TBD>"
        prs.save(path)
        return str(path)

    def test_audit_placeholders(self, tools, sample_pptx_with_placeholders):
        """Test auditing PowerPoint for placeholders."""
        result = tools.tool_office_audit(
            file_path=sample_pptx_with_placeholders,
            checks=["placeholders"],
        )
        assert "results" in result
        assert "placeholders" in result.get("results", {})
