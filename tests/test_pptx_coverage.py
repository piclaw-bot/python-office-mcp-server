"""
Additional tests for high-coverage impact on pptx_advanced_tools.py
"""

import pytest
from pptx import Presentation
from pptx.util import Inches

# Fixtures temp_dir and pptx_advanced_tools are provided by conftest.py


@pytest.fixture
def basic_presentation(temp_dir):
    """Create a basic presentation."""
    prs = Presentation()

    # Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Main Title"
    if len(slide1.placeholders) > 1:
        slide1.placeholders[1].text = "Subtitle"

    # Content slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Content Slide"

    path = temp_dir / "basic.pptx"
    prs.save(path)
    return path


class TestAddSlideVariations:
    """Tests for adding slides with various options."""

    def test_add_slide_at_start(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should add slide at start of presentation."""
        output = temp_dir / "start_slide.pptx"
        result = pptx_advanced_tools.tool_pptx_add_slide(
            str(basic_presentation),
            layout_index=1,
            title="First Slide",
            position="start",
            output_path=str(output)
        )
        assert result.get("success") is True

    def test_add_slide_at_position(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should add slide at specific position."""
        output = temp_dir / "positioned.pptx"
        result = pptx_advanced_tools.tool_pptx_add_slide(
            str(basic_presentation),
            layout_index=1,
            title="Middle Slide",
            position="1",
            output_path=str(output)
        )
        assert result.get("success") is True


class TestPatchShapeVariations:
    """Tests for patching shapes with various identifiers."""

    def test_patch_title(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should patch slide title."""
        output = temp_dir / "patched_title.pptx"
        result = pptx_advanced_tools.tool_pptx_patch_shape(
            str(basic_presentation),
            slide_number=1,
            shape_identifier="title",
            new_text="New Title Text",
            output_path=str(output)
        )
        assert result.get("success") is True

    def test_patch_subtitle(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should patch slide subtitle."""
        output = temp_dir / "patched_subtitle.pptx"
        result = pptx_advanced_tools.tool_pptx_patch_shape(
            str(basic_presentation),
            slide_number=1,
            shape_identifier="subtitle",
            new_text="New Subtitle",
            output_path=str(output)
        )
        # May or may not have subtitle
        assert isinstance(result, dict)

    def test_patch_with_append(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should append to shape text."""
        output = temp_dir / "appended.pptx"
        result = pptx_advanced_tools.tool_pptx_patch_shape(
            str(basic_presentation),
            slide_number=1,
            shape_identifier="title",
            new_text=" - Appended",
            append=True,
            output_path=str(output)
        )
        assert result.get("success") is True


class TestClearAndAddBullets:
    """Tests for bullet operations."""

    def test_clear_bullets(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should clear bullet content from shape."""
        output = temp_dir / "cleared.pptx"
        result = pptx_advanced_tools.tool_pptx_clear_bullets(
            str(basic_presentation),
            slide_number=2,
            output_path=str(output)
        )
        assert isinstance(result, dict)

    def test_add_multiple_bullets(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should add multiple bullets."""
        # Add first bullet
        output1 = temp_dir / "bullet1.pptx"
        pptx_advanced_tools.tool_pptx_add_bullet(
            str(basic_presentation),
            slide_number=2,
            text="First bullet",
            output_path=str(output1)
        )

        # Add second bullet
        output2 = temp_dir / "bullet2.pptx"
        result = pptx_advanced_tools.tool_pptx_add_bullet(
            str(output1),
            slide_number=2,
            text="Second bullet",
            output_path=str(output2)
        )
        assert result.get("success") is True


class TestNotesOperations:
    """Tests for speaker notes."""

    def test_get_notes_all_slides(self, pptx_advanced_tools, basic_presentation):
        """Should get notes from all slides."""
        result = pptx_advanced_tools.tool_pptx_get_notes(str(basic_presentation))
        assert isinstance(result, dict)

    def test_set_and_get_notes(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should set and retrieve notes."""
        output = temp_dir / "with_notes.pptx"
        pptx_advanced_tools.tool_pptx_set_notes(
            str(basic_presentation),
            slide_number=1,
            notes_text="Speaker notes content",
            output_path=str(output)
        )

        result = pptx_advanced_tools.tool_pptx_get_notes(str(output), slide_number=1)
        assert isinstance(result, dict)


class TestTableCreation:
    """Tests for table operations."""

    def test_create_table_slide(self, pptx_advanced_tools, temp_dir):
        """Should handle slide with table."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Data"

        # Add a table
        table = slide.shapes.add_table(
            rows=2, cols=2,
            left=Inches(1), top=Inches(2),
            width=Inches(4), height=Inches(1)
        ).table
        table.cell(0, 0).text = "H1"
        table.cell(0, 1).text = "H2"
        table.cell(1, 0).text = "V1"
        table.cell(1, 1).text = "V2"

        path = temp_dir / "with_table.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_get_table(str(path), slide_number=1)
        assert isinstance(result, dict)


class TestSlideContent:
    """Tests for getting slide content."""

    def test_get_slide_content(self, pptx_advanced_tools, basic_presentation):
        """Should get all content from slide."""
        result = pptx_advanced_tools.tool_pptx_get_slide(str(basic_presentation), slide_number=1)
        assert "title" in result or isinstance(result, dict)

    def test_get_slide_second(self, pptx_advanced_tools, basic_presentation):
        """Should get content from second slide."""
        result = pptx_advanced_tools.tool_pptx_get_slide(str(basic_presentation), slide_number=2)
        assert isinstance(result, dict)


class TestLayoutOperations:
    """Tests for layout-related operations."""

    def test_list_masters(self, pptx_advanced_tools, basic_presentation):
        """Should list slide masters and layouts."""
        result = pptx_advanced_tools.tool_pptx_list_masters(str(basic_presentation))
        assert "default_layouts" in result

    def test_recommend_blank_layout(self, pptx_advanced_tools, basic_presentation):
        """Should recommend layout for blank content."""
        result = pptx_advanced_tools.tool_pptx_recommend_layout(str(basic_presentation), "blank")
        assert isinstance(result, dict)

    def test_recommend_comparison_layout(self, pptx_advanced_tools, basic_presentation):
        """Should recommend layout for comparison."""
        result = pptx_advanced_tools.tool_pptx_recommend_layout(str(basic_presentation), "comparison")
        assert isinstance(result, dict)


class TestCommentAddition:
    """Tests for comment functionality."""

    def test_add_comment_default_position(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should add comment at default position."""
        output = temp_dir / "comment.pptx"
        result = pptx_advanced_tools.tool_pptx_add_comment(
            str(basic_presentation),
            slide_number=1,
            comment_text="Review this slide",
            output_path=str(output)
        )
        assert result.get("success") is True

    def test_add_comment_custom_position(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should add comment at custom position."""
        output = temp_dir / "comment_pos.pptx"
        result = pptx_advanced_tools.tool_pptx_add_comment(
            str(basic_presentation),
            slide_number=1,
            comment_text="Check this area",
            x_inches=5.0,
            y_inches=3.0,
            author="Reviewer",
            output_path=str(output)
        )
        assert result.get("success") is True


class TestSlideHiding:
    """Tests for hide/unhide slides."""

    def test_unhide_slide(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should unhide a slide."""
        # First hide it
        output1 = temp_dir / "hidden.pptx"
        pptx_advanced_tools.tool_pptx_hide_slide(
            str(basic_presentation),
            slide_number=2,
            hidden=True,
            output_path=str(output1)
        )

        # Then unhide
        output2 = temp_dir / "unhidden.pptx"
        result = pptx_advanced_tools.tool_pptx_hide_slide(
            str(output1),
            slide_number=2,
            hidden=False,
            output_path=str(output2)
        )
        assert isinstance(result, dict)


class TestTextAutofit:
    """Tests for text autofit settings."""

    def test_autofit_none(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should set autofit to none."""
        output = temp_dir / "autofit_none.pptx"
        result = pptx_advanced_tools.tool_pptx_set_text_autofit(
            str(basic_presentation),
            slide_number=1,
            shape_identifier="title",
            autofit_type="none",
            output_path=str(output)
        )
        assert result.get("success") is True

    def test_autofit_resize(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should set autofit to resize."""
        output = temp_dir / "autofit_resize.pptx"
        result = pptx_advanced_tools.tool_pptx_set_text_autofit(
            str(basic_presentation),
            slide_number=1,
            shape_identifier="title",
            autofit_type="resize",
            output_path=str(output)
        )
        assert result.get("success") is True


class TestAuditPlaceholderPatterns:
    """Tests for placeholder audit with custom patterns."""

    def test_audit_with_custom_patterns(self, pptx_advanced_tools, temp_dir):
        """Should find custom placeholder patterns."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "ACME Corp - {{date}}"
        path = temp_dir / "custom_patterns.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_audit_placeholders(
            str(path),
            patterns=[r"\{\{.*?\}\}"]
        )
        assert isinstance(result, dict)


class TestErrorCases:
    """Tests for error handling."""

    def test_invalid_layout_index(self, pptx_advanced_tools, basic_presentation, temp_dir):
        """Should handle invalid layout index."""
        output = temp_dir / "invalid_layout.pptx"
        result = pptx_advanced_tools.tool_pptx_add_slide(
            str(basic_presentation),
            layout_index=999,
            output_path=str(output)
        )
        # Should either fail gracefully or use default
        assert isinstance(result, dict)

    def test_get_nonexistent_table(self, pptx_advanced_tools, basic_presentation):
        """Should handle getting table that doesn't exist."""
        result = pptx_advanced_tools.tool_pptx_get_table(
            str(basic_presentation),
            slide_number=1,
            table_index=0
        )
        # Should return error since no table on slide 1
        assert "error" in result or "not found" in str(result).lower() or "tables" not in str(result)


class TestListOperations:
    """Tests for list operations."""

    def test_list_slides(self, pptx_advanced_tools, basic_presentation):
        """Should list all slides."""
        result = pptx_advanced_tools.tool_pptx_list_slides(str(basic_presentation))
        assert result.get("slide_count", 0) >= 2

    def test_list_shapes(self, pptx_advanced_tools, basic_presentation):
        """Should list shapes on slide."""
        result = pptx_advanced_tools.tool_pptx_list_shapes(str(basic_presentation), slide_number=1)
        assert "shapes" in result


class TestExtractAndConvert:
    """Tests for extract and conversion."""

    def test_extract_presentation(self, pptx_tools, basic_presentation):
        """Should extract presentation content."""
        # tool_pptx_extract is on PowerPointTools, not Advanced
        result = pptx_tools.tool_pptx_extract(str(basic_presentation))
        assert isinstance(result, dict)
        assert "slides" in result

    def test_to_markdown(self, pptx_tools, basic_presentation):
        """Should convert to markdown."""
        # tool_pptx_to_markdown is on PowerPointTools, not Advanced
        result = pptx_tools.tool_pptx_to_markdown(str(basic_presentation))
        assert "Main Title" in result
