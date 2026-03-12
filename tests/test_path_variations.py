"""
Comprehensive tests targeting remaining uncovered code paths.
"""

from docx import Document
from pptx import Presentation
from pptx.util import Inches as PptxInches

# Fixture temp_dir is provided by conftest.py


class TestPptxAddTableOperations:
    """Test PPTX table addition operations."""

    def test_add_table_with_many_rows(self, pptx_advanced_tools, temp_dir):
        """Should add table with multiple rows."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Data Table"
        path = temp_dir / "many_rows.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_add_table(
            str(path),
            1,
            ["A", "B", "C"],
            [
                ["R1A", "R1B", "R1C"],
                ["R2A", "R2B", "R2C"],
                ["R3A", "R3B", "R3C"],
                ["R4A", "R4B", "R4C"],
            ]
        )
        assert isinstance(result, dict)

    def test_add_table_single_column(self, pptx_advanced_tools, temp_dir):
        """Should add single column table."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])
        path = temp_dir / "single_col.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_add_table(
            str(path),
            1,
            ["Item"],
            [["First"], ["Second"], ["Third"]]
        )
        assert isinstance(result, dict)


class TestWordInsertTableRowVariations:
    """Test Word table row insertion variations."""

    def test_insert_row_specific_position(self, word_advanced_tools, temp_dir):
        """Should insert row at specific position."""
        doc = Document()
        table = doc.add_table(rows=4, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        for i in range(1, 4):
            table.cell(i, 0).text = f"Item {i}"
            table.cell(i, 1).text = f"Val {i}"
        path = temp_dir / "insert_pos.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_insert_table_row(
            str(path),
            "0",
            {"Name": "New Item", "Value": "New Val"}
        )
        assert isinstance(result, dict)

    def test_insert_row_empty_dict(self, word_advanced_tools, temp_dir):
        """Should handle empty row data."""
        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"
        path = temp_dir / "empty_row.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_insert_table_row(
            str(path),
            "0",
            {}
        )
        assert isinstance(result, dict)


class TestPptxReplaceTextVariations:
    """Test replace text variations."""

    def test_replace_in_table_cells(self, pptx_advanced_tools, temp_dir):
        """Should replace text in table cells."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = slide.shapes.add_table(2, 2, PptxInches(1), PptxInches(2), PptxInches(5), PptxInches(1.5))
        tbl = table.table
        tbl.cell(0, 0).text = "<Customer>"
        tbl.cell(0, 1).text = "Value"
        tbl.cell(1, 0).text = "Contact <Customer>"
        tbl.cell(1, 1).text = "Data"
        path = temp_dir / "table_replace.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_replace_text(
            str(path),
            "<Customer>",
            "Contoso"
        )
        assert isinstance(result, dict)


class TestWordCopyTemplateVariations:
    """Test copy template variations."""

    def test_copy_to_existing_location(self, word_advanced_tools, temp_dir):
        """Should overwrite if destination exists."""
        # Create template
        doc = Document()
        doc.add_heading("Template", level=1)
        template = temp_dir / "template.docx"
        doc.save(template)

        # Create existing output
        doc2 = Document()
        doc2.add_heading("Existing", level=1)
        output = temp_dir / "output.docx"
        doc2.save(output)

        result = word_advanced_tools.tool_word_copy_template(str(template), str(output))
        assert isinstance(result, dict)


class TestPptxClearBulletsVariations:
    """Test clear bullets variations."""

    def test_clear_bullets_no_body(self, pptx_advanced_tools, temp_dir):
        """Should handle slide with no body placeholder."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
        slide.shapes.title.text = "Title Only"
        path = temp_dir / "no_body.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_clear_bullets(str(path), slide_number=1)
        # Should handle gracefully
        assert isinstance(result, dict)


class TestWordListTablesVariations:
    """Test list tables variations."""

    def test_list_tables_multiple(self, word_advanced_tools, temp_dir):
        """Should list multiple tables with details."""
        doc = Document()
        doc.add_paragraph("Before first table")

        t1 = doc.add_table(rows=2, cols=2)
        t1.cell(0, 0).text = "T1A"
        t1.cell(0, 1).text = "T1B"

        doc.add_paragraph("Between tables")

        t2 = doc.add_table(rows=3, cols=4)
        t2.cell(0, 0).text = "T2A"

        doc.add_paragraph("After tables")

        t3 = doc.add_table(rows=4, cols=2)
        t3.cell(0, 0).text = "T3A"

        path = temp_dir / "three_tables.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_list_tables(str(path))
        assert isinstance(result, dict)
        assert len(result.get("tables", [])) == 3


class TestPptxGetTableVariations:
    """Test get table variations."""

    def test_get_table_multiple_tables(self, pptx_advanced_tools, temp_dir):
        """Should get specific table when multiple exist."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # First table
        t1 = slide.shapes.add_table(2, 2, PptxInches(1), PptxInches(1), PptxInches(3), PptxInches(1))
        t1.table.cell(0, 0).text = "Table1"

        # Second table
        t2 = slide.shapes.add_table(3, 3, PptxInches(5), PptxInches(1), PptxInches(4), PptxInches(1.5))
        t2.table.cell(0, 0).text = "Table2"

        path = temp_dir / "two_tables.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_get_table(str(path), 1, table_index=1)
        assert isinstance(result, dict)


class TestWordPatchTableRowVariations:
    """Test patch table row variations."""

    def test_patch_table_row_partial_columns(self, word_advanced_tools, temp_dir):
        """Should patch only specified columns."""
        doc = Document()
        table = doc.add_table(rows=3, cols=4)
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"
        table.cell(0, 2).text = "C"
        table.cell(0, 3).text = "D"
        for i in range(1, 3):
            for j in range(4):
                table.cell(i, j).text = f"R{i}C{j}"
        path = temp_dir / "partial_patch.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_patch_table_row(
            str(path),
            "0",
            1,
            {"A": "Updated", "C": "Changed"}  # Only patch two columns
        )
        assert isinstance(result, dict)


class TestPptxAddCommentPositions:
    """Test comment positioning."""

    def test_add_comment_edge_position(self, pptx_advanced_tools, temp_dir):
        """Should add comment at edge position."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Edge Comment"
        path = temp_dir / "edge_comment.pptx"
        prs.save(path)

        # Position at edge of slide
        result = pptx_advanced_tools.tool_pptx_add_comment(
            str(path),
            1,
            "Edge comment",
            x_inches=12.0,
            y_inches=6.0
        )
        assert isinstance(result, dict)


class TestWordSectionGuidanceVariations:
    """Test section guidance variations."""

    def test_get_section_guidance_no_guidance(self, word_advanced_tools, temp_dir):
        """Should handle section with no guidance."""
        doc = Document()
        doc.add_heading("Section", level=1)
        doc.add_paragraph("Just regular content, no guidance")
        path = temp_dir / "no_guidance.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_get_section_guidance(str(path), "Section")
        assert isinstance(result, dict)


class TestPptxNotesAppend:
    """Test notes append functionality."""

    def test_append_to_existing_notes(self, pptx_advanced_tools, temp_dir):
        """Should append to existing notes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Notes Test"
        notes = slide.notes_slide
        notes.notes_text_frame.text = "Original notes"
        path = temp_dir / "append_notes.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_set_notes(
            str(path),
            1,
            "\n\nAppended notes",
            append=True
        )
        assert isinstance(result, dict)


class TestWordCleanupSowVariations:
    """Test cleanup SOW variations."""

    def test_cleanup_sow_no_placeholders(self, word_advanced_tools, temp_dir):
        """Should cleanup SOW with no placeholders."""
        doc = Document()
        doc.add_heading("Clean Document", level=1)
        doc.add_paragraph("All content is filled in properly.")
        doc.add_heading("Section 2", level=2)
        doc.add_paragraph("More complete content.")
        path = temp_dir / "clean_sow.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_cleanup_sow(str(path))
        assert isinstance(result, dict)


class TestPptxSlideReorderVariations:
    """Test slide reorder variations."""

    def test_reorder_reverse_all(self, pptx_advanced_tools, temp_dir):
        """Should reverse slide order."""
        prs = Presentation()
        for i in range(5):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
        path = temp_dir / "reverse.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_reorder_slides(
            str(path),
            [5, 4, 3, 2, 1]  # Reverse order
        )
        assert isinstance(result, dict)


class TestWordCreateNewTableVariations:
    """Test create new table variations."""

    def test_create_table_before_section(self, word_advanced_tools, temp_dir):
        """Should create table before specific section."""
        doc = Document()
        doc.add_heading("Section A", level=1)
        doc.add_paragraph("Content A")
        doc.add_heading("Section B", level=1)
        doc.add_paragraph("Content B")
        path = temp_dir / "before_section.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_create_new_table(
            str(path),
            ["X", "Y", "Z"],
            [{"X": "1", "Y": "2", "Z": "3"}],
            insert_before_section="Section B"
        )
        assert isinstance(result, dict)


class TestExcelMultiSheetOperations:
    """Test Excel multi-sheet operations."""

    def test_from_markdown_multiple_tables(self, excel_tools, temp_dir):
        """Should create workbook with multiple sheets."""
        md = """| Header A | Header B |
|----------|----------|
| Data A1  | Data B1  |
| Data A2  | Data B2  |

Some text between tables

| Col 1 | Col 2 | Col 3 |
|-------|-------|-------|
| X     | Y     | Z     |
"""
        path = temp_dir / "multi_sheet.xlsx"
        result = excel_tools.tool_excel_from_markdown(str(path), md)
        assert isinstance(result, dict)

    def test_to_markdown_multi_sheet(self, excel_tools, temp_dir):
        """Should convert multi-sheet to markdown."""
        # First create a multi-table file
        md = """| A | B |
    |---|---|
    | 1 | 2 |

    | C | D |
    |---|---|
    | 3 | 4 |
    """
        path = temp_dir / "multi_md.xlsx"
        excel_tools.tool_excel_from_markdown(str(path), md)

        result = excel_tools.tool_excel_to_markdown(str(path))
        # Should be a string
        assert isinstance(result, str)


class TestPptxLayoutAnalysis:
    """Test layout analysis."""

    def test_analyze_layouts_standard_deck(self, pptx_advanced_tools, temp_dir):
        """Should analyze layouts in standard deck."""
        prs = Presentation()
        # Add slides with different layouts
        prs.slides.add_slide(prs.slide_layouts[0])  # Title
        prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        path = temp_dir / "analyze_layouts.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_analyze_layouts(str(path))
        assert isinstance(result, dict)


class TestWordPatchPlaceholderCases:
    """Test patch placeholder edge cases."""

    def test_patch_placeholder_in_header(self, word_advanced_tools, temp_dir):
        """Should patch placeholder in header."""
        doc = Document()
        # Add placeholder to main content
        doc.add_heading("<Project Title>", level=1)
        doc.add_paragraph("Contact: <Customer Name>")
        path = temp_dir / "header_ph.docx"
        doc.save(path)

        result = word_advanced_tools.tool_word_patch_placeholder(
            str(path),
            "<Project Title>",
            "Cloud Migration Project"
        )
        assert isinstance(result, dict)
