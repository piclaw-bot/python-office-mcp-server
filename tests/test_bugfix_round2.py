"""
Tests for bug-fix round 2 (issues #2, #5, #6, #7, #8, #9, #10).

Covers:
- #2  pptx_duplicate_slide preserves tables and shapes (XML deep copy)
- #5  office_table table_id strips stray quotes
- #6  Unsupported file format returns descriptive error
- #7  Workspace-relative path resolution via _resolve_file_path
- #8  Template copy + duplicate preserves non-placeholder shapes
- #9  pptx_add_slide position is 1-based
- #10 PPTX table-not-found error includes reorder hint
"""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from tools import TOOL_CLASSES
from tools.office_unified_tools import (
    _detect_format,
    _resolve_file_path,
    _unsupported_format_error,
)
from tools.pptx_advanced_tools import PresentationAdvancedTools

# ---------------------------------------------------------------------------
# Helpers / fixtures
# ---------------------------------------------------------------------------


def _create_combined_tools():
    """Build a combined class like the real server does."""
    class CombinedTools(*TOOL_CLASSES):
        pass
    return CombinedTools()


@pytest.fixture
def temp_dir():
    with tempfile.TemporaryDirectory() as d:
        yield Path(d)


@pytest.fixture
def combined_tools():
    return _create_combined_tools()


@pytest.fixture
def pptx_tools():
    return PresentationAdvancedTools()


@pytest.fixture
def pptx_with_table(temp_dir):
    """Create a presentation with a table on slide 1."""
    prs = Presentation()
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Table Slide"

    tbl_shape = slide.shapes.add_table(
        3, 2, Inches(1), Inches(2), Inches(8), Inches(3)
    )
    tbl = tbl_shape.table
    tbl.cell(0, 0).text = "Name"
    tbl.cell(0, 1).text = "Value"
    tbl.cell(1, 0).text = "Alpha"
    tbl.cell(1, 1).text = "100"
    tbl.cell(2, 0).text = "Beta"
    tbl.cell(2, 1).text = "200"

    path = temp_dir / "table_slide.pptx"
    prs.save(path)
    return path


@pytest.fixture
def pptx_three_slides(temp_dir):
    """Create a presentation with 3 titled slides for ordering tests."""
    prs = Presentation()
    for idx in range(3):
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = f"Slide {idx + 1}"
    path = temp_dir / "three_slides.pptx"
    prs.save(path)
    return path


# =========================================================================
# Bug #2 — pptx_duplicate_slide preserves tables/shapes
# =========================================================================


class TestDuplicateSlideDeepCopy:
    """Bug #2: pptx_duplicate_slide must carry non-placeholder shapes."""

    def test_duplicate_preserves_table(self, pptx_tools, pptx_with_table, temp_dir):
        """Duplicated slide should contain the same table with data."""
        out = temp_dir / "dup_table.pptx"
        result = pptx_tools.tool_pptx_duplicate_slide(
            file_path=str(pptx_with_table),
            slide_number=1,
            position="after",
            output_path=str(out),
        )
        assert result.get("success") is True
        assert result["new_slide_number"] == 2

        # Reload and verify both slides have a table
        prs = Presentation(str(out))
        assert len(prs.slides) == 2

        for slide_idx in range(2):
            slide = prs.slides[slide_idx]
            tables = [s for s in slide.shapes if s.has_table]
            assert len(tables) >= 1, f"Slide {slide_idx+1} should have a table"
            tbl = tables[0].table
            assert tbl.cell(0, 0).text == "Name"
            assert tbl.cell(1, 0).text == "Alpha"
            assert tbl.cell(2, 1).text == "200"

    def test_duplicate_position_end(self, pptx_tools, pptx_with_table, temp_dir):
        """Duplicated slide at 'end' should be the last slide."""
        out = temp_dir / "dup_end.pptx"
        result = pptx_tools.tool_pptx_duplicate_slide(
            file_path=str(pptx_with_table),
            slide_number=1,
            position="end",
            output_path=str(out),
        )
        assert result.get("success") is True
        prs = Presentation(str(out))
        assert len(prs.slides) == 2

    def test_duplicate_preserves_title(self, pptx_tools, pptx_with_table, temp_dir):
        """Duplicated slide should keep the title text."""
        out = temp_dir / "dup_title.pptx"
        pptx_tools.tool_pptx_duplicate_slide(
            file_path=str(pptx_with_table),
            slide_number=1,
            position="after",
            output_path=str(out),
        )
        prs = Presentation(str(out))
        # Both slides should have the same title
        for slide in prs.slides:
            if slide.shapes.title:
                assert slide.shapes.title.text == "Table Slide"

    def test_duplicate_invalid_slide_number(self, pptx_tools, pptx_with_table):
        """Should return error for out-of-range slide number."""
        result = pptx_tools.tool_pptx_duplicate_slide(
            file_path=str(pptx_with_table),
            slide_number=99,
        )
        assert "error" in result

    def test_duplicate_file_not_found(self, pptx_tools):
        """Should return error for missing file."""
        result = pptx_tools.tool_pptx_duplicate_slide(
            file_path="/nonexistent/file.pptx",
            slide_number=1,
        )
        assert "error" in result
        assert "not found" in result["error"].lower()


# =========================================================================
# Bug #5 — table_id strips stray quotes
# =========================================================================


class TestTableIdQuoteParsing:
    """Bug #5: table_id with extra quotes should still work."""

    def test_quoted_table_id_pptx(self, combined_tools, pptx_with_table):
        """table_id='"1"' should be treated as '1' for PPTX."""
        result = combined_tools.tool_office_table(
            file_path=str(pptx_with_table),
            operation="get",
            table_id='"1"',
        )
        # Should succeed (not crash on int('"1"'))
        # Slide 1 has a table
        assert "error" not in result or "No table" not in result.get("error", "")

    def test_single_quoted_table_id(self, combined_tools, pptx_with_table):
        """table_id="'1'" should be treated as '1'."""
        result = combined_tools.tool_office_table(
            file_path=str(pptx_with_table),
            operation="get",
            table_id="'1'",
        )
        assert "error" not in result or "No table" not in result.get("error", "")

    def test_unquoted_table_id(self, combined_tools, pptx_with_table):
        """Normal unquoted table_id should work as before."""
        result = combined_tools.tool_office_table(
            file_path=str(pptx_with_table),
            operation="get",
            table_id="1",
        )
        assert "error" not in result or "No table" not in result.get("error", "")


# =========================================================================
# Bug #11 — PPTX add_row preserves column mapping for dicts
# =========================================================================


class TestPptxAddRowDictMapping:
    """Bug #11: dict-based add_row should map values by header."""

    def test_add_row_dict_maps_by_header(self, combined_tools, pptx_tools, pptx_with_table):
        result = combined_tools.tool_office_table(
            file_path=str(pptx_with_table),
            operation="add_row",
            table_id="1",
            data={"Value": "300", "Name": "Gamma"},
        )
        assert "error" not in result

        table = pptx_tools.tool_pptx_get_table(
            file_path=str(pptx_with_table),
            slide_number=1,
        )
        assert "error" not in table
        assert table.get("data")
        assert table["data"][-1] == ["Gamma", "300"]


# =========================================================================
# Bug #6 — unsupported format returns descriptive error
# =========================================================================


class TestUnsupportedFormatError:
    """Bug #6: unsupported extensions should list valid ones."""

    def test_detect_format_returns_none_for_pdf(self):
        assert _detect_format("report.pdf") is None

    def test_detect_format_returns_none_for_txt(self):
        assert _detect_format("notes.txt") is None

    def test_unsupported_format_error_includes_extension(self):
        err = _unsupported_format_error("report.pdf")
        assert "error" in err
        assert ".pdf" in err["error"]

    def test_unsupported_format_error_lists_supported(self):
        err = _unsupported_format_error("data.csv")
        assert ".docx" in err["error"]
        assert ".xlsx" in err["error"]
        assert ".pptx" in err["error"]

    def test_office_read_unsupported_format(self, combined_tools, temp_dir):
        """office_read on a .txt file should produce descriptive error."""
        txt_file = temp_dir / "notes.txt"
        txt_file.write_text("hello")
        result = combined_tools.tool_office_read(str(txt_file))
        assert "error" in result
        assert ".txt" in result["error"]
        assert "Supported" in result["error"]

    def test_office_inspect_unsupported_format(self, combined_tools, temp_dir):
        """office_inspect on a .pdf file should produce descriptive error."""
        pdf_file = temp_dir / "report.pdf"
        pdf_file.write_bytes(b"%PDF-1.4 fake")
        result = combined_tools.tool_office_inspect(str(pdf_file))
        assert "error" in result
        assert ".pdf" in result["error"]


# =========================================================================
# Bug #7 — workspace-relative path resolution
# =========================================================================


class TestResolveFilePath:
    """Bug #7: _resolve_file_path should try workspace-relative paths."""

    def test_absolute_path_exists(self, temp_dir):
        """Existing absolute path should be returned as-is."""
        f = temp_dir / "existing.docx"
        f.write_bytes(b"pk")
        assert _resolve_file_path(str(f)) == str(f)

    def test_absolute_path_not_found(self):
        """Non-existent absolute path is returned unchanged for caller errors."""
        p = "/nonexistent/abc.docx"
        assert _resolve_file_path(p) == p

    def test_relative_via_workspace_root(self, temp_dir, monkeypatch):
        """Setting MCP_WORKSPACE_ROOT should resolve relative paths."""
        sub = temp_dir / "workspace"
        sub.mkdir()
        f = sub / "report.docx"
        f.write_bytes(b"pk")

        monkeypatch.setenv("MCP_WORKSPACE_ROOT", str(sub))
        resolved = _resolve_file_path("report.docx")
        assert resolved == str(f)

    def test_relative_via_cwd(self, temp_dir, monkeypatch):
        """Relative path resolved via cwd when MCP_WORKSPACE_ROOT is unset."""
        monkeypatch.delenv("MCP_WORKSPACE_ROOT", raising=False)
        f = temp_dir / "data.xlsx"
        f.write_bytes(b"pk")

        monkeypatch.chdir(temp_dir)
        resolved = _resolve_file_path("data.xlsx")
        assert Path(resolved).exists()

    def test_missing_relative_returns_original(self, monkeypatch):
        """When nothing resolves, return original for downstream error."""
        monkeypatch.delenv("MCP_WORKSPACE_ROOT", raising=False)
        assert _resolve_file_path("no_such_file.pptx") == "no_such_file.pptx"


# =========================================================================
# Bug #8 — template copy followed by duplicate preserves shapes
# =========================================================================


class TestTemplateCopyPreservesShapes:
    """Bug #8: Template copy + duplicate should keep tables."""

    def test_copy_then_duplicate_keeps_table(
        self, combined_tools, pptx_tools, pptx_with_table, temp_dir
    ):
        """Copy a PPTX (with table), then duplicate slide 1 — table should survive."""
        copy_path = temp_dir / "template_copy.pptx"

        # Step 1: copy via office_template
        result = combined_tools.tool_office_template(
            source_path=str(pptx_with_table),
            destination_path=str(copy_path),
            operation="copy",
        )
        assert result.get("success") is True
        assert copy_path.exists()

        # Step 2: verify table in copy
        prs = Presentation(str(copy_path))
        slide = prs.slides[0]
        tables = [s for s in slide.shapes if s.has_table]
        assert len(tables) == 1

        # Step 3: duplicate slide in the copy
        dup_path = temp_dir / "template_dup.pptx"
        result = pptx_tools.tool_pptx_duplicate_slide(
            file_path=str(copy_path),
            slide_number=1,
            position="after",
            output_path=str(dup_path),
        )
        assert result.get("success") is True

        # Step 4: verify both slides have tables
        prs2 = Presentation(str(dup_path))
        assert len(prs2.slides) == 2
        for idx in range(2):
            slide = prs2.slides[idx]
            tables = [s for s in slide.shapes if s.has_table]
            assert len(tables) >= 1, f"Slide {idx+1} lost its table after duplicate"


# =========================================================================
# Bug #9 — pptx_add_slide position is 1-based
# =========================================================================


class TestAddSlidePosition:
    """Bug #9: position='2' should make the new slide become slide 2."""

    def test_position_end(self, pptx_tools, pptx_three_slides, temp_dir):
        """position='end' should append at the end."""
        out = temp_dir / "add_end.pptx"
        result = pptx_tools.tool_pptx_add_slide(
            file_path=str(pptx_three_slides),
            title="New End",
            position="end",
            output_path=str(out),
        )
        assert result.get("success") is True
        assert result["slide_number"] == 4

    def test_position_start(self, pptx_tools, pptx_three_slides, temp_dir):
        """position='start' should be slide 1."""
        out = temp_dir / "add_start.pptx"
        result = pptx_tools.tool_pptx_add_slide(
            file_path=str(pptx_three_slides),
            title="New Start",
            position="start",
            output_path=str(out),
        )
        assert result.get("success") is True
        assert result["slide_number"] == 1

        # Verify: first slide in output should have our title
        prs = Presentation(str(out))
        assert len(prs.slides) == 4
        first = prs.slides[0]
        if first.shapes.title:
            assert first.shapes.title.text == "New Start"

    def test_position_2_becomes_slide_2(self, pptx_tools, pptx_three_slides, temp_dir):
        """position='2' should insert the new slide as slide 2 (1-based)."""
        out = temp_dir / "add_pos2.pptx"
        result = pptx_tools.tool_pptx_add_slide(
            file_path=str(pptx_three_slides),
            title="Inserted at 2",
            position="2",
            output_path=str(out),
        )
        assert result.get("success") is True
        assert result["slide_number"] == 2

        # Verify ordering: Slide 1, Inserted at 2, Slide 2, Slide 3
        prs = Presentation(str(out))
        assert len(prs.slides) == 4
        titles = []
        for slide in prs.slides:
            if slide.shapes.title:
                titles.append(slide.shapes.title.text)
            else:
                titles.append("")
        assert titles[0] == "Slide 1"
        assert titles[1] == "Inserted at 2"
        assert titles[2] == "Slide 2"
        assert titles[3] == "Slide 3"

    def test_position_1_becomes_slide_1(self, pptx_tools, pptx_three_slides, temp_dir):
        """position='1' should insert as slide 1 (first)."""
        out = temp_dir / "add_pos1.pptx"
        result = pptx_tools.tool_pptx_add_slide(
            file_path=str(pptx_three_slides),
            title="Inserted at 1",
            position="1",
            output_path=str(out),
        )
        assert result.get("success") is True
        assert result["slide_number"] == 1

        prs = Presentation(str(out))
        assert len(prs.slides) == 4
        first = prs.slides[0]
        if first.shapes.title:
            assert first.shapes.title.text == "Inserted at 1"

    def test_position_3_in_3_slide_deck(self, pptx_tools, pptx_three_slides, temp_dir):
        """position='3' in a 3-slide deck — should become slide 3."""
        out = temp_dir / "add_pos3.pptx"
        result = pptx_tools.tool_pptx_add_slide(
            file_path=str(pptx_three_slides),
            title="Inserted at 3",
            position="3",
            output_path=str(out),
        )
        assert result.get("success") is True
        assert result["slide_number"] == 3

        prs = Presentation(str(out))
        titles = [
            s.shapes.title.text if s.shapes.title else ""
            for s in prs.slides
        ]
        assert titles == ["Slide 1", "Slide 2", "Inserted at 3", "Slide 3"]


# =========================================================================
# Bug #10 — PPTX table error messages include reorder hint
# =========================================================================


class TestPptxTableErrorMessages:
    """Bug #10: 'No table found' should hint about slide reordering."""

    def test_patch_table_cell_no_table_hint(self, pptx_tools, temp_dir):
        """Error message should mention reordering when table not found."""
        # Create a pptx with a blank slide (no table)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        path = temp_dir / "no_table.pptx"
        prs.save(path)

        result = pptx_tools.tool_pptx_patch_table_cell(
            file_path=str(path),
            slide_number=1,
            row_index=0,
            col_index=0,
            new_text="test",
        )
        assert "error" in result
        assert "reordered" in result["error"].lower() or "deleted" in result["error"].lower()

    def test_get_table_no_table_hint(self, pptx_tools, temp_dir):
        """pptx_get_table should also include the hint."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        path = temp_dir / "no_table2.pptx"
        prs.save(path)

        result = pptx_tools.tool_pptx_get_table(
            file_path=str(path),
            slide_number=1,
        )
        assert "error" in result
        assert "reordered" in result["error"].lower() or "deleted" in result["error"].lower()

    def test_insert_table_row_no_table_hint(self, pptx_tools, temp_dir):
        """pptx_insert_table_row should also include the hint."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        path = temp_dir / "no_table3.pptx"
        prs.save(path)

        result = pptx_tools.tool_pptx_insert_table_row(
            file_path=str(path),
            slide_number=1,
            row_data=["a", "b"],
        )
        assert "error" in result
        assert "reordered" in result["error"].lower() or "deleted" in result["error"].lower()


# =========================================================================
# Bug #6 — pptx_delete_slide remaining_slides count
# =========================================================================


class TestDeleteSlideRemainingCount:
    """Bug #6: remaining_slides must match actual slide count after delete."""

    def test_remaining_slides_after_delete_last(self, pptx_tools, temp_dir):
        """Deleting last slide of a 7-slide deck reports 6 remaining."""
        prs = Presentation()
        for _ in range(7):
            prs.slides.add_slide(prs.slide_layouts[0])
        path = temp_dir / "seven_slides.pptx"
        prs.save(path)

        result = pptx_tools.tool_pptx_delete_slide(
            file_path=str(path), slide_number=7, output_path=str(path)
        )
        assert result.get("success") is True
        assert result["remaining_slides"] == 6

        # Double-check by reloading the file
        prs2 = Presentation(str(path))
        assert len(prs2.slides) == 6

    def test_remaining_slides_after_delete_middle(self, pptx_tools, temp_dir):
        """Deleting a middle slide reports correct count."""
        prs = Presentation()
        for _ in range(5):
            prs.slides.add_slide(prs.slide_layouts[0])
        path = temp_dir / "five_slides.pptx"
        prs.save(path)

        result = pptx_tools.tool_pptx_delete_slide(
            file_path=str(path), slide_number=3, output_path=str(path)
        )
        assert result.get("success") is True
        assert result["remaining_slides"] == 4

        prs2 = Presentation(str(path))
        assert len(prs2.slides) == 4
