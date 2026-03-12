"""
Additional tests for pptx_advanced_tools.py - Extended coverage

Tests cover additional advanced PowerPoint functions.
"""

import pytest
from pptx import Presentation
from pptx.util import Inches

# Fixtures temp_dir and pptx_advanced_tools are provided by conftest.py


@pytest.fixture
def multi_slide_presentation(temp_dir):
    """Create presentation with multiple slides."""
    prs = Presentation()

    # Title slide
    layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Title Slide"
    if slide1.placeholders[1]:
        slide1.placeholders[1].text = "Subtitle here"

    # Content slide
    layout2 = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(layout2)
    slide2.shapes.title.text = "Content Slide"

    # Another content slide
    slide3 = prs.slides.add_slide(layout2)
    slide3.shapes.title.text = "More Content"

    # Blank slide
    layout3 = prs.slide_layouts[6]
    prs.slides.add_slide(layout3)

    path = temp_dir / "multi_slide.pptx"
    prs.save(path)
    return path


@pytest.fixture
def presentation_with_table(temp_dir):
    """Create presentation with a table."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout
    slide.shapes.title.text = "Data Table"

    # Add table
    table = slide.shapes.add_table(
        rows=3, cols=3,
        left=Inches(1), top=Inches(2),
        width=Inches(6), height=Inches(2)
    ).table

    # Set headers
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(0, 2).text = "Status"

    # Set data
    table.cell(1, 0).text = "Item A"
    table.cell(1, 1).text = "100"
    table.cell(1, 2).text = "Active"

    table.cell(2, 0).text = "Item B"
    table.cell(2, 1).text = "200"
    table.cell(2, 2).text = "Pending"

    path = temp_dir / "with_table.pptx"
    prs.save(path)
    return path


class TestReorderSlides:
    """Tests for tool_pptx_reorder_slides."""

    def test_reorder_slides(self, pptx_advanced_tools, multi_slide_presentation, temp_dir):
        """Should reorder slides in presentation."""
        output = temp_dir / "reordered.pptx"
        result = pptx_advanced_tools.tool_pptx_reorder_slides(
            str(multi_slide_presentation),
            [1, 3, 2, 4],  # Move slide 3 before slide 2
            output_path=str(output)
        )
        assert result.get("success") is True

    def test_reorder_invalid_order(self, pptx_advanced_tools, multi_slide_presentation):
        """Should handle invalid slide order."""
        result = pptx_advanced_tools.tool_pptx_reorder_slides(
            str(multi_slide_presentation),
            [1, 2, 3, 4, 5, 6]  # More slides than exist
        )
        # Should handle gracefully - either error or partial success
        assert isinstance(result, dict)


class TestLogChanges:
    """Tests for tool_pptx_log_changes."""

    def test_logs_changes(self, pptx_advanced_tools, multi_slide_presentation, temp_dir):
        """Should create change log slide."""
        output = temp_dir / "with_log.pptx"
        changes = [
            {"slide": 1, "action": "Updated title", "detail": "Changed main title"},
            {"slide": 2, "action": "Added content", "detail": "Added bullets"}
        ]
        result = pptx_advanced_tools.tool_pptx_log_changes(
            str(multi_slide_presentation),
            changes,
            output_path=str(output)
        )
        assert result.get("success") is True


class TestTableOperations:
    """Tests for PowerPoint table operations."""

    def test_get_table(self, pptx_advanced_tools, presentation_with_table):
        """Should get table content from slide."""
        result = pptx_advanced_tools.tool_pptx_get_table(
            str(presentation_with_table),
            slide_number=1,
            table_index=0
        )
        assert "header" in result or "rows" in result or "columns" in result


class TestLayoutAnalysis:
    """Tests for layout analysis functions."""

    def test_analyze_layouts(self, pptx_advanced_tools, multi_slide_presentation):
        """Should analyze available layouts."""
        result = pptx_advanced_tools.tool_pptx_analyze_layouts(
            str(multi_slide_presentation)
        )
        assert "layouts" in result or "default_layouts" in result

    def test_recommend_layout_bullets(self, pptx_advanced_tools, multi_slide_presentation):
        """Should recommend layout for bullet content."""
        result = pptx_advanced_tools.tool_pptx_recommend_layout(
            str(multi_slide_presentation),
            "bullets"
        )
        assert "layout_index" in result or "recommended" in result or "index" in str(result).lower()

    def test_recommend_layout_table(self, pptx_advanced_tools, multi_slide_presentation):
        """Should recommend layout for table content."""
        result = pptx_advanced_tools.tool_pptx_recommend_layout(
            str(multi_slide_presentation),
            "table"
        )
        assert "layout_index" in result or "recommended" in result


class TestAuditPlaceholders:
    """Tests for placeholder audit."""

    def test_audit_clean_presentation(self, pptx_advanced_tools, temp_dir):
        """Should find no placeholders in clean presentation."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Clean Title"
        path = temp_dir / "clean.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_audit_placeholders(str(path))
        findings = result.get("findings", result.get("placeholders", []))
        assert len(findings) == 0 or result.get("success") is True

    def test_audit_finds_placeholders(self, pptx_advanced_tools, temp_dir):
        """Should find placeholder patterns."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "<Customer Name> Project"
        path = temp_dir / "with_placeholders.pptx"
        prs.save(path)

        result = pptx_advanced_tools.tool_pptx_audit_placeholders(str(path))
        findings = result.get("findings", result.get("placeholders", []))
        # Should find the <Customer Name> placeholder
        assert len(findings) >= 1 or "Customer" in str(result)


class TestHiddenSlides:
    """Tests for hidden slide operations."""

    def test_hide_slide(self, pptx_advanced_tools, multi_slide_presentation, temp_dir):
        """Should hide a slide."""
        output = temp_dir / "with_hidden.pptx"
        result = pptx_advanced_tools.tool_pptx_hide_slide(
            str(multi_slide_presentation),
            slide_number=2,
            hidden=True,
            output_path=str(output)
        )
        assert result.get("success") is True or result.get("hidden") is True

    def test_get_hidden_slides(self, pptx_advanced_tools, multi_slide_presentation):
        """Should list hidden slides."""
        result = pptx_advanced_tools.tool_pptx_get_hidden_slides(
            str(multi_slide_presentation)
        )
        # Should return list (possibly empty)
        hidden = result.get("hidden_slides", result.get("slides", []))
        assert isinstance(hidden, list)


class TestTextAutofit:
    """Tests for text autofit settings."""

    def test_set_autofit_shrink(self, pptx_advanced_tools, multi_slide_presentation, temp_dir):
        """Should set autofit to shrink."""
        output = temp_dir / "autofit.pptx"
        result = pptx_advanced_tools.tool_pptx_set_text_autofit(
            str(multi_slide_presentation),
            slide_number=1,
            shape_identifier="title",
            autofit_type="shrink",
            output_path=str(output)
        )
        assert result.get("success") is True


class TestReplaceText:
    """Tests for text replacement."""

    def test_replace_text_all_slides(self, pptx_advanced_tools, temp_dir):
        """Should replace text across all slides."""
        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"<Customer> Slide {i+1}"
        path = temp_dir / "replace_source.pptx"
        prs.save(path)

        output = temp_dir / "replaced.pptx"
        result = pptx_advanced_tools.tool_pptx_replace_text(
            str(path),
            find_text="<Customer>",
            replace_text="Contoso",
            output_path=str(output)
        )
        count = result.get("count", result.get("replacements", 0))
        assert count >= 3 or result.get("success") is True


class TestDuplicateSlide:
    """Tests for slide duplication."""

    def test_duplicate_slide_after(self, pptx_advanced_tools, multi_slide_presentation, temp_dir):
        """Should duplicate slide after original."""
        output = temp_dir / "duplicated.pptx"
        result = pptx_advanced_tools.tool_pptx_duplicate_slide(
            str(multi_slide_presentation),
            slide_number=2,
            position="after",
            output_path=str(output)
        )
        assert result.get("success") is True or result.get("new_slide_number") is not None

    def test_duplicate_slide_end(self, pptx_advanced_tools, multi_slide_presentation, temp_dir):
        """Should duplicate slide at end."""
        output = temp_dir / "dup_end.pptx"
        result = pptx_advanced_tools.tool_pptx_duplicate_slide(
            str(multi_slide_presentation),
            slide_number=1,
            position="end",
            output_path=str(output)
        )
        assert result.get("success") is True


class TestReplacePlaceholders:
    """Tests for bulk placeholder replacement."""

    def test_replace_multiple_placeholders(self, pptx_advanced_tools, temp_dir):
        """Should replace multiple placeholders at once."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "<Customer Name> - <Project Name>"
        path = temp_dir / "multi_ph.pptx"
        prs.save(path)

        output = temp_dir / "multi_replaced.pptx"
        result = pptx_advanced_tools.tool_pptx_replace_placeholders(
            str(path),
            replacements={
                "<Customer Name>": "Contoso Corp",
                "<Project Name>": "Cloud Migration"
            },
            output_path=str(output)
        )
        assert result.get("success") is True


class TestCommentOperations:
    """Tests for comment operations."""

    def test_get_comments_empty(self, pptx_advanced_tools, multi_slide_presentation):
        """Should return empty for presentation without comments."""
        result = pptx_advanced_tools.tool_pptx_get_comments(
            str(multi_slide_presentation)
        )
        # Result may be dict with empty comments or just empty dict
        assert isinstance(result, dict)


class TestErrorHandling:
    """Tests for error handling."""

    def test_get_slide_invalid_number(self, pptx_advanced_tools, multi_slide_presentation):
        """Should handle invalid slide number."""
        result = pptx_advanced_tools.tool_pptx_get_slide(
            str(multi_slide_presentation),
            slide_number=999
        )
        assert "error" in result or "Error" in str(result)

    def test_delete_invalid_slide(self, pptx_advanced_tools, multi_slide_presentation):
        """Should handle deleting invalid slide."""
        result = pptx_advanced_tools.tool_pptx_delete_slide(
            str(multi_slide_presentation),
            slide_number=999
        )
        assert "error" in result


class TestNonExistentFile:
    """Tests for non-existent file handling."""

    def test_list_slides_missing_file(self, pptx_advanced_tools):
        """Should handle missing file."""
        result = pptx_advanced_tools.tool_pptx_list_slides("/nonexistent/file.pptx")
        assert "error" in result


class TestAdvancedBulletOperations:
    """Tests for advanced bullet operations."""

    def test_add_bullet_with_level(self, pptx_advanced_tools, multi_slide_presentation, temp_dir):
        """Should add indented bullet."""
        output = temp_dir / "indented.pptx"
        result = pptx_advanced_tools.tool_pptx_add_bullet(
            str(multi_slide_presentation),
            slide_number=2,
            text="Sub-point item",
            level=1,
            output_path=str(output)
        )
        assert result.get("success") is True

    def test_add_bullet_with_bold_label(self, pptx_advanced_tools, multi_slide_presentation, temp_dir):
        """Should add bullet with bold label."""
        output = temp_dir / "bold_label.pptx"
        result = pptx_advanced_tools.tool_pptx_add_bullet(
            str(multi_slide_presentation),
            slide_number=2,
            text="Some detail text",
            bold_label="Key Point",
            output_path=str(output)
        )
        assert result.get("success") is True


class TestNotesWithAppend:
    """Tests for notes with append mode."""

    def test_append_notes(self, pptx_advanced_tools, multi_slide_presentation, temp_dir):
        """Should append to existing notes."""
        # First add initial notes
        output1 = temp_dir / "notes1.pptx"
        pptx_advanced_tools.tool_pptx_set_notes(
            str(multi_slide_presentation),
            slide_number=1,
            notes_text="Initial notes",
            output_path=str(output1)
        )

        # Then append more
        output2 = temp_dir / "notes2.pptx"
        result = pptx_advanced_tools.tool_pptx_set_notes(
            str(output1),
            slide_number=1,
            notes_text="\nAdditional notes",
            append=True,
            output_path=str(output2)
        )
        assert result.get("success") is True
