"""Pytest configuration and fixtures for MCP tools testing."""

import tempfile
from pathlib import Path

import pytest
from docx import Document

# Check for openpyxl availability
try:
    from openpyxl import Workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

# Check for python-pptx availability
try:
    from pptx import Presentation
    from pptx.util import Mm
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def _ensure_default_docx_template(templates_dir: Path) -> None:
    """Ensure a usable default.docx exists for python-docx."""
    try:
        import docx.api
    except ImportError:
        return

    default_path = templates_dir / "default.docx"
    if default_path.exists():
        docx.api._default_docx_path = lambda: str(default_path)
        return

    # Create a minimal valid document using python-docx's built-in default
    # This ensures standard styles (List Bullet, List Number, etc.) are available
    templates_dir.mkdir(parents=True, exist_ok=True)
    doc = Document()  # Uses python-docx built-in template
    doc.save(default_path)
    docx.api._default_docx_path = lambda: str(default_path)


def _ensure_default_pptx_template(templates_dir: Path) -> None:
    """Ensure a usable default.pptx exists for python-pptx."""
    if not HAS_PPTX:
        return

    try:
        import pptx.api
    except ImportError:
        return

    default_path = templates_dir / "default.pptx"
    if default_path.exists():
        pptx.api._default_pptx_path = lambda: str(default_path)
        return

    # Create a minimal valid presentation using python-pptx's built-in default
    # This ensures standard placeholder indexes (0=title, 1=subtitle)
    templates_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()  # Uses python-pptx's built-in template
    prs.save(default_path)
    pptx.api._default_pptx_path = lambda: str(default_path)

    prs.save(default_path)
    pptx.api._default_pptx_path = lambda: str(default_path)


def pytest_configure():
    """Configure fallback templates for Office libraries."""
    templates_dir = Path(__file__).resolve().parent / "_templates"
    _ensure_default_docx_template(templates_dir)
    _ensure_default_pptx_template(templates_dir)


# =============================================================================
# Core fixtures
# =============================================================================


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as tmpdir:
        yield Path(tmpdir)


# =============================================================================
# Tool instance fixtures
# =============================================================================


@pytest.fixture
def word_tools():
    """Create an instance of WordTools."""
    from tools.word_tools import WordTools
    return WordTools()


@pytest.fixture
def word_advanced_tools():
    """Create an instance of WordAdvancedTools."""
    from tools.word_advanced_tools import WordAdvancedTools
    return WordAdvancedTools()


@pytest.fixture
def excel_tools():
    """Create an instance of ExcelTools."""
    from tools.excel_tools import ExcelTools
    return ExcelTools()


@pytest.fixture
def excel_advanced_tools():
    """Create an instance of ExcelAdvancedTools."""
    from tools.excel_advanced_tools import ExcelAdvancedTools
    return ExcelAdvancedTools()


@pytest.fixture
def pptx_tools():
    """Create an instance of PowerPointTools."""
    from tools.pptx_tools import PowerPointTools
    return PowerPointTools()


@pytest.fixture
def pptx_advanced_tools():
    """Create an instance of PresentationAdvancedTools."""
    from tools.pptx_advanced_tools import PresentationAdvancedTools
    return PresentationAdvancedTools()


@pytest.fixture
def web_tools():
    """Create an instance of WebTools."""
    from tools.web_tools import WebTools
    return WebTools()


# =============================================================================
# Sample document fixtures
# =============================================================================


@pytest.fixture
def simple_docx(temp_dir):
    """Create a simple test document with placeholder text."""
    doc = Document()
    doc.add_paragraph("Hello <Customer Name>, welcome to <Project Name>.")
    doc.add_paragraph("This is a test document for <Customer Name>.")

    # Add a table
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Header 1"
    table.rows[0].cells[1].text = "Header 2"
    table.rows[1].cells[0].text = "<Customer Name>"
    table.rows[1].cells[1].text = "Value"

    filepath = temp_dir / "test_document.docx"
    doc.save(filepath)
    return filepath


@pytest.fixture
def multi_paragraph_docx(temp_dir):
    """Create a document with multiple paragraphs for track changes testing."""
    doc = Document()
    doc.add_heading("Test Document", 0)
    doc.add_paragraph("First paragraph with PLACEHOLDER text.")
    doc.add_paragraph("Second paragraph with PLACEHOLDER and more PLACEHOLDER content.")
    doc.add_paragraph("Third paragraph without any placeholders.")
    doc.add_paragraph("Fourth paragraph with single PLACEHOLDER.")

    filepath = temp_dir / "multi_paragraph.docx"
    doc.save(filepath)
    return filepath


@pytest.fixture
def empty_docx(temp_dir):
    """Create an empty Word document."""
    doc = Document()
    path = temp_dir / "empty.docx"
    doc.save(path)
    return path


@pytest.fixture
def empty_pptx(temp_dir):
    """Create an empty PowerPoint presentation."""
    if not HAS_PPTX:
        pytest.skip("python-pptx not installed")

    prs = Presentation()
    path = temp_dir / "empty.pptx"
    prs.save(path)
    return path


@pytest.fixture
def sample_docx(temp_dir):
    """Create a test Word document with sections."""
    doc = Document()
    doc.add_heading("Test Document", level=0)
    doc.add_paragraph("This is a test paragraph.")
    doc.add_heading("Section 1", level=1)
    doc.add_paragraph("Content for section 1.")

    # Add a table
    table = doc.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(0, 2).text = "Header 3"
    table.cell(1, 0).text = "Data 1"
    table.cell(1, 1).text = "Data 2"
    table.cell(1, 2).text = "Data 3"

    path = temp_dir / "test.docx"
    doc.save(path)
    return path


@pytest.fixture
def sow_template(temp_dir):
    """Create a SOW-style template document."""
    doc = Document()
    doc.add_heading("Statement of Work", level=0)
    doc.add_paragraph("<Customer Name>")
    doc.add_paragraph("<Project Name>")

    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph("[Template Guidance: Describe the project here]")

    doc.add_heading("Scope", level=1)
    doc.add_paragraph("<Customer Name> requires the following services.")

    doc.add_heading("Staffing", level=1)
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Role"
    table.cell(0, 1).text = "Hours"

    path = temp_dir / "template.docx"
    doc.save(path)
    return path


@pytest.fixture
def sample_xlsx(temp_dir):
    """Create a simple test Excel file."""
    if not HAS_OPENPYXL:
        pytest.skip("openpyxl not installed")

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"

    # Headers
    ws['A1'] = "Name"
    ws['B1'] = "Value"
    ws['C1'] = "Status"

    # Data
    ws['A2'] = "Item 1"
    ws['B2'] = 100
    ws['C2'] = "Active"

    ws['A3'] = "Item 2"
    ws['B3'] = 200
    ws['C3'] = "Inactive"

    path = temp_dir / "test.xlsx"
    wb.save(path)
    return path


@pytest.fixture
def sample_pptx(temp_dir):
    """Create a simple test presentation."""
    if not HAS_PPTX:
        pytest.skip("python-pptx not installed")

    prs = Presentation()
    prs.slide_width = Mm(338.67)  # 16:9
    prs.slide_height = Mm(190.5)

    # Add a title slide
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Test Title"

    # Add content slide
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Content Slide"

    path = temp_dir / "test.pptx"
    prs.save(path)
    return path
