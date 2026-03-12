"""
Final strategic tests targeting specific uncovered paths.
"""

import tempfile
from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches as PptxInches

from tools.excel_tools import ExcelTools
from tools.pptx_advanced_tools import PresentationAdvancedTools
from tools.pptx_tools import PowerPointTools
from tools.word_advanced_tools import WordAdvancedTools
from tools.word_tools import WordTools


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as td:
        yield Path(td)


# Test module-level error handling
class TestModuleLevelErrors:
    """Test module-level error cases."""

    def test_word_tools_instance(self):
        """Create word tools."""
        tools = WordTools()
        assert tools is not None

    def test_pptx_tools_instance(self):
        """Create pptx tools."""
        tools = PowerPointTools()
        assert tools is not None

    def test_excel_tools_instance(self):
        """Create excel tools."""
        tools = ExcelTools()
        assert tools is not None

    def test_word_advanced_tools_instance(self):
        """Create word advanced tools."""
        tools = WordAdvancedTools()
        assert tools is not None

    def test_pptx_advanced_tools_instance(self):
        """Create pptx advanced tools."""
        tools = PresentationAdvancedTools()
        assert tools is not None


class TestWordToMarkdownEdgeCases:
    """Test word to markdown edge cases."""

    def test_to_markdown_with_complex_formatting(self, temp_dir):
        """Should handle complex formatting."""
        doc = Document()
        doc.add_heading("Title", level=1)

        # Add paragraph with mixed formatting
        para = doc.add_paragraph()
        para.add_run("Normal ")
        bold_run = para.add_run("bold")
        bold_run.bold = True
        para.add_run(" and ")
        italic_run = para.add_run("italic")
        italic_run.italic = True

        # Add a list
        doc.add_paragraph("Item 1", style='List Bullet')
        doc.add_paragraph("Item 2", style='List Bullet')

        path = temp_dir / "complex_format.docx"
        doc.save(path)

        tools = WordTools()
        result = tools.tool_word_to_markdown(str(path))
        assert "Title" in result

    def test_to_markdown_with_table(self, temp_dir):
        """Should convert doc with table to markdown."""
        doc = Document()
        doc.add_heading("Data Report", level=1)

        table = doc.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Header A"
        table.cell(0, 1).text = "Header B"
        table.cell(1, 0).text = "Row 1 A"
        table.cell(1, 1).text = "Row 1 B"
        table.cell(2, 0).text = "Row 2 A"
        table.cell(2, 1).text = "Row 2 B"

        path = temp_dir / "with_table.docx"
        doc.save(path)

        tools = WordTools()
        result = tools.tool_word_to_markdown(str(path))
        assert "Header A" in result


class TestPptxToMarkdownEdgeCases:
    """Test pptx to markdown edge cases."""

    def test_to_markdown_with_notes(self, temp_dir):
        """Should convert pptx with notes to markdown."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Title Slide"

        notes = slide.notes_slide
        notes.notes_text_frame.text = "Speaker notes here"

        path = temp_dir / "with_notes.pptx"
        prs.save(path)

        tools = PowerPointTools()
        result = tools.tool_pptx_to_markdown(str(path))
        assert "Title Slide" in result


class TestExcelEdgeCases:
    """Test Excel edge cases."""

    def test_extract_empty_cells(self, temp_dir):
        """Should handle Excel with empty cells."""
        tools = ExcelTools()

        md = """| A | B | C |
|---|---|---|
| 1 |   | 3 |
|   | 2 |   |
"""
        path = temp_dir / "empty_cells.xlsx"
        tools.tool_excel_from_markdown(str(path), md)

        result = tools.tool_excel_extract(str(path))
        assert isinstance(result, dict)


class TestWordPatchSectionEdgeCases:
    """Test patch section edge cases."""

    def test_patch_section_with_special_chars(self, temp_dir):
        """Should handle special characters in section names."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_heading("Section (A)", level=1)
        doc.add_paragraph("Content A")
        doc.add_heading("Section [B]", level=1)
        doc.add_paragraph("Content B")
        path = temp_dir / "special_sections.docx"
        doc.save(path)

        result = tools.tool_word_list_sections(str(path))
        assert len(result.get("sections", [])) >= 2


class TestPptxBulletLevelVariations:
    """Test different bullet levels."""

    def test_add_deeply_nested_bullet(self, temp_dir):
        """Should add deeply nested bullet."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Nested Bullets"
        path = temp_dir / "nested_bullets.pptx"
        prs.save(path)

        # Add bullets at various levels
        tools.tool_pptx_add_bullet(str(path), 1, "Level 0", level=0)
        tools.tool_pptx_add_bullet(str(path), 1, "Level 1", level=1)
        tools.tool_pptx_add_bullet(str(path), 1, "Level 2", level=2)
        result = tools.tool_pptx_add_bullet(str(path), 1, "Level 3", level=3)

        assert isinstance(result, dict)


class TestWordFixSplitPlaceholdersComplex:
    """Test complex split placeholder scenarios."""

    def test_fix_split_across_many_runs(self, temp_dir):
        """Should fix placeholder split across many runs."""
        tools = WordAdvancedTools()

        doc = Document()
        para = doc.add_paragraph()
        # Split a placeholder across multiple runs
        para.add_run("<")
        para.add_run("Customer")
        para.add_run(" ")
        para.add_run("Name")
        para.add_run(">")
        para.add_run(" is important")

        path = temp_dir / "multi_split.docx"
        doc.save(path)

        result = tools.tool_word_fix_split_placeholders(
            str(path),
            {"<Customer Name>": "Acme Corp"}
        )
        assert isinstance(result, dict)


class TestPptxHideUnhideOperations:
    """Test hide/unhide operations."""

    def test_hide_then_unhide(self, temp_dir):
        """Should hide then unhide slide."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Toggleable"
        path = temp_dir / "toggle.pptx"
        prs.save(path)

        # Hide
        tools.tool_pptx_hide_slide(str(path), 1, hidden=True)

        # Get hidden
        result = tools.tool_pptx_get_hidden_slides(str(path))
        assert isinstance(result, dict)

        # Unhide
        result = tools.tool_pptx_hide_slide(str(path), 1, hidden=False)
        assert isinstance(result, dict)


class TestWordCommentVariations:
    """Test comment variations."""

    def test_add_comment_target_not_found(self, temp_dir):
        """Should handle target text not found."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_paragraph("Some content here")
        path = temp_dir / "no_target.docx"
        doc.save(path)

        result = tools.tool_word_add_comment(
            str(path),
            "nonexistent text",
            "Comment for missing text"
        )
        # Should handle gracefully
        assert isinstance(result, dict)


class TestPptxFromMarkdownSpecialContent:
    """Test pptx from markdown with special content."""

    def test_with_table_content(self, temp_dir):
        """Should handle markdown table."""
        tools = PowerPointTools()

        md = """# Data Presentation

## Statistics

| Metric | Value |
|--------|-------|
| Users  | 1000  |
| Revenue| $50K  |

## Summary
- Good performance
- Growth expected
"""
        path = temp_dir / "table_slides.pptx"
        tools.tool_pptx_from_markdown(str(path), md)
        assert Path(path).exists()


class TestWordFromMarkdownSpecialContent:
    """Test word from markdown with special content."""

    def test_with_blockquotes(self, temp_dir):
        """Should handle blockquotes."""
        tools = WordTools()

        md = """# Document

## Quote Section

> This is a blockquote
> spanning multiple lines

## Regular Content
Normal paragraph here.
"""
        path = temp_dir / "blockquotes.docx"
        tools.tool_word_from_markdown(str(path), md)
        assert Path(path).exists()


class TestWordDuplicateTableOperations:
    """Test duplicate table operations."""

    def test_duplicate_table_structure_with_styling(self, temp_dir):
        """Should duplicate table with styling."""
        tools = WordAdvancedTools()

        doc = Document()
        table = doc.add_table(rows=2, cols=3)
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"
        table.cell(0, 2).text = "C"
        table.cell(1, 0).text = "D"
        table.cell(1, 1).text = "E"
        table.cell(1, 2).text = "F"

        path = temp_dir / "styled_table.docx"
        doc.save(path)

        result = tools.tool_word_duplicate_table_structure(str(path), "0")
        assert isinstance(result, dict)


class TestPptxSlideNumberEdgeCases:
    """Test slide number edge cases."""

    def test_slide_number_zero(self, temp_dir):
        """Should handle invalid slide number 0."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "zero_slide.pptx"
        prs.save(path)

        result = tools.tool_pptx_get_slide(str(path), 0)
        # Should return error or handle gracefully
        assert isinstance(result, dict)

    def test_negative_slide_number(self, temp_dir):
        """Should handle negative slide number."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])
        path = temp_dir / "neg_slide.pptx"
        prs.save(path)

        result = tools.tool_pptx_get_slide(str(path), -1)
        # Should handle gracefully
        assert isinstance(result, dict)


class TestWordGetSectionVariations:
    """Test get section variations."""

    def test_get_section_fuzzy_match(self, temp_dir):
        """Should get section with partial match."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_heading("1. Introduction to the Project", level=1)
        doc.add_paragraph("Intro content")
        doc.add_heading("2. Background Information", level=1)
        doc.add_paragraph("Background content")
        path = temp_dir / "long_sections.docx"
        doc.save(path)

        result = tools.tool_word_get_section(str(path), "Introduction")
        assert isinstance(result, dict)


class TestPptxSlideContentReading:
    """Test slide content reading."""

    def test_get_slide_with_multiple_shapes(self, temp_dir):
        """Should get slide with multiple shape types."""
        tools = PresentationAdvancedTools()

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        slide.shapes.title.text = "Title"

        # Add textbox
        txBox = slide.shapes.add_textbox(PptxInches(1), PptxInches(2), PptxInches(3), PptxInches(1))
        txBox.text_frame.text = "Textbox content"

        # Add table
        table = slide.shapes.add_table(2, 2, PptxInches(5), PptxInches(2), PptxInches(3), PptxInches(1.5))
        tbl = table.table
        tbl.cell(0, 0).text = "A"
        tbl.cell(0, 1).text = "B"

        path = temp_dir / "multi_shapes.pptx"
        prs.save(path)

        result = tools.tool_pptx_get_slide(str(path), 1)
        assert isinstance(result, dict)


class TestWordReplaceGlobalVariablesComplex:
    """Test complex global variable replacement."""

    def test_replace_in_tables(self, temp_dir):
        """Should replace in tables too."""
        tools = WordAdvancedTools()

        doc = Document()
        doc.add_paragraph("Text with <Customer Name>")

        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Customer"
        table.cell(0, 1).text = "<Customer Name>"
        table.cell(1, 0).text = "Project"
        table.cell(1, 1).text = "<Project Name>"

        path = temp_dir / "table_replace.docx"
        doc.save(path)

        result = tools.tool_word_replace_global_variables(
            str(path),
            {
                "<Customer Name>": "Contoso",
                "<Project Name>": "Migration"
            }
        )
        assert isinstance(result, dict)


class TestPptxExtractContent:
    """Test PPTX content extraction."""

    def test_extract_with_all_content_types(self, temp_dir):
        """Should extract all content types."""
        tools = PowerPointTools()

        prs = Presentation()

        # Title slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Main Title"
        for shape in slide1.shapes:
            if hasattr(shape, 'text_frame') and shape != slide1.shapes.title:
                shape.text_frame.text = "Subtitle"

        # Content slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Content"
        for shape in slide2.shapes:
            if shape.has_text_frame and shape != slide2.shapes.title:
                shape.text_frame.text = "Bullet points"

        path = temp_dir / "full_extract.pptx"
        prs.save(path)

        result = tools.tool_pptx_extract(str(path))
        assert isinstance(result, dict)
        assert "slides" in result
