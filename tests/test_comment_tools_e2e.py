"""End-to-end comment tool tests across Excel, Word, and PowerPoint.

These tests validate read/write/delete behavior through:
- Unified MCP comment tool: office_comment
- Format-specific comment tools for each document type
"""

from __future__ import annotations

import pytest

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import lxml  # noqa: F401
    HAS_LXML = True
except ImportError:
    HAS_LXML = False

try:
    import pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from tools import TOOL_CLASSES


def create_combined_tools_class():
    """Create a class that combines all MCP tool mixins."""

    class CombinedTools(*TOOL_CLASSES):
        """Combined tool class for E2E testing."""

    return CombinedTools


@pytest.fixture
def tools():
    """Provide a full tool instance similar to the dynamic server composition."""
    CombinedTools = create_combined_tools_class()
    return CombinedTools()


@pytest.fixture
def excel_comment_file(temp_dir):
    """Create an Excel workbook prepared for comment operations."""
    if not HAS_OPENPYXL:
        pytest.skip("openpyxl not installed")
    path = temp_dir / "comments.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Data"
    ws["A1"] = "Revenue"
    ws["B1"] = 1000
    wb.save(path)
    return str(path)


@pytest.fixture
def word_comment_file(temp_dir):
    """Create a Word document prepared for comment operations."""
    if not HAS_DOCX:
        pytest.skip("python-docx not installed")
    path = temp_dir / "comments.docx"
    doc = docx.Document()
    doc.add_paragraph("This sentence is the comment target.")
    doc.add_paragraph("Another paragraph for context.")
    doc.save(path)
    return str(path)


@pytest.fixture
def pptx_comment_file(temp_dir):
    """Create a PowerPoint deck prepared for comment operations."""
    if not HAS_PPTX:
        pytest.skip("python-pptx not installed")
    path = temp_dir / "comments.pptx"
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "Comment Test Slide"
    prs.save(path)
    return str(path)


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
def test_excel_comments_unified_read_write_delete(tools, excel_comment_file):
    """office_comment should add, read, and delete Excel comments."""
    add = tools.tool_office_comment(
        file_path=excel_comment_file,
        operation="add",
        target="A1",
        text="Validate this number with Finance",
    )
    assert add.get("success") is True

    got = tools.tool_office_comment(file_path=excel_comment_file, operation="get")
    assert "error" not in got
    assert got.get("total_comments", 0) == 1

    sheet_comments = got.get("by_sheet", {}).get("Data", [])
    assert len(sheet_comments) == 1
    assert sheet_comments[0].get("cell") == "A1"
    assert "Finance" in sheet_comments[0].get("text", "")

    deleted = tools.tool_office_comment(
        file_path=excel_comment_file,
        operation="delete",
        target="A1",
    )
    assert deleted.get("success") is True

    after = tools.tool_office_comment(file_path=excel_comment_file, operation="get")
    assert after.get("total_comments", 0) == 0


@pytest.mark.skipif(not HAS_OPENPYXL, reason="openpyxl not installed")
def test_excel_comments_direct_read_write_delete(tools, excel_comment_file):
    """excel_* comment tools should add, read, and delete Excel comments."""
    add = tools.tool_excel_add_comment(
        file_path=excel_comment_file,
        cell_ref="A1",
        text="Direct Excel comment",
    )
    assert add.get("success") is True

    got = tools.tool_excel_get_comments(file_path=excel_comment_file)
    assert got.get("total_comments", 0) == 1

    deleted = tools.tool_excel_delete_comment(file_path=excel_comment_file, cell_ref="A1")
    assert deleted.get("success") is True

    after = tools.tool_excel_get_comments(file_path=excel_comment_file)
    assert after.get("total_comments", 0) == 0


@pytest.mark.skipif(not (HAS_DOCX and HAS_LXML), reason="python-docx or lxml not installed")
def test_word_comments_unified_read_write_delete(tools, word_comment_file):
    """office_comment should add, read, and delete Word comments."""
    add = tools.tool_office_comment(
        file_path=word_comment_file,
        operation="add",
        target="comment target",
        text="Please verify this claim",
    )
    assert add.get("success") is True

    got = tools.tool_office_comment(file_path=word_comment_file, operation="get")
    assert "error" not in got
    assert got.get("comment_count", 0) >= 1

    comments = got.get("comments", [])
    assert comments
    assert "verify" in comments[0].get("text", "").lower()
    comment_id = comments[0]["id"]

    deleted = tools.tool_office_comment(
        file_path=word_comment_file,
        operation="delete",
        target=str(comment_id),
    )
    assert deleted.get("success") is True

    after = tools.tool_office_comment(file_path=word_comment_file, operation="get")
    assert after.get("comment_count", 0) == 0


@pytest.mark.skipif(not (HAS_DOCX and HAS_LXML), reason="python-docx or lxml not installed")
def test_word_comments_direct_read_write_delete(tools, word_comment_file):
    """word_* comment tools should add, read, and delete Word comments."""
    add = tools.tool_word_add_comment(
        file_path=word_comment_file,
        target_text="comment target",
        comment_text="Direct Word comment",
    )
    assert add.get("success") is True

    got = tools.tool_word_get_comments(file_path=word_comment_file)
    assert got.get("comment_count", 0) >= 1
    comment_id = got["comments"][0]["id"]

    deleted = tools.tool_word_delete_comment(file_path=word_comment_file, comment_id=str(comment_id))
    assert deleted.get("success") is True

    after = tools.tool_word_get_comments(file_path=word_comment_file)
    assert after.get("comment_count", 0) == 0


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
def test_pptx_comments_unified_read_write_delete(tools, pptx_comment_file):
    """office_comment should add, read, and delete PowerPoint comments."""
    add = tools.tool_office_comment(
        file_path=pptx_comment_file,
        operation="add",
        target="slide:1",
        text="Please update this title",
    )
    assert add.get("success") is True

    got = tools.tool_office_comment(file_path=pptx_comment_file, operation="get")
    assert "error" not in got
    assert got.get("total_comments", 0) >= 1

    slide_comments = got.get("comments", {}).get(1, [])
    assert slide_comments
    assert "update" in slide_comments[0].get("text", "").lower()
    idx = slide_comments[0]["index"]

    deleted = tools.tool_office_comment(
        file_path=pptx_comment_file,
        operation="delete",
        target=f"slide:1/comment:{idx}",
    )
    assert deleted.get("success") is True

    after = tools.tool_office_comment(file_path=pptx_comment_file, operation="get")
    assert after.get("total_comments", 0) == 0


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
def test_pptx_comments_direct_read_write_delete(tools, pptx_comment_file):
    """pptx_* comment tools should add, read, and delete PowerPoint comments."""
    add = tools.tool_pptx_add_comment(
        file_path=pptx_comment_file,
        slide_number=1,
        comment_text="Direct PPTX comment",
    )
    assert add.get("success") is True

    got = tools.tool_pptx_get_comments(file_path=pptx_comment_file, slide_number=1)
    assert got.get("total_comments", 0) >= 1
    idx = got["comments"][1][0]["index"]

    deleted = tools.tool_pptx_delete_comment(
        file_path=pptx_comment_file,
        slide_number=1,
        comment_index=int(idx),
    )
    assert deleted.get("success") is True

    after = tools.tool_pptx_get_comments(file_path=pptx_comment_file, slide_number=1)
    assert after.get("total_comments", 0) == 0
