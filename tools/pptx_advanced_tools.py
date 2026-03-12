#!/usr/bin/env python3
"""
pptx_advanced_tools.py - Advanced MCP tools for PowerPoint presentation manipulation

Provides comprehensive tools to:
- Introspect presentations (masters, layouts, slides, shapes)
- Manipulate slide order and structure
- Patch text, tables, and content
- Add speaker notes and comments
- Work with existing templates while preserving formatting

Quality bar: All tools preserve presentation structure and formatting
by editing templates rather than creating documents from scratch.
"""

import os
import re
import shutil
from contextlib import suppress
from copy import deepcopy
from pathlib import Path
from typing import Any, TypedDict

try:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml.ns import qn
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from .save_utils import open_pptx_with_retries, safe_save_pptx


class ChangeLogEntry(TypedDict, total=False):
    """A single change log entry for pptx_log_changes."""
    slide: int     # Slide number that was modified
    action: str    # What was done (e.g., 'Updated title')
    detail: str    # Specific details about the change


# Unit conversion notes:
# - All public API parameters use inches
# - Widescreen 16:9 = 13.333" × 7.5"
# - Standard margin = 0.5"
# - Content width = 12.333"


# Author name for comments (from environment or default)
DEFAULT_AUTHOR = os.environ.get("MCP_AUTHOR", "Solution Architect Agent")


def _derive_initials(author: str) -> str:
    """Derive initials from an author display name."""
    tokens = [part for part in str(author).split() if part]
    return "".join(part[0].upper() for part in tokens[:2]) if tokens else "SA"


def _get_shape_info(shape) -> dict:
    """Extract comprehensive information about a shape."""
    info = {
        "shape_type": str(shape.shape_type),
        "name": shape.name,
        "shape_id": shape.shape_id,
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }

    if shape.is_placeholder:
        ph = shape.placeholder_format
        info["is_placeholder"] = True
        info["placeholder_idx"] = ph.idx
        info["placeholder_type"] = str(ph.type)
    else:
        info["is_placeholder"] = False

    if shape.has_text_frame:
        info["has_text_frame"] = True
        info["text_preview"] = shape.text_frame.text[:100] if shape.text_frame.text else ""
    else:
        info["has_text_frame"] = False

    if shape.has_table:
        info["has_table"] = True
        info["table_rows"] = len(shape.table.rows)
        info["table_cols"] = len(shape.table.columns)
    else:
        info["has_table"] = False

    return info


def _find_shape_by_identifier(slide, identifier: str):
    """Find a shape by name, index, or placeholder type.

    Args:
        slide: pptx Slide object
        identifier: Shape name, index (as string), or placeholder type

    Returns:
        Shape object or None
    """
    # Try as index
    if identifier.isdigit():
        idx = int(identifier)
        shapes_list = list(slide.shapes)
        if 0 <= idx < len(shapes_list):
            return shapes_list[idx]

    # Try as placeholder type
    # Note: "body" matches both BODY and OBJECT placeholder types since
    # OBJECT is commonly used for content placeholders in many layouts
    placeholder_types = {
        "title": [PP_PLACEHOLDER.TITLE],
        "body": [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT],
        "content": [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT],
        "subtitle": [PP_PLACEHOLDER.SUBTITLE],
        "center_title": [PP_PLACEHOLDER.CENTER_TITLE],
        "footer": [PP_PLACEHOLDER.FOOTER],
        "date": [PP_PLACEHOLDER.DATE],
        "slide_number": [PP_PLACEHOLDER.SLIDE_NUMBER],
    }

    if identifier.lower() in placeholder_types:
        target_types = placeholder_types[identifier.lower()]
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type in target_types:
                return shape

    # Try as shape name (exact or partial match)
    for shape in slide.shapes:
        if shape.name.lower() == identifier.lower():
            return shape
        if identifier.lower() in shape.name.lower():
            return shape

    return None


def _find_table_in_slide(slide, identifier: str = None):
    """Find a table shape in a slide.

    Args:
        slide: pptx Slide object
        identifier: Optional table index (as string) or None for first table

    Returns:
        (table_shape, table) tuple or (None, None)
    """
    tables = [(s, s.table) for s in slide.shapes if s.has_table]

    if not tables:
        return None, None

    if identifier and identifier.isdigit():
        idx = int(identifier)
        if 0 <= idx < len(tables):
            return tables[idx]

    return tables[0] if tables else (None, None)


def _load_presentation(file_path: str):
    """Load presentation with retry and path resolution diagnostics."""
    prs, resolved_path, error = open_pptx_with_retries(file_path)
    if error:
        return None, resolved_path, {"error": error}
    return prs, resolved_path, None


def _capture_text_frame_props(text_frame) -> dict[str, Any]:
    """Capture text frame properties for later restoration.

    Saves autofit settings and first paragraph/run formatting
    that should be preserved during text operations.
    """
    props = {
        "auto_size": text_frame.auto_size,
        "word_wrap": text_frame.word_wrap,
    }

    # Capture first paragraph properties if available
    if text_frame.paragraphs:
        first_para = text_frame.paragraphs[0]
        props["para_level"] = first_para.level
        props["para_alignment"] = first_para.alignment

        # Capture first run font properties if available
        if first_para.runs:
            first_run = first_para.runs[0]
            props["font_size"] = first_run.font.size
            props["font_bold"] = first_run.font.bold
            props["font_italic"] = first_run.font.italic
            props["font_name"] = first_run.font.name
            # Only capture RGB color if explicitly set
            if first_run.font.color.type is not None:
                with suppress(AttributeError):
                    props["font_color_rgb"] = first_run.font.color.rgb

    return props


def _restore_text_frame_props(text_frame, props: dict[str, Any], restore_autofit: bool = True) -> None:
    """Restore text frame properties after text modification.

    Args:
        text_frame: The text frame to restore properties to
        props: Dictionary from _capture_text_frame_props
        restore_autofit: If True, restore autofit settings (default True)
    """
    if restore_autofit and props.get("auto_size") is not None:
        text_frame.auto_size = props["auto_size"]

    if props.get("word_wrap") is not None:
        text_frame.word_wrap = props["word_wrap"]


def _set_text_preserving_format(text_frame, new_text: str, props: dict[str, Any] = None) -> None:
    """Set text while preserving formatting from captured properties.

    Clears existing text and sets new text, then restores formatting
    from the captured properties dictionary.

    Args:
        text_frame: The text frame to modify
        new_text: The new text to set
        props: Dictionary from _capture_text_frame_props (optional)
    """
    # Clear and set new text
    text_frame.clear()

    if not new_text:
        return

    # Add new text
    para = text_frame.paragraphs[0]
    run = para.add_run()
    run.text = new_text

    # Restore properties if provided
    if props:
        _restore_text_frame_props(text_frame, props)

        # Restore paragraph properties
        if props.get("para_level") is not None:
            para.level = props["para_level"]
        if props.get("para_alignment") is not None:
            para.alignment = props["para_alignment"]

        # Restore font properties
        if props.get("font_size") is not None:
            run.font.size = props["font_size"]
        if props.get("font_bold") is not None:
            run.font.bold = props["font_bold"]
        if props.get("font_italic") is not None:
            run.font.italic = props["font_italic"]
        if props.get("font_name") is not None:
            run.font.name = props["font_name"]
        if props.get("font_color_rgb") is not None:
            run.font.color.rgb = props["font_color_rgb"]


class PresentationAdvancedTools:
    """MCP tool mixin for advanced PowerPoint presentation manipulation."""

    # =========================================================================
    # INTROSPECTION TOOLS
    # =========================================================================

    def tool_pptx_list_slides(self, file_path: str) -> dict[str, Any]:
        """Get an overview of all slides in a PowerPoint file.

        USE THIS FIRST when working with an existing presentation to understand
        its structure before making changes.

        Returns for each slide:
        - number: Slide position (1-based)
        - title: The slide's title text
        - layout_name: Which layout template is used
        - shape_count: Number of shapes on the slide
        - has_notes: Whether speaker notes exist
        - hidden: Whether the slide is hidden

        Args:
            file_path: Path to the .pptx file

        Returns:
            Dictionary with slide_count and slides array
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        try:
            prs = Presentation(file_path)
        except Exception as e:
            return {"error": f"Failed to open presentation: {e}"}

        slides = []

        for slide_num, slide in enumerate(prs.slides, 1):
            # Check if slide is hidden
            slide_elem = slide._element
            is_hidden = slide_elem.get('show') == '0'

            slide_info = {
                "number": slide_num,
                "slide_id": slide.slide_id,
                "title": None,
                "layout_name": slide.slide_layout.name,
                "shape_count": len(slide.shapes),
                "has_notes": slide.has_notes_slide and bool(slide.notes_slide.notes_text_frame.text.strip()),
                "hidden": is_hidden,
                "shape_types": [],
            }

            for shape in slide.shapes:
                shape_type_str = str(shape.shape_type).replace("MSO_SHAPE_TYPE.", "")
                if shape_type_str not in slide_info["shape_types"]:
                    slide_info["shape_types"].append(shape_type_str)

                # Get title
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    if (ph_type == PP_PLACEHOLDER.TITLE or ph_type == PP_PLACEHOLDER.CENTER_TITLE) and shape.has_text_frame:
                        slide_info["title"] = shape.text_frame.text.strip()

            slides.append(slide_info)

        return {
            "file": path.name,
            "slide_count": len(slides),
            "slides": slides,
            "next_tools": ["pptx_get_slide", "pptx_list_shapes", "pptx_replace_placeholders"]
        }

    def tool_pptx_list_masters(self, file_path: str) -> dict[str, Any]:
        """List available slide layouts in a presentation template.

        USE THIS to find layout_index values needed for pptx_add_slide.
        Each layout defines placeholder positions for titles, content, etc.

        Common layout indices (may vary by template):
        - 0: Title Slide (centered title + subtitle)
        - 1: Title and Content (title + bullet area)
        - 2: Section Header
        - 5: Title Only (good for custom content/tables)
        - 6: Blank

        Args:
            file_path: Path to the .pptx file

        Returns:
            Dictionary with default_layouts array containing index and name
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)

        masters = []
        for master_idx, master in enumerate(prs.slide_masters):
            master_info = {
                "index": master_idx,
                "layouts": [],
                "placeholders": [],
            }

            # Get master placeholders
            for ph in master.placeholders:
                master_info["placeholders"].append({
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type).replace("PP_PLACEHOLDER.", ""),
                })

            # Get layouts
            for layout_idx, layout in enumerate(master.slide_layouts):
                layout_info = {
                    "index": layout_idx,
                    "name": layout.name,
                    "placeholders": [],
                }

                for ph in layout.placeholders:
                    layout_info["placeholders"].append({
                        "idx": ph.placeholder_format.idx,
                        "type": str(ph.placeholder_format.type).replace("PP_PLACEHOLDER.", ""),
                    })

                master_info["layouts"].append(layout_info)

            masters.append(master_info)

        # Also list the default slide_layouts for convenience
        default_layouts = []
        for idx, layout in enumerate(prs.slide_layouts):
            default_layouts.append({
                "index": idx,
                "name": layout.name,
            })

        return {
            "file": path.name,
            "master_count": len(masters),
            "masters": masters,
            "default_layouts": default_layouts,
            "next_tools": ["pptx_add_slide", "pptx_copy_template"]
        }

    def tool_pptx_list_shapes(self, file_path: str, slide_number: int) -> dict[str, Any]:
        """List all shapes on a slide to find shape names for editing.

        USE THIS to find shape_identifier values for pptx_patch_shape.
        Shape names like 'Title 1', 'Content Placeholder 2', 'Table 3' can
        be used to target specific elements.

        Returns for each shape:
        - name: Use this as shape_identifier in other tools
        - is_placeholder: True for title/content placeholders
        - placeholder_type: 'TITLE', 'BODY', 'SUBTITLE', etc.
        - text_preview: First 100 chars of text content
        - has_table: Whether this shape contains a table

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number (first slide = 1)

        Returns:
            Dictionary with shapes array
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]
        shapes = []

        for idx, shape in enumerate(slide.shapes):
            shape_info = _get_shape_info(shape)
            shape_info["index"] = idx
            shapes.append(shape_info)

        return {
            "file": path.name,
            "slide_number": slide_number,
            "layout_name": slide.slide_layout.name,
            "shape_count": len(shapes),
            "shapes": shapes,
            "next_tools": ["pptx_patch_shape", "pptx_clear_bullets", "pptx_add_bullet"]
        }

    def tool_pptx_get_slide(self, file_path: str, slide_number: int) -> dict[str, Any]:
        """Read all content from a specific slide.

        USE THIS to see what's currently on a slide before editing.
        Returns structured content including title, body bullets, tables, and notes.

        Returns:
        - title: The slide title text
        - subtitle: Subtitle text (if present)
        - body_content: Array of {text, level} for each bullet
        - tables: Array of tables with header and data rows
        - notes: Speaker notes text

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number (first slide = 1)

        Returns:
            Dictionary with slide content
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]

        content = {
            "slide_number": slide_number,
            "slide_id": slide.slide_id,
            "layout_name": slide.slide_layout.name,
            "title": None,
            "subtitle": None,
            "body_content": [],
            "tables": [],
            "other_text": [],
            "notes": None,
        }

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type

                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE) and shape.has_text_frame:
                    content["title"] = shape.text_frame.text.strip()
                elif ph_type == PP_PLACEHOLDER.SUBTITLE and shape.has_text_frame:
                    content["subtitle"] = shape.text_frame.text.strip()
                elif ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) and shape.has_text_frame:
                    paragraphs = []
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            paragraphs.append({
                                "text": para.text,
                                "level": para.level,
                            })
                    content["body_content"].extend(paragraphs)

            if shape.has_table:
                table_data = {
                    "shape_name": shape.name,
                    "rows": len(shape.table.rows),
                    "columns": len(shape.table.columns),
                    "header": [],
                    "data": [],
                }

                for row_idx, row in enumerate(shape.table.rows):
                    row_data = [cell.text.strip() for cell in row.cells]
                    if row_idx == 0:
                        table_data["header"] = row_data
                    else:
                        table_data["data"].append(row_data)

                content["tables"].append(table_data)

            elif shape.has_text_frame and not shape.is_placeholder:
                text = shape.text_frame.text.strip()
                if text:
                    content["other_text"].append({
                        "shape_name": shape.name,
                        "text": text,
                    })

        # Get notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                content["notes"] = notes_text

        return {
            "file": path.name,
            "slide": content,
            "next_tools": ["pptx_patch_shape", "pptx_add_bullet", "pptx_add_comment"]
        }

    # =========================================================================
    # SLIDE MANAGEMENT TOOLS
    # =========================================================================

    def tool_pptx_copy_template(
        self,
        template_path: str,
        output_path: str,
        template_dir: str = None
    ) -> dict[str, Any]:
        """Copy a PowerPoint template to create a new presentation.

        START HERE when creating a presentation from a template.
        This preserves the template's formatting, layouts, and theme.

        After copying, use:
        1. pptx_list_slides - See the template structure
        2. pptx_replace_placeholders - Fill in <placeholder> values
        3. pptx_patch_shape - Update specific content
        4. pptx_add_bullet - Add bullet points

        Args:
            template_path: Template filename (e.g., 'SOW-Presentation.pptx')
                          or full path to template
            output_path: Where to save the new presentation
            template_dir: Optional directory containing templates

        Returns:
            Status with template and output paths
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        # Resolve template path
        template = Path(template_path)
        if not template.is_absolute() and template_dir:
            template = Path(template_dir) / template_path
        elif not template.is_absolute():
            # Check common locations
            skill_templates = Path(__file__).parent.parent.parent / "skills" / "present-sow" / "templates" / template_path
            if skill_templates.exists():
                template = skill_templates

        if not template.exists():
            return {"error": f"Template not found: {template_path}"}

        # Create output directory if needed
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)

        # Copy template
        shutil.copy2(template, output)

        return {
            "success": True,
            "template": str(template),
            "output": str(output),
            "message": f"Copied template to '{output_path}'",
            "next_tools": ["pptx_list_slides", "pptx_replace_placeholders", "pptx_list_masters"]
        }

    def tool_pptx_add_slide(
        self,
        file_path: str,
        layout_index: int = 1,
        title: str = None,
        position: str = "end",
        output_path: str = None
    ) -> dict[str, Any]:
        """Add a new slide to an existing presentation.

        Common layout_index values (use pptx_list_masters to see all):
        - 0: Title Slide
        - 1: Title and Content (default - has title + bullet area)
        - 5: Title Only (good for tables or custom content)
        - 6: Blank

        Args:
            file_path: Path to the .pptx file
            layout_index: Which layout to use (default: 1 = Title and Content)
            title: Optional title text for the new slide
            position: Where to insert - 'end' (default), 'start', or 1-based slide number (e.g. '2' to make it slide 2)
            output_path: Optional output path (defaults to overwriting input)

        Returns:
            Status with new slide_number
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, load_error = _load_presentation(file_path)
        if load_error:
            return load_error

        if layout_index < 0 or layout_index >= len(prs.slide_layouts):
            return {"error": f"Layout index {layout_index} invalid. Available: 0-{len(prs.slide_layouts)-1}"}

        position_text = str(position).strip().strip('"').strip("'")
        if position_text not in {"end", "start"} and not position_text.isdigit():
            return {
                "error": "Invalid position. Use 'end', 'start', or a 1-based slide number as string/integer."
            }

        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)

        # Set title if provided
        if title and slide.shapes.title:
            slide.shapes.title.text = title

        new_slide_num = len(prs.slides)

        # Handle positioning via XML manipulation
        if position_text != "end":
            sldIdLst = prs.part._element.find(qn('p:sldIdLst'))
            if sldIdLst is not None:
                children = list(sldIdLst)
                last_slide = children[-1]
                sldIdLst.remove(last_slide)

                if position_text == "start":
                    sldIdLst.insert(0, last_slide)
                    new_slide_num = 1
                elif position_text.isdigit():
                    # position is 1-based slide number
                    target_pos = int(position_text) - 1
                    num_existing = len(children) - 1  # minus the one we removed
                    if 0 <= target_pos <= num_existing:
                        sldIdLst.insert(target_pos, last_slide)
                        new_slide_num = target_pos + 1
                    else:
                        return {
                            "error": f"Position {position_text} out of range. Valid range: 1-{num_existing + 1}."
                        }

        # Save
        save_path = output_path or resolved_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "slide_number": new_slide_num,
            "layout_used": layout.name,
            "title": title,
            "message": f"Added slide at position {new_slide_num}",
            "next_tools": ["pptx_patch_shape", "pptx_add_bullet", "pptx_add_comment"]
        }

    def tool_pptx_delete_slide(
        self,
        file_path: str,
        slide_number: int,
        output_path: str = None
    ) -> dict[str, Any]:
        """Remove a slide from the presentation.

        USE THIS to remove template slides you don't need or to clean up
        unwanted content.

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number to delete
            output_path: Optional output path (defaults to overwriting input)

        Returns:
            Status with remaining slide count
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, load_error = _load_presentation(file_path)
        if load_error:
            return load_error

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        # Delete slide via XML manipulation
        slide_idx = slide_number - 1

        # Remove the slide's relationship from the presentation part
        slide_part = prs.slides[slide_idx].part
        for rId, rel in list(prs.part.rels.items()):
            if rel.target_part is slide_part:
                prs.part.drop_rel(rId)
                break

        # Remove from sldIdLst
        sldIdLst = prs.part._element.find(qn('p:sldIdLst'))
        if sldIdLst is not None:
            children = list(sldIdLst)
            if slide_idx < len(children):
                sldIdLst.remove(children[slide_idx])

        # Count remaining slides from sldIdLst (len(prs.slides) may be stale)
        remaining = len(list(sldIdLst)) if sldIdLst is not None else 0

        # Save
        save_path = output_path or resolved_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "deleted_slide": slide_number,
            "remaining_slides": remaining,
            "message": f"Deleted slide {slide_number}",
            "next_tools": ["pptx_list_slides"]
        }

    def tool_pptx_reorder_slides(
        self,
        file_path: str,
        new_order: list[int],
        output_path: str = None
    ) -> dict[str, Any]:
        """Change the order of slides in a presentation.

        Provide the complete slide order as a list. For example,
        [1, 3, 2, 4, 5] moves slide 3 before slide 2.

        Args:
            file_path: Path to the .pptx file
            new_order: Complete list of slide numbers in desired order (1-based)
            output_path: Optional output path (defaults to overwriting input)

        Returns:
            Status with new slide order
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, load_error = _load_presentation(file_path)
        if load_error:
            return load_error
        total_slides = len(prs.slides)

        # Validate new_order
        if len(new_order) != total_slides:
            return {"error": f"new_order must contain {total_slides} slide numbers"}

        if set(new_order) != set(range(1, total_slides + 1)):
            return {"error": f"new_order must contain each slide number 1-{total_slides} exactly once"}

        # Reorder via XML manipulation
        sldIdLst = prs.part._element.find(qn('p:sldIdLst'))
        if sldIdLst is not None:
            children = list(sldIdLst)

            # Clear and re-add in new order
            for child in children:
                sldIdLst.remove(child)

            for slide_num in new_order:
                sldIdLst.append(children[slide_num - 1])

        # Save
        save_path = output_path or resolved_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "new_order": new_order,
            "message": f"Reordered {total_slides} slides",
            "next_tools": ["pptx_list_slides"]
        }

    def tool_pptx_duplicate_slide(
        self,
        file_path: str,
        slide_number: int,
        position: str = "after",
        output_path: str = None
    ) -> dict[str, Any]:
        """Copy a slide including all shapes, tables, and formatting.

        USE THIS when you need multiple slides based on a template slide.
        Performs a full XML deep copy so tables, images, and other shapes
        are faithfully duplicated.

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number to copy
            position: 'after' (right after original) or 'end'
            output_path: Optional output path (defaults to overwriting input)

        Returns:
            Status with new slide number
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, load_error = _load_presentation(file_path)
        if load_error:
            return load_error

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        source_slide = prs.slides[slide_number - 1]
        position_text = str(position).strip().strip('"').strip("'").lower()
        if position_text not in {"after", "end"}:
            return {"error": "Invalid position. Use 'after' or 'end'."}

        target_slide = prs.slides.add_slide(source_slide.slide_layout)

        for shape in list(target_slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)

        for shape in source_slide.shapes:
            target_slide.shapes._spTree.insert_element_before(deepcopy(shape._element), 'p:extLst')

        notes_reltype = getattr(RT, "NOTES_SLIDE", None)
        for rel in source_slide.part.rels.values():
            if rel.reltype == RT.SLIDE_LAYOUT:
                continue
            if notes_reltype and rel.reltype == notes_reltype:
                continue
            try:
                target_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.rId, rel.is_external)
            except Exception:
                continue

        sldIdLst = prs.part._element.find(qn('p:sldIdLst'))
        new_slide_num = len(prs.slides)

        if position_text == "after" and sldIdLst is not None:
            children = list(sldIdLst)
            last_slide = children[-1]
            sldIdLst.remove(last_slide)
            sldIdLst.insert(slide_number, last_slide)
            new_slide_num = slide_number + 1

        # Save
        save_path = output_path or resolved_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "source_slide": slide_number,
            "new_slide_number": new_slide_num,
            "message": f"Duplicated slide {slide_number} to position {new_slide_num}",
            "next_tools": ["pptx_patch_shape", "pptx_get_slide"]
        }

    def tool_pptx_hide_slide(
        self,
        file_path: str,
        slide_number: int,
        hidden: bool = True,
        output_path: str = None
    ) -> dict[str, Any]:
        """Hide or unhide a slide in the presentation.

        Hidden slides are skipped during slideshow but remain editable.
        Use for backup content or speaker-only material.

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            hidden: True to hide, False to show again
            output_path: Optional output path (defaults to overwriting input)

        Returns:
            Status with visibility state
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]

        # Access the slide's XML element and set the show attribute
        # In OOXML, show="0" means hidden, show="1" or absent means visible
        slide_elem = slide._element
        if hidden:
            slide_elem.set('show', '0')
        else:
            # Remove the show attribute to make visible (default is shown)
            if 'show' in slide_elem.attrib:
                del slide_elem.attrib['show']

        # Save
        save_path = output_path or file_path
        safe_save_pptx(prs, save_path)

        action = "Hidden" if hidden else "Unhidden"
        return {
            "success": True,
            "file": save_path,
            "slide_number": slide_number,
            "hidden": hidden,
            "message": f"{action} slide {slide_number}",
            "next_tools": ["pptx_list_slides", "pptx_get_slide"]
        }

    def tool_pptx_get_hidden_slides(
        self,
        file_path: str
    ) -> dict[str, Any]:
        """Find all hidden slides in a presentation.

        Returns slide numbers and titles of slides that are hidden.
        Use pptx_hide_slide with hidden=False to unhide them.

        Args:
            file_path: Path to the .pptx file

        Returns:
            List of hidden slides with numbers and titles
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)
        hidden_slides = []

        for idx, slide in enumerate(prs.slides, 1):
            slide_elem = slide._element
            show_attr = slide_elem.get('show')
            if show_attr == '0':
                title = slide.shapes.title.text if slide.shapes.title else None
                hidden_slides.append({
                    "slide_number": idx,
                    "title": title,
                })

        return {
            "file": path.name,
            "total_slides": len(prs.slides),
            "hidden_count": len(hidden_slides),
            "hidden_slides": hidden_slides,
            "next_tools": ["pptx_hide_slide"] if hidden_slides else []
        }

    # =========================================================================
    # CONTENT PATCHING TOOLS
    # =========================================================================

    def tool_pptx_patch_shape(
        self,
        file_path: str,
        slide_number: int,
        shape_identifier: str,
        new_text: str = None,
        append: bool = False,
        autofit: bool = True,
        output_path: str = None
    ) -> dict[str, Any]:
        """Replace or update text in a shape on a slide.

        USE THIS to update titles, subtitles, or any text content.

        IMPORTANT: When patching content placeholders (body, Content Placeholder),
        do NOT include bullet characters (•, -, *, etc.) in text lines.
        PowerPoint content placeholders automatically render each line as a
        bullet point. Including bullet characters will cause duplication.

        For indented sub-bullets, use 4 spaces at the start of a line:
            "Top level item\\n    Indented sub-item\\nAnother top level"

        Each 4 spaces of indentation = 1 bullet level (max 8 levels).

        Common shape_identifier values:
        - 'title' - The slide title
        - 'subtitle' - Subtitle (on title slides)
        - 'body' - Main content area
        - Shape name from pptx_list_shapes (e.g., 'Content Placeholder 2')

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            shape_identifier: 'title', 'body', 'subtitle', or shape name
            new_text: The text to set (without bullet characters for body content)
            append: If True, add to existing text; if False, replace all
            autofit: If True (default), enable text auto-sizing to fit shape

        Returns:
            Status with old_text and new_text
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, load_error = _load_presentation(file_path)
        if load_error:
            return load_error

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]
        shape = _find_shape_by_identifier(slide, shape_identifier)

        if not shape:
            return {"error": f"Shape '{shape_identifier}' not found on slide {slide_number}"}

        if not shape.has_text_frame:
            return {"error": f"Shape '{shape_identifier}' does not have a text frame"}

        tf = shape.text_frame
        old_text = tf.text

        if new_text is not None:
            if append:
                # Append to existing text
                if tf.paragraphs:
                    last_para = tf.paragraphs[-1]
                    last_para.add_run().text = "\n" + new_text
                else:
                    tf.text = new_text
            else:
                # Replace text with proper paragraph handling
                # Split into lines and detect indentation levels
                lines = new_text.split('\n')

                # Clear existing paragraphs (keep the first one for structure)
                for p in list(tf.paragraphs)[1:]:
                    p._p.getparent().remove(p._p)

                # Process each line
                for i, line in enumerate(lines):
                    # Skip empty lines
                    if not line.strip():
                        continue

                    # Calculate indent level: 4 spaces = 1 level
                    stripped = line.lstrip(' ')
                    leading_spaces = len(line) - len(stripped)
                    level = min(leading_spaces // 4, 8)  # Max 8 levels

                    # Strip any bullet characters from the text
                    cleaned_text = stripped.lstrip('•●○◦▪▸►-–—* \t')

                    if i == 0:
                        # Use existing first paragraph
                        p = tf.paragraphs[0]
                        p.clear()
                        p.level = level
                        run = p.add_run()
                        run.text = cleaned_text
                    else:
                        # Add new paragraph
                        p = tf.add_paragraph()
                        p.level = level
                        run = p.add_run()
                        run.text = cleaned_text

            # Enable auto-fit to prevent text overflow
            if autofit:
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Save
        save_path = output_path or resolved_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "slide_number": slide_number,
            "shape_name": shape.name,
            "old_text": old_text[:100] + "..." if len(old_text) > 100 else old_text,
            "new_text": new_text[:100] + "..." if new_text and len(new_text) > 100 else new_text,
            "mode": "append" if append else "replace",
            "autofit": autofit,
            "message": f"Updated shape '{shape.name}' on slide {slide_number}",
            "next_tools": ["pptx_add_bullet", "pptx_add_comment", "pptx_get_slide"]
        }

    def tool_pptx_add_bullet(
        self,
        file_path: str,
        slide_number: int,
        text: str,
        shape_identifier: str = "body",
        level: int = 0,
        bold_label: str = None,
        output_path: str = None
    ) -> dict[str, Any]:
        """Add a bullet point to a slide's content area.

        USE THIS to add content bullets one at a time. For slides with
        existing template content, use pptx_clear_bullets first.

        IMPORTANT: Do NOT include bullet characters (•, -, *, etc.) in the text.
        PowerPoint content placeholders automatically render each line/paragraph
        as a bullet point. Including bullet characters will cause duplication.

        The bold_label parameter creates 'Label: text' formatting,
        useful for key-value style bullets like 'Duration: 18 months'.

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            text: The bullet text content (without bullet characters)
            shape_identifier: Target shape (default: 'body' for content area)
            level: Indentation level 0-8 (0=top level, 1=sub-bullet, etc.)
            bold_label: Optional bold prefix (e.g., 'Phase 1' → 'Phase 1: text')

        Returns:
            Status with bullet_count
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]
        shape = _find_shape_by_identifier(slide, shape_identifier)

        if not shape:
            return {"error": f"Shape '{shape_identifier}' not found on slide {slide_number}"}

        if not shape.has_text_frame:
            return {"error": f"Shape '{shape_identifier}' does not have a text frame"}

        tf = shape.text_frame

        # Check if first paragraph is empty
        if tf.paragraphs and not tf.paragraphs[0].text.strip():
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.level = min(max(level, 0), 8)  # Clamp to valid range

        # Strip any leading bullet characters from the text to avoid duplication
        # PowerPoint adds bullets automatically based on paragraph formatting
        cleaned_text = text.lstrip('•●○◦▪▸►-–—* \t')

        if bold_label:
            run = p.add_run()
            run.text = bold_label + ": "
            run.font.bold = True
            run = p.add_run()
            run.text = cleaned_text
        else:
            # Use add_run instead of setting p.text to preserve bullet formatting
            run = p.add_run()
            run.text = cleaned_text

        # Save
        save_path = output_path or file_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "slide_number": slide_number,
            "shape_name": shape.name,
            "bullet_count": len(tf.paragraphs),
            "level": level,
            "message": f"Added bullet to slide {slide_number}",
            "next_tools": ["pptx_add_bullet", "pptx_add_comment", "pptx_get_slide"]
        }

    def tool_pptx_clear_bullets(
        self,
        file_path: str,
        slide_number: int,
        shape_identifier: str = "body",
        output_path: str = None
    ) -> dict[str, Any]:
        """Clear all bullet content from a shape to start fresh.

        USE THIS before pptx_add_bullet when the template has placeholder
        content you want to replace entirely.

        Preserves text frame formatting properties (autofit, word wrap) while
        clearing the text content.

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            shape_identifier: Shape to clear (default: 'body')

        Returns:
            Status with number of paragraphs cleared
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]
        shape = _find_shape_by_identifier(slide, shape_identifier)

        if not shape:
            return {"error": f"Shape '{shape_identifier}' not found on slide {slide_number}"}

        if not shape.has_text_frame:
            return {"error": f"Shape '{shape_identifier}' does not have a text frame"}

        tf = shape.text_frame
        cleared_count = len(tf.paragraphs)

        # Capture properties before clearing
        props = _capture_text_frame_props(tf)

        # Clear the text frame
        tf.clear()

        # Restore text frame properties (autofit, word wrap)
        _restore_text_frame_props(tf, props)

        # Save
        save_path = output_path or file_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "slide_number": slide_number,
            "shape_name": shape.name,
            "cleared_paragraphs": cleared_count,
            "autofit_preserved": props.get("auto_size") is not None,
            "message": f"Cleared {cleared_count} paragraphs from slide {slide_number}",
            "next_tools": ["pptx_add_bullet", "pptx_patch_shape"]
        }

    # =========================================================================
    # TABLE TOOLS
    # =========================================================================

    def tool_pptx_get_table(
        self,
        file_path: str,
        slide_number: int,
        table_index: int = 0
    ) -> dict[str, Any]:
        """Read content from a table on a slide.

        Returns the table's headers and all data rows.
        USE THIS before modifying table content.

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            table_index: Which table if slide has multiple (0-based)

        Returns:
            Table with headers and rows as lists
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, load_error = _load_presentation(file_path)
        if load_error:
            return load_error

        path = Path(resolved_path)

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]
        table_shape, table = _find_table_in_slide(slide, str(table_index))

        if not table:
            return {
                "error": f"No table found on slide {slide_number}. "
                         f"Looked for table index {table_index}. "
                         "Hint: slide numbers may have changed if slides were reordered or deleted."
            }

        result = {
            "file": path.name,
            "slide_number": slide_number,
            "table_index": table_index,
            "shape_name": table_shape.name,
            "rows": len(table.rows),
            "columns": len(table.columns),
            "header": [],
            "data": [],
        }

        for row_idx, row in enumerate(table.rows):
            row_data = [cell.text.strip() for cell in row.cells]
            if row_idx == 0:
                result["header"] = row_data
            else:
                result["data"].append(row_data)

        return result

    def tool_pptx_add_table(
        self,
        file_path: str,
        slide_number: int,
        headers: list[str],
        rows: list[list[str]] = None,
        left: float = 1.0,
        top: float = 2.0,
        width: float = 11.0,
        height: float = 3.0,
        output_path: str = None
    ) -> dict[str, Any]:
        """Add a data table to a slide.

        Creates a table with column headers and optional data rows.
        Header row is bold. Position/size in inches (16:9 slide = 13.3" × 7.5").

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            headers: Column header names (e.g., ['Phase', 'Duration', 'Deliverables'])
            rows: Data rows as list of lists (optional)
            left, top: Position from top-left (default: 1.0", 2.0")
            width, height: Table size (default: 11.0" × 3.0")

        Returns:
            Status with table dimensions
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, load_error = _load_presentation(file_path)
        if load_error:
            return load_error

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]

        num_rows = 1 + (len(rows) if rows else 0)
        num_cols = len(headers)

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        table = table_shape.table

        # Set headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            # Make header bold
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True

        # Set row data
        if rows:
            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < num_cols:
                        table.cell(row_idx + 1, col_idx).text = str(cell_text)

        # Save
        save_path = output_path or resolved_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "slide_number": slide_number,
            "rows": num_rows,
            "columns": num_cols,
            "message": f"Added table with {num_rows} rows and {num_cols} columns",
            "next_tools": ["pptx_insert_table_row", "pptx_patch_table_cell", "pptx_add_comment"]
        }

    def tool_pptx_insert_table_row(
        self,
        file_path: str,
        slide_number: int,
        row_data: list[str],
        table_index: int = 0,
        position: str = "end",
        output_path: str = None
    ) -> dict[str, Any]:
        """Add a row to an existing table.

        Provide cell values matching the table columns.
        (Internally recreates the table - original formatting preserved.)

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            row_data: Cell values as list (e.g., ['Phase 3', '6 months', 'Optimization'])
            table_index: Which table if slide has multiple (0-based)
            position: 'end' (default), 'start', or row number

        Returns:
            Status with row position
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, load_error = _load_presentation(file_path)
        if load_error:
            return load_error

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]
        table_shape, table = _find_table_in_slide(slide, str(table_index))

        if not table:
            return {
                "error": f"No table found on slide {slide_number}. "
                         f"Looked for table index {table_index}. "
                         "Hint: slide numbers may have changed if slides were reordered or deleted."
            }

        # Extract current table data
        current_data = []
        for row in table.rows:
            current_data.append([cell.text for cell in row.cells])

        # Determine insert position
        if position == "end":
            insert_idx = len(current_data)
        else:
            insert_idx = min(int(position), len(current_data))

        # Insert new row
        current_data.insert(insert_idx, row_data)

        # Get table position and size
        left = table_shape.left
        top = table_shape.top
        width = table_shape.width
        height = table_shape.height

        # Remove old table shape
        sp = table_shape._element
        sp.getparent().remove(sp)

        # Create new table with updated data
        num_rows = len(current_data)
        num_cols = len(current_data[0]) if current_data else len(row_data)

        new_table_shape = slide.shapes.add_table(
            num_rows, num_cols, left, top, width, height
        )
        new_table = new_table_shape.table

        # Populate new table
        for row_idx, row in enumerate(current_data):
            for col_idx, cell_text in enumerate(row):
                if col_idx < num_cols:
                    new_table.cell(row_idx, col_idx).text = str(cell_text)
                    # Make header bold
                    if row_idx == 0:
                        for para in new_table.cell(row_idx, col_idx).text_frame.paragraphs:
                            for run in para.runs:
                                run.font.bold = True

        # Save
        save_path = output_path or resolved_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "slide_number": slide_number,
            "row_inserted_at": insert_idx,
            "total_rows": num_rows,
            "message": f"Inserted row at position {insert_idx}",
            "next_tools": ["pptx_insert_table_row", "pptx_patch_table_cell", "pptx_add_comment"]
        }

    def tool_pptx_patch_table_cell(
        self,
        file_path: str,
        slide_number: int,
        row_index: int,
        col_index: int,
        new_text: str,
        table_index: int = 0,
        output_path: str = None
    ) -> dict[str, Any]:
        """Update a single cell in a table.

        USE THIS for targeted cell edits. Row/column indices are 0-based
        (row 0 = header row).

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            row_index: 0-based row (0 = header)
            col_index: 0-based column
            new_text: Text to put in the cell
            table_index: Which table if slide has multiple (0-based)

        Returns:
            Status with old and new cell values
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, load_error = _load_presentation(file_path)
        if load_error:
            return load_error

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]
        table_shape, table = _find_table_in_slide(slide, str(table_index))

        if not table:
            return {
                "error": f"No table found on slide {slide_number}. "
                         f"Looked for table index {table_index}. "
                         "Hint: slide numbers may have changed if slides were reordered or deleted."
            }

        if row_index < 0 or row_index >= len(table.rows):
            return {"error": f"Row index {row_index} out of range (0-{len(table.rows)-1})"}

        if col_index < 0 or col_index >= len(table.columns):
            return {"error": f"Column index {col_index} out of range (0-{len(table.columns)-1})"}

        cell = table.cell(row_index, col_index)
        old_text = cell.text
        cell.text = new_text

        # Save
        save_path = output_path or resolved_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "slide_number": slide_number,
            "row_index": row_index,
            "col_index": col_index,
            "old_text": old_text,
            "new_text": new_text,
            "message": f"Updated cell [{row_index}, {col_index}]",
            "next_tools": ["pptx_get_table", "pptx_add_comment"]
        }

    # =========================================================================
    # NOTES AND COMMENTS TOOLS
    # =========================================================================

    def tool_pptx_set_notes(
        self,
        file_path: str,
        slide_number: int,
        notes_text: str,
        append: bool = False,
        output_path: str = None
    ) -> dict[str, Any]:
        """Set speaker notes for a slide.

        Speaker notes appear below the slide in Presenter View and can
        be printed as handouts. Use for talking points and context.

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            notes_text: The notes content (supports newlines)
            append: If True, add to existing notes; if False, replace
            output_path: Optional output path (defaults to overwriting input)

        Returns:
            Status with notes preview
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]
        notes_slide = slide.notes_slide

        # Try to get notes text frame - may be None if notes structure is incomplete
        notes_tf = notes_slide.notes_text_frame
        if notes_tf is None:
            # Look for any shape with a text frame
            for shape in notes_slide.shapes:
                if shape.has_text_frame:
                    notes_tf = shape.text_frame
                    break

        if notes_tf is None:
            # Notes slide has no text placeholder - we need to add one via XML
            # This happens with some templates that don't properly initialize notes
            from pptx.oxml import parse_xml
            from pptx.oxml.ns import qn

            # Create a body placeholder shape XML (mimics PowerPoint's notes structure)
            body_placeholder_xml = '''
<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
      xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="2" name="Notes Placeholder 2"/>
    <p:cNvSpPr>
      <a:spLocks noGrp="1"/>
    </p:cNvSpPr>
    <p:nvPr>
      <p:ph type="body" idx="1"/>
    </p:nvPr>
  </p:nvSpPr>
  <p:spPr/>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p>
      <a:endParaRPr/>
    </a:p>
  </p:txBody>
</p:sp>
'''
            sp_tree = notes_slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
            sp_elem = parse_xml(body_placeholder_xml)
            sp_tree.append(sp_elem)

            # Re-fetch the notes text frame after adding the placeholder
            notes_tf = notes_slide.notes_text_frame

        if notes_tf is None:
            return {"error": f"Could not create notes text frame on slide {slide_number}"}

        old_notes = notes_tf.text if notes_tf.text else ""

        if append and old_notes:
            notes_tf.text = old_notes + "\n\n" + notes_text
        else:
            notes_tf.text = notes_text

        # Save
        save_path = output_path or file_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "slide_number": slide_number,
            "mode": "append" if append else "replace",
            "notes_preview": notes_tf.text[:200] + "..." if len(notes_tf.text) > 200 else notes_tf.text,
            "message": f"Updated notes on slide {slide_number}",
            "next_tools": ["pptx_add_comment", "pptx_get_slide"]
        }

    def tool_pptx_get_notes(
        self,
        file_path: str,
        slide_number: int = None
    ) -> dict[str, Any]:
        """Read speaker notes from slides.

        USE THIS to see existing notes before updating them.
        Without slide_number, returns notes from ALL slides.

        Args:
            file_path: Path to the .pptx file
            slide_number: Specific slide (omit for all slides)

        Returns:
            Notes content for requested slide(s)
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)

        if slide_number is not None:
            if slide_number < 1 or slide_number > len(prs.slides):
                return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

            slide = prs.slides[slide_number - 1]
            notes_text = ""
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            return {
                "file": path.name,
                "slide_number": slide_number,
                "has_notes": bool(notes_text),
                "notes": notes_text,
            }

        # Get all notes
        all_notes = []
        for idx, slide in enumerate(prs.slides, 1):
            notes_text = ""
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            if notes_text:
                all_notes.append({
                    "slide_number": idx,
                    "notes": notes_text,
                })

        return {
            "file": path.name,
            "slides_with_notes": len(all_notes),
            "total_slides": len(prs.slides),
            "notes": all_notes,
        }

    # =========================================================================
    # PLACEHOLDER REPLACEMENT TOOLS
    # =========================================================================

    def tool_pptx_replace_text(
        self,
        file_path: str,
        find_text: str,
        replace_text: str,
        slide_number: int = None,
        output_path: str = None
    ) -> dict[str, Any]:
        """Find and replace text in shapes, tables, and notes.

        Replaces ALL occurrences of find_text. Use for single text
        replacements. For multiple placeholders, use pptx_replace_placeholders.

        Args:
            file_path: Path to the .pptx file
            find_text: Text to find (e.g., '<Customer Name>')
            replace_text: Replacement value (e.g., 'Contoso Ltd')
            slide_number: Specific slide only (omit for all slides)

        Returns:
            Status with replacement count
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, load_error = _load_presentation(file_path)
        if load_error:
            return load_error

        replacement_count = 0
        slides_modified = []

        slides_to_process = [prs.slides[slide_number - 1]] if slide_number else prs.slides

        for slide_idx, slide in enumerate(slides_to_process):
            actual_slide_num = slide_number if slide_number else slide_idx + 1
            slide_replaced = False

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if find_text in run.text:
                                run.text = run.text.replace(find_text, replace_text)
                                replacement_count += 1
                                slide_replaced = True

                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            # Use run-level replacement to preserve formatting
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    if find_text in run.text:
                                        run.text = run.text.replace(find_text, replace_text)
                                        replacement_count += 1
                                        slide_replaced = True

            # Also check notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                for para in notes_tf.paragraphs:
                    for run in para.runs:
                        if find_text in run.text:
                            run.text = run.text.replace(find_text, replace_text)
                            replacement_count += 1
                            slide_replaced = True

            if slide_replaced:
                slides_modified.append(actual_slide_num)

        # Save
        save_path = output_path or resolved_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "find_text": find_text,
            "replace_text": replace_text,
            "replacement_count": replacement_count,
            "slides_modified": slides_modified,
            "message": f"Replaced {replacement_count} occurrences across {len(slides_modified)} slide(s)",
            "next_tools": ["pptx_replace_placeholders", "pptx_audit_placeholders"]
        }

    def tool_pptx_replace_placeholders(
        self,
        file_path: str,
        replacements: dict[str, str],
        output_path: str = None
    ) -> dict[str, Any]:
        """Replace placeholder text throughout a presentation.

        USE THIS after copying a template to fill in standard placeholders
        like <Customer Name>, <Project Name>, [Date], etc.

        Searches ALL slides, shapes, tables, and notes for each placeholder
        and replaces with the specified value.

        Args:
            file_path: Path to the .pptx file
            replacements: Dictionary of placeholder → value mappings:
                {
                    '<Customer Name>': 'Contoso Ltd',
                    '<Project Name>': 'Customer Care Transformation',
                    '[Date]': 'January 2026'
                }

        Returns:
            Status with counts of replacements per placeholder
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, load_error = _load_presentation(file_path)
        if load_error:
            return load_error

        counts = dict.fromkeys(replacements.keys(), 0)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            for find_text, replace_text in replacements.items():
                                if find_text in run.text:
                                    run.text = run.text.replace(find_text, replace_text)
                                    counts[find_text] += 1

                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            # Use run-level replacement to preserve formatting
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    for find_text, replace_text in replacements.items():
                                        if find_text in run.text:
                                            run.text = run.text.replace(find_text, replace_text)
                                            counts[find_text] += 1

            # Also process notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                for para in notes_tf.paragraphs:
                    for run in para.runs:
                        for find_text, replace_text in replacements.items():
                            if find_text in run.text:
                                run.text = run.text.replace(find_text, replace_text)
                                counts[find_text] += 1

        total_replacements = sum(counts.values())

        # Save
        save_path = output_path or resolved_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "total_replacements": total_replacements,
            "by_placeholder": counts,
            "message": f"Replaced {total_replacements} placeholders",
            "next_tools": ["pptx_audit_placeholders", "pptx_add_comment", "pptx_list_slides"]
        }

    # =========================================================================
    # FORMATTING TOOLS
    # =========================================================================

    def tool_pptx_set_text_autofit(
        self,
        file_path: str,
        slide_number: int,
        shape_identifier: str,
        autofit_type: str = "shrink",
        output_path: str = None
    ) -> dict[str, Any]:
        """Configure how text fits within a shape.

        Options:
        - 'shrink': Shrink text to fit (recommended for templates)
        - 'none': No auto-fit (text may overflow)
        - 'resize': Resize shape to fit text

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            shape_identifier: 'title', 'body', or shape name
            autofit_type: 'shrink' (default), 'none', or 'resize'

        Returns:
            Status with applied setting
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)

        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides."}

        slide = prs.slides[slide_number - 1]
        shape = _find_shape_by_identifier(slide, shape_identifier)

        if not shape:
            return {"error": f"Shape '{shape_identifier}' not found on slide {slide_number}"}

        if not shape.has_text_frame:
            return {"error": f"Shape '{shape_identifier}' does not have a text frame"}

        autofit_map = {
            "none": MSO_AUTO_SIZE.NONE,
            "shrink": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
            "resize": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
        }

        if autofit_type.lower() not in autofit_map:
            return {"error": "Invalid autofit_type. Use: none, shrink, or resize"}

        shape.text_frame.auto_size = autofit_map[autofit_type.lower()]

        # Save
        save_path = output_path or file_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "slide_number": slide_number,
            "shape_name": shape.name,
            "autofit_type": autofit_type,
            "message": f"Set auto-fit to '{autofit_type}' on slide {slide_number}",
            "next_tools": ["pptx_patch_shape", "pptx_get_slide"]
        }

    # =========================================================================
    # AUDIT TOOLS
    # =========================================================================

    def tool_pptx_audit_placeholders(
        self,
        file_path: str,
        patterns: list[str] = None
    ) -> dict[str, Any]:
        """Scan presentation for unfilled placeholders.

        Finds <...>, [...], [TBD], and similar patterns in all slides.
        USE THIS before delivery to ensure no placeholders remain.

        Args:
            file_path: Path to the .pptx file
            patterns: Custom regex patterns (optional)

        Returns:
            Audit report with placeholders found and locations
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        default_patterns = [
            r'<[^>]+>',           # <Customer Name>
            r'\[TBD\]',           # [TBD]
            r'\[insert[^\]]*\]',  # [insert ...]
            r'\[TODO[^\]]*\]',    # [TODO ...]
        ]

        search_patterns = patterns or default_patterns

        prs = Presentation(file_path)

        findings = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    for pattern in search_patterns:
                        matches = re.findall(pattern, text, re.IGNORECASE)
                        for match in matches:
                            findings.append({
                                "slide": slide_idx,
                                "shape": shape.name,
                                "placeholder": match,
                            })

                if shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            for pattern in search_patterns:
                                matches = re.findall(pattern, cell.text, re.IGNORECASE)
                                for match in matches:
                                    findings.append({
                                        "slide": slide_idx,
                                        "shape": shape.name,
                                        "location": f"table[{row_idx},{col_idx}]",
                                        "placeholder": match,
                                    })

        # Group by placeholder
        by_placeholder = {}
        for f in findings:
            ph = f["placeholder"]
            if ph not in by_placeholder:
                by_placeholder[ph] = []
            by_placeholder[ph].append(f)

        status = "CLEAN" if not findings else "PLACEHOLDERS_FOUND"

        return {
            "file": path.name,
            "status": status,
            "total_placeholders": len(findings),
            "unique_placeholders": list(by_placeholder.keys()),
            "findings": findings[:50],  # Limit output
            "by_placeholder": {k: len(v) for k, v in by_placeholder.items()},
            "next_tools": ["pptx_replace_placeholders", "pptx_replace_text"] if findings else ["pptx_add_comment"]
        }

    # =========================================================================
    # COMMENT TOOLS
    # =========================================================================

    def tool_pptx_add_comment(
        self,
        file_path: str,
        slide_number: int,
        comment_text: str,
        x_inches: float = None,
        y_inches: float = None,
        author: str = None,
        output_path: str = None
    ) -> dict[str, Any]:
        """Add a review comment to a slide.

        Comments appear as markers that can be viewed in PowerPoint's
        Comments pane. Use for review notes, context, or feedback.

        Position (x_inches, y_inches) is from top-left corner.
        For 16:9 slides: width ~13.3", height ~7.5".

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            comment_text: The comment content
            x_inches: Position from left edge (default: 1.0)
            y_inches: Position from top edge (default: 1.0)
            author: Comment author name (default: 'Solution Architect Agent')

        Returns:
            Status with comment details
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        author = author or getattr(self, "_comment_author", DEFAULT_AUTHOR)
        initials = getattr(self, "_comment_initials", _derive_initials(author))
        save_path = output_path or file_path

        # Use default positions if not specified
        # 16:9 slide is 13.333" x 7.5"
        if x_inches is None:
            x_inches = 1.0  # 1 inch from left
        if y_inches is None:
            y_inches = 1.0  # 1 inch from top

        try:
            result = _add_pptx_comment(
                file_path, save_path, slide_number,
                comment_text, x_inches, y_inches, author, initials
            )
            return result
        except Exception as e:
            return {"error": f"Failed to add comment: {str(e)}"}

    def tool_pptx_get_comments(
        self,
        file_path: str,
        slide_number: int = None
    ) -> dict[str, Any]:
        """Read review comments from slides.

        Returns all comments with author, position, and text.
        Without slide_number, returns comments from ALL slides.

        Args:
            file_path: Path to the .pptx file
            slide_number: Specific slide only (omit for all)

        Returns:
            Comments grouped by slide
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        try:
            result = _get_pptx_comments(file_path, slide_number)
            return result
        except Exception as e:
            return {"error": f"Failed to get comments: {str(e)}"}

    def tool_pptx_delete_comment(
        self,
        file_path: str,
        slide_number: int,
        comment_index: int | None = None,
        output_path: str | None = None,
    ) -> dict[str, Any]:
        """Delete comments from a slide.

        Deletes one comment (by index) or all comments on a slide when
        comment_index is not provided.

        Args:
            file_path: Path to the .pptx file
            slide_number: 1-based slide number
            comment_index: Optional comment index on that slide
            output_path: Optional output path (defaults to overwriting input)

        Returns:
            Status with deletion details
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        save_path = output_path or file_path
        try:
            return _delete_pptx_comment(
                file_path=file_path,
                output_path=save_path,
                slide_number=slide_number,
                comment_index=comment_index,
            )
        except Exception as e:
            return {"error": f"Failed to delete comment: {str(e)}"}

    # =========================================================================
    # LAYOUT RECOMMENDATION TOOLS
    # =========================================================================

    def tool_pptx_analyze_layouts(
        self,
        file_path: str
    ) -> dict[str, Any]:
        """Get detailed information about all slide layouts.

        Returns placeholder types, positions, and recommended uses for
        each layout in the presentation template.

        USE THIS when you need to understand what layouts are available
        before creating new slides.

        Args:
            file_path: Path to the .pptx file

        Returns:
            Layout analysis with classifications and recommendations
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)
        layouts = []

        for idx, layout in enumerate(prs.slide_layouts):
            placeholders = []
            for ph in layout.placeholders:
                ph_type = str(ph.placeholder_format.type).replace("PP_PLACEHOLDER.", "")
                placeholders.append({
                    "idx": ph.placeholder_format.idx,
                    "type": ph_type,
                    "name": ph.name,
                    "left": ph.left,
                    "top": ph.top,
                    "width": ph.width,
                    "height": ph.height,
                })

            # Classify the layout
            ph_types = [p["type"] for p in placeholders]
            classification = _classify_layout(ph_types, layout.name)

            layouts.append({
                "index": idx,
                "name": layout.name,
                "classification": classification,
                "placeholder_count": len(placeholders),
                "placeholders": placeholders,
                "recommended_for": _get_layout_recommendations(classification),
            })

        return {
            "file": path.name,
            "layout_count": len(layouts),
            "layouts": layouts,
            "next_tools": ["pptx_add_slide", "pptx_recommend_layout", "pptx_copy_template"]
        }

    def tool_pptx_recommend_layout(
        self,
        file_path: str,
        content_type: str
    ) -> dict[str, Any]:
        """Get the best layout for a specific content type.

        Content types:
        - 'title': Section title slide
        - 'bullets': Bullet point list
        - 'two_column': Side-by-side content
        - 'comparison': Comparison with headers
        - 'image': Image with caption
        - 'table': Table-heavy content
        - 'blank': Custom content

        Args:
            file_path: Path to the .pptx file
            content_type: What you want to show on the slide

        Returns:
            Recommended layout_index and alternatives
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        prs = Presentation(file_path)

        # Map content types to ideal layout characteristics
        content_preferences = {
            "title": ["title_slide", "section_header"],
            "bullets": ["title_and_content", "title_and_vertical_text"],
            "two_column": ["two_content", "comparison"],
            "comparison": ["comparison", "two_content"],
            "image": ["picture_with_caption", "content_with_caption", "title_only"],
            "table": ["title_and_content", "title_only", "blank"],
            "blank": ["blank", "title_only"],
        }

        preferred = content_preferences.get(content_type.lower(), ["title_and_content"])

        # Score each layout
        scored_layouts = []
        for idx, layout in enumerate(prs.slide_layouts):
            ph_types = [str(ph.placeholder_format.type).replace("PP_PLACEHOLDER.", "")
                       for ph in layout.placeholders]
            classification = _classify_layout(ph_types, layout.name)

            score = 0
            if classification in preferred:
                score = len(preferred) - preferred.index(classification)

            scored_layouts.append({
                "index": idx,
                "name": layout.name,
                "classification": classification,
                "score": score,
            })

        # Sort by score descending
        scored_layouts.sort(key=lambda x: x["score"], reverse=True)

        recommended = scored_layouts[0] if scored_layouts else None
        alternatives = scored_layouts[1:4] if len(scored_layouts) > 1 else []

        return {
            "file": path.name,
            "content_type": content_type,
            "recommended": recommended,
            "alternatives": alternatives,
            "usage": f"pptx_add_slide(file_path='{file_path}', layout_index={recommended['index'] if recommended else 1}, title='Your Title')",
            "next_tools": ["pptx_add_slide"]
        }

    # =========================================================================
    # CHANGE LOG TOOLS (for auditability without track changes)
    # =========================================================================

    def tool_pptx_log_changes(
        self,
        file_path: str,
        changes: list[ChangeLogEntry],
        output_path: str = None
    ) -> dict[str, Any]:
        """Append change log entries to the first slide's notes.

        Use for audit trail since PowerPoint doesn't have track changes.
        Changes are appended to the notes of slide 1 with timestamps.

        Args:
            file_path: Path to the .pptx file
            changes: List of change entries, each with slide (number), action (what was done), and detail (specifics)
            output_path: Optional output path (defaults to overwriting input)

        Returns:
            Status with changes logged count
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        path = Path(file_path)
        if not path.exists():
            return {"error": f"File not found: {file_path}"}

        from datetime import datetime

        prs = Presentation(file_path)

        if len(prs.slides) == 0:
            return {"error": "Presentation has no slides"}

        # Get first slide
        first_slide = prs.slides[0]

        # Get or create notes slide
        if not first_slide.has_notes_slide:
            _ = first_slide.notes_slide  # This creates the notes slide

        notes_slide = first_slide.notes_slide
        notes_tf = notes_slide.notes_text_frame

        # Add timestamp header
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
        if notes_tf.paragraphs and not notes_tf.paragraphs[0].text.strip():
            p = notes_tf.paragraphs[0]
        else:
            p = notes_tf.add_paragraph()
        p.text = f"--- Changes recorded {timestamp} ---"

        # Add each change
        for change in changes:
            p = notes_tf.add_paragraph()
            slide_ref = f"Slide {change.get('slide', '?')}"
            action = change.get('action', 'Modified')
            detail = change.get('detail', '')
            p.text = f"{slide_ref}: {action} - {detail}"

        # Save
        save_path = output_path or file_path
        safe_save_pptx(prs, save_path)

        return {
            "success": True,
            "file": save_path,
            "slide": 1,
            "changes_logged": len(changes),
            "message": f"Logged {len(changes)} changes to slide 1 notes",
            "next_tools": ["pptx_get_notes", "pptx_list_slides", "pptx_audit_placeholders"]
        }


# =============================================================================
# HELPER FUNCTIONS FOR COMMENTS
# =============================================================================

def _add_pptx_comment(
    file_path,
    output_path,
    slide_number,
    comment_text,
    x_inches,
    y_inches,
    author,
    initials=None,
    parent_id=None,
):
    """Add a modern (Office 365) comment to a PPTX file."""
    import tempfile
    import uuid
    import zipfile
    from datetime import datetime
    from posixpath import dirname, join, normpath

    from lxml import etree

    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    P188_NS = 'http://schemas.microsoft.com/office/powerpoint/2018/8/main'
    AC_NS = 'http://schemas.microsoft.com/office/drawing/2013/main/command'
    PC_NS = 'http://schemas.microsoft.com/office/powerpoint/2013/main/command'
    MODERN_COMMENTS_REL = 'http://schemas.microsoft.com/office/2018/10/relationships/comments'
    MODERN_COMMENTS_CT = 'application/vnd.ms-powerpoint.comments+xml'

    def inches_to_emu(inches):
        return int(inches * 914400)

    def _to_int(value):
        try:
            return int(str(value))
        except (TypeError, ValueError):
            return None

    def _resolve_relationship_target(rels_part, target):
        rels = str(rels_part)
        if '/_rels/' in rels and rels.endswith('.rels'):
            source_part = rels.replace('/_rels/', '/', 1)[:-5]
            base_dir = dirname(source_part)
            return normpath(join(base_dir, str(target or '')))
        return normpath(join(dirname(rels), str(target or '')))

    def _next_rel_id(rels_root):
        max_id = 0
        for rel in rels_root.findall(f'.//{{*}}Relationship'):
            rid = rel.get('Id', '')
            if rid.startswith('rId'):
                with suppress(ValueError):
                    max_id = max(max_id, int(rid[3:]))
        return f'rId{max_id + 1}'

    def _first_shape_id(slide_xml_bytes):
        try:
            root = etree.fromstring(slide_xml_bytes)
            node = root.find(f'.//{{{P_NS}}}sp/{{{P_NS}}}nvSpPr/{{{P_NS}}}cNvPr')
            if node is not None and node.get('id') is not None:
                return str(node.get('id'))
        except Exception:
            pass
        return '1'

    author = author or DEFAULT_AUTHOR
    initials = (initials or '').strip().upper() or _derive_initials(author)

    with zipfile.ZipFile(file_path, 'r') as zf_in:
        names = set(zf_in.namelist())
        slide_part = f'ppt/slides/slide{slide_number}.xml'
        slide_rels_part = f'ppt/slides/_rels/slide{slide_number}.xml.rels'
        if slide_part not in names:
            return {'error': f'Invalid slide_number: {slide_number}'}
        if slide_rels_part not in names:
            return {'error': f'Missing relationships for slide {slide_number}'}

        slide_id = None
        if 'ppt/presentation.xml' in names and 'ppt/_rels/presentation.xml.rels' in names:
            pres_root = etree.fromstring(zf_in.read('ppt/presentation.xml'))
            pres_rels_root = etree.fromstring(zf_in.read('ppt/_rels/presentation.xml.rels'))
            rel_target_by_id = {}
            for rel in pres_rels_root.findall(f'.//{{*}}Relationship'):
                rel_id = rel.get('Id')
                rel_type = rel.get('Type', '')
                if rel_id and rel_type.endswith('/slide'):
                    rel_target_by_id[rel_id] = normpath(join('ppt', rel.get('Target', '')))
            sld_ids = pres_root.findall(f'.//{{{P_NS}}}sldIdLst/{{{P_NS}}}sldId')
            for sld in sld_ids:
                rel_id = sld.get(f'{{{R_NS}}}id')
                if rel_id and rel_target_by_id.get(rel_id) == slide_part:
                    slide_id = _to_int(sld.get('id'))
                    break
        if slide_id is None:
            slide_id = 0x7FFFFFFF + slide_number

        slide_rels_root = etree.fromstring(zf_in.read(slide_rels_part))
        modern_part = None
        for rel in slide_rels_root.findall(f'.//{{*}}Relationship'):
            if rel.get('Type') == MODERN_COMMENTS_REL:
                modern_part = _resolve_relationship_target(slide_rels_part, rel.get('Target'))
                break
        if not modern_part:
            modern_part = f'ppt/comments/modernComment_{slide_id:08X}_0.xml'
            rel_ns = slide_rels_root.nsmap.get(None, 'http://schemas.openxmlformats.org/package/2006/relationships')
            etree.SubElement(
                slide_rels_root,
                f'{{{rel_ns}}}Relationship',
                Id=_next_rel_id(slide_rels_root),
                Type=MODERN_COMMENTS_REL,
                Target=f'../comments/{Path(modern_part).name}',
            )

        if 'ppt/authors.xml' in names:
            authors_root = etree.fromstring(zf_in.read('ppt/authors.xml'))
        else:
            authors_root = etree.Element(f'{{{P188_NS}}}authorLst', nsmap={'p188': P188_NS})

        author_id = None
        for node in authors_root.findall(f'.//{{{P188_NS}}}author'):
            if (node.get('name') or '').strip() == author:
                author_id = node.get('id')
                break
        if not author_id:
            author_id = '{' + str(uuid.uuid4()).upper() + '}'
            etree.SubElement(
                authors_root,
                f'{{{P188_NS}}}author',
                id=author_id,
                name=author,
                initials=initials,
                userId=os.environ.get('MCP_AUTHOR_IDENTITY') or f'S::{author.lower().replace(" ", ".")}',
                providerId='AD',
            )

        if modern_part in names:
            cm_root = etree.fromstring(zf_in.read(modern_part))
        else:
            cm_root = etree.Element(
                f'{{{P188_NS}}}cmLst',
                nsmap={'a': A_NS, 'ac': AC_NS, 'pc': PC_NS, 'p188': P188_NS},
            )

        cm = etree.SubElement(
            cm_root,
            f'{{{P188_NS}}}cm',
            id='{' + str(uuid.uuid4()).upper() + '}',
            authorId=author_id,
            created=datetime.now().isoformat(timespec='milliseconds'),
        )
        de = etree.SubElement(cm, f'{{{AC_NS}}}deMkLst')
        etree.SubElement(de, f'{{{PC_NS}}}docMk')
        etree.SubElement(de, f'{{{PC_NS}}}sldMk', cId='0', sldId=str(slide_id))
        etree.SubElement(de, f'{{{AC_NS}}}spMk', id=_first_shape_id(zf_in.read(slide_part)))

        pos = etree.SubElement(cm, f'{{{P188_NS}}}pos')
        pos.set('x', str(inches_to_emu(x_inches)))
        pos.set('y', str(inches_to_emu(y_inches)))

        tx_body = etree.SubElement(cm, f'{{{P188_NS}}}txBody')
        etree.SubElement(tx_body, f'{{{A_NS}}}bodyPr')
        etree.SubElement(tx_body, f'{{{A_NS}}}lstStyle')
        para = etree.SubElement(tx_body, f'{{{A_NS}}}p')
        run = etree.SubElement(para, f'{{{A_NS}}}r')
        etree.SubElement(run, f'{{{A_NS}}}rPr', lang='en-US')
        text_elem = etree.SubElement(run, f'{{{A_NS}}}t')
        text_elem.text = comment_text

        ct_root = etree.fromstring(zf_in.read('[Content_Types].xml'))
        part_name = f'/{modern_part}'
        if not any(node.get('PartName') == part_name for node in ct_root):
            etree.SubElement(
                ct_root,
                'Override',
                PartName=part_name,
                ContentType=MODERN_COMMENTS_CT,
            )

        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            tmp_path = tmp.name
        with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zf_out:
            for item in zf_in.namelist():
                if item == '[Content_Types].xml':
                    zf_out.writestr(item, etree.tostring(ct_root, xml_declaration=True, encoding='UTF-8', standalone='yes'))
                elif item == slide_rels_part:
                    zf_out.writestr(item, etree.tostring(slide_rels_root, xml_declaration=True, encoding='UTF-8', standalone='yes'))
                elif item == 'ppt/authors.xml':
                    zf_out.writestr(item, etree.tostring(authors_root, xml_declaration=True, encoding='UTF-8', standalone='yes'))
                elif item == modern_part:
                    zf_out.writestr(item, etree.tostring(cm_root, xml_declaration=True, encoding='UTF-8', standalone='yes'))
                else:
                    zf_out.writestr(item, zf_in.read(item))
            if 'ppt/authors.xml' not in names:
                zf_out.writestr('ppt/authors.xml', etree.tostring(authors_root, xml_declaration=True, encoding='UTF-8', standalone='yes'))
            if modern_part not in names:
                zf_out.writestr(modern_part, etree.tostring(cm_root, xml_declaration=True, encoding='UTF-8', standalone='yes'))

    # Move temp file to output
    shutil.move(tmp_path, output_path)

    return {
        "success": True,
        "file": output_path,
        "slide_number": slide_number,
        "comment": comment_text[:100] + "..." if len(comment_text) > 100 else comment_text,
        "position": {"x_inches": x_inches, "y_inches": y_inches},
        "author": author,
        "format": "modern",
        "message": f"Added comment to slide {slide_number}",
        "next_tools": ["pptx_add_comment", "pptx_get_comments", "pptx_list_slides"]
    }


def _get_pptx_comments(file_path, slide_number=None):
    """Get comments from a PPTX file by reading the ZIP package.

    Supports both legacy PresentationML comments (`commentN.xml`) and
    Office 365 modern comments (`modernComment_*.xml`).
    """
    import zipfile
    from posixpath import dirname, join, normpath

    from lxml import etree

    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    P188_NS = 'http://schemas.microsoft.com/office/powerpoint/2018/8/main'
    AC_NS = 'http://schemas.microsoft.com/office/drawing/2013/main/command'
    PC_NS = 'http://schemas.microsoft.com/office/powerpoint/2013/main/command'

    LEGACY_COMMENTS_REL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments'
    MODERN_COMMENTS_REL = 'http://schemas.microsoft.com/office/2018/10/relationships/comments'

    def emu_to_inches(emu):
        # 1 inch = 914400 EMU
        return round(emu / 914400, 2)

    def _to_int(value):
        try:
            return int(str(value))
        except (TypeError, ValueError):
            return None

    def _resolve_relationship_target(rels_part, target):
        # OOXML targets are relative to the source part, not the _rels folder.
        rels = str(rels_part)
        if '/_rels/' in rels and rels.endswith('.rels'):
            source_part = rels.replace('/_rels/', '/', 1)[:-5]
            base_dir = dirname(source_part)
            return normpath(join(base_dir, str(target or '')))
        return normpath(join(dirname(rels), str(target or '')))

    def _extract_modern_text(cm_elem):
        parts = []
        for text_elem in cm_elem.findall(f'.//{{{A_NS}}}t'):
            if text_elem.text:
                parts.append(text_elem.text)
        return ''.join(parts)

    def _extract_modern_anchor(cm_elem):
        anchor = {}

        sld_mk = cm_elem.find(f'.//{{{PC_NS}}}sldMk')
        if sld_mk is not None and sld_mk.get('sldId') is not None:
            anchor['slide_id'] = sld_mk.get('sldId')

        tc_mk = cm_elem.find(f'.//{{{AC_NS}}}tcMk')
        if tc_mk is not None:
            anchor['type'] = 'table_cell'
            if tc_mk.get('rowId') is not None:
                anchor['rowId'] = tc_mk.get('rowId')
            if tc_mk.get('colId') is not None:
                anchor['colId'] = tc_mk.get('colId')

            frame_mk = cm_elem.find(f'.//{{{AC_NS}}}graphicFrameMk')
            if frame_mk is not None and frame_mk.get('id') is not None:
                anchor['shapeId'] = frame_mk.get('id')
            return anchor

        sp_mk = cm_elem.find(f'.//{{{AC_NS}}}spMk')
        if sp_mk is not None:
            anchor['type'] = 'shape'
            if sp_mk.get('id') is not None:
                anchor['shapeId'] = sp_mk.get('id')
            return anchor

        if anchor:
            anchor['type'] = 'unknown'
        return anchor

    comments_by_slide = {}
    authors = {}
    modern_authors = {}
    slide_part_to_number = {}
    slide_id_to_number = {}
    modern_part_to_slide = {}

    with zipfile.ZipFile(file_path, 'r') as zf:
        # Build slide part -> slide number map from presentation relationships.
        if 'ppt/presentation.xml' in zf.namelist() and 'ppt/_rels/presentation.xml.rels' in zf.namelist():
            presentation_root = etree.fromstring(zf.read('ppt/presentation.xml'))
            presentation_rels_root = etree.fromstring(zf.read('ppt/_rels/presentation.xml.rels'))

            rel_target_by_id = {}
            for rel in presentation_rels_root.findall(f'.//{{*}}Relationship'):
                rel_id = rel.get('Id')
                rel_type = rel.get('Type', '')
                target = rel.get('Target', '')
                if rel_id and rel_type.endswith('/slide'):
                    rel_target_by_id[rel_id] = normpath(join('ppt', target))

            sld_id_elements = presentation_root.findall(f'.//{{{P_NS}}}sldIdLst/{{{P_NS}}}sldId')
            for index, sld_id_elem in enumerate(sld_id_elements, start=1):
                rel_id = sld_id_elem.get(f'{{{R_NS}}}id')
                sld_id = sld_id_elem.get('id')
                if rel_id and rel_id in rel_target_by_id:
                    slide_part_to_number[rel_target_by_id[rel_id]] = index
                sld_id_int = _to_int(sld_id)
                if sld_id_int is not None:
                    slide_id_to_number[sld_id_int] = index

        # Build comment-part -> slide number map from slide relationship files.
        for item in zf.namelist():
            if not item.startswith('ppt/slides/_rels/slide') or not item.endswith('.xml.rels'):
                continue

            match = re.search(r'slide(\d+)\.xml\.rels$', item)
            if not match:
                continue
            rel_slide_number = int(match.group(1))

            rels_root = etree.fromstring(zf.read(item))
            for rel in rels_root.findall(f'.//{{*}}Relationship'):
                rel_type = rel.get('Type', '')
                if rel_type not in {LEGACY_COMMENTS_REL, MODERN_COMMENTS_REL}:
                    continue
                resolved = _resolve_relationship_target(item, rel.get('Target'))
                modern_part_to_slide[resolved] = rel_slide_number

        # Read authors
        if 'ppt/commentAuthors.xml' in zf.namelist():
            authors_xml = zf.read('ppt/commentAuthors.xml')
            root = etree.fromstring(authors_xml)
            for author in root:
                author_id = author.get('id')
                author_name = author.get('name')
                authors[author_id] = author_name

        # Read modern authors
        if 'ppt/authors.xml' in zf.namelist():
            authors_xml = zf.read('ppt/authors.xml')
            root = etree.fromstring(authors_xml)
            for author in root.findall(f'.//{{{P188_NS}}}author'):
                author_id = author.get('id')
                if not author_id:
                    continue
                modern_authors[author_id] = {
                    'name': author.get('name') or author_id,
                    'initials': author.get('initials', ''),
                    'userId': author.get('userId'),
                    'providerId': author.get('providerId'),
                }

        # Read legacy comments
        for item in zf.namelist():
            if item.startswith('ppt/comments/comment') and item.endswith('.xml'):
                # Extract slide number from filename
                match = re.search(r'comment(\d+)\.xml', item)
                if match:
                    slide_num = int(match.group(1))

                    if slide_number is not None and slide_num != slide_number:
                        continue

                    comments_xml = zf.read(item)
                    root = etree.fromstring(comments_xml)

                    slide_comments = []
                    for cm in root:
                        author_id = cm.get('authorId', '0')
                        dt = cm.get('dt', '')
                        idx = cm.get('idx', '')

                        pos = cm.find(f'{{{P_NS}}}pos')
                        x = emu_to_inches(int(pos.get('x', 0))) if pos is not None else 0
                        y = emu_to_inches(int(pos.get('y', 0))) if pos is not None else 0

                        text = cm.find(f'{{{P_NS}}}text')
                        comment_text = text.text if text is not None else ''

                        slide_comments.append({
                            "index": idx,
                            "author": authors.get(author_id, f"Author {author_id}"),
                            "datetime": dt,
                            "position": {"x_inches": x, "y_inches": y},
                            "text": comment_text,
                        })

                    comments_by_slide[slide_num] = slide_comments

        # Read modern Office 365 comments
        for item in zf.namelist():
            if not item.startswith('ppt/comments/modernComment_') or not item.endswith('.xml'):
                continue

            modern_root = etree.fromstring(zf.read(item))
            modern_comments = modern_root.findall(f'.//{{{P188_NS}}}cm')
            if not modern_comments:
                continue

            # Resolve slide number from rels first.
            slide_num = modern_part_to_slide.get(item)

            # Fallback: infer from first slide marker in this comments part.
            if slide_num is None:
                marker = modern_root.find(f'.//{{{PC_NS}}}sldMk')
                if marker is not None:
                    marker_id = _to_int(marker.get('sldId'))
                    if marker_id is not None:
                        slide_num = slide_id_to_number.get(marker_id)

            # Last resort: infer from filename hex token (maps to slide id in some decks).
            if slide_num is None:
                m_hex = re.search(r'modernComment_([0-9A-Fa-f]+)_\d+\.xml$', item)
                if m_hex:
                    with suppress(ValueError):
                        slide_num = slide_id_to_number.get(int(m_hex.group(1), 16))

            # If slide cannot be resolved, keep it in bucket 0 for visibility.
            if slide_num is None:
                slide_num = 0

            if slide_number is not None and slide_num != slide_number:
                continue

            slide_comments = comments_by_slide.setdefault(slide_num, [])
            for cm in modern_comments:
                author_id = cm.get('authorId')
                author_info = modern_authors.get(author_id, {})

                pos = cm.find(f'{{{P188_NS}}}pos')
                x = emu_to_inches(int(pos.get('x', 0))) if pos is not None else 0
                y = emu_to_inches(int(pos.get('y', 0))) if pos is not None else 0

                comment_text = _extract_modern_text(cm)
                anchor = _extract_modern_anchor(cm)
                comment_id = cm.get('id', '')
                created = cm.get('created', '')

                slide_comments.append({
                    'index': str(len(slide_comments) + 1),
                    'comment_id': comment_id,
                    'author': author_info.get('name', author_id or 'Unknown'),
                    'author_id': author_id,
                    'initials': author_info.get('initials', ''),
                    'user_id': author_info.get('userId'),
                    'datetime': created,
                    'created': created,
                    'position': {'x_inches': x, 'y_inches': y},
                    'anchor': anchor,
                    'text': comment_text,
                    'format': 'modern',
                    'source_part': item,
                })

        # Keep slide order stable in output.
        comments_by_slide = dict(sorted(comments_by_slide.items(), key=lambda kv: kv[0]))

    total_comments = sum(len(c) for c in comments_by_slide.values())

    return {
        "file": Path(file_path).name,
        "total_comments": total_comments,
        "slides_with_comments": len(comments_by_slide),
        "comments": comments_by_slide,
    }


def _delete_pptx_comment(file_path, output_path, slide_number, comment_index=None):
    """Delete comments from a PPTX slide (modern first, legacy fallback)."""
    import tempfile
    import zipfile
    from posixpath import dirname, join, normpath

    from lxml import etree

    slide_comments_part = f'ppt/comments/comment{slide_number}.xml'
    slide_comment_part_name = f'/ppt/comments/comment{slide_number}.xml'
    slide_rels_part = f'ppt/slides/_rels/slide{slide_number}.xml.rels'
    MODERN_COMMENTS_REL = 'http://schemas.microsoft.com/office/2018/10/relationships/comments'
    MODERN_COMMENTS_CT = 'application/vnd.ms-powerpoint.comments+xml'

    def _resolve_relationship_target(rels_part, target):
        rels = str(rels_part)
        if '/_rels/' in rels and rels.endswith('.rels'):
            source_part = rels.replace('/_rels/', '/', 1)[:-5]
            base_dir = dirname(source_part)
            return normpath(join(base_dir, str(target or '')))
        return normpath(join(dirname(rels), str(target or '')))

    with zipfile.ZipFile(file_path, 'r') as zf_in:
        names = set(zf_in.namelist())

        modern_comment_part = None
        if slide_rels_part in names:
            rels_root = etree.fromstring(zf_in.read(slide_rels_part))
            for rel in rels_root.findall(f'.//{{*}}Relationship'):
                if rel.get('Type') == MODERN_COMMENTS_REL:
                    modern_comment_part = _resolve_relationship_target(slide_rels_part, rel.get('Target'))
                    break

        if modern_comment_part and modern_comment_part in names:
            comments_root = etree.fromstring(zf_in.read(modern_comment_part))
            comments = list(comments_root)
            removed = 0
            if comment_index is None:
                removed = len(comments)
                for cm in comments:
                    comments_root.remove(cm)
            else:
                idx = int(comment_index)
                if 1 <= idx <= len(comments):
                    comments_root.remove(comments[idx - 1])
                    removed = 1

            if removed == 0:
                return {"error": f"Comment index not found on slide {slide_number}: {comment_index}"}

            remove_comment_part = len(comments_root) == 0
            updated_comments_xml = etree.tostring(comments_root, xml_declaration=True, encoding='UTF-8', standalone='yes')

            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
                tmp_path = tmp.name
            with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zf_out:
                for item in zf_in.namelist():
                    if item == '[Content_Types].xml':
                        root = etree.fromstring(zf_in.read(item))
                        if remove_comment_part:
                            for node in list(root):
                                if node.get('PartName') == f'/{modern_comment_part}' and node.get('ContentType') == MODERN_COMMENTS_CT:
                                    root.remove(node)
                        zf_out.writestr(item, etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone='yes'))
                    elif item == slide_rels_part:
                        root = etree.fromstring(zf_in.read(item))
                        if remove_comment_part:
                            for rel in list(root):
                                if rel.get('Type') == MODERN_COMMENTS_REL and _resolve_relationship_target(slide_rels_part, rel.get('Target')) == modern_comment_part:
                                    root.remove(rel)
                        zf_out.writestr(item, etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone='yes'))
                    elif item == modern_comment_part:
                        if not remove_comment_part:
                            zf_out.writestr(item, updated_comments_xml)
                    else:
                        zf_out.writestr(item, zf_in.read(item))

            shutil.move(tmp_path, output_path)

            return {
                "success": True,
                "file": output_path,
                "slide_number": slide_number,
                "deleted_comments": removed,
                "deleted_all_on_slide": comment_index is None,
                "removed_comment_part": remove_comment_part,
                "format": "modern",
                "message": f"Deleted {removed} comment(s) from slide {slide_number}",
                "next_tools": ["pptx_get_comments", "pptx_add_comment", "pptx_list_slides"],
            }

        if slide_comments_part not in zf_in.namelist():
            return {"error": f"No comments found for slide {slide_number}"}

        comments_root = etree.fromstring(zf_in.read(slide_comments_part))
        comments = list(comments_root)
        removed = 0

        if comment_index is None:
            removed = len(comments)
            for cm in comments:
                comments_root.remove(cm)
        else:
            target = str(comment_index)
            for cm in comments:
                if cm.get('idx') == target:
                    comments_root.remove(cm)
                    removed = 1
                    break

        if removed == 0:
            return {
                "error": f"Comment index not found on slide {slide_number}: {comment_index}",
            }

        remove_comment_part = len(comments_root) == 0
        updated_comments_xml = etree.tostring(
            comments_root,
            xml_declaration=True,
            encoding='UTF-8',
            standalone='yes',
        )

        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            tmp_path = tmp.name

        with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zf_out:
            for item in zf_in.namelist():
                if item == '[Content_Types].xml':
                    root = etree.fromstring(zf_in.read(item))
                    if remove_comment_part:
                        for node in list(root):
                            if node.get('PartName') == slide_comment_part_name:
                                root.remove(node)
                    zf_out.writestr(
                        item,
                        etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone='yes'),
                    )
                elif item == f'ppt/slides/_rels/slide{slide_number}.xml.rels':
                    rels_root = etree.fromstring(zf_in.read(item))
                    if remove_comment_part:
                        for rel in list(rels_root):
                            target = rel.get('Target', '')
                            rel_type = rel.get('Type', '')
                            if target == f'../comments/comment{slide_number}.xml' or rel_type.endswith('/comments'):
                                rels_root.remove(rel)
                    zf_out.writestr(
                        item,
                        etree.tostring(rels_root, xml_declaration=True, encoding='UTF-8', standalone='yes'),
                    )
                elif item == slide_comments_part:
                    if not remove_comment_part:
                        zf_out.writestr(item, updated_comments_xml)
                else:
                    zf_out.writestr(item, zf_in.read(item))

    shutil.move(tmp_path, output_path)

    return {
        "success": True,
        "file": output_path,
        "slide_number": slide_number,
        "deleted_comments": removed,
        "deleted_all_on_slide": comment_index is None,
        "removed_comment_part": remove_comment_part,
        "message": f"Deleted {removed} comment(s) from slide {slide_number}",
        "next_tools": ["pptx_get_comments", "pptx_add_comment", "pptx_list_slides"],
    }


def _classify_layout(ph_types, layout_name):
    """Classify a layout based on its placeholder types."""
    name_lower = layout_name.lower()

    # Try to classify by name first
    if 'title slide' in name_lower or 'title master' in name_lower:
        return 'title_slide'
    if 'section' in name_lower:
        return 'section_header'
    if 'two content' in name_lower:
        return 'two_content'
    if 'comparison' in name_lower:
        return 'comparison'
    if 'picture' in name_lower:
        return 'picture_with_caption'
    if 'content with caption' in name_lower:
        return 'content_with_caption'
    if 'vertical' in name_lower:
        return 'title_and_vertical_text'
    if 'blank' in name_lower:
        return 'blank'
    if 'title only' in name_lower:
        return 'title_only'
    if 'title and content' in name_lower:
        return 'title_and_content'

    # Classify by placeholder composition
    has_center_title = 'CENTER_TITLE (3)' in ph_types
    has_title = 'TITLE (1)' in ph_types or has_center_title
    has_subtitle = 'SUBTITLE (4)' in ph_types
    object_count = ph_types.count('OBJECT (7)')
    body_count = ph_types.count('BODY (2)')
    has_picture = 'PICTURE (18)' in ph_types

    if has_center_title and has_subtitle:
        return 'title_slide'
    if has_picture:
        return 'picture_with_caption'
    if object_count >= 2 or (object_count >= 1 and body_count >= 2):
        return 'comparison' if body_count >= 2 else 'two_content'
    if has_title and (object_count >= 1 or body_count >= 1):
        return 'title_and_content'
    if has_title and not object_count and not body_count:
        return 'title_only'
    if not has_title:
        return 'blank'

    return 'title_and_content'  # Default


def _get_layout_recommendations(classification):
    """Get content recommendations for a layout classification."""
    recommendations = {
        'title_slide': ['Opening/closing slides', 'Section dividers', 'Cover pages'],
        'section_header': ['Chapter starts', 'Major transitions', 'Agenda items'],
        'title_and_content': ['Bullet lists', 'Tables', 'Single charts', 'Key points'],
        'two_content': ['Side-by-side comparisons', 'Before/after', 'Dual charts'],
        'comparison': ['Feature comparisons', 'Options analysis', 'Pros/cons'],
        'title_only': ['Large images', 'Full-slide charts', 'Custom layouts'],
        'blank': ['Complex diagrams', 'Free-form content', 'Custom designs'],
        'picture_with_caption': ['Hero images', 'Screenshots', 'Photo highlights'],
        'content_with_caption': ['Annotated visuals', 'Diagrams with notes'],
        'title_and_vertical_text': ['Vertical lists', 'Asian language content'],
    }
    return recommendations.get(classification, ['General content'])

