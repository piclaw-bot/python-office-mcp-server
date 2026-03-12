#!/usr/bin/env python3
"""
pptx_tools.py - MCP tools for PowerPoint presentation processing

Provides tools to parse, extract content from, and generate PowerPoint (.pptx) files.
Supports GitHub Flavored Markdown slide generation with labeled lines and bullet points.
"""

import re
from pathlib import Path
from typing import Any

try:
    from lxml import etree
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
    from pptx.oxml.ns import qn
    from pptx.util import Mm
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from .markdown_parser import (
    HorizontalRule,
    Paragraph,
    Table,
    parse_markdown_to_nodes,
)
from .save_utils import open_pptx_with_retries, safe_save_pptx

# Default theme fonts (can be overridden)
DEFAULT_TITLE_FONT = "Segoe UI Semibold"
DEFAULT_BODY_FONT = "Segoe UI"

# Layout indices in default PowerPoint template
LAYOUT_TITLE = 0          # Title Slide
LAYOUT_TITLE_CONTENT = 1  # Title and Content (bullets)
LAYOUT_SECTION = 2        # Section Header
LAYOUT_TWO_CONTENT = 3    # Two Content
LAYOUT_COMPARISON = 4     # Comparison
LAYOUT_TITLE_ONLY = 5     # Title Only (good for tables)
LAYOUT_BLANK = 6          # Blank


def _clear_all_slides(prs) -> None:
    """Remove all slides from a presentation in-place."""
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def _build_blank_pptx_bytes():
    """Build a minimal blank 16:9 PPTX in memory with 7 standard slide layouts.

    Produces a valid OOXML package with a slide master, theme, and the
    standard layout set (Title, Title+Content, Section Header, Two Content,
    Comparison, Title Only, Blank) so the rest of the generation code can
    reference layouts by index without depending on any external template
    file — critical for PyInstaller-frozen binaries where python-pptx's
    bundled default.pptx is unavailable.
    """
    import io
    import zipfile

    # 16:9 slide dimensions in EMU (English Metric Units)
    CX, CY = 12192000, 6858000  # 13.333" x 7.5"

    # ── package scaffolding ───────────────────────────────────────────
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
        '<Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>'
        '<Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>'
        + ''.join(
            f'<Override PartName="/ppt/slideLayouts/slideLayout{i}.xml" '
            f'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>'
            for i in range(1, 8)
        )
        + '</Types>'
    )

    top_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>'
        '</Relationships>'
    )

    pres_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>'
        f'<p:sldSz cx="{CX}" cy="{CY}"/>'
        f'<p:notesSz cx="{CY}" cy="{CX}"/>'
        '</p:presentation>'
    )

    pres_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>'
        '<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="theme/theme1.xml"/>'
        '</Relationships>'
    )

    # ── slide master ──────────────────────────────────────────────────
    layout_id_list = ''.join(
        f'<p:sldLayoutId id="{2147483649 + i}" r:id="rId{i + 2}"/>'
        for i in range(7)
    )
    master_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:cSld><p:spTree>'
        '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        '<p:grpSpPr/></p:spTree></p:cSld>'
        '<p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2"'
        ' accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4"'
        ' accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>'
        f'<p:sldLayoutIdLst>{layout_id_list}</p:sldLayoutIdLst>'
        '</p:sldMaster>'
    )

    master_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>'
        + ''.join(
            f'<Relationship Id="rId{i + 2}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"'
            f' Target="../slideLayouts/slideLayout{i + 1}.xml"/>'
            for i in range(7)
        )
        + '</Relationships>'
    )

    # ── slide layouts (7 standard types) ──────────────────────────────
    # Each entry: (OOXML type, display name, [(ph_type, idx, l, t, w, h)])
    LAYOUTS = [
        ("title", "Title Slide", [
            ("ctrTitle", 0, 685800, 2286000, 10820400, 1325563),
            ("subTitle", 1, 1371600, 3886200, 9448800, 1752600),
        ]),
        ("obj", "Title and Content", [
            ("title", 0, 457200, 274638, 10515600, 1143000),
            ("body",  1, 457200, 1600200, 10515600, 4525963),
        ]),
        ("secHead", "Section Header", [
            ("title", 0, 685800, 2286000, 10820400, 1325563),
            ("body",  1, 685800, 3886200, 10820400, 1325563),
        ]),
        ("twoObj", "Two Content", [
            ("title", 0, 457200, 274638, 10515600, 1143000),
            ("body",  1, 457200, 1600200, 5156200, 4525963),
            ("body",  2, 5816600, 1600200, 5156200, 4525963),
        ]),
        ("obj", "Comparison", [
            ("title", 0, 457200, 274638, 10515600, 1143000),
            ("body",  1, 457200, 1600200, 5156200, 4525963),
            ("body",  2, 5816600, 1600200, 5156200, 4525963),
        ]),
        ("titleOnly", "Title Only", [
            ("title", 0, 457200, 274638, 10515600, 1143000),
        ]),
        ("blank", "Blank", []),
    ]

    layout_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster"'
        ' Target="../slideMasters/slideMaster1.xml"/>'
        '</Relationships>'
    )

    def _layout_xml(ltype, name, placeholders):
        shapes = ''
        for sid, (ph_type, ph_idx, x_emu, y_emu, cx_emu, cy_emu) in enumerate(placeholders, start=2):
            shapes += (
                f'<p:sp><p:nvSpPr><p:cNvPr id="{sid}" name="{ph_type} {ph_idx}"/>'
                '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
                f'<p:nvPr><p:ph type="{ph_type}" idx="{ph_idx}"/></p:nvPr></p:nvSpPr>'
                f'<p:spPr><a:xfrm><a:off x="{x_emu}" y="{y_emu}"/><a:ext cx="{cx_emu}" cy="{cy_emu}"/>'
                '</a:xfrm></p:spPr>'
                '<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/>'
                '</a:p></p:txBody></p:sp>'
            )
        return (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
            ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            f' type="{ltype}">'
            f'<p:cSld name="{name}"><p:spTree>'
            '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            f'<p:grpSpPr/>{shapes}</p:spTree></p:cSld>'
            '</p:sldLayout>'
        )

    # ── theme (Office-standard colours, Calibri fonts) ────────────────
    theme_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="Blank">'
        '<a:themeElements>'
        '<a:clrScheme name="Office">'
        '<a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1>'
        '<a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1>'
        '<a:dk2><a:srgbClr val="44546A"/></a:dk2>'
        '<a:lt2><a:srgbClr val="E7E6E6"/></a:lt2>'
        '<a:accent1><a:srgbClr val="4472C4"/></a:accent1>'
        '<a:accent2><a:srgbClr val="ED7D31"/></a:accent2>'
        '<a:accent3><a:srgbClr val="A5A5A5"/></a:accent3>'
        '<a:accent4><a:srgbClr val="FFC000"/></a:accent4>'
        '<a:accent5><a:srgbClr val="5B9BD5"/></a:accent5>'
        '<a:accent6><a:srgbClr val="70AD47"/></a:accent6>'
        '<a:hlink><a:srgbClr val="0563C1"/></a:hlink>'
        '<a:folHlink><a:srgbClr val="954F72"/></a:folHlink>'
        '</a:clrScheme>'
        '<a:fontScheme name="Office">'
        '<a:majorFont><a:latin typeface="Calibri Light"/><a:ea typeface=""/><a:cs typeface=""/></a:majorFont>'
        '<a:minorFont><a:latin typeface="Calibri"/><a:ea typeface=""/><a:cs typeface=""/></a:minorFont>'
        '</a:fontScheme>'
        '<a:fmtScheme name="Office">'
        '<a:fillStyleLst>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '</a:fillStyleLst>'
        '<a:lnStyleLst>'
        '<a:ln w="6350"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
        '<a:ln w="12700"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
        '<a:ln w="19050"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
        '</a:lnStyleLst>'
        '<a:effectStyleLst>'
        '<a:effectStyle><a:effectLst/></a:effectStyle>'
        '<a:effectStyle><a:effectLst/></a:effectStyle>'
        '<a:effectStyle><a:effectLst/></a:effectStyle>'
        '</a:effectStyleLst>'
        '<a:bgFillStyleLst>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
        '</a:bgFillStyleLst>'
        '</a:fmtScheme>'
        '</a:themeElements>'
        '</a:theme>'
    )

    # ── assemble ZIP ──────────────────────────────────────────────────
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        zf.writestr('[Content_Types].xml', content_types)
        zf.writestr('_rels/.rels', top_rels)
        zf.writestr('ppt/presentation.xml', pres_xml)
        zf.writestr('ppt/_rels/presentation.xml.rels', pres_rels)
        zf.writestr('ppt/slideMasters/slideMaster1.xml', master_xml)
        zf.writestr('ppt/slideMasters/_rels/slideMaster1.xml.rels', master_rels)
        zf.writestr('ppt/theme/theme1.xml', theme_xml)
        for idx, (ltype, name, phs) in enumerate(LAYOUTS):
            zf.writestr(f'ppt/slideLayouts/slideLayout{idx + 1}.xml',
                        _layout_xml(ltype, name, phs))
            zf.writestr(f'ppt/slideLayouts/_rels/slideLayout{idx + 1}.xml.rels',
                        layout_rels)
    buf.seek(0)
    return buf


def _create_new_presentation_with_fallback() -> Any:
    """Create a new blank 16:9 presentation.

    Tries python-pptx's bundled template first (works in normal Python
    installs).  Falls back to an entirely in-memory blank PPTX so the
    function works inside PyInstaller-frozen binaries where the bundled
    template is unavailable.
    """
    try:
        return Presentation()
    except (PackageNotFoundError, Exception):
        return Presentation(_build_blank_pptx_bytes())


def _set_theme_fonts(prs, title_font: str = None, body_font: str = None):
    """Set the theme fonts in the presentation's theme.

    Args:
        prs: Presentation object
        title_font: Font name for titles (major font)
        body_font: Font name for body text (minor font)
    """
    title_font = title_font or DEFAULT_TITLE_FONT
    body_font = body_font or DEFAULT_BODY_FONT

    # Find theme part through slide master relationships
    for rel in prs.slide_master.part.rels.values():
        if 'theme' in rel.reltype.lower():
            theme_part = rel.target_part
            theme_xml = etree.fromstring(theme_part.blob)

            # Find font scheme
            font_scheme = theme_xml.find('.//' + qn('a:fontScheme'))
            if font_scheme is not None:
                # Set major font (titles)
                major = font_scheme.find(qn('a:majorFont'))
                if major is not None:
                    latin = major.find(qn('a:latin'))
                    if latin is not None:
                        latin.set('typeface', title_font)

                # Set minor font (body)
                minor = font_scheme.find(qn('a:minorFont'))
                if minor is not None:
                    latin = minor.find(qn('a:latin'))
                    if latin is not None:
                        latin.set('typeface', body_font)

                # Write back modified theme
                theme_part._blob = etree.tostring(theme_xml)
            break


def _analyze_markdown_for_layouts(markdown: str) -> dict:
    """Analyze markdown content to determine needed layouts.

    Uses the GFM parser to analyze content structure.

    Returns:
        Dict with layout analysis:
        - has_title_slide: bool
        - slides: list of {type, title, has_table, has_bullets}
    """
    nodes = parse_markdown_to_nodes(markdown)
    analysis = {
        "has_title_slide": False,
        "slides": []
    }

    current_slide = None
    first_heading = True

    for node in nodes:
        # Horizontal rule - new slide boundary
        if isinstance(node, HorizontalRule):
            if current_slide:
                analysis["slides"].append(current_slide)
                current_slide = None
            continue

        # Headings start new slides
        if isinstance(node, Paragraph) and node.level > 0:
            if current_slide:
                analysis["slides"].append(current_slide)

            title = node.text
            if node.level == 1 and first_heading:
                analysis["has_title_slide"] = True
                current_slide = {"type": "title", "title": title,
                                "has_table": False, "has_bullets": False}
                first_heading = False
            else:
                current_slide = {"type": "content", "title": title,
                                "has_table": False, "has_bullets": False}
            continue

        # Tables
        if isinstance(node, Table) and current_slide:
            current_slide["has_table"] = True
            continue

        # Bullets
        if isinstance(node, Paragraph) and node.list_type and current_slide:
            current_slide["has_bullets"] = True
            continue

    # Don't forget last slide
    if current_slide:
        analysis["slides"].append(current_slide)

    return analysis


def _get_layout_for_slide(prs, slide_info: dict) -> Any:
    """Select the best layout for a slide based on its content.

    Args:
        prs: Presentation object
        slide_info: Dict with type, has_table, has_bullets

    Returns:
        SlideLayout object
    """
    layouts = prs.slide_layouts

    if slide_info["type"] == "title":
        return layouts[LAYOUT_TITLE]

    # For content slides, choose based on content type
    if slide_info.get("has_table") and not slide_info.get("has_bullets"):
        # Table-only slide - use Title Only layout (cleaner for tables)
        return layouts[LAYOUT_TITLE_ONLY] if len(layouts) > LAYOUT_TITLE_ONLY else layouts[LAYOUT_TITLE_CONTENT]

    # Default to Title and Content for bullets
    return layouts[LAYOUT_TITLE_CONTENT] if len(layouts) > LAYOUT_TITLE_CONTENT else layouts[0]


class PowerPointTools:
    """MCP tool mixin for PowerPoint presentation processing.

    Note: The extract/to_markdown methods are kept for internal use by
    office_unified_tools but are not exposed as separate MCP tools.
    They are filtered out in office_server.py via DEPRECATED_TOOLS.

    Public MCP tools:
    - pptx_from_markdown: Convert Markdown to PowerPoint presentation

    Internal methods (used by unified tools):
    - tool_pptx_extract: Extract slides and content
    - tool_pptx_to_markdown: Convert to Markdown format

    See PresentationAdvancedTools for slide management and specialized tools.
    """

    def tool_pptx_extract(self, file_path: str) -> dict[str, Any]:
        """Extract text and structure from a PowerPoint presentation.

        NOTE: This method is internal - use office_read() instead.

        Args:
            file_path: Path to the .pptx file

        Returns:
            Dictionary with slides, their titles, and text content
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        prs, resolved_path, error = open_pptx_with_retries(file_path)
        if error:
            return {"error": error}

        path = Path(resolved_path)
        slides = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                "number": slide_num,
                "title": None,
                "content": [],
                "notes": None
            }

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                    if text:
                        # Check if this is the title
                        if (hasattr(shape, "is_placeholder") and shape.is_placeholder
                                and shape.placeholder_format.type == 1):  # Title
                            slide_data["title"] = text
                            continue
                        slide_data["content"].append(text)

                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    slide_data["content"].append({"table": table_data})

            # Extract notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    slide_data["notes"] = notes_frame.text.strip()

            slides.append(slide_data)

        return {
            "file": path.name,
            "slide_count": len(slides),
            "slides": slides
        }

    def tool_pptx_to_markdown(self, file_path: str) -> str:
        """Convert PowerPoint to Markdown format.

        NOTE: This method is internal - use office_read(output_format="markdown") instead.

        Args:
            file_path: Path to the .pptx file

        Returns:
            Markdown string representation
        """
        data = self.tool_pptx_extract(file_path)
        if "error" in data:
            return f"Error: {data['error']}"

        lines = []
        for slide in data["slides"]:
            title = slide["title"] or f"Slide {slide['number']}"
            lines.append(f"## {title}")
            lines.append("")

            for content in slide["content"]:
                if isinstance(content, dict) and "table" in content:
                    # Render table
                    table = content["table"]
                    if table:
                        lines.append("| " + " | ".join(table[0]) + " |")
                        lines.append("| " + " | ".join(["---"] * len(table[0])) + " |")
                        for row in table[1:]:
                            lines.append("| " + " | ".join(row) + " |")
                    lines.append("")
                else:
                    # Bullet points from content
                    for line in content.split("\n"):
                        if line.strip():
                            lines.append(f"- {line.strip()}")
                    lines.append("")

            if slide["notes"]:
                lines.append(f"> **Notes:** {slide['notes']}")
                lines.append("")

            lines.append("---")
            lines.append("")

        return "\n".join(lines)

    def tool_pptx_from_markdown(self, output_path: str, markdown: str,
                                   title_font: str = None, body_font: str = None) -> dict[str, Any]:
        """Convert Markdown content to a PowerPoint presentation.

        This is the primary tool for creating PowerPoint decks from text content.

        Features:
        - Analyzes content to select appropriate layouts for each slide
        - Sets theme fonts (default: Segoe UI Semibold for titles, Segoe UI for body)
        - Uses millimeters internally for precise positioning
        - Fonts inherit from theme - no hardcoded font overrides

        Slide Mapping:
        - First # heading becomes title slide with large centered title
        - Subsequent # or ## headings start new content slides
        - --- (horizontal rule) also starts a new slide context
        - Bullet points (- or *) become slide body content
        - **Label:** patterns are rendered with bold labels (great for key points)
        - **Context (assumptions):** after title becomes subtitle on title slide
        - Tables (| col | col |) are rendered as PowerPoint tables
        - Non-heading, non-bullet paragraphs become plain text

        Example:
            pptx_from_markdown(
                output_path="04. Artifacts/proposal.pptx",
                markdown='''
# Cloud Migration Proposal
**Context:** Enterprise transformation for ACME Corp

---

## Executive Summary
- **Objective:** Migrate 15 legacy applications to Azure
- **Timeline:** 12 months with phased approach

---

## Investment

| Phase | Cost | Timeline |
|-------|------|----------|
| Phase 1 | $400K | Q1 |
| Phase 2 | $800K | Q2-Q3 |

---

## Next Steps
- Approve project charter
- Schedule kickoff workshop
'''
            )

        Args:
            output_path: Path for the output .pptx file
            markdown: Markdown content following the slide pattern
            title_font: Font for titles (default: Segoe UI Semibold)
            body_font: Font for body text (default: Segoe UI)

        Returns:
            Status dictionary with file path and slide count
        """
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

        from pptx.enum.text import PP_ALIGN

        # Step 1: Analyze markdown to understand content structure
        content_analysis = _analyze_markdown_for_layouts(markdown)

        # Step 2: Create presentation and configure theme
        prs = _create_new_presentation_with_fallback()
        prs.slide_width = Mm(338.67)   # Widescreen 16:9 (13.333")
        prs.slide_height = Mm(190.5)   # (7.5")

        # Step 3: Set theme fonts BEFORE creating any slides
        _set_theme_fonts(prs, title_font or DEFAULT_TITLE_FONT, body_font or DEFAULT_BODY_FONT)

        # Standard positioning for 16:9 widescreen (all dimensions in mm)
        MARGIN_LEFT = Mm(12.7)         # 0.5"
        CONTENT_WIDTH = Mm(313.27)     # 12.333" = slide width - margins
        TITLE_TOP = Mm(7.62)           # 0.3"
        TITLE_HEIGHT = Mm(25.4)        # 1.0"
        CONTENT_TOP = Mm(38.1)         # 1.5"
        CONTENT_HEIGHT = Mm(139.7)     # 5.5"
        TABLE_TOP = Mm(45.72)          # 1.8"

        lines = markdown.split("\n")

        # Track state
        current_slide = None
        title_slide_active = False
        has_title_slide = False
        body_segments = []  # Accumulate body content for current slide
        table_buffer = []   # Accumulate table rows
        slide_index = 0     # Track which slide we're creating for layout lookup

        def strip_markdown(text: str) -> str:
            """Remove bold/italic markers from text."""
            t = text or ""
            t = re.sub(r'\*\*(.+?)\*\*', r'\1', t)
            t = re.sub(r'\*(.+?)\*', r'\1', t)
            return t.strip()

        def parse_table_row(line: str) -> list[str] | None:
            """Parse a markdown table row, return cells or None if not a table row."""
            if not line.strip().startswith('|'):
                return None
            # Remove leading/trailing pipes and split
            cells = [c.strip() for c in line.strip().strip('|').split('|')]
            return cells

        def is_separator_row(cells: list[str]) -> bool:
            """Check if this is a table separator row (---|---|---)."""
            return all(re.match(r'^[-:]+$', c.strip()) for c in cells if c.strip())

        def add_title_slide(title: str):
            """Add a title slide with large title and optional subtitle."""
            nonlocal current_slide, title_slide_active, slide_index
            # Use layout from analysis
            layout = prs.slide_layouts[LAYOUT_TITLE]
            slide = prs.slides.add_slide(layout)

            # Only adjust width for widescreen - preserve vertical positioning from layout
            if slide.shapes.title:
                # Store original vertical position before modifying
                original_top = slide.shapes.title.top
                original_height = slide.shapes.title.height
                slide.shapes.title.left = MARGIN_LEFT
                slide.shapes.title.width = CONTENT_WIDTH
                # Restore vertical positioning (layout has it centered)
                slide.shapes.title.top = original_top
                slide.shapes.title.height = original_height
                slide.shapes.title.text = title

            # Resize subtitle placeholder - preserve vertical position
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Subtitle
                    original_top = shape.top
                    original_height = shape.height
                    shape.left = MARGIN_LEFT
                    shape.width = CONTENT_WIDTH
                    shape.top = original_top
                    shape.height = original_height

            title_slide_active = True
            current_slide = slide
            slide_index += 1
            return slide

        def add_content_slide(title: str, for_table: bool = False):
            """Add a content slide with appropriate layout based on content type."""
            nonlocal current_slide, title_slide_active, slide_index

            # Select layout based on content type
            if for_table:
                # Title Only layout is cleaner for table-only slides
                layout_idx = LAYOUT_TITLE_ONLY if len(prs.slide_layouts) > LAYOUT_TITLE_ONLY else LAYOUT_TITLE_CONTENT
            else:
                layout_idx = LAYOUT_TITLE_CONTENT

            layout = prs.slide_layouts[layout_idx] if len(prs.slide_layouts) > layout_idx else prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)

            # Resize title placeholder for widescreen
            if slide.shapes.title:
                slide.shapes.title.left = MARGIN_LEFT
                slide.shapes.title.top = TITLE_TOP
                slide.shapes.title.width = CONTENT_WIDTH
                slide.shapes.title.height = TITLE_HEIGHT
                slide.shapes.title.text = title

            # Resize content placeholder for widescreen
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Body placeholder
                    shape.left = MARGIN_LEFT
                    shape.top = CONTENT_TOP
                    shape.width = CONTENT_WIDTH
                    shape.height = CONTENT_HEIGHT

            title_slide_active = False
            current_slide = slide
            slide_index += 1
            return slide

        def add_table_to_slide(table_data: list[list[str]]):
            """Add a table to the current slide with smart column sizing."""
            if not current_slide or not table_data:
                return

            rows = len(table_data)
            cols = max(len(row) for row in table_data)

            # Position table below title area, full width for widescreen
            left = MARGIN_LEFT
            top = TABLE_TOP
            width = CONTENT_WIDTH

            # Calculate smart column widths based on content
            col_max_chars = [0] * cols
            for row_data in table_data:
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < cols:
                        text = strip_markdown(cell_text)
                        col_max_chars[col_idx] = max(col_max_chars[col_idx], len(text))

            # Determine column width ratios based on content length
            # Use sqrt to dampen extreme differences (long text should wrap)
            total_weight = sum(max(5, c ** 0.6) for c in col_max_chars)  # min weight of 5
            col_weights = [max(5, c ** 0.6) / total_weight for c in col_max_chars]

            # Apply minimum column width (1 inch) and maximum (6 inches)
            MIN_COL_WIDTH = Mm(25.4)   # 1"
            MAX_COL_WIDTH = Mm(152.4)  # 6"

            col_widths = []
            for weight in col_weights:
                w = int(width * weight)
                w = max(MIN_COL_WIDTH, min(MAX_COL_WIDTH, w))
                col_widths.append(w)

            # Adjust to fit total width exactly
            total_calculated = sum(col_widths)
            if total_calculated != width:
                diff = width - total_calculated
                # Distribute difference to columns proportionally
                for i in range(cols):
                    adjustment = int(diff * col_weights[i])
                    col_widths[i] += adjustment
                # Handle any remaining rounding error
                remaining = width - sum(col_widths)
                if remaining != 0:
                    col_widths[-1] += remaining

            # Calculate row height - use taller rows for better readability
            base_row_height = Mm(12.7)  # 0.5" base height
            # For rows with longer content, estimate if wrapping needed
            max_content_len = max(max(len(strip_markdown(c)) for c in row) for row in table_data)
            if max_content_len > 50:
                row_height = Mm(15.24)  # 0.6" for tables with long content
            else:
                row_height = base_row_height

            # Calculate total height, cap to fit slide
            height = min(row_height * rows, Mm(132.08))  # Cap at 5.2"

            table_shape = current_slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table

            # Set calculated column widths
            for col_idx in range(cols):
                table.columns[col_idx].width = col_widths[col_idx]

            # Import alignment constants
            from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

            # Fill cells with proper text wrapping and alignment
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < cols:
                        cell = table.cell(row_idx, col_idx)
                        cell.text = strip_markdown(cell_text)

                        # Set text frame properties
                        tf = cell.text_frame
                        tf.word_wrap = True

                        # Set cell margins for breathing room
                        tf.margin_left = Mm(2.54)   # 0.1"
                        tf.margin_right = Mm(2.54)  # 0.1"
                        tf.margin_top = Mm(1.27)    # 0.05"
                        tf.margin_bottom = Mm(1.27) # 0.05"

                        # Vertical alignment: center text in cell
                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                        # Paragraph formatting
                        para = tf.paragraphs[0]

                        # Header row styling
                        if row_idx == 0:
                            para.font.bold = True
                            para.alignment = PP_ALIGN.CENTER  # Center header text
                        else:
                            # Data rows: left align text, right align numbers
                            text = strip_markdown(cell_text)
                            # Check if content looks like a number/currency/percentage
                            if re.match(r'^[\$€£]?[\d,]+\.?\d*%?$', text.strip().replace(',', '')) or text.startswith('$') or text.endswith('%'):
                                para.alignment = PP_ALIGN.RIGHT
                            else:
                                para.alignment = PP_ALIGN.LEFT

        def flush_table():
            """Flush accumulated table rows to the current slide."""
            nonlocal table_buffer
            if not table_buffer:
                return

            # Filter out separator rows
            data_rows = [row for row in table_buffer if not is_separator_row(row)]
            if data_rows:
                add_table_to_slide(data_rows)
            table_buffer = []

        def flush_body():
            """Flush accumulated body content to the current slide."""
            nonlocal body_segments, current_slide
            if not current_slide or not body_segments:
                return

            # Find the content placeholder
            content_placeholder = None
            for shape in current_slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Body placeholder
                    content_placeholder = shape
                    break

            if content_placeholder and hasattr(content_placeholder, 'text_frame'):
                tf = content_placeholder.text_frame
                tf.clear()

                for i, segment in enumerate(body_segments):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

                    # Handle segments with label/text split
                    if isinstance(segment, dict):
                        if segment.get("label"):
                            run = p.add_run()
                            run.text = segment["label"] + ": "
                            run.font.bold = True
                        if segment.get("text"):
                            run = p.add_run()
                            run.text = segment["text"]
                    else:
                        # Plain text - use add_run to work with bullet formatting
                        run = p.add_run()
                        run.text = str(segment)

                    p.level = 0

            body_segments = []

        for line in lines:
            stripped = line.strip()

            if not stripped:
                # Empty line might end a table
                if table_buffer:
                    flush_table()
                continue

            # Horizontal rule - reset context
            if re.match(r'^[-*_]{3,}$', stripped):
                flush_table()
                flush_body()
                current_slide = None
                continue

            # Table row detection
            table_cells = parse_table_row(stripped)
            if table_cells is not None:
                # If we have body content, flush it first (table comes after bullets)
                if body_segments:
                    flush_body()
                table_buffer.append(table_cells)
                continue
            elif table_buffer:
                # Non-table line after table - flush the table
                flush_table()

            # Heading 1 - title slide or content slide
            h1_match = re.match(r'^#\s+(.+)', stripped)
            if h1_match:
                flush_body()
                title = h1_match.group(1).strip()

                if not has_title_slide:
                    add_title_slide(title)
                    has_title_slide = True
                else:
                    # Check content analysis for this slide's content type
                    slide_info = content_analysis["slides"][slide_index] if slide_index < len(content_analysis["slides"]) else {}
                    for_table = slide_info.get("has_table", False) and not slide_info.get("has_bullets", False)
                    add_content_slide(title, for_table=for_table)
                continue

            # Heading 2 - content slide
            h2_match = re.match(r'^##\s+(.+)', stripped)
            if h2_match:
                flush_body()
                title = h2_match.group(1).strip()
                # Check content analysis for this slide's content type
                slide_info = content_analysis["slides"][slide_index] if slide_index < len(content_analysis["slides"]) else {}
                for_table = slide_info.get("has_table", False) and not slide_info.get("has_bullets", False)
                add_content_slide(title, for_table=for_table)
                continue

            # **Label:** Text pattern (standalone, not bullet)
            labeled_match = re.match(r'^\*\*(.+?):\*\*\s*(.*)', stripped)
            if labeled_match:
                label = strip_markdown(labeled_match.group(1)).rstrip(':')
                text = strip_markdown(labeled_match.group(2))

                # On title slide, add as subtitle using the subtitle placeholder
                if current_slide and title_slide_active and label.lower().startswith('context'):
                    # Find subtitle placeholder
                    subtitle_shape = None
                    for shape in current_slide.placeholders:
                        if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                            subtitle_shape = shape
                            break

                    if subtitle_shape and hasattr(subtitle_shape, 'text_frame'):
                        tf = subtitle_shape.text_frame
                        p = tf.paragraphs[0]
                        run = p.add_run()
                        run.text = text
                        # Font size/color inherits from subtitle placeholder in master
                    else:
                        # Fallback: add text box in consistent position for widescreen
                        left = MARGIN_LEFT
                        top = Mm(114.3)   # 4.5"
                        width = CONTENT_WIDTH
                        height = Mm(25.4)  # 1"
                        textbox = current_slide.shapes.add_textbox(left, top, width, height)
                        tf = textbox.text_frame
                        p = tf.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        run = p.add_run()
                        run.text = text
                        # Font inherits from theme; only text boxes outside
                        # placeholders may need explicit styling
                else:
                    body_segments.append({"label": label, "text": text})
                continue

            # Bullet with **Label:** pattern
            bullet_labeled = re.match(r'^[-*]\s+\*\*(.+?):\*\*\s*(.*)', stripped)
            if bullet_labeled:
                label = strip_markdown(bullet_labeled.group(1)).rstrip(':')
                text = strip_markdown(bullet_labeled.group(2))
                body_segments.append({"label": label, "text": text})
                continue

            # Plain bullet - DO NOT add bullet character, PowerPoint handles this
            bullet_match = re.match(r'^[-*]\s+(.*)', stripped)
            if bullet_match:
                text = strip_markdown(bullet_match.group(1))
                # Strip any existing bullet chars from text to avoid duplication
                text = text.lstrip('•●○◦▪▸►-–—* \t')
                body_segments.append(text)
                continue

            # Fallback - plain text
            body_segments.append(strip_markdown(stripped))

        # Flush any remaining content
        flush_table()
        flush_body()

        slide_count = len(prs.slides)
        safe_save_pptx(prs, output_path)

        return {
            "success": True,
            "file": output_path,
            "slides": slide_count,
            "message": f"Created PowerPoint with {slide_count} slide(s)"
        }
