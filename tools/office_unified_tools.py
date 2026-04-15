"""
Unified Office Tools - Consolidated MCP tools for Word, Excel, and PowerPoint.

This module provides a unified interface for document operations, reducing
the total tool count from 80 to ~15 tools. Document format is auto-detected
from the file extension.

Consolidated tools:
- office_read: Read document content (replaces 7 extract/markdown tools)
- office_inspect: Get document structure (replaces 18 list/get tools)
- office_patch: Edit content (replaces 13 patch/replace tools)
- office_comment: Manage comments (replaces 6 comment tools)
- office_table: Table operations (replaces 10 table tools)
- office_template: Template operations (replaces 5 template tools)
- office_audit: Audit documents (replaces 6 audit tools)
- office_image: Insert images into documents

Note: This class is designed as a mixin that inherits alongside other tool classes.
It delegates to sibling tool methods via self, not by creating new instances.
"""

import contextlib
import os
import re
from datetime import date, datetime
from pathlib import Path
from typing import Any, Literal, TypedDict

try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

from .diagnostics import build_mutation_diagnostics
from .excel_advanced_tools import (
    DEFAULT_AUTHOR,
    _auto_row_height,
    _coerce_value,
    _get_range_bounds,
    _parse_cell_reference,
    _set_cell_with_coercion,
)
from .save_utils import merge_xlsx_preserving_package, resolve_office_path, safe_save_pptx


class PatchChange(TypedDict, total=False):
    """A single change to apply to a document via office_patch."""
    target: str   # What to edit: cell ref ("B5"), placeholder ("<Name>"), or shape path ("slide:1/Title 1")
    value: Any    # New value to set


SUPPORTED_EXTENSIONS = {
    ".docx": "word", ".doc": "word",
    ".xlsx": "excel", ".xlsm": "excel", ".xls": "excel",
    ".pptx": "powerpoint", ".ppt": "powerpoint",
}


def _detect_format(file_path: str) -> str | None:
    """Detect document format from file extension."""
    path = Path(file_path)
    suffix = path.suffix.lower()
    return SUPPORTED_EXTENSIONS.get(suffix)


def _resolve_file_path(file_path: str) -> str:
    """Resolve a file path, trying workspace-relative if absolute doesn't exist.

    Checks in order:
    1. Path as-is (absolute or cwd-relative)
    2. Relative to MCP_WORKSPACE_ROOT environment variable
    3. Relative to current working directory (explicit join)

    Returns the first path that exists, or the original path if none found.
    """
    return resolve_office_path(file_path)


def _unsupported_format_error(file_path: str) -> dict[str, str]:
    """Return a descriptive error for unsupported file formats."""
    ext = Path(file_path).suffix or "(none)"
    supported = ".docx, .xlsx, .xlsm, .pptx"
    return {
        "error": (
            f"Unsupported file format '{ext}' for: {file_path}. "
            f"Supported extensions: {supported}."
        )
    }


def _has_tool(obj: Any, tool_name: str) -> bool:
    """Check if object has a specific tool method."""
    method_name = f"tool_{tool_name}"
    return hasattr(obj, method_name) and callable(getattr(obj, method_name))


def _is_excel_cell_or_range(ref: str) -> bool:
    """Check if a string looks like an Excel cell or range reference."""
    if not ref:
        return False
    trimmed = ref.strip().strip("'").strip('"')
    return bool(re.match(r"^[A-Z]+\d+(?::[A-Z]+\d+)?$", trimmed, re.IGNORECASE))


def _preview_value(value: Any, limit: int = 100) -> str:
    """Return a truncated preview string for a value."""
    preview = str(value)
    if len(preview) > limit:
        return f"{preview[:limit]}..."
    return preview


def _normalize_pptx_comment_target(target: str) -> tuple[int | None, str | None]:
    """Normalize PowerPoint comment target to a slide number."""
    if not target:
        return None, "Missing target slide number"

    raw = str(target).strip()
    if raw.startswith("slide:"):
        raw = raw[6:]

    if "/" in raw:
        return None, "PowerPoint comments target must be a slide number (e.g., '3' or 'slide:3')."

    if not raw.isdigit():
        return None, "PowerPoint comments target must be a slide number (e.g., '3' or 'slide:3')."

    return int(raw), None


def _normalize_pptx_comment_delete_target(target: str | None) -> tuple[int | None, int | None, str | None]:
    """Parse PowerPoint comment delete target.

    Accepted formats:
    - "3" or "slide:3"                  -> delete all comments on slide 3
    - "slide:3/comment:2"                -> delete comment index 2 on slide 3
    """
    if not target:
        return None, None, "Missing target for PowerPoint delete. Use 'slide:N' or 'slide:N/comment:I'."

    raw = str(target).strip()
    if "/" in raw:
        left, right = raw.split("/", 1)
        slide_num, err = _normalize_pptx_comment_target(left)
        if err:
            return None, None, err
        right = right.strip().lower()
        if not right.startswith("comment:"):
            return None, None, "Invalid PowerPoint delete target. Use 'slide:N/comment:I'."
        idx_raw = right.split(":", 1)[1].strip()
        if not idx_raw.isdigit():
            return None, None, "Comment index must be a positive integer."
        return slide_num, int(idx_raw), None

    slide_num, err = _normalize_pptx_comment_target(raw)
    if err:
        return None, None, err
    return slide_num, None, None


def _close_openpyxl_workbook(wb) -> None:
    """Close an openpyxl workbook and any VBA archive safely."""
    vba_archive = getattr(wb, "vba_archive", None)
    if vba_archive is not None:
        with contextlib.suppress(Exception):
            vba_archive.close()
        with contextlib.suppress(Exception):
            wb.vba_archive = None

    archive = getattr(wb, "_archive", None)
    if archive is not None:
        with contextlib.suppress(Exception):
            archive.close()
        with contextlib.suppress(Exception):
            wb._archive = None

    with contextlib.suppress(Exception):
        wb.close()


def _derive_initials(name: str) -> str:
    """Derive initials from a display name."""
    tokens = [part for part in str(name).strip().split() if part]
    if not tokens:
        return "SA"
    return "".join(token[0].upper() for token in tokens[:2])


class OfficeUnifiedTools:
    """Unified tools for Office document processing.

    These consolidated tools auto-detect document format and delegate
    to the appropriate format-specific implementation via self methods.

    This class is designed as a mixin - it expects to be used alongside
    other tool classes (WordTools, ExcelTools, etc.) through multiple
    inheritance. It delegates to sibling tool methods via self.
    """

    # =========================================================================
    # office_set_comment_identity - Runtime comment identity configuration
    # =========================================================================

    def tool_office_set_comment_identity(
        self,
        name: str,
        identity: str | None = None,
        initials: str | None = None,
    ) -> dict[str, Any]:
        """Set default commenter identity for subsequent comment operations.

        This updates in-memory defaults used by office_comment and format-specific
        add-comment tools when the author argument is omitted.

        Args:
            name: Display name for comments (for example, "Jane Doe")
            identity: Optional identity string (for example, email or alias)
            initials: Optional initials override for formats that support it

        Returns:
            Updated identity configuration
        """
        normalized_name = str(name or "").strip()
        if not normalized_name:
            return {"error": "name is required"}

        normalized_identity = str(identity).strip() if identity is not None else None
        if normalized_identity == "":
            normalized_identity = None

        normalized_initials = str(initials).strip().upper() if initials is not None else None
        if not normalized_initials:
            normalized_initials = _derive_initials(normalized_name)

        self._comment_author = normalized_name
        self._comment_identity = normalized_identity
        self._comment_initials = normalized_initials

        return {
            "success": True,
            "comment_author": self._comment_author,
            "comment_identity": self._comment_identity,
            "comment_initials": self._comment_initials,
            "message": "Default comment identity updated for this server session.",
        }

    # =========================================================================
    # office_read - Consolidated read/extract tool
    # =========================================================================

    def tool_office_read(
        self,
        file_path: str,
        output_format: Literal["json", "markdown"] = "json",
        scope: str | None = None,
        include_formulas: bool = False,
    ) -> dict[str, Any] | str:
        """Read content from Word, Excel, or PowerPoint documents.

        Auto-detects document format from file extension and extracts content
        in the requested output format.

        Replaces: word_extract, word_to_markdown, excel_extract, excel_to_markdown,
        excel_get_range, pptx_extract, pptx_to_markdown

        Examples:
            # Read entire Excel file as JSON
            office_read(file_path="data.xlsx")

            # Read Excel file as markdown
            office_read(file_path="data.xlsx", output_format="markdown")

            # Read specific range from Excel
            office_read(file_path="data.xlsx", scope="Sheet1!A1:D10")

            # Read a single worksheet
            office_read(file_path="data.xlsx", scope="Sheet1")

            # Read Word document as markdown
            office_read(file_path="report.docx", output_format="markdown")

            # Read PowerPoint as JSON
            office_read(file_path="deck.pptx")

        Args:
            file_path: Path to the document (.docx, .xlsx, .xlsm, .pptx)
            output_format: Output format - "json" for structured data,
                          "markdown" for text representation
            scope: Optional scope limiter:
                     - Excel: sheet name like "Sheet1" or range like "A1:D10" or "Sheet1!B2:C5"
                   - Word: section title (not yet implemented)
                   - PowerPoint: slide number like "slide:3" (not yet implemented)
            include_formulas: For Excel, return formulas instead of values

        Returns:
            dict for JSON format, str for markdown format
        """
        file_path = _resolve_file_path(file_path)
        doc_format = _detect_format(file_path)

        if doc_format is None:
            return _unsupported_format_error(file_path)

        # Excel
        if doc_format == "excel":
            if not _has_tool(self, "excel_extract"):
                return {"error": "Excel support not available"}

            # If scope is provided, use get_range
            if scope and _has_tool(self, "excel_get_range"):
                parsed_sheet, cell_ref = _parse_cell_reference(scope)
                if _is_excel_cell_or_range(cell_ref):
                    return self.tool_excel_get_range(
                        file_path=file_path,
                        range_ref=scope,
                        include_formulas=include_formulas,
                    )

                sheet_name = (parsed_sheet or cell_ref).strip("'\"")
                if output_format == "markdown":
                    return self.tool_excel_to_markdown(
                        file_path=file_path,
                        sheet_name=sheet_name,
                        include_formulas=include_formulas,
                    )
                return self.tool_excel_extract(
                    file_path=file_path,
                    sheet_name=sheet_name,
                    include_formulas=include_formulas,
                )

            # Full document
            if output_format == "markdown":
                return self.tool_excel_to_markdown(file_path, include_formulas=include_formulas)
            else:
                return self.tool_excel_extract(file_path, include_formulas=include_formulas)

        # Word
        elif doc_format == "word":
            if not _has_tool(self, "word_extract"):
                return {"error": "Word support not available"}

            if output_format == "markdown":
                return self.tool_word_to_markdown(file_path)
            else:
                return self.tool_word_extract(file_path)

        # PowerPoint
        elif doc_format == "powerpoint":
            if not _has_tool(self, "pptx_extract"):
                return {"error": "PowerPoint support not available"}

            if output_format == "markdown":
                return self.tool_pptx_to_markdown(file_path)
            else:
                return self.tool_pptx_extract(file_path)

        return {"error": f"Unhandled format: {doc_format}"}

    # =========================================================================
    # office_inspect - Consolidated inspection tool
    # =========================================================================

    def tool_office_inspect(
        self,
        file_path: str,
        what: Literal[
            "structure",
            "sheets",
            "slides",
            "sections",
            "tables",
            "named_ranges",
            "merged_cells",
            "comments",
            "tracking",
            "shapes",
            "masters",
        ] = "structure",
        target: str | None = None,
    ) -> dict[str, Any]:
        """Inspect document structure and metadata.

        Auto-detects document format and returns requested structural information.

        Replaces: excel_list_sheets, excel_list_tables, excel_list_named_ranges,
        excel_list_merged_cells, excel_get_comments, excel_get_change_log,
        word_list_sections, word_list_tables, word_check_tracking,
        pptx_list_slides, pptx_list_shapes, pptx_list_masters, pptx_get_notes,
        pptx_get_comments, pptx_get_hidden_slides

        Examples:
            # List all sheets in Excel workbook
            office_inspect(file_path="data.xlsx", what="sheets")

            # List tables in Excel
            office_inspect(file_path="data.xlsx", what="tables")

            # Get comments from Excel
            office_inspect(file_path="data.xlsx", what="comments")

            # List slides in PowerPoint
            office_inspect(file_path="deck.pptx", what="slides")

            # List sections in Word
            office_inspect(file_path="report.docx", what="sections")

            # Get shapes on a specific slide
            office_inspect(file_path="deck.pptx", what="shapes", target="3")

        Args:
            file_path: Path to the document
            what: What to inspect:
                  - "structure": Overview of document structure
                  - "sheets": Excel sheets list
                  - "slides": PowerPoint slides list
                  - "sections": Word sections list
                  - "tables": Tables in document
                  - "named_ranges": Excel named ranges
                  - "merged_cells": Excel merged cell regions
                  - "comments": Comments/notes in document
                  - "tracking": Track changes status (Word)
                  - "shapes": Shapes on a slide (PowerPoint)
                  - "masters": Slide masters (PowerPoint)
            target: Optional target for scoped inspection:
                    - Sheet name for Excel
                    - Slide number for PowerPoint

        Returns:
            Dictionary with inspection results
        """
        file_path = _resolve_file_path(file_path)
        doc_format = _detect_format(file_path)

        if doc_format is None:
            return _unsupported_format_error(file_path)

        # Excel inspections
        if doc_format == "excel":
            if not _has_tool(self, "excel_list_sheets"):
                return {"error": "Excel support not available"}

            if what == "structure" or what == "sheets":
                return self.tool_excel_list_sheets(file_path)
            elif what == "tables":
                return self.tool_excel_list_tables(file_path, sheet_name=target)
            elif what == "named_ranges":
                return self.tool_excel_list_named_ranges(file_path)
            elif what == "merged_cells":
                return self.tool_excel_list_merged_cells(file_path, sheet_name=target)
            elif what == "comments":
                return self.tool_excel_get_comments(file_path, sheet_name=target)
            elif what == "tracking":
                return self.tool_excel_get_change_log(file_path)
            else:
                return {"error": f"Inspection '{what}' not supported for Excel"}

        # Word inspections
        elif doc_format == "word":
            if not _has_tool(self, "word_list_sections"):
                return {"error": "Word support not available"}

            if what == "structure" or what == "sections":
                return self.tool_word_list_sections(file_path)
            elif what == "tables":
                return self.tool_word_list_tables(file_path)
            elif what == "tracking":
                return self.tool_word_check_tracking(file_path)
            else:
                return {"error": f"Inspection '{what}' not supported for Word"}

        # PowerPoint inspections
        elif doc_format == "powerpoint":
            if not _has_tool(self, "pptx_list_slides"):
                return {"error": "PowerPoint support not available"}

            if what == "structure" or what == "slides":
                return self.tool_pptx_list_slides(file_path)
            elif what == "shapes":
                if not target:
                    return {"error": "Specify slide number in 'target' for shapes"}
                return self.tool_pptx_list_shapes(file_path, slide_number=int(target))
            elif what == "masters":
                return self.tool_pptx_list_masters(file_path)
            elif what == "tables":
                if not target:
                    return {"error": "Specify slide number in 'target' for tables"}
                return self.tool_pptx_get_table(file_path, slide_number=int(target))
            elif what == "comments":
                return self.tool_pptx_get_comments(file_path)
            else:
                return {"error": f"Inspection '{what}' not supported for PowerPoint"}

        return {"error": f"Unhandled format: {doc_format}"}

    # =========================================================================
    # office_comment - Consolidated comment tool
    # =========================================================================

    def tool_office_comment(
        self,
        file_path: str,
        operation: Literal["add", "get", "reply", "resolve", "reopen", "delete"] = "get",
        target: str | None = None,
        text: str | None = None,
        author: str | None = None,
        output_path: str | None = None,
        format: Literal["flat", "threaded"] = "flat",
        filter: Literal["open", "resolved", "mine", "all"] = "all",
    ) -> dict[str, Any]:
        """Manage comments in Word, Excel, or PowerPoint documents.

        Replaces: excel_add_comment, excel_get_comments, word_add_comment,
        pptx_add_comment, pptx_get_comments

        Examples:
            # Get all comments from Excel
            office_comment(file_path="data.xlsx", operation="get")

            # Add comment to Excel cell
            office_comment(
                file_path="data.xlsx",
                operation="add",
                target="B5",
                text="Review this value"
            )

            # Add comment to Word text
            office_comment(
                file_path="report.docx",
                operation="add",
                target="project timeline",
                text="Verify dates with PM"
            )

            # Reply to an existing Word comment by comment ID
            office_comment(
                file_path="report.docx",
                operation="reply",
                target="12",  # comment ID from office_comment(..., operation="get")
                text="Acknowledged - updated in v2"
            )

            # Add comment to PowerPoint slide
            office_comment(
                file_path="deck.pptx",
                operation="add",
                target="3",  # or "slide:3"
                text="Update chart data"
            )

        Args:
            file_path: Path to the document
            operation: add/get/reply/delete plus resolve/reopen for Word
            target: Target location/ID
                - Excel add/delete: cell reference (e.g., "B5")
                - Excel get: optional sheet name filter
                - Word add: text span to annotate
                - Word reply/resolve/reopen/delete: comment ID
                - PowerPoint add: slide number ("3" or "slide:3")
                - PowerPoint delete: "slide:N" or "slide:N/comment:I"
            text: Comment text (required for add/reply)
            author: Author display name. Also used by Word get(filter="mine")
            output_path: Optional output path (defaults to overwriting input)
            format: For Word get only: "flat" (default) or "threaded"
            filter: For Word get only: "all" (default), "open", "resolved", "mine"

        Returns:
            Dictionary with operation results
        """
        file_path = _resolve_file_path(file_path)
        doc_format = _detect_format(file_path)

        if doc_format is None:
            return _unsupported_format_error(file_path)

        operation = str(operation or "get").lower()

        # Excel comments
        if doc_format == "excel":
            if not _has_tool(self, "excel_get_comments"):
                return {"error": "Excel support not available"}

            if operation == "get":
                if text is not None:
                    return {"error": "Excel comment retrieval does not accept text parameter"}
                if format != "flat":
                    return {"error": "Excel comment retrieval supports format='flat' only"}
                if filter != "all":
                    return {"error": "Excel comment retrieval supports filter='all' only"}
                # target can be sheet_name for filtering
                return self.tool_excel_get_comments(file_path, sheet_name=target)
            elif operation == "reply":
                return {"error": "Excel comment reply is not supported; use operation='add' with a cell target"}
            elif operation in {"resolve", "reopen"}:
                return {"error": "Excel comment resolve/reopen is not supported by this server"}
            elif operation == "add":
                if format != "flat" or filter != "all":
                    return {"error": "format/filter are only supported for operation='get'"}
                if not target or not text:
                    return {"error": "Both 'target' (cell ref) and 'text' required"}
                return self.tool_excel_add_comment(
                    file_path=file_path,
                    cell_ref=target,
                    text=text,
                    author=author,
                    output_path=output_path,
                )
            elif operation == "delete":
                if format != "flat" or filter != "all":
                    return {"error": "format/filter are only supported for operation='get'"}
                if text is not None:
                    return {"error": "Excel comment delete does not accept text parameter"}
                if not target:
                    return {"error": "'target' (cell ref) required for Excel delete"}
                if not _has_tool(self, "excel_delete_comment"):
                    return {"error": "Excel comment delete support not available"}
                return self.tool_excel_delete_comment(
                    file_path=file_path,
                    cell_ref=target,
                    output_path=output_path,
                )

            return {"error": "Unsupported operation for Excel comments"}

        # Word comments
        elif doc_format == "word":
            if operation == "get":
                if text is not None or target is not None:
                    return {"error": "Word comment retrieval does not accept target or text parameters"}
                if not _has_tool(self, "word_get_comments"):
                    return {"error": "Word comment support not available"}
                return self.tool_word_get_comments(
                    file_path=file_path,
                    filter=filter,
                    author=author,
                    format=format,
                )
            elif operation == "add":
                if format != "flat" or filter != "all":
                    return {"error": "format/filter are only supported for operation='get'"}
                if not _has_tool(self, "word_add_comment"):
                    return {"error": "Word comment support not available"}
                if not target or not text:
                    return {"error": "Both 'target' (text to find) and 'text' required"}
                return self.tool_word_add_comment(
                    file_path=file_path,
                    target_text=target,
                    comment_text=text,
                    author=author,
                    output_path=output_path,
                )
            elif operation == "reply":
                if format != "flat" or filter != "all":
                    return {"error": "format/filter are only supported for operation='get'"}
                if not _has_tool(self, "word_reply_to_comment"):
                    return {"error": "Word comment reply support not available"}
                if not target or not text:
                    return {"error": "Both 'target' (comment ID) and 'text' required for Word reply"}
                return self.tool_word_reply_to_comment(
                    file_path=file_path,
                    comment_id=str(target),
                    text=text,
                    author=author,
                    output_path=output_path,
                )
            elif operation in {"resolve", "reopen"}:
                if format != "flat" or filter != "all":
                    return {"error": "format/filter are only supported for operation='get'"}
                if text is not None:
                    return {"error": "Word comment resolve/reopen does not accept text parameter"}
                if not target:
                    return {"error": "'target' (comment ID) required for Word resolve/reopen"}
                if not _has_tool(self, "word_resolve_comment"):
                    return {"error": "Word comment resolve support not available"}
                return self.tool_word_resolve_comment(
                    file_path=file_path,
                    comment_id=str(target),
                    resolved=(operation == "resolve"),
                    output_path=output_path,
                )
            elif operation == "delete":
                if format != "flat" or filter != "all":
                    return {"error": "format/filter are only supported for operation='get'"}
                if text is not None:
                    return {"error": "Word comment delete does not accept text parameter"}
                if not target:
                    return {"error": "'target' (comment ID) required for Word delete"}
                if not _has_tool(self, "word_delete_comment"):
                    return {"error": "Word comment delete support not available"}
                return self.tool_word_delete_comment(
                    file_path=file_path,
                    comment_id=str(target),
                    output_path=output_path,
                )

            return {"error": "Unsupported operation for Word comments"}

        # PowerPoint comments
        elif doc_format == "powerpoint":
            if not _has_tool(self, "pptx_get_comments"):
                return {"error": "PowerPoint support not available"}

            if operation == "get":
                if text is not None or target is not None:
                    return {"error": "PowerPoint comment retrieval does not accept target or text parameters"}
                if format != "flat":
                    return {"error": "PowerPoint comment retrieval supports format='flat' only"}
                if filter != "all":
                    return {"error": "PowerPoint comment retrieval supports filter='all' only"}
                return self.tool_pptx_get_comments(file_path)
            elif operation == "reply":
                return {"error": "PowerPoint comment reply is not supported; use operation='add' with a slide target"}
            elif operation in {"resolve", "reopen"}:
                return {"error": "PowerPoint comment resolve/reopen is not supported"}
            elif operation == "add":
                if format != "flat" or filter != "all":
                    return {"error": "format/filter are only supported for operation='get'"}
                if not target or not text:
                    return {"error": "Both 'target' (slide number) and 'text' required"}
                slide_number, error = _normalize_pptx_comment_target(target)
                if error:
                    return {"error": error}
                return self.tool_pptx_add_comment(
                    file_path=file_path,
                    slide_number=slide_number,
                    comment_text=text,
                    author=author,
                )
            elif operation == "delete":
                if format != "flat" or filter != "all":
                    return {"error": "format/filter are only supported for operation='get'"}
                if text is not None:
                    return {"error": "PowerPoint comment delete does not accept text parameter"}
                if not _has_tool(self, "pptx_delete_comment"):
                    return {"error": "PowerPoint comment delete support not available"}
                slide_number, comment_index, error = _normalize_pptx_comment_delete_target(target)
                if error:
                    return {"error": error}
                return self.tool_pptx_delete_comment(
                    file_path=file_path,
                    slide_number=slide_number,
                    comment_index=comment_index,
                    output_path=output_path,
                )

            return {"error": "Unsupported operation for PowerPoint comments"}

        return {"error": f"Unhandled format: {doc_format}"}

    # =========================================================================
    # office_patch - Consolidated editing tool
    # =========================================================================

    def tool_office_patch(
        self,
        file_path: str,
        changes: list[PatchChange],
        track_changes: bool = True,
        output_path: str | None = None,
        mode: Literal["best_effort", "safe", "strict", "dry_run"] = "best_effort",
    ) -> dict[str, Any]:
        """Apply edits to Word, Excel, or PowerPoint documents.

        Accepts a list of changes and applies them to the document. Each change
        specifies a target (cell, shape, placeholder, section) and new value.

        Replaces: excel_patch_cell, excel_patch_range, excel_replace_placeholders,
        word_patch_section, word_patch_placeholder, word_fix_split_placeholders,
        word_replace_global_variables, pptx_patch_shape, pptx_replace_text,
        pptx_replace_placeholders

        Examples:
            # Patch Excel cells
            office_patch(
                file_path="data.xlsx",
                changes=[
                    {"target": "B5", "value": "New Value"},
                    {"target": "C10", "value": 42},
                    {"target": "D1", "value": "=SUM(A1:A10)"},
                ]
            )

            # Patch Excel range (multiple cells at once)
            office_patch(
                file_path="data.xlsx",
                changes=[{"target": "A1:B3", "value": [["H1", "H2"], ["A", 1], ["B", 2]]}]
            )

            # Patch cells on a specific sheet (quote sheet names with special chars)
            office_patch(
                file_path="form.xlsm",
                changes=[
                    {"target": "'ECIF Work Scope (E)'!B5", "value": "Contoso Ltd"},
                    {"target": "'ECIF Work Scope (E)'!B28", "value": "02/01/2026"},
                ]
            )

            # Replace placeholders in Word
            office_patch(
                file_path="template.docx",
                changes=[
                    {"target": "<Customer Name>", "value": "Acme Corp"},
                    {"target": "<Date>", "value": "2026-01-23"},
                ]
            )

            # Patch PowerPoint shape
            office_patch(
                file_path="deck.pptx",
                changes=[{"target": "slide:1/Title 1", "value": "New Title"}]
            )

            # Patch PowerPoint with soft return
            office_patch(
                file_path="deck.pptx",
                changes=[{"target": "slide:1/Title 2", "value": "Contoso{br}Project"}]
            )

        IMPORTANT for PowerPoint: When patching content placeholders (body,
        Content Placeholder), do NOT include bullet characters (•, -, *, etc.)
        in text lines. PowerPoint placeholders automatically render each line
        as a bullet. Including bullet characters causes duplication like '- • text'.
        Use newlines to separate items, and leading spaces (4 spaces) for indentation.

        Args:
            file_path: Path to the document
            changes: List of changes, each a dict with \"target\" (cell ref, placeholder, or shape path) and \"value\" (new content, no bullet chars for PPTX body)
            track_changes: Log changes for audit trail (default True)
            output_path: Optional output path (defaults to overwriting input)

        Returns:
            Dictionary with results of all changes. Each successful result includes
            a "value_preview" field with a truncated preview of the value applied.
        """
        file_path = _resolve_file_path(file_path)
        doc_format = _detect_format(file_path)

        if doc_format is None:
            return _unsupported_format_error(file_path)

        if not changes:
            return {"error": "No changes provided"}

        if mode == "safe" and (output_path is None or Path(output_path).resolve() == Path(file_path).resolve()):
            return {
                "success": False,
                "mode": mode,
                "status": "failed",
                "warnings": ["safe mode requires an explicit output_path different from the source file."],
                "matched_targets": [],
                "unmatched_targets": [{"target": "document", "reason": "safe_mode_requires_distinct_output_path"}],
                "skipped_targets": [],
                "diagnostics": {"changes_requested": len(changes)},
                "next_tools": ["office_help", "office_template", "office_inspect"],
            }

        results = []
        errors = []

        # Excel patching
        if doc_format == "excel":
            if not HAS_OPENPYXL or not _has_tool(self, "excel_patch_cell"):
                return {"error": "Excel support not available"}

            try:
                wb = load_workbook(file_path, data_only=False, keep_vba=True)
            except Exception as e:
                return {"error": f"Failed to load workbook: {e}"}

            edited_sheets: set[str] = set()
            try:
                any_applied = False

                for change in changes:
                    target = change.get("target")
                    value = change.get("value")

                    if not target:
                        errors.append({"error": "Missing 'target' in change"})
                        continue

                    if ":" in target and not target.startswith("="):
                        parsed_sheet, cell_ref = _parse_cell_reference(target)
                        target_sheet = parsed_sheet or wb.active.title

                        if target_sheet not in wb.sheetnames:
                            errors.append({"target": target, "error": f"Sheet not found: {target_sheet}"})
                            continue

                        range_error = False

                        try:
                            min_row, min_col, max_row, max_col = _get_range_bounds(cell_ref)
                        except ValueError as e:
                            errors.append({"target": target, "error": str(e)})
                            continue

                        if not isinstance(value, list):
                            errors.append({"target": target, "error": "Range value must be a 2D list"})
                            continue

                        expected_rows = max_row - min_row + 1
                        expected_cols = max_col - min_col + 1

                        if len(value) != expected_rows:
                            errors.append({
                                "target": target,
                                "error": f"Row count mismatch: expected {expected_rows}, got {len(value)}",
                            })
                            continue

                        rows_to_adjust = set()
                        ws = wb[target_sheet]

                        for row_idx, row_values in enumerate(value):
                            if len(row_values) != expected_cols:
                                errors.append({
                                    "target": target,
                                    "error": f"Column count mismatch in row {row_idx + 1}: "
                                             f"expected {expected_cols}, got {len(row_values)}",
                                })
                                range_error = True
                                break

                            for col_idx, new_value in enumerate(row_values):
                                cell = ws.cell(row=min_row + row_idx, column=min_col + col_idx)
                                old_value = cell.value
                                coerced_value = _coerce_value(new_value)

                                if old_value != coerced_value:
                                    if mode != "dry_run":
                                        _set_cell_with_coercion(cell, new_value, auto_height=True)
                                        rows_to_adjust.add(min_row + row_idx)
                                    any_applied = True
                                    edited_sheets.add(target_sheet)

                        for row_num in rows_to_adjust:
                            _auto_row_height(ws, row_num)

                        if not range_error:
                            results.append({
                                "target": target,
                                "success": True,
                                "value_preview": _preview_value(value),
                            })
                        continue

                    parsed_sheet, cell_address = _parse_cell_reference(target)
                    target_sheet = parsed_sheet or wb.active.title

                    if target_sheet not in wb.sheetnames:
                        errors.append({"target": target, "error": f"Sheet not found: {target_sheet}"})
                        continue

                    ws = wb[target_sheet]
                    cell = ws[cell_address]
                    old_value = cell.value
                    coerced_value = _coerce_value(value) if mode == "dry_run" else _set_cell_with_coercion(cell, value, auto_height=True)

                    if mode != "dry_run":
                        _auto_row_height(ws, cell.row, cell=cell)
                    any_applied = True
                    edited_sheets.add(target_sheet)
                    results.append({
                        "target": target,
                        "success": True,
                        "value_preview": _preview_value(value),
                    })

                if any_applied and mode != "dry_run":
                    save_path = output_path or file_path
                    staged_save_path = save_path
                    cleanup_path: str | None = None
                    if Path(save_path).resolve() == Path(file_path).resolve():
                        staged_save_path = str(Path(save_path).with_suffix(Path(save_path).suffix + ".openpyxl.tmp"))
                        cleanup_path = staged_save_path
                    try:
                        wb.save(staged_save_path)
                        merge_xlsx_preserving_package(
                            source_path=file_path,
                            staged_path=staged_save_path,
                            output_path=save_path,
                            edited_sheets=edited_sheets,
                        )
                    except Exception as e:
                        return {"error": f"Failed to save workbook: {e}"}
                    finally:
                        if cleanup_path:
                            with contextlib.suppress(FileNotFoundError):
                                Path(cleanup_path).unlink()
            finally:
                _close_openpyxl_workbook(wb)

            matched_targets = [item for item in results if item.get("success")]
            unmatched_targets = [item for item in errors if item.get("target")]
            skipped_targets = [item for item in results if not item.get("success")]
            warnings = []
            if not matched_targets:
                warnings.append("No Excel changes were applied; inspect sheets/ranges before retrying.")
            diag = build_mutation_diagnostics(
                matched_targets=matched_targets,
                unmatched_targets=unmatched_targets,
                skipped_targets=skipped_targets,
                warnings=warnings,
                diagnostics={
                    "changes_requested": len(changes),
                    "changes_applied": len(results),
                    "edited_sheets": sorted(edited_sheets),
                    "errors": errors,
                    "preserved_parts_summary": {
                        "strategy": "merge_original_package_with_edited_sheets",
                        "edited_sheets": sorted(edited_sheets),
                    },
                },
                next_tools=["office_read", "office_inspect", "office_audit", "office_help"],
            )
            if mode == "strict" and (unmatched_targets or skipped_targets or not matched_targets):
                diag["success"] = False
                diag["status"] = "failed"
                diag["warnings"] = list(dict.fromkeys(diag.get("warnings", []) + ["strict mode requires all requested Excel targets to match and apply cleanly."]))
            return {
                **diag,
                "mode": mode,
                "file": file_path,
                "changes_applied": len(results),
                "errors": len(errors),
                "results": results,
                "error_details": errors if errors else None,
                "edited_sheets": sorted(edited_sheets),
                "preserved_parts_summary": {
                    "strategy": "merge_original_package_with_edited_sheets",
                    "edited_sheets": sorted(edited_sheets),
                },
            }

        # Word patching
        elif doc_format == "word":
            if not _has_tool(self, "word_fix_split_placeholders"):
                return {"error": "Word support not available"}

            placeholders = {}
            for change in changes:
                target = change.get("target")
                value = change.get("value")

                if not target:
                    errors.append({"error": "Missing 'target' in change"})
                    continue

                if str(target).lower().startswith("section:"):
                    if not _has_tool(self, "word_patch_section"):
                        errors.append({"target": target, "error": "Word section patching support not available"})
                        continue

                    section_title = str(target).split(":", 1)[1].strip()
                    if not section_title:
                        errors.append({"target": target, "error": "Section title is required after 'section:'"})
                        continue

                    if isinstance(value, list):
                        new_content = [str(item) for item in value if str(item).strip()]
                    elif value is None:
                        new_content = []
                    else:
                        raw_text = str(value)
                        new_content = [part.strip() for part in re.split(r"\n\s*\n", raw_text) if part.strip()]
                        if not new_content and raw_text.strip():
                            new_content = [raw_text.strip()]

                    section_result = self.tool_word_patch_section(
                        file_path=file_path,
                        section_title=section_title,
                        new_content=new_content,
                        output_path=output_path,
                        mode=mode,
                    )
                    if "error" in section_result:
                        errors.append({"target": target, "error": section_result["error"]})
                    else:
                        results.append({
                            "target": target,
                            "success": True,
                            "value_preview": _preview_value(value),
                        })
                    continue

                placeholders[str(target)] = str(value) if value is not None else ""

            if placeholders:
                if mode == "dry_run":
                    read_result = self.tool_word_extract(file_path) if _has_tool(self, "word_extract") else {"content": ""}
                    content = str(read_result.get("content", ""))
                    for target in placeholders:
                        results.append({
                            "target": target,
                            "success": target in content,
                            "value_preview": _preview_value(placeholders.get(target, "")),
                        })
                else:
                    result = self.tool_word_fix_split_placeholders(
                        file_path=file_path,
                        replacements=placeholders,
                        output_path=output_path,
                    )
                    if "error" in result:
                        errors.append(result)
                    else:
                        total_replacements = result.get("total_replacements", 0)
                        for target in placeholders:
                            results.append({
                                "target": target,
                                "success": total_replacements > 0,
                                "value_preview": _preview_value(placeholders.get(target, "")),
                            })

        # PowerPoint patching
        elif doc_format == "powerpoint":
            if not _has_tool(self, "pptx_patch_shape"):
                return {"error": "PowerPoint support not available"}

            for change in changes:
                target = change.get("target")
                value = change.get("value")

                if not target:
                    errors.append({"error": "Missing 'target' in change"})
                    continue

                if isinstance(value, str):
                    value = (
                        value.replace("{br}", "\v")
                        .replace("<br/>", "\v")
                        .replace("<br>", "\v")
                    )

                # Check if target specifies slide/shape (e.g., "slide:1/Title 1")
                if target.startswith("slide:"):
                    # Parse slide:N/shape_name format
                    parts = target[6:].split("/", 1)
                    slide_num = int(parts[0])
                    shape_identifier = parts[1] if len(parts) > 1 else None

                    if shape_identifier:
                        result = self.tool_pptx_patch_shape(
                            file_path=file_path,
                            slide_number=slide_num,
                            shape_identifier=shape_identifier,
                            new_text=str(value),
                            output_path=output_path,
                        )
                    else:
                        errors.append({
                            "target": target,
                            "error": "Shape name required (format: slide:N/ShapeName)"
                        })
                        continue
                else:
                    # Assume it's a placeholder to replace globally
                    result = self.tool_pptx_replace_text(
                        file_path=file_path,
                        find_text=target,
                        replace_text=str(value),
                        output_path=output_path,
                    )

                if "error" in result:
                    errors.append({"target": target, "error": result["error"]})
                else:
                    results.append({
                        "target": target,
                        "success": True,
                        "value_preview": _preview_value(value),
                    })

        matched_targets = [item for item in results if item.get("success")]
        unmatched_targets = [item for item in errors if item.get("target")]
        skipped_targets = [item for item in results if not item.get("success")]
        warnings = []
        if not matched_targets:
            warnings.append(f"No {doc_format} changes were applied; inspect targets before retrying.")
        diag = build_mutation_diagnostics(
            matched_targets=matched_targets,
            unmatched_targets=unmatched_targets,
            skipped_targets=skipped_targets,
            warnings=warnings,
            diagnostics={
                "changes_requested": len(changes),
                "changes_applied": len(results),
                "errors": errors,
            },
            next_tools=["office_read", "office_inspect", "office_audit", "office_help"],
        )
        if mode == "strict" and (unmatched_targets or skipped_targets or not matched_targets):
            diag["success"] = False
            diag["status"] = "failed"
            diag["warnings"] = list(dict.fromkeys(diag.get("warnings", []) + [f"strict mode requires all requested {doc_format} targets to match and apply cleanly."]))
        return {
            **diag,
            "mode": mode,
            "file": file_path,
            "changes_applied": len(results),
            "errors": len(errors),
            "results": results,
            "error_details": errors if errors else None,
        }

    # =========================================================================
    # office_table - Consolidated table operations
    # =========================================================================

    def tool_office_table(
        self,
        file_path: str,
        operation: Literal["get", "add_row", "update_row", "create"] = "get",
        table_id: str | None = None,
        data: dict[str, Any] | None = None,
        row_index: int | None = None,
        output_path: str | None = None,
        mode: Literal["best_effort", "safe", "strict", "dry_run"] = "best_effort",
    ) -> dict[str, Any]:
        """Manage tables in Word, Excel, or PowerPoint documents.

        Replaces: excel_get_table, excel_append_table_row, excel_update_table_row,
        word_get_table, word_insert_table_row, word_patch_table_row, word_create_new_table,
        pptx_get_table, pptx_insert_table_row, pptx_patch_table_cell

        Examples:
            # Get Excel table data
            office_table(file_path="data.xlsx", operation="get", table_id="Sales")

            # Add row to Excel table
            office_table(
                file_path="data.xlsx",
                operation="add_row",
                table_id="Sales",
                data={"Product": "Widget", "Amount": 100}
            )

            # Update Excel table row
            office_table(
                file_path="data.xlsx",
                operation="update_row",
                table_id="Sales",
                row_index=2,
                data={"Amount": 150}
            )

            # Get Word table (by index, passed as string)
            office_table(file_path="report.docx", operation="get", table_id="0")

            # Create Word table
            office_table(
                file_path="report.docx",
                operation="create",
                data={
                    "headers": ["Phase", "Owner", "Target Date"],
                    "rows": [{"Phase": "Discovery", "Owner": "PM", "Target Date": "2026-04-01"}],
                    "insert_after_section": "Delivery Plan"
                }
            )

            # Get PowerPoint table (slide number as string)
            office_table(file_path="deck.pptx", operation="get", table_id="3")

        Args:
            file_path: Path to the document
            operation: "get" to retrieve table data, "add_row" to append a row,
                "update_row" to modify an existing row, or "create" to create a
                new table (Word and PowerPoint)
            table_id: Table identifier as a string. For Excel pass the table name
                (e.g. "Sales"). For Word pass the 0-based table index for get/add/update
                (e.g. "0"). For PowerPoint pass the slide number (e.g. "3").
                For Word create, table_id is optional.
            data: Row data as a dict with column names or indices as keys and cell
                values as values. Required for add_row and update_row. For Word create,
                provide an object with "headers" and optional "rows",
                "insert_after_section", "insert_before_section", "output_path", and
                "author". For PowerPoint update_row include "row", "col", and "value" keys.
            row_index: 1-based row index for update_row operations

        Returns:
            Dictionary with table data or operation result
        """
        file_path = _resolve_file_path(file_path)
        doc_format = _detect_format(file_path)

        if doc_format is None:
            return _unsupported_format_error(file_path)

        # Strip stray quotes from table_id (some MCP clients wrap values)
        if table_id is not None:
            table_id = str(table_id).strip('"').strip("'")

        # Excel tables
        if doc_format == "excel":
            if not _has_tool(self, "excel_get_table"):
                return {"error": "Excel support not available"}

            if not table_id:
                return {"error": "table_id (table name) required for Excel"}

            table_name = str(table_id)

            if operation == "get":
                return self.tool_excel_get_table(
                    file_path=file_path,
                    table_name=table_name,
                )
            elif operation == "add_row":
                if not data:
                    return {"error": "data required for add_row operation"}
                return self.tool_excel_append_table_row(
                    file_path=file_path,
                    table_name=table_name,
                    row_data=data,
                    output_path=output_path,
                    mode=mode,
                )
            elif operation == "update_row":
                if not data or row_index is None:
                    return {"error": "data and row_index required for update_row"}
                return self.tool_excel_update_table_row(
                    file_path=file_path,
                    table_name=table_name,
                    row_index=row_index,
                    row_data=data,
                    output_path=output_path,
                    mode=mode,
                )

        # Word tables
        elif doc_format == "word":
            if operation != "get" and mode != "best_effort":
                return {"error": "mode support for Word table mutations is not implemented yet; use best_effort for this path"}
            if not _has_tool(self, "word_get_table"):
                return {"error": "Word support not available"}

            if operation == "create":
                if not _has_tool(self, "word_create_new_table"):
                    return {"error": "Word table creation support not available"}
                if not isinstance(data, dict):
                    return {"error": "For Word create, data must be an object with headers and optional rows"}

                headers = data.get("headers")
                rows = data.get("rows", [])
                if not isinstance(headers, list) or not headers:
                    return {"error": "For Word create, data.headers must be a non-empty list"}
                if rows is not None and not isinstance(rows, list):
                    return {"error": "For Word create, data.rows must be a list of row objects"}

                return self.tool_word_create_new_table(
                    file_path=file_path,
                    headers=[str(h) for h in headers],
                    rows=rows,
                    insert_after_section=data.get("insert_after_section"),
                    insert_before_section=data.get("insert_before_section"),
                    output_path=data.get("output_path"),
                    author=data.get("author", DEFAULT_AUTHOR),
                )

            if table_id is None:
                return {"error": "table_id (table index or identifier) required"}

            if operation == "get":
                return self.tool_word_get_table(
                    file_path=file_path,
                    table_identifier=str(table_id),
                )
            elif operation == "add_row":
                if not data:
                    return {"error": "data required for add_row operation"}
                return self.tool_word_insert_table_row(
                    file_path=file_path,
                    table_identifier=str(table_id),
                    row_data=data,
                )
            elif operation == "update_row":
                if not data or row_index is None:
                    return {"error": "data and row_index required for update_row"}
                return self.tool_word_patch_table_row(
                    file_path=file_path,
                    table_identifier=str(table_id),
                    row_index=int(row_index),
                    updates=data if isinstance(data, dict) else {
                        str(i): v for i, v in enumerate(data)
                    },
                )

        # PowerPoint tables
        elif doc_format == "powerpoint":
            if operation != "get" and mode != "best_effort":
                return {"error": "mode support for PowerPoint table mutations is not implemented yet; use best_effort for this path"}
            if not _has_tool(self, "pptx_get_table"):
                return {"error": "PowerPoint support not available"}

            if operation == "create":
                if not _has_tool(self, "pptx_add_table"):
                    return {"error": "PowerPoint table creation support not available"}
                if not isinstance(data, dict):
                    return {"error": "For PowerPoint create, data must be an object with headers and optional rows"}

                slide_value = table_id if table_id is not None else data.get("slide")
                if slide_value is None:
                    return {"error": "For PowerPoint create, provide slide number via table_id or data.slide"}

                headers = data.get("headers")
                rows = data.get("rows", [])
                if not isinstance(headers, list) or not headers:
                    return {"error": "For PowerPoint create, data.headers must be a non-empty list"}
                if rows is not None and not isinstance(rows, list):
                    return {"error": "For PowerPoint create, data.rows must be a list of rows"}

                return self.tool_pptx_add_table(
                    file_path=file_path,
                    slide_number=int(slide_value),
                    headers=[str(h) for h in headers],
                    rows=[[str(cell) for cell in row] for row in rows],
                    left=float(data.get("left", 1.0)),
                    top=float(data.get("top", 2.0)),
                    width=float(data.get("width", 11.0)),
                    height=float(data.get("height", 3.0)),
                )

            if table_id is None:
                return {"error": "table_id (slide number) required for PowerPoint"}

            slide_num = int(table_id)

            if operation == "get":
                return self.tool_pptx_get_table(
                    file_path=file_path,
                    slide_number=slide_num,
                )
            elif operation == "add_row":
                if not data:
                    return {"error": "data required for add_row operation"}
                row_data = data
                if isinstance(data, dict):
                    table_result = self.tool_pptx_get_table(
                        file_path=file_path,
                        slide_number=slide_num,
                    )
                    if "error" in table_result:
                        return table_result

                    header = table_result.get("header", [])
                    if header:
                        ordered = []
                        for col_name in header:
                            value = ""
                            if col_name in data:
                                value = data[col_name]
                            else:
                                for key, val in data.items():
                                    if str(key).strip().lower() == str(col_name).strip().lower():
                                        value = val
                                        break
                            ordered.append(value)
                        row_data = ordered
                    else:
                        row_data = list(data.values())

                return self.tool_pptx_insert_table_row(
                    file_path=file_path,
                    slide_number=slide_num,
                    row_data=row_data,
                )
            elif operation == "update_row":
                # PowerPoint uses patch_table_cell for individual cells
                if not isinstance(data, dict) or "row" not in data or "col" not in data:
                    return {
                        "error": "For PowerPoint update_row, data must include "
                                "'row', 'col', and 'value' keys"
                    }
                return self.tool_pptx_patch_table_cell(
                    file_path=file_path,
                    slide_number=slide_num,
                    row_index=int(data["row"]),
                    col_index=int(data["col"]),
                    new_text=str(data.get("value", "")),
                )

        return {"error": f"Unhandled format: {doc_format}"}

    # =========================================================================
    # office_template - Consolidated template operations
    # =========================================================================

    def tool_office_template(
        self,
        source_path: str,
        destination_path: str,
        operation: Literal["copy", "analyze"] = "copy",
    ) -> dict[str, Any]:
        """Copy templates or analyze template structure.

        Replaces: excel_copy_template, word_copy_template, pptx_copy_template,
        word_analyze_template_formatting

        Examples:
            # Copy Excel template
            office_template(
                source_path="templates/budget.xlsx",
                destination_path="output/q1-budget.xlsx"
            )

            # Copy Word template
            office_template(
                source_path="templates/sow.docx",
                destination_path="output/acme-sow.docx"
            )

            # Analyze Word template formatting
            office_template(
                source_path="templates/sow.docx",
                destination_path="",  # Not used for analyze
                operation="analyze"
            )

        Args:
            source_path: Path to the template file
            destination_path: Path for the copy (ignored for analyze)
            operation: "copy" to copy template, "analyze" to inspect formatting

        Returns:
            Dictionary with operation results
        """
        source_path = _resolve_file_path(source_path)
        doc_format = _detect_format(source_path)

        if doc_format is None:
            return _unsupported_format_error(source_path)

        # Excel template
        if doc_format == "excel":
            if not _has_tool(self, "excel_copy_template"):
                return {"error": "Excel support not available"}

            if operation == "copy":
                return self.tool_excel_copy_template(
                    template_path=source_path,
                    output_path=destination_path,
                )
            elif operation == "analyze":
                return {"error": "Analyze not yet supported for Excel templates"}

        # Word template
        elif doc_format == "word":
            if operation == "copy":
                # For direct path copy, use shutil
                import shutil
                from pathlib import Path as PathLib
                src = PathLib(source_path)
                if not src.exists():
                    return {"error": f"Source file not found: {source_path}"}
                try:
                    shutil.copy2(source_path, destination_path)
                    return {
                        "success": True,
                        "source": source_path,
                        "destination": destination_path,
                    }
                except Exception as e:
                    return {"error": f"Failed to copy template: {e}"}
            elif operation == "analyze":
                if not _has_tool(self, "word_analyze_template_formatting"):
                    return {"error": "Word support not available"}
                formatting = self.tool_word_analyze_template_formatting(
                    file_path=source_path,
                )
                if "error" in formatting:
                    return formatting
                if _has_tool(self, "word_parse_sow_template"):
                    formatting["template_metadata"] = self.tool_word_parse_sow_template(
                        template_path=source_path,
                    )
                return formatting

        # PowerPoint template
        elif doc_format == "powerpoint":
            if not _has_tool(self, "pptx_copy_template"):
                return {"error": "PowerPoint support not available"}

            if operation == "copy":
                return self.tool_pptx_copy_template(
                    template_path=source_path,
                    output_path=destination_path,
                )
            elif operation == "analyze":
                return self.tool_pptx_analyze_layouts(
                    file_path=source_path,
                )

        return {"error": f"Unhandled format: {doc_format}"}

    # =========================================================================
    # office_audit - Consolidated audit operations
    # =========================================================================

    def tool_office_audit(
        self,
        file_path: str,
        checks: list[Literal[
            "placeholders", "completion", "tracking", "formatting",
            "empty_cells", "totals", "dates"
        ]] | None = None,
        audit_config: dict[str, Any] | None = None,
    ) -> dict[str, Any]:
        """Audit documents for completeness, placeholders, and issues.

        Replaces: excel_audit_placeholders, word_audit_completion, word_audit_sow,
        pptx_audit_placeholders

        Examples:
            # Check for unfilled placeholders
            office_audit(file_path="contract.docx", checks=["placeholders"])

            # Full completion audit
            office_audit(file_path="sow.docx", checks=["completion"])

            # Check Excel for placeholders
            office_audit(file_path="estimate.xlsx", checks=["placeholders"])

            # Multiple checks
            office_audit(
                file_path="document.docx",
                checks=["placeholders", "tracking"]
            )

        Args:
            file_path: Path to the document
            checks: List of checks to perform:
                    - "placeholders": Find unfilled <...>, [...], [TBD] patterns
                    - "completion": Full completion audit (Word SOW)
                    - "tracking": Check for pending track changes
                    - "formatting": Check for formatting issues
                    - "empty_cells": Check required Excel cells for empty values
                    - "totals": Verify Excel totals based on configured ranges
                    - "dates": Validate Excel date formats (MM/DD/YYYY)
                audit_config: Optional configuration for Excel checks:
                      - required_cells: list of cell refs to check for empty values
                      - date_cells: list of cell refs to validate as MM/DD/YYYY
                      - totals: list of dicts with sum_range, target, and optional tolerance

        Returns:
            Dictionary with audit findings
        """
        file_path = _resolve_file_path(file_path)
        doc_format = _detect_format(file_path)

        if doc_format is None:
            return _unsupported_format_error(file_path)

        if not checks:
            checks = ["placeholders"]  # Default check

        results = {}

        # Excel audit
        if doc_format == "excel":
            if not _has_tool(self, "excel_audit_placeholders"):
                return {"error": "Excel support not available"}

            if "placeholders" in checks:
                results["placeholders"] = self.tool_excel_audit_placeholders(
                    file_path=file_path,
                )
            if "tracking" in checks:
                results["tracking"] = self.tool_excel_get_change_log(
                    file_path=file_path,
                )

            if any(check in checks for check in ["empty_cells", "totals", "dates"]):
                results.update(
                    _audit_excel_additional(
                        file_path,
                        audit_config or {},
                        checks,
                    )
                )

        # Word audit
        elif doc_format == "word":
            if not _has_tool(self, "word_audit_sow"):
                return {"error": "Word support not available"}

            if "placeholders" in checks:
                results["placeholders"] = self.tool_word_audit_sow(
                    file_path=file_path,
                )
            if "completion" in checks:
                results["completion"] = self.tool_word_audit_completion(
                    file_path=file_path,
                )
            if "tracking" in checks:
                results["tracking"] = self.tool_word_check_tracking(
                    file_path=file_path,
                )

        # PowerPoint audit
        elif doc_format == "powerpoint":
            if not _has_tool(self, "pptx_audit_placeholders"):
                return {"error": "PowerPoint support not available"}

            if "placeholders" in checks:
                results["placeholders"] = self.tool_pptx_audit_placeholders(
                    file_path=file_path,
                )

        return {
            "file": file_path,
            "checks_performed": checks,
            "results": results,
        }


def _audit_excel_additional(
    file_path: str,
    audit_config: dict[str, Any],
    checks: list[str],
) -> dict[str, Any]:
    """Run additional Excel checks beyond placeholder auditing."""
    if not HAS_OPENPYXL:
        return {"error": "openpyxl not installed. Run: pip install openpyxl"}

    required_cells = audit_config.get("required_cells", [])
    date_cells = audit_config.get("date_cells", [])
    totals = audit_config.get("totals", [])

    results: dict[str, Any] = {}

    try:
        wb = load_workbook(file_path, data_only=True, keep_vba=True)
    except Exception as e:
        return {"error": f"Failed to load workbook: {e}"}

    try:
        if "empty_cells" in checks:
            missing = []
            for ref in required_cells:
                sheet_name, cell_ref = _parse_cell_reference(ref)
                target_sheet = sheet_name or wb.active.title
                if target_sheet not in wb.sheetnames:
                    missing.append({"cell": ref, "error": f"Sheet not found: {target_sheet}"})
                    continue
                value = wb[target_sheet][cell_ref].value
                if value is None or str(value).strip() == "":
                    missing.append({"cell": ref, "value": value})
            results["empty_cells"] = {
                "missing": missing,
                "count": len(missing),
            }

        if "dates" in checks:
            invalid = []
            date_pattern = re.compile(r"^(0[1-9]|1[0-2])/(0[1-9]|[12]\d|3[01])/\d{4}$")

            for ref in date_cells:
                sheet_name, cell_ref = _parse_cell_reference(ref)
                target_sheet = sheet_name or wb.active.title
                if target_sheet not in wb.sheetnames:
                    invalid.append({"cell": ref, "error": f"Sheet not found: {target_sheet}"})
                    continue
                value = wb[target_sheet][cell_ref].value
                if isinstance(value, (datetime, date)):
                    continue
                value_text = "" if value is None else str(value).strip()
                if not date_pattern.match(value_text):
                    invalid.append({"cell": ref, "value": value})
            results["dates"] = {
                "invalid": invalid,
                "count": len(invalid),
            }

        if "totals" in checks:
            mismatches = []
            for rule in totals:
                sum_range = rule.get("sum_range")
                target = rule.get("target")
                tolerance = float(rule.get("tolerance", 0.01))

                if not sum_range or not target:
                    mismatches.append({"error": "Totals rule requires sum_range and target"})
                    continue

                sum_sheet, sum_ref = _parse_cell_reference(sum_range)
                target_sheet, target_ref = _parse_cell_reference(target)

                sum_sheet = sum_sheet or wb.active.title
                target_sheet = target_sheet or wb.active.title

                if sum_sheet not in wb.sheetnames:
                    mismatches.append({"sum_range": sum_range, "error": f"Sheet not found: {sum_sheet}"})
                    continue
                if target_sheet not in wb.sheetnames:
                    mismatches.append({"target": target, "error": f"Sheet not found: {target_sheet}"})
                    continue

                try:
                    min_row, min_col, max_row, max_col = _get_range_bounds(sum_ref)
                except ValueError as e:
                    mismatches.append({"sum_range": sum_range, "error": str(e)})
                    continue

                ws = wb[sum_sheet]
                total_value = 0.0

                for row in ws.iter_rows(
                    min_row=min_row,
                    max_row=max_row,
                    min_col=min_col,
                    max_col=max_col,
                ):
                    for cell in row:
                        value = cell.value
                        if isinstance(value, (int, float)):
                            total_value += float(value)

                target_value = wb[target_sheet][target_ref].value
                if not isinstance(target_value, (int, float)):
                    mismatches.append({
                        "sum_range": sum_range,
                        "target": target,
                        "expected": total_value,
                        "actual": target_value,
                        "error": "Target cell is not numeric",
                    })
                    continue

                delta = abs(total_value - float(target_value))
                if delta > tolerance:
                    mismatches.append({
                        "sum_range": sum_range,
                        "target": target,
                        "expected": total_value,
                        "actual": target_value,
                        "delta": delta,
                        "tolerance": tolerance,
                    })

            results["totals"] = {
                "mismatches": mismatches,
                "count": len(mismatches),
            }
    finally:
        _close_openpyxl_workbook(wb)

    return results


class OfficeImageTools:
    """Image insertion tools for Office documents."""

    # =========================================================================
    # office_image - Insert images into documents
    # =========================================================================

    def tool_office_image(
        self,
        file_path: str,
        image_path: str,
        target: str | None = None,
        width_inches: float | None = None,
        height_inches: float | None = None,
        output_path: str | None = None,
    ) -> dict[str, Any]:
        """Insert an image into a Word, Excel, or PowerPoint document.

        Auto-detects document format from file extension and inserts the image
        at the specified location. Supports PNG, JPG/JPEG, and GIF formats.

        Examples:
            # Insert image at end of Word document
            office_image(
                file_path="report.docx",
                image_path="logo.png",
                width_inches=2.0
            )

            # Insert image at specific paragraph in Word
            office_image(
                file_path="report.docx",
                image_path="chart.png",
                target="after:Executive Summary",
                width_inches=5.0
            )

            # Insert image in Excel cell
            office_image(
                file_path="data.xlsx",
                image_path="logo.png",
                target="A1",
                width_inches=1.5
            )

            # Insert image on specific Excel sheet
            office_image(
                file_path="data.xlsx",
                image_path="chart.png",
                target="Sheet2!B5"
            )

            # Insert image on PowerPoint slide
            office_image(
                file_path="deck.pptx",
                image_path="diagram.png",
                target="slide:2",
                width_inches=4.0,
                height_inches=3.0
            )

        Args:
            file_path: Path to the document (.docx, .xlsx, .pptx)
            image_path: Path to the image file (.png, .jpg, .jpeg, .gif)
            target: Where to insert the image:
                   - Word: "after:Section Title" or "end" (default)
                   - Excel: cell reference like "A1" or "Sheet1!B5"
                   - PowerPoint: "slide:N" where N is slide number (1-based)
            width_inches: Image width in inches (height auto-scales if not set)
            height_inches: Image height in inches (width auto-scales if not set)
            output_path: Optional output path (defaults to overwriting input)

        Returns:
            Dictionary with insertion result
        """
        file_path = _resolve_file_path(file_path)
        doc_format = _detect_format(file_path)

        if doc_format is None:
            return _unsupported_format_error(file_path)

        # Validate image file
        image_path = _resolve_file_path(image_path)
        if not os.path.exists(image_path):
            return {"error": f"Image file not found: {image_path}"}

        image_ext = Path(image_path).suffix.lower()
        if image_ext not in (".png", ".jpg", ".jpeg", ".gif", ".svg"):
            return {"error": f"Unsupported image format: {image_ext}. Use PNG, JPG, GIF, or SVG."}

        # SVG is not supported in Excel (openpyxl cannot embed SVG)
        if image_ext == ".svg" and doc_format == "excel":
            return {"error": "SVG images are not supported in Excel documents. Use PNG or JPG."}

        final_path = output_path or file_path

        try:
            if doc_format == "word":
                return self._insert_image_word(
                    file_path, image_path, target, width_inches, height_inches, final_path
                )
            elif doc_format == "excel":
                return self._insert_image_excel(
                    file_path, image_path, target, width_inches, height_inches, final_path
                )
            elif doc_format == "powerpoint":
                return self._insert_image_pptx(
                    file_path, image_path, target, width_inches, height_inches, final_path
                )
            else:
                return {"error": f"Image insertion not supported for {doc_format}"}
        except Exception as e:
            return {"error": str(e)}

    @staticmethod
    def _make_fallback_png() -> bytes:
        """Create a minimal 1x1 transparent PNG for SVG fallback."""
        import struct
        import zlib

        def _chunk(ctype: bytes, data: bytes) -> bytes:
            c = ctype + data
            return (
                struct.pack(">I", len(data))
                + c
                + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
            )

        sig = b"\x89PNG\r\n\x1a\n"
        ihdr = _chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 6, 0, 0, 0))
        idat = _chunk(b"IDAT", zlib.compress(b"\x00\x00\x00\x00\x00"))
        iend = _chunk(b"IEND", b"")
        return sig + ihdr + idat + iend

    @staticmethod
    def _add_svg_extension(blip_element, svg_rId: str) -> None:
        """Add SVG extension to an a:blip element (OOXML 2016 SVG spec)."""
        from lxml import etree

        NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        NS_ASVG = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"

        ext_lst = etree.SubElement(blip_element, f"{{{NS_A}}}extLst")
        ext = etree.SubElement(ext_lst, f"{{{NS_A}}}ext")
        ext.set("uri", "{96DAC541-7B7A-43D3-8B79-37D633B846F1}")
        svg_blip = etree.SubElement(ext, f"{{{NS_ASVG}}}svgBlip")
        svg_blip.set(f"{{{NS_R}}}embed", svg_rId)

    def _insert_image_word(
        self,
        file_path: str,
        image_path: str,
        target: str | None,
        width_inches: float | None,
        height_inches: float | None,
        output_path: str,
    ) -> dict[str, Any]:
        """Insert image into Word document.

        For SVG files, embeds the SVG via the OOXML 2016 SVG extension with
        a 1px PNG fallback for older readers.
        """
        try:
            from docx import Document
            from docx.shared import Inches
        except ImportError:
            return {"error": "python-docx not installed"}

        doc = Document(file_path)
        is_svg = Path(image_path).suffix.lower() == ".svg"

        # Convert dimensions
        width = Inches(width_inches) if width_inches else None
        height = Inches(height_inches) if height_inches else None

        # For SVG: compute dimensions from SVG metadata if not specified
        if is_svg and not (width or height):
            w_in, h_in = self._get_image_dimensions(image_path)
            width = Inches(w_in)
            height = Inches(h_in)

        def _add_picture_to_run(run):
            """Add picture (raster or SVG) to a run and return inline shape."""
            if is_svg:
                import io

                from docx.opc.constants import RELATIONSHIP_TYPE as RT
                from docx.opc.packuri import PackURI
                from docx.opc.part import Part as DocxPart

                # Insert 1px fallback PNG — creates the drawing XML
                fallback = self._make_fallback_png()
                inline_shape = run.add_picture(
                    io.BytesIO(fallback), width=width, height=height
                )

                # Add SVG as an OPC part
                svg_bytes = Path(image_path).read_bytes()
                idx = len([r for r in doc.part.rels.values() if "image" in r.reltype]) + 1
                svg_part = DocxPart(
                    PackURI(f"/word/media/image_svg{idx}.svg"),
                    "image/svg+xml",
                    svg_bytes,
                    doc.part.package,
                )
                svg_rId = doc.part.relate_to(svg_part, RT.IMAGE)

                # Patch the blip with the SVG extension
                NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
                blip = inline_shape._inline.graphic.graphicData.find(
                    f".//{{{NS_A}}}blip"
                )
                self._add_svg_extension(blip, svg_rId)
            else:
                run.add_picture(image_path, width=width, height=height)

        # Parse target
        if target and target.startswith("after:"):
            section_title = target[6:]  # Remove "after:" prefix
            inserted = False
            for _i, para in enumerate(doc.paragraphs):
                if section_title.lower() in para.text.lower():
                    new_para = doc.add_paragraph()
                    para._element.addnext(new_para._element)
                    _add_picture_to_run(new_para.add_run())
                    inserted = True
                    break
            if not inserted:
                return {"error": f"Section not found: {section_title}"}
        else:
            # Insert at end of document
            para = doc.add_paragraph()
            _add_picture_to_run(para.add_run())

        doc.save(output_path)

        return {
            "status": "success",
            "file": output_path,
            "image": image_path,
            "location": target or "end",
        }

    def _insert_image_excel(
        self,
        file_path: str,
        image_path: str,
        target: str | None,
        width_inches: float | None,
        height_inches: float | None,
        output_path: str,
    ) -> dict[str, Any]:
        """Insert image into Excel workbook."""
        try:
            from openpyxl import load_workbook
            from openpyxl.drawing.image import Image
        except ImportError:
            return {"error": "openpyxl not installed"}

        wb = load_workbook(file_path)

        # Parse target for sheet and cell
        sheet_name = None
        cell_ref = "A1"

        if target:
            if "!" in target:
                sheet_name, cell_ref = target.split("!", 1)
            else:
                cell_ref = target

        ws = wb[sheet_name] if sheet_name else wb.active

        # Create and configure image
        img = Image(image_path)

        # Set dimensions (openpyxl uses pixels, 96 DPI standard)
        if width_inches:
            img.width = int(width_inches * 96)
        if height_inches:
            img.height = int(height_inches * 96)

        # Anchor image to cell
        img.anchor = cell_ref

        ws.add_image(img)
        try:
            wb.save(output_path)
        finally:
            _close_openpyxl_workbook(wb)

        return {
            "status": "success",
            "file": output_path,
            "image": image_path,
            "sheet": ws.title,
            "cell": cell_ref,
        }

    def _get_image_dimensions(self, image_path: str) -> tuple[float, float]:
        """Get image dimensions in inches.

        Handles both raster images (via PIL) and SVG files (via XML parsing).

        Args:
            image_path: Path to the image file

        Returns:
            Tuple of (width_inches, height_inches)
        """
        import re
        from pathlib import Path

        path = Path(image_path)

        if path.suffix.lower() == ".svg":
            # Parse SVG dimensions from XML
            import xml.etree.ElementTree as ET

            tree = ET.parse(image_path)
            root = tree.getroot()

            # Try width/height attributes first
            width_attr = root.get("width", "")
            height_attr = root.get("height", "")

            def parse_svg_length(value: str, default_dpi: float = 96.0) -> float | None:
                """Parse SVG length value to inches."""
                if not value:
                    return None
                value = value.strip()
                # Match number with optional unit
                match = re.match(r"^([\d.]+)\s*(px|pt|in|mm|cm|em|%)?$", value, re.I)
                if not match:
                    return None
                num = float(match.group(1))
                unit = (match.group(2) or "px").lower()

                if unit == "in":
                    return num
                elif unit == "px":
                    return num / default_dpi
                elif unit == "pt":
                    return num / 72.0
                elif unit == "mm":
                    return num / 25.4
                elif unit == "cm":
                    return num / 2.54
                elif unit == "em":
                    return num * 16 / default_dpi  # Assume 16px em
                elif unit == "%":
                    return None  # Can't resolve percentage without context
                return num / default_dpi

            width_in = parse_svg_length(width_attr)
            height_in = parse_svg_length(height_attr)

            # Fall back to viewBox if width/height not usable
            if width_in is None or height_in is None:
                viewbox = root.get("viewBox", "")
                if viewbox:
                    parts = viewbox.split()
                    if len(parts) == 4:
                        vb_width = float(parts[2])
                        vb_height = float(parts[3])
                        # viewBox is in user units, assume 96 DPI
                        if width_in is None:
                            width_in = vb_width / 96.0
                        if height_in is None:
                            height_in = vb_height / 96.0

            # Default fallback
            if width_in is None:
                width_in = 6.0  # Default 6 inches
            if height_in is None:
                height_in = 4.0  # Default 4 inches

            return (width_in, height_in)

        else:
            # Use PIL for raster images
            from PIL import Image as PILImage

            with PILImage.open(image_path) as img:
                img_width_px, img_height_px = img.size
                # Check for DPI info in image metadata
                dpi = img.info.get("dpi", (96, 96))
                if isinstance(dpi, tuple):
                    dpi_x, dpi_y = dpi
                else:
                    dpi_x = dpi_y = dpi
                return (img_width_px / dpi_x, img_height_px / dpi_y)

    def _insert_image_pptx(
        self,
        file_path: str,
        image_path: str,
        target: str | None,
        width_inches: float | None,
        height_inches: float | None,
        output_path: str,
    ) -> dict[str, Any]:
        """Insert image into PowerPoint presentation.

        Centers the image horizontally and vertically on the slide.
        If dimensions are provided, uses those; otherwise scales to fit
        within the slide with reasonable margins.

        Supports raster images (PNG, JPG, GIF) and SVG files.
        """
        try:
            from pptx import Presentation
            from pptx.util import Emu, Inches
        except ImportError as e:
            return {"error": f"Required package not installed: {e}"}

        prs = Presentation(file_path)

        # Parse target slide
        slide_num = 1  # Default to first slide
        if target and target.startswith("slide:"):
            try:
                slide_num = int(target[6:])
            except ValueError:
                return {"error": f"Invalid slide number: {target}"}

        if slide_num < 1 or slide_num > len(prs.slides):
            return {"error": f"Slide {slide_num} out of range (1-{len(prs.slides)})"}

        slide = prs.slides[slide_num - 1]

        # Get slide dimensions (in EMUs)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Get image native dimensions (handles both raster and SVG)
        try:
            img_width_native, img_height_native = self._get_image_dimensions(image_path)
        except Exception as e:
            return {"error": f"Could not read image dimensions: {e}"}

        # Determine final image dimensions
        if width_inches and height_inches:
            # Both specified - use as-is
            final_width = Inches(width_inches)
            final_height = Inches(height_inches)
        elif width_inches:
            # Width specified - scale height proportionally
            final_width = Inches(width_inches)
            scale = width_inches / img_width_native
            final_height = Inches(img_height_native * scale)
        elif height_inches:
            # Height specified - scale width proportionally
            final_height = Inches(height_inches)
            scale = height_inches / img_height_native
            final_width = Inches(img_width_native * scale)
        else:
            # No dimensions specified - fit within slide with margins
            margin = Inches(0.5)
            max_width = slide_width - (2 * margin)
            max_height = slide_height - (2 * margin)

            # Calculate scale to fit within bounds
            width_scale = max_width / Inches(img_width_native)
            height_scale = max_height / Inches(img_height_native)
            scale = min(width_scale, height_scale, 1.0)  # Don't upscale

            final_width = Emu(int(Inches(img_width_native) * scale))
            final_height = Emu(int(Inches(img_height_native) * scale))

        # Center the image on the slide
        left = (slide_width - final_width) // 2
        top = (slide_height - final_height) // 2

        # Add picture to slide
        is_svg = Path(image_path).suffix.lower() == ".svg"
        if is_svg:
            import io

            from pptx.opc.constants import RELATIONSHIP_TYPE as RT
            from pptx.opc.package import Part as PptxPart
            from pptx.opc.packuri import PackURI
            from pptx.oxml import parse_xml

            # 1. Add fallback 1px PNG via normal mechanism
            fallback = self._make_fallback_png()
            png_image_part, png_rId = slide.part.get_or_add_image_part(
                io.BytesIO(fallback)
            )

            # 2. Add SVG as an OPC part
            svg_bytes = Path(image_path).read_bytes()
            idx = len([r for r in slide.part.rels.values() if "image" in r.reltype]) + 1
            svg_part = PptxPart(
                PackURI(f"/ppt/media/image_svg{idx}.svg"),
                "image/svg+xml",
                prs.part.package,
                svg_bytes,
            )
            svg_rId = slide.part.relate_to(svg_part, RT.IMAGE)

            # 3. Build p:pic element with SVG extension in a:blip
            NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
            NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
            NS_ASVG = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"

            shape_id = max(
                (int(sp.get("id", 0)) for sp in slide.shapes._spTree.iter()),
                default=0,
            ) + 1

            pic_xml = (
                f'<p:pic xmlns:a="{NS_A}" xmlns:r="{NS_R}" xmlns:p="{NS_P}">'
                f"  <p:nvPicPr>"
                f'    <p:cNvPr id="{shape_id}" name="Picture {shape_id - 1}"'
                f'     descr="SVG Image"/>'
                f"    <p:cNvPicPr>"
                f'      <a:picLocks noChangeAspect="1"/>'
                f"    </p:cNvPicPr>"
                f"    <p:nvPr/>"
                f"  </p:nvPicPr>"
                f"  <p:blipFill>"
                f'    <a:blip r:embed="{png_rId}">'
                f"      <a:extLst>"
                f'        <a:ext uri="{{96DAC541-7B7A-43D3-8B79-37D633B846F1}}">'
                f'          <asvg:svgBlip xmlns:asvg="{NS_ASVG}"'
                f'           r:embed="{svg_rId}"/>'
                f"        </a:ext>"
                f"      </a:extLst>"
                f"    </a:blip>"
                f"    <a:stretch><a:fillRect/></a:stretch>"
                f"  </p:blipFill>"
                f"  <p:spPr>"
                f"    <a:xfrm>"
                f'      <a:off x="{left}" y="{top}"/>'
                f'      <a:ext cx="{final_width}" cy="{final_height}"/>'
                f"    </a:xfrm>"
                f'    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
                f"  </p:spPr>"
                f"</p:pic>"
            )
            slide.shapes._spTree.append(parse_xml(pic_xml))
        else:
            slide.shapes.add_picture(
                image_path, left, top, width=final_width, height=final_height
            )

        safe_save_pptx(prs, output_path)

        return {
            "status": "success",
            "file": output_path,
            "image": image_path,
            "slide": slide_num,
            "position": {
                "left_inches": left / 914400,
                "top_inches": top / 914400,
                "width_inches": final_width / 914400,
                "height_inches": final_height / 914400,
            },
            "centered": True,
        }
